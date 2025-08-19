

#views.py original
from django.contrib.auth.models import User
from django.contrib.auth import authenticate, login
from rest_framework import status
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.permissions import IsAuthenticated, AllowAny  
from notebook.models import Document as NotebookDocument, ProcessedIndex as NotebookProcessedIndex

from moviepy import VideoFileClip, AudioFileClip
import time
import faiss
import numpy as np
import os
import pickle
import json
import re
from datetime import datetime
from django.core.files.storage import default_storage
from django.conf import settings
from django.utils.deprecation import MiddlewareMixin
import google.generativeai as genai  
from core.utils import update_project_timestamp
from .models import (
    ChatHistory,
    ChatMessage,
    Document,
    ProcessedIndex,
    UserAPITokens,
    ConversationMemoryBuffer,
    UserModulePermissions,
    UserUploadPermissions,
    ConversationTransaction,
    UserTransaction,
    DocumentTransaction,
    TransactionType,
    DocumentProcessingStatus,

)
import uuid
from rest_framework.authtoken.models import Token
from django.utils.safestring import mark_safe
import logging
from .models import UserProfile
logger = logging.getLogger(__name__)
from core.models import Project
from openai import OpenAI
from dotenv import load_dotenv
import os

from io import StringIO, BytesIO
from django.http import HttpResponse, FileResponse, JsonResponse
from django.utils import timezone
from django.utils.html import strip_tags
from django.shortcuts import get_object_or_404
from django.contrib.auth import update_session_auth_hash
from llama_parse import LlamaParse
import requests
from bs4 import BeautifulSoup
import re
from duckduckgo_search import DDGS
from urllib.parse import urlparse
import nltk
from nltk.tokenize import sent_tokenize
from django.utils.html import strip_tags
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

from moviepy import VideoFileClip

import tempfile
from django.core.files import File

from chat.models import UserAPITokens
from io import BytesIO
from datetime import datetime, timedelta
from docx import Document as DocxDocument
from docx.shared import Pt, Inches, RGBColor
from docx.enum.style import WD_STYLE_TYPE
from bs4 import BeautifulSoup, NavigableString
from django.http import HttpResponse
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated
from rest_framework import status
from rest_framework.response import Response
import logging
from django.db.models import F
from django.shortcuts import redirect
from django.conf import settings
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import AllowAny
from rest_framework.authtoken.models import Token
from django.contrib.auth.models import User
from msal import ConfidentialClientApplication
import uuid
import requests
 

logger = logging.getLogger(__name__)



load_dotenv()

OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

if not OPENAI_API_KEY:
    raise ValueError("Missing required API keys in environment variables")



client = OpenAI(api_key=OPENAI_API_KEY)

# # Configure Google Generative AI
# GOOGLE_API_KEY = "AIzaSyDOKm5KYY6LjLa20IbZg027fQauwyMOKWQ"
# genai.configure(api_key=GOOGLE_API_KEY)
# # model = genai.GenerativeModel('gemini-1.5-flash')
# GENERATIVE_MODEL = genai.GenerativeModel('gemini-1.5-flash', 
#     generation_config={
#         'temperature': 0.7,
#         'max_output_tokens': 1024
#     },
#     safety_settings={
#         genai.types.HarmCategory.HARM_CATEGORY_HATE_SPEECH: genai.types.HarmBlockThreshold.BLOCK_NONE,
#         genai.types.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: genai.types.HarmBlockThreshold.BLOCK_NONE,
#         genai.types.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: genai.types.HarmBlockThreshold.BLOCK_NONE,
#         genai.types.HarmCategory.HARM_CATEGORY_HARASSMENT: genai.types.HarmBlockThreshold.BLOCK_NONE
#     }
# )

class SignupView(APIView):
    permission_classes = [AllowAny]
    authentication_classes = []
 
    def post(self, request):
        from dotenv import load_dotenv
        import os
 
        load_dotenv()
        # Extract form data
        username = request.data.get('username')
        email = request.data.get('email')
        password = request.data.get('password')
 
        # Optional API tokens provided by user during signup
        nebius_token = request.data.get('nebius_token', '')
        gemini_token = request.data.get('gemini_token', '')
        llama_token = request.data.get('llama_token', '')
        huggingface_token = request.data.get('huggingface_token', '')
 
        if not username or not email or not password:
            return Response({
                'error': 'Please provide username, email, and password'
            }, status=status.HTTP_400_BAD_REQUEST)
 
        if User.objects.filter(username=username).exists():
            return Response({
                'error': 'Username already exists'
            }, status=status.HTTP_400_BAD_REQUEST)
 
        # Create new user
        try:
            # Create the user
            user = User.objects.create_user(
                username=username,
                email=email,
                password=password
            )
 
            # Set user tokens: Use user input if available, otherwise fall back to env defaults
            # ✅ Update the existing UserAPITokens created by signal
            user.api_tokens.huggingface_token = huggingface_token or os.getenv("HUGGINGFACE_TOKEN", "")
            user.api_tokens.gemini_token = gemini_token or os.getenv("GOOGLE_API_KEY", "")
            user.api_tokens.nebius_token = nebius_token or os.getenv("NEBIUS_TOKEN", "")
            user.api_tokens.llama_token = llama_token or os.getenv("LLAMA_TOKEN", "")
            user.api_tokens.save()
 
 
            # Generate token
            token, _ = Token.objects.get_or_create(user=user)
            return Response({
                'message': 'User created successfully',
                'token': token.key,
                'username': user.username
            }, status=status.HTTP_201_CREATED)
 
        except Exception as e:
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class LoginView(APIView):
    # Explicitly set permission to allow any user (including unauthenticated)
    permission_classes = [AllowAny]
    authentication_classes = []  # Disable authentication checks

    def post(self, request):
        username = request.data.get('username')
        password = request.data.get('password')

        # Validate input
        if not username or not password:
            return Response({
                'error': 'Please provide username and password'
            }, status=status.HTTP_400_BAD_REQUEST)

        # Authenticate user
        user = authenticate(username=username, password=password)

        if user:
            # Generate or get existing token
            token, _ = Token.objects.get_or_create(user=user)
            
            return Response({
                'token': token.key,
                'username': user.username
            }, status=status.HTTP_200_OK)
        
        return Response({
            'error': 'Invalid credentials'
        }, status=status.HTTP_401_UNAUTHORIZED)
    

class MicrosoftSSOView(APIView):
    """Single view to handle Microsoft SSO - direct redirect approach"""
    permission_classes = [AllowAny]
    authentication_classes = []
   
    def get_msal_app(self):
        return ConfidentialClientApplication(
            client_id=settings.AZURE_AD['CLIENT_ID'],
            authority=settings.AZURE_AD['AUTHORITY'],
            client_credential=settings.AZURE_AD['CLIENT_SECRET']
        )
   
    def get(self, request):
        """Initiate Microsoft SSO - redirect to Microsoft login"""
        try:
            msal_app = self.get_msal_app()
           
            # Generate state for security
            state = str(uuid.uuid4())
            request.session['oauth_state'] = state
           
            # Use the correct redirect URI
            auth_url = msal_app.get_authorization_request_url(
                scopes=settings.AZURE_AD['SCOPE'],
                state=state,
                redirect_uri=settings.AZURE_AD['REDIRECT_URI']
            )
           
            logger.info(f"Redirecting to Microsoft auth URL: {auth_url}")
            return redirect(auth_url)
           
        except Exception as e:
            logger.error(f"Microsoft SSO initiation error: {str(e)}")
            return redirect(f"{settings.FRONTEND_URL}/login?error=sso_failed")
 
 
class MicrosoftSSOCallbackView(APIView):
    """Handle Microsoft SSO callback after authentication"""
    permission_classes = [AllowAny]
    authentication_classes = []
   
    def get_msal_app(self):
        return ConfidentialClientApplication(
            client_id=settings.AZURE_AD['CLIENT_ID'],
            authority=settings.AZURE_AD['AUTHORITY'],
            client_credential=settings.AZURE_AD['CLIENT_SECRET']
        )
   
    def get_user_info(self, access_token):
        """Get user information from Microsoft Graph API"""
        headers = {
            'Authorization': f'Bearer {access_token}',
            'Content-Type': 'application/json'
        }
       
        try:
            response = requests.get(
                'https://graph.microsoft.com/v1.0/me',
                headers=headers,
                timeout=10
            )
            if response.status_code == 200:
                return response.json()
            else:
                logger.error(f"Failed to get user info: {response.status_code} - {response.text}")
            return None
        except Exception as e:
            logger.error(f"Error getting user info: {str(e)}")
            return None
   
    def create_or_update_user(self, user_info):
        """Create or update user - they're already validated by Microsoft"""
 
        from dotenv import load_dotenv
        import os
 
        load_dotenv()
 
 
        email = user_info.get('mail') or user_info.get('userPrincipalName')
        display_name = user_info.get('displayName', '')
        given_name = user_info.get('givenName', '')
        surname = user_info.get('surname', '')
 
        if not email:
            logger.error("No email found in user info")
            return None
 
        try:
            # Get existing user
            user = User.objects.get(email=email)
            # Update user info in case it changed
            user.first_name = given_name
            user.last_name = surname
            user.save()
            logger.info(f"Updated existing user: {email}")
        except User.DoesNotExist:
            # Create new user - Microsoft already validated them
            username = email.split('@')[0]
 
            # Ensure unique username
            counter = 1
            original_username = username
            while User.objects.filter(username=username).exists():
                username = f"{original_username}_{counter}"
                counter += 1
 
            user = User.objects.create_user(
                username=username,
                email=email,
                first_name=given_name,
                last_name=surname
            )
            user.set_unusable_password()
            user.save()
 
            logger.info(f"Created new user: {email}")
 
 
            # ✅ Update the existing UserAPITokens created by signal
            user.api_tokens.huggingface_token = os.getenv("HUGGINGFACE_TOKEN", "")
            user.api_tokens.gemini_token =  os.getenv("GOOGLE_API_KEY", "")
            user.api_tokens.nebius_token =  os.getenv("NEBIUS_TOKEN", "")
            user.api_tokens.llama_token =   os.getenv("LLAMA_TOKEN", "")
            user.api_tokens.save()
 
            # 🔒 Disable all modules for the new user
            try:
                module_perm = user.module_permissions
            except UserModulePermissions.DoesNotExist:
                module_perm = UserModulePermissions.objects.create(user=user)
 
            # Assuming these are your modules (customize as needed)
            all_modules = ["document-qa", "ad-campaign-generator"]
            module_perm.disabled_modules = {module: True for module in all_modules}
            module_perm.save()
 
        return user
   
    def get(self, request):
        """Handle callback from Microsoft with authorization code"""
        code = request.GET.get('code')
        state = request.GET.get('state')
        error = request.GET.get('error')
        error_description = request.GET.get('error_description', '')
       
        logger.info(f"SSO Callback received - Code: {'Yes' if code else 'No'}, State: {state}, Error: {error}")
       
        # Check for errors from Microsoft
        if error:
            logger.error(f"Microsoft SSO error: {error} - {error_description}")
            return redirect(f"{settings.FRONTEND_URL}/login?error=microsoft_{error}")
       
        if not code:
            logger.error("No authorization code received")
            return redirect(f"{settings.FRONTEND_URL}/login?error=no_code")
       
        # Validate state for security
        stored_state = request.session.get('oauth_state')
        if stored_state and stored_state != state:
            logger.error(f"State mismatch - Stored: {stored_state}, Received: {state}")
            return redirect(f"{settings.FRONTEND_URL}/login?error=invalid_state")
       
        try:
            msal_app = self.get_msal_app()
           
            # Exchange authorization code for tokens
            result = msal_app.acquire_token_by_authorization_code(
                code,
                scopes=settings.AZURE_AD['SCOPE'],
                redirect_uri=settings.AZURE_AD['REDIRECT_URI']
            )
           
            if 'error' in result:
                logger.error(f"Token acquisition error: {result.get('error')} - {result.get('error_description')}")
                return redirect(f"{settings.FRONTEND_URL}/login?error=token_failed")
           
            # Get user info from Microsoft
            access_token = result.get('access_token')
            if not access_token:
                logger.error("No access token received")
                return redirect(f"{settings.FRONTEND_URL}/login?error=no_access_token")
               
            user_info = self.get_user_info(access_token)
           
            if not user_info:
                logger.error("Failed to get user info from Microsoft")
                return redirect(f"{settings.FRONTEND_URL}/login?error=user_info_failed")
           
            # Create/update user in our system
            user = self.create_or_update_user(user_info)
           
            if not user:
                logger.error("Failed to create or update user")
                return redirect(f"{settings.FRONTEND_URL}/login?error=user_creation_failed")
           
            # Generate auth token
            token, created = Token.objects.get_or_create(user=user)
            logger.info(f"Generated token for user: {user.username}")
           
            # Clean up session
            if 'oauth_state' in request.session:
                del request.session['oauth_state']
           
            # REDIRECT TO FRONTEND WITH TOKEN - This is the key change!
            # Instead of returning JSON, redirect to frontend with token and username
            redirect_url = f"{settings.FRONTEND_URL}/sso-callback?token={token.key}&username={user.username}"
            return redirect(redirect_url)
       
        except Exception as e:
            logger.error(f"SSO callback exception: {str(e)}")
            return redirect(f"{settings.FRONTEND_URL}/login?error=sso_exception")

class ChangePasswordView(APIView):
    permission_classes = [IsAuthenticated]
 
    def post(self, request):
        user = request.user
        current_password = request.data.get('current_password')
        new_password = request.data.get('new_password')
 
        if not current_password or not new_password:
            return Response({'message': 'Both current and new password are required'}, status=status.HTTP_400_BAD_REQUEST)
 
        if not user.check_password(current_password):  # Use user.check_password()
            return Response({'message': 'Current password is incorrect'}, status=status.HTTP_400_BAD_REQUEST)
 
        user.set_password(new_password)
        user.save()
        update_session_auth_hash(request, user)  # Important: Keep user logged in
 
        return Response({'message': 'Password updated successfully'}, status=status.HTTP_200_OK)

#new
# Fix for UserProfileView in views.py

#new

 
#new
from django.db.models import Sum,Count
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import IsAuthenticated
from rest_framework.parsers import MultiPartParser, FormParser
import os
import uuid
 
class UserProfileView(APIView):
    permission_classes = [IsAuthenticated]
    parser_classes = (MultiPartParser, FormParser)
 
    def get(self, request):
        user = request.user
        try:
            profile = UserProfile.objects.get(user=user)
            if profile.profile_picture:
                # Get the full URL for the profile picture
                profile_picture = profile.get_profile_picture_url()
            else:
                profile_picture = f'https://ui-avatars.com/api/?name={user.username}&background=random'
        except UserProfile.DoesNotExist:
            profile_picture = f'https://ui-avatars.com/api/?name={user.username}&background=random'
 
        # Get module permissions data
        try:
            module_permissions = user.module_permissions.disabled_modules
        except (AttributeError, UserModulePermissions.DoesNotExist):
            # Create module permissions if they don't exist
            module_permissions = {}
            UserModulePermissions.objects.create(user=user, disabled_modules={})
 
        # Get token usage and question count statistics
        token_stats = self.get_user_token_stats(user)
       
        # Get document statistics
        document_stats = self.get_user_document_stats(user)
        api_tokens = user.api_tokens if hasattr(user, 'api_tokens') else None
       
        print(token_stats)
        print(document_stats)
       
        return Response({
            'username': user.username,
            'email': user.email,
            'first_name': user.first_name,
            'last_name': user.last_name,
            'profile_picture': profile_picture,
            'date_joined': user.date_joined,
            'disabled_modules': module_permissions,
            'total_tokens_used': token_stats['total_tokens'],
            'total_input_tokens': token_stats['total_input_tokens'],
            'total_output_tokens': token_stats['total_output_tokens'],
            'total_questions_asked': token_stats['total_questions'],
            'total_pages_processed': document_stats['total_pages'],
            'total_documents_uploaded': document_stats['total_documents'],
            'token_limit': api_tokens.token_limit if api_tokens and api_tokens.token_limit is not None else None,
            'page_limit': (api_tokens.page_limit) if api_tokens and api_tokens.page_limit is not None else None,

        }, status=status.HTTP_200_OK)
 
    def get_user_token_stats(self, user):
        """
        Calculate total token usage and question count for the user
        """
        from notebook.models import ConversationTransaction
        # Get all ConversationTransaction records for this user
        # We need to filter through UserTransaction to get to ConversationTransaction
        conversation_stats = ConversationTransaction.objects.filter(
            user_transaction__user=user
        ).aggregate(
            total_tokens=Sum('total_tokens_used') or 0,
            total_input_tokens=Sum('input_api_tokens') or 0,
            total_output_tokens=Sum('output_api_tokens') or 0,
            total_questions=Sum('question_count') or 0,
        )
        print(conversation_stats)
       
        return {
            'total_tokens': conversation_stats['total_tokens'],
            'total_input_tokens': conversation_stats['total_input_tokens'],
            'total_output_tokens': conversation_stats['total_output_tokens'],
            'total_questions': conversation_stats['total_questions'],
        }
 
    def get_user_document_stats(self, user):
        """
        Calculate total document pages and document count for the user
        """
        from notebook.models import DocumentTransaction
        # Get all DocumentTransaction records for this user
        # We need to filter through UserTransaction to get to DocumentTransaction
        document_stats = DocumentTransaction.objects.filter(
            user_transaction__user=user
        ).aggregate(
            total_pages=Sum('no_pages') or 0,
            total_documents=Count('id'),
        )
       
        return {
            'total_pages': document_stats['total_pages'],
            'total_documents': document_stats['total_documents'],
        }
 
    def post(self, request):
        try:
            user = request.user
            profile_picture = request.FILES.get('profile_picture')
           
            if not profile_picture:
                return Response({
                    'error': 'No profile picture provided'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Validate file type
            allowed_types = ['image/jpeg', 'image/png', 'image/gif']
            if profile_picture.content_type not in allowed_types:
                return Response({
                    'error': 'Invalid file type. Only JPG, PNG, and GIF are allowed.'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Validate file size (2MB limit)
            if profile_picture.size > 2 * 1024 * 1024:
                return Response({
                    'error': 'File size too large. Maximum size is 2MB.'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Get or create profile
            profile, created = UserProfile.objects.get_or_create(user=user)
           
            # Delete old profile picture if it exists
            if profile.profile_picture:
                try:
                    old_file_path = profile.profile_picture.path
                    if os.path.exists(old_file_path):
                        os.remove(old_file_path)
                except Exception as e:
                    print(f"Error deleting old profile picture: {e}")
           
            # Generate unique filename
            file_extension = os.path.splitext(profile_picture.name)[1]
            unique_filename = f"{user.username}_{uuid.uuid4().hex[:8]}{file_extension}"
           
            # Save the new profile picture
            profile.profile_picture.save(
                unique_filename,
                profile_picture,
                save=True
            )
           
            # Build the full URL
            profile_picture_url = request.build_absolute_uri(profile.profile_picture.url)
           
            return Response({
                'message': 'Profile picture updated successfully',
                'profile_picture': profile_picture_url
            }, status=status.HTTP_200_OK)
           
        except Exception as e:
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class GetUserDocumentsView(APIView):
    def get(self, request):
        try:
            user = request.user
            main_project_id = request.query_params.get('main_project_id')
            print(f"Getting documents for user {user.username} and project {main_project_id}")
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            documents = Document.objects.filter(user=user, main_project_id=main_project_id).select_related('processedindex')
        
            document_list = []
            for doc in documents:
                try:
                    processed = doc.processedindex
                    document_data = {
                        'id': doc.id,
                        'filename': doc.filename,
                        'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M'),
                        'summary': processed.summary,
                        'follow_up_questions': [
                            processed.follow_up_question_1,
                            processed.follow_up_question_2,
                            processed.follow_up_question_3
                        ] if all([processed.follow_up_question_1,
                                processed.follow_up_question_2,
                                processed.follow_up_question_3]) else []
                    }
                    document_list.append(document_data)

                   
                except ProcessedIndex.DoesNotExist:
                    document_data = {
                        'id': doc.id,
                        'filename': doc.filename,
                        'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M'),
                        'summary': 'Document processing pending',
                        'follow_up_questions': []
                    }
                    document_list.append(document_data)
            
            # Update project timestamp to track activity
            update_project_timestamp(main_project_id, user)
            
            return Response(document_list, status=status.HTTP_200_OK)
        
        except Exception as e:
            print(f"Error in GetUserDocumentsView: {str(e)}")
            return Response(
                {'error': f'Failed to fetch documents: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class ManageConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def update_conversation(self, request, conversation_id):
        try:
            # Log incoming request data for debugging
            print(f"Incoming update request for conversation {conversation_id}")
            print(f"Request data: {request.data}")

            # Validate input
            new_title = request.data.get('title')
            is_active = request.data.get('is_active', True)

            if not new_title or not new_title.strip():
                return Response(
                    {'error': 'Title cannot be empty'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Find the conversation
            try:
                conversation = ChatHistory.objects.get(
                    conversation_id=conversation_id, 
                    user=request.user
                )
            except ChatHistory.DoesNotExist:
                return Response(
                    {'error': 'Conversation not found'},
                    status=status.HTTP_404_NOT_FOUND
                )

            # Update the title and active status
            conversation.title = new_title.strip()
            conversation.is_active = is_active
            conversation.save()

            # Update project timestamp if main_project_id exists
            if hasattr(conversation, 'main_project_id') and conversation.main_project_id:
                update_project_timestamp(conversation.main_project_id, request.user)

            # Log successful update
            print(f"Conversation {conversation_id} updated successfully")
            print(f"New title: {conversation.title}")

            return Response({
                'message': 'Conversation title updated successfully',
                'conversation_id': str(conversation.conversation_id),
                'new_title': conversation.title,
                'is_active': conversation.is_active
            }, status=status.HTTP_200_OK)

        except Exception as e:
            # Comprehensive error logging
            print(f"Error updating conversation title: {str(e)}")
            logger.error(f"Conversation title update error: {str(e)}")

            return Response(
                {'error': 'Failed to update conversation title', 'details': str(e)},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def put(self, request, conversation_id):
        return self.update_conversation(request, conversation_id)

    def patch(self, request, conversation_id):
        return self.update_conversation(request, conversation_id)

class DocumentProcessingMixin:

    def _parse_json_response(self, response_content, fallback_message="Error processing response"):
        """Helper method to safely parse JSON responses from LLM"""
        try:
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing failed: {str(e)}")
            return {
                "content": fallback_message,
                "error": True,
                "format_type": "error"
            }


    def convert_video_to_audio(self, video_path):
        """
        Convert video files to audio format using MoviePy.
        Returns the path to the created audio file.
        
        Args:
            video_path: Path to the video file
        """
        import tempfile
        from moviepy import VideoFileClip
        import logging
        
        logger = logging.getLogger(__name__)
        
        # Set up the audio directory
        AUDIO_DIR = "audio_files"
        os.makedirs(AUDIO_DIR, exist_ok=True)
        
        try:
            # Generate output path for the audio file
            base_name = os.path.splitext(os.path.basename(video_path))[0]
            output_path = os.path.join(AUDIO_DIR, f"{base_name}.mp3")
            
            logger.info(f"Loading video: {video_path}")
            video = VideoFileClip(video_path)
            
            logger.info("Extracting audio...")
            audio = video.audio
            
            if audio is None:
                logger.warning(f"No audio stream found in video: {video_path}")
                return None
            
            logger.info(f"Converting to MP3: {output_path}")
            audio.write_audiofile(output_path, codec='mp3')
            
            # Close files
            audio.close()
            video.close()
            
            logger.info(f"Video-to-audio conversion complete: {output_path}")
            return output_path
        
        except Exception as e:
            logger.error(f"Error during video-to-audio conversion: {str(e)}")
            return None
    def extract_text_from_video(self, file_path, user=None):
        """
        Extract text from video files by first converting to audio, then transcribing.
        Returns the transcribed text.
        
        Args:
            file_path: Path to the video file
            user: Django user object to get API token
        """
        
        logger = logging.getLogger(__name__)
        print(f"Starting video-to-text extraction for: {file_path}")
        
        try:
            # First convert video to audio
            print("Attempting to convert video to audio...")
            audio_path = self.convert_video_to_audio(file_path)
            
            if not audio_path:
                print(f"Failed to convert video to audio: {file_path}")
                logger.error(f"Failed to convert video to audio: {file_path}")
                return f"Error: Failed to extract audio from video file."
            
            print(f"Successfully converted video to audio: {audio_path}")
            
            # Then extract text from the resulting audio file
            print(f"Starting audio transcription...")
            extracted_text = self.extract_text_from_audio(audio_path, user=user)
            
            print(f"Audio transcription completed, text length: {len(extracted_text) if extracted_text else 0}")
            
            #Clean up the temporary audio file if needed
            #Uncomment if you want to delete the audio file after processing
            if os.path.exists(audio_path):
                os.remove(audio_path)
            
            return extracted_text
        
        except Exception as e:
            print(f"Exception in extract_text_from_video: {str(e)}")
            logger.error(f"Error extracting text from video: {str(e)}")
            return f"Error transcribing video: {str(e)}"

    def process_document_from_text(self, text, filename):
        """
        Process document directly from provided text.
        For use with transcripts that are already extracted.
        """
        
        try:
            # Clean the text
            cleaned_text = self.clean_text(text)
            
            # Create document record for processing
            doc = {
                'text': cleaned_text,
                'name': filename
            }
            
            all_chunks = []
            # Split text into chunks, with smaller size for transcripts
            chunks = self.split_text_into_chunks(doc['text'], chunk_size=800, chunk_overlap=200)
            
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    'text': chunk,
                    'source': doc['name'],
                    'source_file': doc['name'],
                    'chunk_id': i,
                    'is_transcript': True  # Mark this as a transcript
                })
            
            # Get embeddings for all chunks
            text_chunks = [chunk['text'] for chunk in all_chunks]
            print(f"Getting embeddings for {len(text_chunks)} text chunks")
            embeddings = self.get_embeddings(text_chunks)
            
            if not embeddings:
                raise ValueError("Failed to generate embeddings for document chunks")
            
            # Create FAISS index
            index = self.create_faiss_index(embeddings)
            
            # Generate a unique session ID for this document
            session_id = uuid.uuid4().hex
            
            # Save the index and chunks
            index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
            
            print(f"Transcript processed successfully: {len(all_chunks)} chunks created")
            
            return {
                'index_path': index_file,
                'metadata_path': pickle_file,
                'full_text': doc['text']
            }
            
        except Exception as e:
            print(f"Error in process_document_from_text: {str(e)}")
            raise


    def extract_text_from_audio(self, file_path, user=None):
        """
        Extract text from audio files using Google's Gemini API.
        For large files, chunks the audio before transcription.
        Returns the transcribed text.
        
        Args:
            file_path: Path to the audio file
            user: Django user object to get API token
        """

        logger = logging.getLogger(__name__)
        
        # Set up the transcripts directory
        TRANSCRIPT_DIR = "transcripts"
        os.makedirs(TRANSCRIPT_DIR, exist_ok=True)
        
        # Audio chunking settings
        CHUNK_LENGTH_MINUTES = 5
        CHUNK_THRESHOLD_MB = 20
        DURATION_THRESHOLD_MINUTES = 10
        MAX_RETRIES = 3
        
        try:
            # Get the user's Gemini API token
            gemini_api_key = None
            if user:
                try:
                    user_api_tokens = UserAPITokens.objects.get(user=user)
                    gemini_api_key = user_api_tokens.gemini_token
                    
                    # If no token is saved for the user, use a fallback mechanism or raise an error
                    if not gemini_api_key:
                        logger.warning(f"No Gemini API token found for user {user.username}")
                        # Optional: You could implement a fallback or use an environment variable
                        gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                        if not gemini_api_key:
                            raise ValueError("Gemini API token is required for processing audio files")
                        
                except UserAPITokens.DoesNotExist:
                    logger.error(f"No API tokens record found for user {user.username}")
                    raise ValueError("User API tokens not configured")
            else:
                # Fallback to environment variable if no user provided
                gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                if not gemini_api_key:
                    raise ValueError("Gemini API token is required for processing audio files")
            
            # Configure the Gemini API with the new library
            from google import genai
            client = genai.Client(api_key=gemini_api_key)
            model = "gemini-2.0-flash-exp"
            
            # Generate a unique base filename for this transcription
            base_filename = f"audio_transcript_{uuid.uuid4().hex}"
            full_transcript = ""
            temp_files = []
            
            # Check if audio file needs chunking
            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            
            # Load audio file to check duration using MoviePy
            audio_clip = AudioFileClip(file_path)
            duration_minutes = audio_clip.duration / 60  # Convert seconds to minutes
            
            logger.info(f"Audio file: {file_path}, Size: {file_size_mb:.2f}MB, Duration: {duration_minutes:.2f} minutes")
            
            # Determine if chunking is needed
            need_chunking = False
            if file_size_mb > CHUNK_THRESHOLD_MB or duration_minutes > DURATION_THRESHOLD_MINUTES:
                need_chunking = True
                logger.info(f"File exceeds size/duration thresholds, will chunk for processing")
            
            if need_chunking:
                # Calculate parameters for chunking
                chunk_length_seconds = CHUNK_LENGTH_MINUTES * 60
                overlap_seconds = 5  # 5 second overlap
                total_chunks = int(audio_clip.duration / chunk_length_seconds) + (1 if audio_clip.duration % chunk_length_seconds > 0 else 0)
                
                logger.info(f"Splitting audio into {total_chunks} chunks of {CHUNK_LENGTH_MINUTES} minutes each")
                
                for i in range(total_chunks):
                    # Calculate start and end times with overlap
                    start_time = max(0, i * chunk_length_seconds - overlap_seconds) if i > 0 else 0
                    end_time = min(audio_clip.duration, (i + 1) * chunk_length_seconds + overlap_seconds)
                    
                    # Create a subclip
                    chunk = audio_clip.subclipped(start_time, end_time)
                    
                    # Export the chunk to a temporary file
                    chunk_filename = f"{os.path.splitext(file_path)[0]}_chunk_{i}.mp3"
                    chunk.write_audiofile(chunk_filename, codec='mp3')
                    temp_files.append(chunk_filename)
                    chunk.close()  # Close the subclip to free resources
                    
                    # Process the chunk
                    logger.info(f"Processing chunk {i + 1}/{total_chunks}")
                    
                    # Transcribe each chunk with retry logic
                    retry_count = 0
                    success = False
                    
                    while retry_count < MAX_RETRIES and not success:
                        try:
                            # Upload file using new genai library
                            with open(chunk_filename, 'rb') as audio_file:
                                uploaded_file = client.files.upload(file=chunk_filename)
                            
                            # Define the prompt for transcription
                            prompt = """You are a transcription and translation assistant. Your job is to:
                            1. Transcribe spoken content from audio.
                            2. Translate everything fully into English.
                            3. Output only the translated English version.
                            4. If other languages are spoken, translate them fully into English.
                            5. Identify intelligently the Moderator and Respondent(s) in the conversation.
                            6. If a speaker is continously speaking, do not repeat the speaker label.
 
                            **Formatting Instructions:**
                            - Format each line like this:
                            M: [Moderator's translated English speech]
                            R: [Respondent's translated English speech]
                            - Use “M:” for the moderator and “R:” for any respondent. Do not attempt to label or distinguish individual respondents.
                            - Do not include names or speaker numbers.
                            - Do not include the original language (e.g., Hindi).
                            - Do not paraphrase or summarize. Translate every spoken sentence accurately.
                            - Output should be clean, easy-to-read translated English dialogue.
                            - Do not use JSON or structured formats.
 
                            **Example Output:**
                            M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
                            R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
                            R: I live with my parents and younger brother.
 
                            Only output the above format. Nothing else.
                            """
                            
                            # Generate the transcription using new API
                            response = client.models.generate_content(
                                model=model,
                                contents=[prompt, uploaded_file]
                            )
                            chunk_transcription = response.text
                            
                            # Add a separator between chunks
                            if i > 0:
                                full_transcript += "\n\n--- NEXT SEGMENT ---\n\n"
                            
                            full_transcript += chunk_transcription
                            success = True
                            
                        except Exception as e:
                            retry_count += 1
                            logger.error(f"Error in chunk transcription (attempt {retry_count}): {str(e)}")
                            if retry_count >= MAX_RETRIES:
                                full_transcript += f"\n\n[Error transcribing segment {i + 1}: {str(e)}]\n\n"
                            else:
                                time.sleep(2)  # Wait before retry
                
                # Close the main audio clip
                audio_clip.close()
                    
            else:
                # Process the entire file at once if it's small enough
                logger.info("Processing audio file in a single pass")
                audio_clip.close()  # Close the clip since we're not using it for chunking
                
                # Upload file using new genai library
                uploaded_file = client.files.upload(file=file_path)
                
                # Define the prompt for transcription
                prompt = """You are a transcription and translation assistant. Your job is to:
                            1. Transcribe spoken content from audio.
                            2. Translate everything fully into English.
                            3. Output only the translated English version.
                            4. If other languages are spoken, translate them fully into English.
                            5. Identify intelligently the Moderator and Respondent(s) in the conversation.
                            6. If a speaker is continously speaking, do not repeat the speaker label.
 
                            **Formatting Instructions:**
                            - Format each line like this:
                            M: [Moderator's translated English speech]
                            R: [Respondent's translated English speech]
                            - Use “M:” for the moderator and “R:” for any respondent. Do not attempt to label or distinguish individual respondents.
                            - Do not include names or speaker numbers.
                            - Do not include the original language (e.g., Hindi).
                            - Do not paraphrase or summarize. Translate every spoken sentence accurately.
                            - Output should be clean, easy-to-read translated English dialogue.
                            - Do not use JSON or structured formats.
 
                            **Example Output:**
                            M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
                            R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
                            R: I live with my parents and younger brother.
 
                            Only output the above format. Nothing else.
                            """
                
                # Generate the transcription using new API
                response = client.models.generate_content(
                    model=model,
                    contents=[prompt, uploaded_file]
                )
                full_transcript = response.text
            
            # Save the full transcription as a text file
            txt_path = os.path.join(TRANSCRIPT_DIR, f"{base_filename}.txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(full_transcript)
            
            # Clean up temporary chunk files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                        logger.info(f"Removed temporary file: {temp_file}")
                except Exception as e:
                    logger.error(f"Error removing file {temp_file}: {str(e)}")
            
            # Extract the text from the saved text file using the existing method
            extracted_text = self.extract_text_from_txt(txt_path)
            
            return extracted_text
        
        except Exception as e:
            logger.error(f"Error transcribing audio file: {str(e)}")
            return f"Error transcribing audio: {str(e)}"
            
    def extract_text_from_file(self, file_path, user=None):
        """Extract text from different file types."""
        from pathlib import Path
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif ext == '.txt':
            return self.extract_text_from_txt(file_path)
        
        elif ext in ['.mp3', '.wav', '.mpeg']:
            return self.extract_text_from_audio(file_path, user)
        elif ext in ['.mp4', '.avi', '.mov', '.wmv', '.mkv', '.flv', '.webm', '.mts']:
            return self.extract_text_from_video(file_path, user)
        else:
            return ""

    def extract_text_from_pdf(self, file_path):
        import fitz
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    def extract_text_from_docx(self, file_path):
        import docx
        doc = docx.Document(file_path)
        text = [para.text for para in doc.paragraphs]
        return "\n".join(text)

    def extract_text_from_pptx(self, file_path):
        import pptx
        pres = pptx.Presentation(file_path)
        text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_text_from_xlsx(self, file_path):
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append(" | ".join(row_text))
        return "\n".join(text)
    
    def extract_text_from_txt(self, file_path):
        """Extract text from a plain text file."""
        try:
            # Try multiple encodings in case UTF-8 fails
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    return text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try binary mode and decode with errors='replace'
            with open(file_path, 'rb') as file:
                binary_data = file.read()
                return binary_data.decode('utf-8', errors='replace')
                
        except Exception as e:
            print(f"Error extracting text from txt file: {str(e)}")
            return f"Error extracting text: {str(e)}"

    def clean_text(self, text):
        """Clean and normalize extracted text."""
        import re
        import unicodedata
        
        # Remove excessive whitespace while preserving paragraph breaks
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()
        
        # Remove any control characters
        text = ''.join(char for char in text if unicodedata.category(char)[0] != 'C')
        
        return text

    def generate_summary(self, text, file_name):
        """Generates a summary covering the entire text using GPT-4 with formatted HTML-like output."""
        import tiktoken
        import re
        
        if not text.strip():
            return "No extractable text found in the document.", []
        
        # Create the detailed prompt with formatting instructions
        format_prompt = f"""
        Please analyze this document '{file_name}' and provide:
        1. A concise summary of the main points and key findings
        
        Content: {{content}}

        ### Instructions:
        - Use semantic HTML-like tags for structure
        - Provide a clear, organized summary
        - Highlight key insights with <b> tags
        - Use <p> tags for paragraphs
        - Use <ul> and <li> for list-based information

        ### Expected Response Format:
        <b>Summary Overview</b>
        <p>High-level introduction to the document's main theme</p>

        <b>Key Highlights</b>
        <ul>
            <li>First major insight</li>
            <li>Second major insight</li>
            <li>Third major insight</li>
        </ul>

        <b>Detailed Insights</b>
        <p>Expanded explanation of the document's core content and significance</p>
        """
        
        encoding = tiktoken.get_encoding("cl100k_base")
        tokens = encoding.encode(text)
        MAX_TOKENS_ALLOWED = 3500
        
        if len(tokens) > MAX_TOKENS_ALLOWED:
            # Split into chunks that do not exceed the token limit
            chunk_size = MAX_TOKENS_ALLOWED
            chunks = []
            for i in range(0, len(tokens), chunk_size):
                chunk_tokens = tokens[i:i+chunk_size]
                chunk_text = encoding.decode(chunk_tokens)
                chunks.append(chunk_text)
                
            partial_summaries = []
            for idx, chunk in enumerate(chunks):
                try:
                    # Updated system message for JSON response
                    system_message = """You are a document summarization expert.
                    You must respond in JSON format with the following structure:
                    {
                        "summary": "Your formatted summary with HTML tags",
                        "key_points": ["point1", "point2", "point3"],
                        "document_section": "chunk_number"
                    }"""

                    chunk_prompt = format_prompt.replace("{content}", chunk)
                    chunk_prompt += f"""
                    
                    CRITICAL: Your response MUST be a valid JSON object with this structure:
                    {{
                        "summary": "Your summary with HTML formatting",
                        "key_points": ["list", "of", "key", "points"],
                        "document_section": "chunk_{idx}"
                    }}
                    """

                    completion = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": system_message},
                            {"role": "user", "content": chunk_prompt}
                        ],
                        response_format={"type": "json_object"},  # Force JSON response
                        temperature=0.3,
                        max_tokens=500
                    )

                    # Parse JSON response
                    json_response = self._parse_json_response(completion.choices[0].message.content)
                    partial_summary = json_response.get("summary", f"Summary generation error for chunk {idx}.")
                    partial_summaries.append(partial_summary)
                except Exception as e:
                    print(f"Error generating summary for chunk {idx}: {str(e)}")
                    partial_summaries.append(f"Summary generation error for chunk {idx}.")
            
            combined_text = "\n\n".join(partial_summaries)
            try:
                # Updated system message for final summary
                system_message = """You are a document summarization expert.
                You must respond in JSON format with the following structure:
                {
                    "summary": "Your comprehensive summary with HTML formatting",
                    "document_name": "document_filename",
                    "summary_type": "consolidated"
                }"""

                final_prompt = f"""
                Combine these partial summaries into a cohesive, comprehensive summary for document '{file_name}':

                {combined_text}

                ### Instructions:
                - Use semantic HTML-like tags for structure
                - Provide a clear, organized summary
                - Highlight key insights with <b> tags
                - Use <p> tags for paragraphs
                - Use <ul> and <li> for list-based information

                ### Expected Response Format:
                <b>Summary Overview</b>
                <p>High-level introduction to the document's main theme</p>

                <b>Key Highlights</b>
                <ul>
                    <li>First major insight</li>
                    <li>Second major insight</li>
                    <li>Third major insight</li>
                </ul>

                <b>Detailed Insights</b>
                <p>Expanded explanation of the document's core content and significance</p>
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "summary": "Your comprehensive summary with HTML formatting",
                    "document_name": "{file_name}",
                    "summary_type": "consolidated"
                }}
                """
                
                final_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": final_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.3,
                    max_tokens=1000
                )

                # Parse JSON response
                json_response = self._parse_json_response(final_completion.choices[0].message.content)
                summary_text = json_response.get("summary", "")
                
                # Format the response
                summary_text = self.format_summary_response(summary_text)
                
                return summary_text, []
            except Exception as e:
                print(f"Error generating final summary: {str(e)}")
                return self.format_summary_response("\n".join(partial_summaries)), []
        else:
            try:
                # Updated system message for single summary
                system_message = """You are a document summarization expert.
                You must respond in JSON format with the following structure:
                {
                    "summary": "Your formatted summary with HTML tags",
                    "document_name": "document_filename", 
                    "summary_type": "complete"
                }"""

                full_prompt = format_prompt.replace("{content}", text)
                full_prompt += f"""
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "summary": "Your summary with HTML formatting",
                    "document_name": "{file_name}",
                    "summary_type": "complete"
                }}
                """

                completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": full_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.3,
                    max_tokens=1000
                )

                # Parse JSON response
                json_response = self._parse_json_response(completion.choices[0].message.content)
                summary_text = json_response.get("summary", "")
                
                # Format the response
                summary_text = self.format_summary_response(summary_text)
                
                return summary_text, []
            except Exception as e:
                print(f"Error generating summary: {str(e)}")
                return self.format_error_message(file_name), []
        
    def format_summary_response(self, response_text):
        """Ensures proper HTML-like formatting in the summary response."""
        # Wrap paragraphs in <p> tags if not already wrapped
        response_text = re.sub(r'^(?!<\s*(?:p|b|u|ul|li)\b)(.*)', r'<p>\1</p>', response_text, flags=re.MULTILINE)
        
        # Ensure bold tags for section headers
        section_headers = ['Summary Overview', 'Key Highlights', 'Detailed Insights']
        for header in section_headers:
            response_text = response_text.replace(header, f'<b>{header}</b>')
        
        return response_text

    def format_error_message(self, file_name):
        """Returns a formatted error message."""
        return f"""
        <b>Summary Generation Error</b>
        <p>Unable to generate a comprehensive summary for {file_name}</p>
        
        <b>Possible Reasons</b>
        <ul>
            <li>Document may be too complex</li>
            <li>Parsing issues encountered</li>
            <li>Insufficient context extracted</li>
        </ul>
        """
    
from rest_framework.parsers import JSONParser






class ConsolidatedSummaryView(DocumentProcessingMixin, APIView):
    parser_classes = (JSONParser,)



    def load_faiss_index(self, session_id):
    
        try:
            index_file = f"faiss_indexes/{session_id}_index.faiss"
            pickle_file = f"faiss_indexes/{session_id}_chunks.pkl"
            index = faiss.read_index(index_file)
            with open(pickle_file, "rb") as f:
                chunks = pickle.load(f)
            return index, chunks
        except Exception as e:
            print(f"Error loading FAISS index: {str(e)}")
            return None, []

    def _parse_json_response(self, response_content):
        """Parse JSON response with fallback handling"""
        try:
            import json
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {str(e)}")
            # Return a fallback structure
            return {
                "summary": response_content,  # Use raw content as summary
                "error": "JSON parsing failed"
            }

    def format_summary_response(self, response_text):
        """Ensures proper HTML-like formatting in the summary response."""
        # Wrap paragraphs in <p> tags if not already wrapped
        response_text = re.sub(r'^(?!<[p|b|u|ul|li])(.*?)$', r'<p>\1</p>', response_text, flags=re.MULTILINE)
        
        # Ensure bold tags for section headers
        section_headers = ['Summary Overview', 'Key Highlights', 'Detailed Insights']
        for header in section_headers:
            response_text = response_text.replace(header, f'<b>{header}</b>')
        
        return response_text

    def hierarchical_summary(self, chunks, group_size=10):
        """
        Generates a consolidated summary by grouping the chunk texts, summarizing each group,
        and then summarizing the aggregated group summaries.
        Enhanced to track document sources for better consolidated summaries.
        """
        import tiktoken
        encoding = tiktoken.get_encoding("cl100k_base")
            
        # Organize chunks by source document first
        chunks_by_source = {}
        for chunk in chunks:
            source = chunk.get('source', 'Unknown Source')
            if source not in chunks_by_source:
                chunks_by_source[source] = []
            chunks_by_source[source].append(chunk)
            
        # Process each document separately at first
        all_group_summaries = []
        for source, source_chunks in chunks_by_source.items():
            group_summaries = []
            # Group chunks to avoid exceeding token limits
            for i in range(0, len(source_chunks), group_size):
                group_text = "\n\n".join(chunk['text'] for chunk in source_chunks[i:i+group_size])
                    
                # Ensure we don't exceed token limits for the API call
                tokens = encoding.encode(group_text)
                if len(tokens) > 8000:
                    group_text = encoding.decode(tokens[:8000])
                    
                # Generate summary for this group
                try:
                    # Updated system message for JSON response
                    system_message = """You are a document summarization expert.
                    You must respond in JSON format with the following structure:
                    {
                        "summary": "Your summary with HTML formatting",
                        "source_document": "document_name",
                        "group_number": "group_identifier"
                    }"""

                    group_prompt = f"""
                    Summarize the following text from document '{source}'. Cover all key details:

                    {group_text}
                    
                    CRITICAL: Your response MUST be a valid JSON object with this structure:
                    {{
                        "summary": "Your summary covering all key details",
                        "source_document": "{source}",
                        "group_number": "group_{i//group_size}"
                    }}
                    """

                    completion = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": system_message},
                            {"role": "user", "content": group_prompt}
                        ],
                        response_format={"type": "json_object"},  # Force JSON response
                        temperature=0.3,
                        max_tokens=500
                    )

                    # Parse JSON response
                    json_response = self._parse_json_response(completion.choices[0].message.content)
                    group_summary = json_response.get("summary", f"Summary generation error for group {i//group_size}.")
                    group_summaries.append(f"[From {source}] {group_summary}")
                except Exception as e:
                    print(f"Error generating group summary for {source}: {str(e)}")
                    # Add a placeholder if summarization fails
                    group_summaries.append(f"[From {source}] Summary generation error for chunk group {i//group_size}.")
            
            # Get document-level summary first
            document_summary = "\n\n".join(group_summaries)
            try:
                # Updated system message for document summary
                system_message = """You are a document summarization expert.
                You must respond in JSON format with the following structure:
                {
                    "summary": "Your cohesive document summary",
                    "document_name": "source_document",
                    "summary_type": "document_level"
                }"""

                doc_prompt = f"""
                Combine these partial summaries into a cohesive summary for document '{source}':

                {document_summary}
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "summary": "Your cohesive summary for the document",
                    "document_name": "{source}",
                    "summary_type": "document_level"
                }}
                """

                doc_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": doc_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.3,
                    max_tokens=800
                )

                # Parse JSON response
                json_response = self._parse_json_response(doc_completion.choices[0].message.content)
                doc_final_summary = f"<b>Document: {source}</b>\n{json_response.get('summary', document_summary)}"
                all_group_summaries.append(doc_final_summary)
            except Exception as e:
                print(f"Error generating document summary for {source}: {str(e)}")
                all_group_summaries.append(f"<b>Document: {source}</b>\n{document_summary}")
        
        # Combine all document summaries for a final consolidated summary
        aggregated_summary = "\n\n".join(all_group_summaries)
        
        # Generate final consolidated summary
        try:
            # Updated system message for final consolidated summary
            system_message = """You are a document summarization expert.
            You must respond in JSON format with the following structure:
            {
                "summary": "Your comprehensive consolidated summary with HTML formatting",
                "documents_analyzed": ["list", "of", "documents"],
                "summary_type": "consolidated"
            }"""

            final_prompt = f"""
            Combine these document summaries into a cohesive, comprehensive consolidated summary.
            First provide a high-level overview of all documents, then highlight key points across documents,
            and finally discuss any important similarities, differences, or connections between the documents.
            
            {aggregated_summary}
            
            Format your response with HTML-like tags:
            <b>Consolidated Summary</b>
            <p>Overall summary of all documents together</p>
            
            <b>Key Points Across Documents</b>
            <ul>
                <li>First key point across documents</li>
                <li>Second key point across documents</li>
                <li>Third key point across documents</li>
            </ul>
            
            <b>Document Relationships</b>
            <p>Analysis of how the documents relate to each other, including any contradictions, complementary information, or overlapping themes.</p>
            
            CRITICAL: Your response MUST be a valid JSON object with this structure:
            {{
                "summary": "Your consolidated summary with HTML formatting",
                "documents_analyzed": [list of document names],
                "summary_type": "consolidated"
            }}
            """

            final_completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": final_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
                temperature=0.3,
                max_tokens=1200
            )

            # Parse JSON response
            json_response = self._parse_json_response(final_completion.choices[0].message.content)
            final_summary = json_response.get("summary", "")
        except Exception as e:
            print(f"Error generating final consolidated summary: {str(e)}")
            # Fallback to a simpler message with the aggregated summaries
            final_summary = f"""
            <b>Consolidated Summary</b>
            <p>The following summarizes content from multiple documents:</p>
            
            {aggregated_summary}
            """
        
        return final_summary

    def post(self, request):
        document_ids = request.data.get('document_ids', [])
        main_project_id = request.data.get('main_project_id')
        use_cached = request.data.get('use_cached', True)  # Option to use cached summary
        user = request.user

        if not document_ids or len(document_ids) < 2:
            return Response({
                'error': 'Please select at least two documents for a consolidated summary'
            }, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Verify project access
            main_project = Project.objects.get(id=main_project_id, user=user)
            
            # Get all selected documents
            documents = Document.objects.filter(
                id__in=document_ids,
                user=user,
                main_project=main_project
            )
            
            if not documents or documents.count() < 2:
                return Response({
                    'error': 'Could not find requested documents'
                }, status=status.HTTP_404_NOT_FOUND)
            
            # Check if we have a cached consolidated summary in blob storage
            if use_cached:
                # Create a consistent identifier for this combination of documents
                sorted_doc_ids = sorted(document_ids)
                doc_ids_str = "_".join(map(str, sorted_doc_ids))
                
                # Try to find existing blob (you might want to implement a more sophisticated caching mechanism)
                # For now, we'll generate a new summary each time, but save it to blob
                pass
            
            # Collect document chunks for hierarchical summarization
            all_chunks = []
            
            for document in documents:
                try:
                    processed_index = ProcessedIndex.objects.get(document=document)
                    
                    # Load FAISS index chunks for this document
                    _, chunks = self.load_faiss_index(processed_index.metadata.split('/')[-1].split('_')[0])
                    
                    # Add these chunks to our collection
                    all_chunks.extend(chunks)
                    
                except ProcessedIndex.DoesNotExist:
                    # Skip documents that haven't been processed
                    continue
            
            if not all_chunks:
                return Response({
                    'error': 'No processed content found for the selected documents'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Generate hierarchical summary across all documents
            consolidated_summary = self.hierarchical_summary(all_chunks)
            
            # Format the summary with HTML-like tags
            formatted_summary = self.format_summary_response(consolidated_summary)

            # Update project timestamp to track activity
            update_project_timestamp(main_project_id, user)
            
            return Response({
                'consolidated_summary': formatted_summary
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"Error generating consolidated summary: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while generating the consolidated summary'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class DocumentUploadView(DocumentProcessingMixin, APIView):
    parser_classes = (MultiPartParser, FormParser)
   
    def post(self, request):
        files = request.FILES.getlist('files')
        user = request.user
        main_project_id = request.data.get('main_project_id')
        target_user_id = request.data.get('target_user_id')

        # ...existing code...
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({'error': 'Target user not found'}, status=status.HTTP_404_NOT_FOUND)
            skip_permission_checks = True
        else:
            user = request.user
            skip_permission_checks = False

        if not skip_permission_checks:
            try:
                permissions = UserModulePermissions.objects.get(user=user)
                if permissions.disabled_modules.get('document-upload', False):
                    return Response({'error': 'Document uploads are disabled for this user'}, status=status.HTTP_403_FORBIDDEN)
            except UserModulePermissions.DoesNotExist:
                pass

            try:
                upload_permissions = UserUploadPermissions.objects.get(user=user)
                if not upload_permissions.can_upload:
                    return Response({'error': 'Document uploads are disabled for this user'}, status=status.HTTP_403_FORBIDDEN)
            except UserUploadPermissions.DoesNotExist:
                pass
        # ...existing code...

        try:
            main_project = Project.objects.get(id=main_project_id, user=user)
            uploaded_docs = []
            last_processed_doc_id = None
            processed_data = None

            for file in files:
                file_ext = os.path.splitext(file.name)[1].lower()
                is_audio_file = file_ext in ['.mp3', '.wav', '.mpeg', '.m4a', '.aac', '.flac']
                is_video_file = file_ext in ['.mp4', '.avi', '.mov', '.wmv', '.mkv', '.flv', '.webm', '.mts']
                is_media_file = is_audio_file or is_video_file

                # --- Status: Uploading ---
                update_doc_status(user, file.name, "uploading", 10, "Uploading file...")

                existing_doc = Document.objects.filter(
                    user=user,
                    filename=file.name,
                    main_project=main_project
                ).first()

                # --- Status: Uploaded ---
                update_doc_status(user, file.name, "uploaded", 20, "File uploaded, awaiting processing...")

                try:
                    if is_media_file:
                        # --- Status: Processing Media ---
                        update_doc_status(user, file.name, "processing", 30, "Processing media file...")

                        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                            for chunk in file.chunks():
                                tmp_file.write(chunk)
                            file_path = tmp_file.name

                        try:
                            if is_video_file:
                                update_doc_status(user, file.name, "processing", 35, "Extracting audio from video...")
                                extracted_text = self.extract_text_from_video(file_path, user=user)
                            else:
                                update_doc_status(user, file.name, "processing", 35, "Transcribing audio file...")
                                extracted_text = self.extract_text_from_audio(file_path, user=user)

                            if not extracted_text or (isinstance(extracted_text, str) and extracted_text.startswith("Error")):
                                update_doc_status(user, file.name, "error", 100, "Failed to extract text from media file.")
                                return Response({'error': f'Failed to extract text from media file: {file.name}'}, status=status.HTTP_400_BAD_REQUEST)

                            base_name = os.path.splitext(file.name)[0]
                            file_type = "video" if is_video_file else "audio"
                            transcript_filename = f"{base_name}_{file_type}_transcript.txt"

                            update_doc_status(user, file.name, "processing", 45, "Processing transcript...")

                            processed_data = self.process_document_from_text(extracted_text, transcript_filename)

                            transcript_file_path = os.path.join(tempfile.gettempdir(), transcript_filename)
                            with open(transcript_file_path, 'w', encoding='utf-8') as f:
                                f.write(extracted_text)

                            try:
                                with open(transcript_file_path, 'rb') as f:
                                    from django.core.files import File
                                    django_file = File(f, name=transcript_filename)

                                    if existing_doc:
                                        existing_doc.file = django_file
                                        existing_doc.filename = transcript_filename
                                        existing_doc.save()
                                        document = existing_doc
                                    else:
                                        document = Document.objects.create(
                                            user=user,
                                            file=django_file,
                                            filename=transcript_filename,
                                            main_project=main_project
                                        )
                                        file_size = file.size if hasattr(file, 'size') else None
                                        upload_method = 'video_transcript' if is_video_file else 'audio_transcript'
                                        self.create_document_transaction(user, document, main_project, upload_method=upload_method, file_size=file_size)
                            finally:
                                if os.path.exists(transcript_file_path):
                                    os.unlink(transcript_file_path)

                            # --- Status: Indexing ---
                            update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                            ProcessedIndex.objects.create(
                                document=document,
                                faiss_index=processed_data['index_path'],
                                metadata=processed_data['metadata_path'],
                                summary="",
                                markdown_path=processed_data.get('markdown_path', '')
                            )

                            # Only for admin uploads
                            if target_user_id and request.user.username == 'admin':
                                notebook_doc = NotebookDocument.objects.create(
                                    user=user,
                                    file=document.file,
                                    filename=document.filename,
                                    main_project=main_project
                                )
                                NotebookProcessedIndex.objects.create(
                                    document=notebook_doc,
                                    faiss_index=processed_data['index_path'],
                                    metadata=processed_data['metadata_path'],
                                    summary="",
                                    markdown_path=processed_data.get('markdown_path', '')
                                )

                            uploaded_docs.append({
                                'id': document.id,
                                'filename': transcript_filename,
                                'original_media_type': file_type
                            })
                            last_processed_doc_id = document.id

                            # --- Status: Complete ---
                            update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=document.id)

                        finally:
                            if os.path.exists(file_path):
                                os.unlink(file_path)

                    else:
                        # --- Status: Processing Document ---
                        update_doc_status(user, file.name, "processing", 30, "Processing document...")

                        if existing_doc:
                            try:
                                processed_index = ProcessedIndex.objects.get(document=existing_doc)
                                uploaded_docs.append({
                                    'id': existing_doc.id,
                                    'filename': existing_doc.filename
                                })
                                last_processed_doc_id = existing_doc.id
                                update_doc_status(user, file.name, "complete", 100, "Already processed!", doc_id=existing_doc.id)
                            except ProcessedIndex.DoesNotExist:
                                processed_data = self.process_document(file, user)
                                existing_doc.file = file
                                existing_doc.save()

                                # --- Status: Indexing ---
                                update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                                ProcessedIndex.objects.create(
                                    document=existing_doc,
                                    faiss_index=processed_data['index_path'],
                                    metadata=processed_data['metadata_path'],
                                    summary="",
                                    markdown_path=processed_data.get('markdown_path', '')
                                )
                                # Only for admin uploads
                                if target_user_id and request.user.username == 'admin':
                                    notebook_doc = NotebookDocument.objects.create(
                                        user=user,
                                        file=document.file,
                                        filename=document.filename,
                                        main_project=main_project
                                    )
                                    NotebookProcessedIndex.objects.create(
                                        document=notebook_doc,
                                        faiss_index=processed_data['index_path'],
                                        metadata=processed_data['metadata_path'],
                                        summary="",
                                        markdown_path=processed_data.get('markdown_path', '')
                                    )
                                uploaded_docs.append({
                                    'id': existing_doc.id,
                                    'filename': existing_doc.filename
                                })
                                last_processed_doc_id = existing_doc.id

                                # --- Status: Complete ---
                                update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=existing_doc.id)
                        else:
                            processed_data = self.process_document(file, user)

                            document = Document.objects.create(
                                user=user,
                                file=file,
                                filename=file.name,
                                main_project=main_project
                            )
                            file_size = file.size if hasattr(file, 'size') else None
                            upload_method = 'regular'
                            self.create_document_transaction(user, document, main_project, upload_method='regular', file_size=file_size)

                            # --- Status: Indexing ---
                            update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                            ProcessedIndex.objects.create(
                                document=document,
                                faiss_index=processed_data['index_path'],
                                metadata=processed_data['metadata_path'],
                                summary="",
                                markdown_path=processed_data.get('markdown_path', '')
                            )
                            # Only for admin uploads
                            if target_user_id and request.user.username == 'admin':
                                notebook_doc = NotebookDocument.objects.create(
                                    user=user,
                                    file=document.file,
                                    filename=document.filename,
                                    main_project=main_project
                                )
                                NotebookProcessedIndex.objects.create(
                                    document=notebook_doc,
                                    faiss_index=processed_data['index_path'],
                                    metadata=processed_data['metadata_path'],
                                    summary="",
                                    markdown_path=processed_data.get('markdown_path', '')
                                )
                            uploaded_docs.append({
                                'id': document.id,
                                'filename': document.filename
                            })
                            last_processed_doc_id = document.id

                            # --- Status: Complete ---
                            update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=document.id)

                except Exception as e:
                    update_doc_status(user, file.name, "error", 100, f"Error: {str(e)}")
                    print(f"Error processing document: {str(e)}")
                    return Response({
                        'error': str(e),
                        'detail': 'An error occurred while processing the document'
                    }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

            request.session['active_document_id'] = last_processed_doc_id

            if processed_data is None:
                return Response({'error': 'No documents were successfully processed'}, status=status.HTTP_400_BAD_REQUEST)

            update_project_timestamp(main_project_id, user)

            return Response({
                'message': 'Documents uploaded successfully',
                'documents': uploaded_docs,
                'active_document_id': last_processed_doc_id
            }, status=status.HTTP_201_CREATED)

        except Exception as e:
            print(f"Error processing document: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while processing the document'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
 
 
    def create_document_transaction(self, user, document, main_project, upload_method, file_size=None):
        """Create transaction record for document upload"""
        try:
            print("&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&")
            # Create main transaction record
            user_transaction = UserTransaction.objects.create(
                user=user,
                transaction_type=TransactionType.DOCUMENT_UPLOAD,
                document_name=document.filename,
                document_id=document.id,
                main_project=main_project,
                metadata={
                    'upload_timestamp': timezone.now().isoformat(),
                    'upload_method': upload_method,
                    'file_size': file_size or getattr(document.file, 'size', None)
                }
            )
           
            # Create detailed document transaction
            DocumentTransaction.objects.create(
                user_transaction=user_transaction,
                original_filename=document.filename,
                file_size=file_size or getattr(document.file, 'size', None),
                file_type=os.path.splitext(document.filename)[1].lower(),
                upload_method=upload_method
            )
           
            print(f"Transaction recorded for document: {document.filename}")
           
        except Exception as e:
            print(f"Error creating document transaction: {str(e)}")
    # -------------------------------
    # Keep the complexity detection from original code
    # -------------------------------
    def detect_document_complexity(self, file_path):
        """Detect if PDF contains images."""
        import fitz
        try:
            doc = fitz.open(file_path)
            for page in doc:
                images = page.get_images()
                if images:
                    return True
            return False
        except Exception as e:
            print(f"Error detecting images in PDF: {e}")
            return False
    
    # Add LlamaParse integration from Streamlit code
    def process_complex_document_with_llamaparse(self, file_path, file_name, user):
        """
        Process complex document using LlamaParse with user's API token.
        """
        import tempfile
        import uuid
        import numpy as np
        import faiss
        import pickle
        from llama_parse import LlamaParse
       
        try:
            # Get the user's Llama API token
            try:
                user_api_tokens = UserAPITokens.objects.get(user=user)
                llama_api_key = user_api_tokens.llama_token
               
                # If no token is saved for the user, use a fallback mechanism or raise an error
                if not llama_api_key:
                    logger.warning(f"No Llama API token found for user {user.username}")
                    # Optional: You could implement a fallback or raise an error
                    raise ValueError("Llama API token is required for processing complex documents")
                   
            except UserAPITokens.DoesNotExist:
                logger.error(f"No API tokens record found for user {user.username}")
                raise ValueError("User API tokens not configured")
           
            parser = LlamaParse(
                api_key=llama_api_key,  # Use the user's Llama API token
                result_type="markdown",
                verbose=True,
                images=True,
                premium_mode=True
            )
           
            parsed_documents = parser.load_data(file_path)
            full_text = '\n'.join([doc.text for doc in parsed_documents])
           
            # Check if extracted text is empty or only whitespace
            if not full_text or not full_text.strip():
                logger.error(f"No text could be extracted from document: {file_name}")
                raise ValueError(f"Failed to extract any text content from the document '{file_name}'. The document may be corrupted, password-protected, or in an unsupported format.")
           
            # Additional check for very short content (optional - adjust threshold as needed)
            if len(full_text.strip()) < 5:  # Less than 10 characters
                logger.warning(f"Very little text extracted from document: {file_name} (only {len(full_text.strip())} characters)")
                raise ValueError(f"Insufficient text content extracted from document '{file_name}'. Only {len(full_text.strip())} characters were extracted.")
           
            logger.info(f"Successfully extracted {len(full_text)} characters from document: {file_name}")
           
            # Save markdown (adapted from Streamlit implementation)
            base_name = os.path.splitext(file_name)[0]
            safe_name = re.sub(r'[^\w\-_.]', '_', base_name)
            md_path = os.path.join("markdown_files", f"{safe_name}.md")
            os.makedirs("markdown_files", exist_ok=True)
            with open(md_path, "w", encoding='utf-8') as f:
                f.write(full_text)
           
            # Create FAISS index
            text_embedding_dim = 1536  # OpenAI embedding dimension
            faiss_index = faiss.IndexFlatL2(text_embedding_dim)
            metadata_store = []
           
            # Create chunks with overlap for better retrieval (from Streamlit implementation)
            chunked_texts = []
            for doc in parsed_documents:
                # Skip documents with no meaningful content
                if not doc.text or not doc.text.strip():
                    continue
                   
                # Original document as a chunk
                chunked_texts.append(doc.text)
                metadata_store.append({
                    "content": doc.text,
                    "source_file": file_name
                })
               
                # Create additional smaller chunks for better retrieval
                words = doc.text.split()
                chunk_size = 200  # Smaller chunk size as in Streamlit
                stride = 100      # With overlap
               
                if len(words) > chunk_size:
                    for i in range(0, len(words) - chunk_size, stride):
                        chunk = " ".join(words[i:i+chunk_size])
                        if len(chunk.split()) > 50:  # Ensure chunk has substantial content
                            chunked_texts.append(chunk)
                            metadata_store.append({
                                "content": chunk,
                                "source_file": file_name
                            })
           
            # Final check to ensure we have chunks to process
            if not chunked_texts:
                logger.error(f"No valid text chunks created from document: {file_name}")
                raise ValueError(f"Unable to create text chunks from document '{file_name}'. The document content may be insufficient or invalid.")
           
            # Process in batches to avoid API limitations
            batch_size = 50  # Adjust based on API limits
            for i in range(0, len(chunked_texts), batch_size):
                batch = chunked_texts[i:i+batch_size]
                embeddings = self.get_embeddings(batch)
               
                if embeddings:
                    faiss_index.add(np.array(embeddings).astype('float32'))
                else:
                    print(f"Failed to generate embeddings for a batch in {file_name}. Continuing with next batch...")
           
            # Save index and metadata
            session_id = uuid.uuid4().hex
            index_file, pickle_file = self.save_faiss_index(faiss_index, metadata_store, session_id)
           
            return {
                'index_path': index_file,
                'metadata_path': pickle_file,
                'full_text': full_text,
                'markdown_path': md_path  # Include the markdown path
            }
           
        except Exception as e:
            logger.error(f"Error in complex document processing for file '{file_name}': {str(e)}")
            print(f"Error in complex document processing: {str(e)}")
            raise
    def split_text_into_chunks(self, text, chunk_size=1500, chunk_overlap=300):
        """
        Splits text into chunks while attempting to preserve sentence and paragraph structure.
        If the text is very short, a single chunk is returned.
        """
        if len(text) <= chunk_size:
            return [text]
        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = min(start + chunk_size, text_length)
            # Try to find a paragraph break to end the chunk
            if end < text_length:
                break_point = text.rfind("\n\n", start, end)
                if break_point != -1 and break_point > start + chunk_size // 2:
                    end = break_point + 2
                else:
                    # Fallback to sentence break
                    sentence_break = text.rfind(". ", start, end)
                    if sentence_break != -1 and sentence_break > start + chunk_size // 2:
                        end = sentence_break + 2
                    else:
                        # Fallback to space
                        space_break = text.rfind(" ", start, end)
                        if space_break != -1 and space_break > start + chunk_size // 2:
                            end = space_break + 1
            chunks.append(text[start:end])
            start = end - chunk_overlap if end < text_length else text_length
        return chunks
    
    
    # Clean text function (from Streamlit)
    def clean_text(self, text):
        """Clean extracted text by removing extra whitespace and special characters."""
        # Remove extra whitespace
        cleaned = re.sub(r'\s+', ' ', text)
        # Remove special characters that might cause issues
        cleaned = re.sub(r'[^\w\s.,;:!?"\'\-()]', ' ', cleaned)
        return cleaned.strip()
   
    def get_embeddings(self, texts):
        """
        Gets embeddings using OpenAI's text-embedding-3-small model.
        """
        try:
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small"
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            print(f"Error getting embeddings: {str(e)}")
            return []

    # -------------------------------
    # FAISS functions for persistence from Streamlit
    # -------------------------------
    def create_faiss_index(self, embeddings):
        import numpy as np
        import faiss
        
        dimension = len(embeddings[0])
        index = faiss.IndexFlatL2(dimension)
        vectors = np.array(embeddings).astype('float32')
        index.add(vectors)
        return index

    def save_faiss_index(self, index, chunks, session_id):
        import pickle
        import faiss
        
        os.makedirs("faiss_indexes", exist_ok=True)
        index_file = f"faiss_indexes/{session_id}_index.faiss"
        pickle_file = f"faiss_indexes/{session_id}_chunks.pkl"
        faiss.write_index(index, index_file)
        with open(pickle_file, "wb") as f:
            pickle.dump(chunks, f)
        return index_file, pickle_file

    def load_faiss_index(self, session_id):
        import faiss
        import pickle
        
        try:
            index_file = f"faiss_indexes/{session_id}_index.faiss"
            pickle_file = f"faiss_indexes/{session_id}_chunks.pkl"
            index = faiss.read_index(index_file)
            with open(pickle_file, "rb") as f:
                chunks = pickle.load(f)
            return index, chunks
        except Exception as e:
            print(f"Error loading FAISS index: {str(e)}")
            return None, []

    # -------------------------------
    # Process documents function that integrates with Django database
    # -------------------------------
    def process_document(self, file, user):
        """
        Process documents: extraction, chunking, embeddings, FAISS creation
        While maintaining compatibility with the Django database structure.
        Now integrates complex document handling with LlamaParse and user-specific API tokens.
        """
        import tempfile
        import uuid
        import numpy as np
        import faiss
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                for chunk in file.chunks():
                    tmp_file.write(chunk)
                file_path = tmp_file.name


                
            try:

                file_ext = os.path.splitext(file_path)[1].lower()

                if file_ext == '.pptx':
                    return self.process_complex_document_with_llamaparse(file_path, file.name, user)

                if file_ext in ['.jpg', '.jpeg', '.png', '.bmp']:
                    return self.process_complex_document_with_llamaparse(file_path, file.name, user)
                

                # Special handling for audio files
                if file_ext in ['.mp3', '.wav', '.mpeg']:
                    print(f"Audio file detected: {file.name}. Using Gemini for transcription...")
                    extracted_text = self.extract_text_from_audio(file_path, user)
                    
                    if not extracted_text:
                        raise ValueError("No content could be transcribed from the audio file")
                    
                    # Clean the text
                    cleaned_text = self.clean_text(extracted_text)
                    
                    # Create document record for processing
                    doc = {
                        'text': cleaned_text,
                        'name': file.name
                    }
                    
                    # Continue with normal processing flow for the extracted text
                    all_chunks = []
                    # Split text into chunks
                    chunks = self.split_text_into_chunks(doc['text'])
                    for i, chunk in enumerate(chunks):
                        all_chunks.append({
                            'text': chunk,
                            'source': doc['name'],
                            'source_file': doc['name'],
                            'chunk_id': i
                        })
                    
                    # Get embeddings for all chunks
                    text_chunks = [chunk['text'] for chunk in all_chunks]
                    print(f"Getting embeddings for {len(text_chunks)} text chunks")
                    embeddings = self.get_embeddings(text_chunks)
                    
                    if not embeddings:
                        raise ValueError("Failed to generate embeddings for document chunks")
                    
                    # Create FAISS index
                    index = self.create_faiss_index(embeddings)
                    
                    # Generate a unique session ID for this document
                    session_id = uuid.uuid4().hex
                    
                    # Save the index and chunks
                    index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
                    
                    print(f"Audio document processed successfully: {len(all_chunks)} chunks created")
                    
                    return {
                        'index_path': index_file,
                        'metadata_path': pickle_file,
                        'full_text': doc['text']
                    }

                # Check document complexity (keep this from original code)
                is_complex = self.detect_document_complexity(file_path)

                # Process document based on complexity
                if is_complex:
                    # Now using the LlamaParse approach for complex documents
                    print(f"Complex document detected: {file.name}. Using LlamaParse...")
                    return self.process_complex_document_with_llamaparse(file_path, file.name, user)
                else:
                    # Original simple document processing
                    extracted_text = self.extract_text_from_file(file_path, user=user)
                
                if not extracted_text:
                    raise ValueError("No content could be extracted from the document")
                
                # Clean the text
                cleaned_text = self.clean_text(extracted_text)
                
                # Create document record for processing
                doc = {
                    'text': cleaned_text,
                    'name': file.name
                }
                
                all_chunks = []
                # Split text into chunks
                chunks = self.split_text_into_chunks(doc['text'])
                for i, chunk in enumerate(chunks):
                    all_chunks.append({
                        'text': chunk,  # Consistent key name 'text' for all chunks
                        'source': doc['name'],
                        'source_file': doc['name'],  # Add source_file for consistency with LlamaParse
                        'chunk_id': i
                    })
                
                # Get embeddings for all chunks
                text_chunks = [chunk['text'] for chunk in all_chunks]
                print(f"Getting embeddings for {len(text_chunks)} text chunks")
                embeddings = self.get_embeddings(text_chunks)
                
                if not embeddings:
                    raise ValueError("Failed to generate embeddings for document chunks")
                
                # Create FAISS index
                index = self.create_faiss_index(embeddings)
                
                # Generate a unique session ID for this document
                session_id = uuid.uuid4().hex
                
                # Save the index and chunks
                index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
                
                print(f"Document processed successfully: {len(all_chunks)} chunks created")
                
                return {
                    'index_path': index_file,
                    'metadata_path': pickle_file,
                    'full_text': doc['text']
                }
                
            finally:
                # Clean up temporary file
                if os.path.exists(file_path):
                    os.unlink(file_path)
                    
        except Exception as e:
            print(f"Error in process_document: {str(e)}")
            raise


def update_doc_status(user, doc_name, status, progress=0, message="", doc_id=None):
    obj, _ = DocumentProcessingStatus.objects.get_or_create(
        user=user, document_name=doc_name,
        defaults={"status": status, "progress": progress, "message": message, "document_id": doc_id}
    )
    obj.status = status
    obj.progress = progress
    obj.message = message
    if doc_id:
        obj.document_id = doc_id
    obj.save()

class DocumentProcessingStatusView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        statuses = DocumentProcessingStatus.objects.filter(user=request.user)
        return Response([
            {
                "filename": s.document_name,
                "status": s.status,
                "progress": s.progress,
                "message": s.message,
                "document_id": s.document_id,
            }
            for s in statuses.order_by("created_at")
        ])

class GenerateDocumentSummaryView(DocumentProcessingMixin, APIView):
    def post(self, request):
        document_ids = request.data.get('document_ids', [])
        user = request.user
        main_project_id = request.data.get('main_project_id')

        try:
            documents = Document.objects.filter(
                id__in=document_ids,
                user=user,
                main_project=main_project_id
            )

            if not documents:
                return Response({
                    'error': 'No valid documents found'
                }, status=status.HTTP_404_NOT_FOUND)

            all_summaries = []
            for document in documents:
                with open(document.file.path, 'rb') as file:
                    # Using inherited methods from the mixin
                    text = self.extract_text_from_file(file.name)
                    cleaned_text = self.clean_text(text)
                    summary, _ = self.generate_summary(cleaned_text, document.filename)
                    
                    processed_index, created = ProcessedIndex.objects.get_or_create(
                        document=document,
                        defaults={
                            'summary': summary
                        }
                    )
                    if not created:
                        processed_index.summary = summary
                        processed_index.save()

                    all_summaries.append({
                        'document_id': document.id,
                        'filename': document.filename,
                        'summary': summary
                    })
            if main_project_id:
                update_project_timestamp(main_project_id, user)        

            return Response({
                'message': 'Summaries generated successfully',
                'summaries': all_summaries
            }, status=status.HTTP_200_OK)

        except Exception as e:
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)     
class DeleteDocumentView(APIView):
    permission_classes = [IsAuthenticated]
    
    def delete(self, request, document_id):
        try:
            # Find the document and ensure it belongs to the current user
            document = Document.objects.get(
                id=document_id, 
                user=request.user
            )

            # Get main_project_id from the document
            main_project_id = document.main_project_id
            
            # Optional: Delete associated ProcessedIndex
            try:
                processed_index = ProcessedIndex.objects.get(document=document)
                processed_index.delete()
            except ProcessedIndex.DoesNotExist:
                pass
            
            # Delete the document
            document.delete()

            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
            
            return Response(
                {'message': 'Document deleted successfully'}, 
                status=status.HTTP_200_OK
            )
        
        except Document.DoesNotExist:
            return Response(
                {'error': 'Document not found'}, 
                status=status.HTTP_404_NOT_FOUND
            )
        except Exception as e:
            return Response(
                {'error': f'Failed to delete document: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


class GetChatHistoryView(APIView):
    def get(self, request):
        try:
            user = request.user
            main_project_id = request.query_params.get('main_project_id')
            
            print(f"Getting chat history for user {user.username} and project {main_project_id}")
            
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)

            # Filter conversations by both user and main_project_id
            conversations = ChatHistory.objects.filter(
                user=user,
                main_project_id=main_project_id,
                is_active=True
            ).order_by('-created_at')
            
            history = []
            for conversation in conversations:
                messages = conversation.messages.all().order_by('created_at')
                message_list = [
                    {   
                        'id': message.id,
                        'role': message.role,
                        'content': message.content,
                        'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                        'citations': message.citations or []
                    } for message in messages
                ]
                
                history.append({
                    'conversation_id': str(conversation.conversation_id),
                    'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                    'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                    'messages': message_list,
                    'preview': message_list[0]['content'] if message_list else "",
                    'follow_up_questions': conversation.follow_up_questions or [],
                    'selected_documents': [str(doc.id) for doc in conversation.documents.all()]
                })

            return Response(history, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"Error in GetChatHistoryView: {str(e)}")
            return Response(
                {'error': f'Failed to fetch chat history: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        
        
class SetActiveDocumentView(APIView):
    def post(self, request):
        try:
            document_id = request.data.get('document_id')
            
            if not document_id:
                return Response(
                    {'error': 'Document ID is required'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Verify the document exists and belongs to the user
            try:
                document = Document.objects.get(
                    id=document_id, 
                    user=request.user
                )
                
                # Check if document is processed
                ProcessedIndex.objects.get(document=document)
                # Get the main_project_id from the document
                main_project_id = document.main_project_id
                
                # Update project timestamp to track activity
                if main_project_id:
                    update_project_timestamp(main_project_id, request.user)
            except Document.DoesNotExist:
                return Response(
                    {'error': 'Document not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
            except ProcessedIndex.DoesNotExist:
                return Response(
                    {'error': 'Document not processed'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Set the active document in the session
            request.session['active_document_id'] = document_id

            return Response({
                'message': 'Active document set successfully',
                'active_document_id': document_id,
                'filename': document.filename
            }, status=status.HTTP_200_OK)
        
        except Document.DoesNotExist:
            return Response(
                {'error': 'Document not found'},
                status=status.HTTP_404_NOT_FOUND
            )
        
        except Exception as e:
            import traceback
            print(f"Detailed Error: {str(e)}")
            print(traceback.format_exc())  # This will print the full stack trace
            return Response(
                {'error': f'An unexpected error occurred: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
def post_process_response(response_text):
    """
    Post-process the generated response to improve formatting and readability
    but preserve the original meaning and content quality.
    
    Args:
        response_text (str): Raw generated response text
    
    Returns:
        str: Cleaned and formatted response
    """
    try:
        import re
        # Create a copy of the original text before processing
        original_text = response_text
        
        # Only remove explicit section headers - not content
        response_text = re.sub(r'^[IVX]+\.\s*[\w\s]+:', '', response_text, flags=re.MULTILINE)
        response_text = re.sub(r'^\d+\.\s*[\w\s]+:', '', response_text, flags=re.MULTILINE)

        # Convert markdown-style formatting to HTML
        response_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', response_text)
        response_text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', response_text)
        
        # Remove specific section headers but keep their content
        section_headers = [
            'Contextual Insight', 
            'Structured Response', 
            'Analytical Depth', 
            'Interactive Engagement'
        ]
        for header in section_headers:
            response_text = response_text.replace(header + ':', '')
        
        # Clean up extra whitespace while preserving paragraph structure
        response_text = re.sub(r'\n{3,}', '\n\n', response_text)
        
        # If the processed text is significantly shorter or might have lost content,
        # revert to the original
        if len(response_text) < len(original_text) * 0.9:
            logger.warning("Post-processing reduced content significantly, reverting to original")
            return original_text
            
        return response_text.strip()
    
    except Exception as e:
        logger.error(f"Error in post-processing response: {str(e)}", exc_info=True)
        return response_text  # Return original if any error occurs


import uuid
import os
import pickle
import numpy as np
import faiss
import re
from datetime import datetime
import logging
from sklearn.feature_extraction.text import TfidfVectorizer
from agno.agent import Agent
from agno.models.openai import OpenAIChat
from agno.tools.duckduckgo import DuckDuckGoTools

logger = logging.getLogger(__name__)


class ChatView(APIView):
    permission_classes = [IsAuthenticated]
 
    def __init__(self, post_process_func=post_process_response):
        self.post_process_func = post_process_func
        self.agent = Agent(
                model=OpenAIChat(id="gpt-4o"),
                description="""You are a helpful, friendly AI assistant providing informative and thoughtful responses.
                Use semantic HTML tags for structure (<b>, <p>, <ul>, <li>) to enhance readability.
                Maintain a natural, conversational tone.
                Provide comprehensive and clear explanations, breaking down complex topics into easily understandable sections.
                When using web search, prioritize recent, reliable, and relevant information to support your responses.""",
                tools=[DuckDuckGoTools()],
                show_tool_calls=True,
                markdown=True
            )
        
    def normalize_citation_markers(self, text):

        """

        Convert all citation formats to consistent [n] format:

        - [Source: 1, 2] -> [1][2]

        - [Sources: 1, 2] -> [1][2]

        - Source: 1, 2 -> [1][2]

        - (Source 1, 2) -> [1][2]

        """

        if not text:

            return text

        # Replace [Source: n] or [Sources: n] with [n]

        text = re.sub(r'\[(?:Source|Sources?)\s*:?\s*(\d+)\]', r'[\1]', text, flags=re.IGNORECASE)

        # Replace [Source: n, m] or [Sources: n, m] with [n][m]

        def bracketize(match):

            nums = re.findall(r'\d+', match.group(0))

            return ''.join([f'[{n}]' for n in nums])

        text = re.sub(r'\[(?:Source|Sources?)\s*:?\s*[\d,\s]+\]', bracketize, text, flags=re.IGNORECASE)

        # Replace Source: n, m or Sources: n, m (not in brackets) with [n][m]

        text = re.sub(r'(?:Source|Sources?)\s*:?\s*([\d,\s]+)', 

                    lambda m: ''.join([f'[{n}]' for n in re.findall(r'\d+', m.group(1))]), 

                    text, flags=re.IGNORECASE)

        # Replace (Source n, m) with [n][m]

        text = re.sub(r'\((?:Source|Sources?)\s*:?\s*([\d,\s]+)\)', 

                    lambda m: ''.join([f'[{n}]' for n in re.findall(r'\d+', m.group(1))]), 

                    text, flags=re.IGNORECASE)

        # Handle grouped citations [1, 2, 3] -> [1][2][3]

        text = re.sub(r'\[(\d+(?:\s*,\s*\d+)+)\]', 

                    lambda m: ''.join([f'[{n.strip()}]' for n in m.group(1).split(',')]), 

                    text)

        return text
 
    def post(self, request):
        user = request.user
        try:
            # Extract data with more robust handling
            message = request.data.get('message')
            conversation_id = request.data.get('conversation_id')
            main_project_id = request.data.get('main_project_id')
            selected_documents = request.data.get('selected_documents', [])
           
            # Extract web_mode flag from request
            use_web_knowledge = request.data.get('use_web_knowledge', False)
           
            # Extract general_chat_mode flag
            general_chat_mode = request.data.get('general_chat_mode', False)
 
            # Extract response_length preference
            response_length = request.data.get('response_length', 'comprehensive')
           
            # Extract response_format parameter or default to 'natural'
            response_format = request.data.get('response_format', 'natural')
 
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Validate message
            if not message:
                return Response(
                    {'error': 'Message is required'},
                    status=status.HTTP_400_BAD_REQUEST
                )
 
            # Process URLs in the query - new feature
            modified_message, url_context, extracted_urls = self.process_urls_in_query(message)
            has_url_content = bool(url_context)
           
            # Log URL processing results
            if has_url_content:
                print(f"URLs detected in query: {extracted_urls}")
                print(f"Modified query: {modified_message}")
                print(f"URL context size: {len(url_context)} characters")
 
            conversation_context = ""
            if conversation_id:
                try:
                    # Try to get existing conversation
                    conversation = ChatHistory.objects.get(
                        conversation_id=conversation_id,
                        user=user
                    )
                   
                    # Get last 5 conversation messages (max 10 messages = 5 turns)
                    recent_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('-created_at')[:10]
                   
                    # Format for context
                    if recent_messages:
                        context_messages = []
                        for msg in recent_messages:
                            prefix = "User: " if msg.role == 'user' else "Assistant: "
                            context_messages.append(f"{prefix}{msg.content[:400]}...")
                       
                        conversation_context = "Previous conversation:\n" + "\n".join(reversed(context_messages))
                   
                    # Try to get memory buffer
                    try:
                        memory_buffer = ConversationMemoryBuffer.objects.get(conversation=conversation)
                        if memory_buffer.context_summary:
                            conversation_context += f"\n\nConversation Summary: {memory_buffer.context_summary}"
                    except ConversationMemoryBuffer.DoesNotExist:
                        # No memory buffer yet, that's fine
                        pass
                       
                except ChatHistory.DoesNotExist:
                    # No conversation yet, that's fine
                    pass
           
            # Skip document validation if in general chat mode or if URLs are present
            if not general_chat_mode and not has_url_content:
                # First, check for active document in session
                active_document_id = request.session.get('active_document_id')
               
                if not selected_documents:
                    active_document_id = request.session.get('active_document_id')
                    if active_document_id:
                        selected_documents = [active_document_id]
               
                # Validate document selection only when not in general chat mode and no URLs are provided
                if not selected_documents:
                    return Response(
                        {'error': 'Please select at least one document or set an active document'},
                        status=status.HTTP_400_BAD_REQUEST
                    )
 
            # Log the inputs for debugging
            print(f"Chat request received")
            print(f"Message: {message}")
            print(f"General chat mode: {general_chat_mode}")
            print(f"Use web knowledge: {use_web_knowledge}")
            print(f"Selected documents: {selected_documents}")
            print(f"Response format: {response_format}")
            print(f"Response length: {response_length}")
            print(f"Has conversation context: {bool(conversation_context)}")
            print(f"Has URL content: {has_url_content}")
 
            # Process document search early to get context for format detection
            all_chunks = []
            content_sources = []
            similar_contents = []  # Store document content separately
            citation_sources = {}  # Initialize citation sources dict
           
            if not general_chat_mode:
                try:
                    processed_docs = ProcessedIndex.objects.filter(
                        document_id__in=selected_documents,
                        document__user=user
                    )
                   
                    if not processed_docs.exists() and not has_url_content:
                        return Response(
                            {'error': 'No valid documents found'},
                            status=status.HTTP_404_NOT_FOUND
                        )
                       
                    # Log the processed docs for debugging
                    print(f"Found {processed_docs.count()} processed documents")
                    for doc in processed_docs:
                        print(f"Document: {doc.document.filename}")
                        print(f"FAISS index path: {doc.faiss_index}")
                        print(f"Metadata path: {doc.metadata}")
                        print(f"Markdown path: {doc.markdown_path}")
                       
                except Exception as e:
                    print(f"Error fetching documents: {str(e)}")
                    if not has_url_content:
                        return Response(
                            {'error': f'Document retrieval error: {str(e)}'},
                            status=status.HTTP_400_BAD_REQUEST
                        )
 
                # Create a merged metadata store from all selected documents
                all_metadata_store = []
               
                # Load FAISS indices and chunks for all selected documents
                for proc_doc in processed_docs:
                    if not proc_doc.faiss_index or not proc_doc.metadata:
                        print(f"Missing index or metadata for {proc_doc.document.filename}")
                        continue
                   
                    if not os.path.exists(proc_doc.faiss_index) or not os.path.exists(proc_doc.metadata):
                        print(f"Index or metadata file does not exist for {proc_doc.document.filename}")
                        continue
                   
                    # Check if this is a LlamaParse document with direct markdown access
                    if proc_doc.markdown_path and os.path.exists(proc_doc.markdown_path):
                        print(f"Found LlamaParse markdown for {proc_doc.document.filename}")
                        # Add dummy entry that will be handled by search_similar_content
                        all_metadata_store.append({
                            'is_llamaparse': True,
                            'source_file': proc_doc.document.filename,
                            'document_id': proc_doc.document.id
                        })
                   
                    # Load metadata - this is needed for both LlamaParse and simple documents
                    try:
                        with open(proc_doc.metadata, 'rb') as f:
                            chunks = pickle.load(f)
                       
                        # Validate the chunks format
                        if not isinstance(chunks, list):
                            print(f"Warning: chunks is not a list for {proc_doc.document.filename}, type: {type(chunks)}")
                            if chunks:  # If it's a non-empty object, try to adapt it
                                chunks = [chunks]
                            else:
                                chunks = []
                       
                        # Add document source information to each chunk
                        for chunk in chunks:
                            if isinstance(chunk, dict):
                                if 'source_file' not in chunk:
                                    chunk['source_file'] = proc_doc.document.filename
                               
                                # Ensure 'text' key exists
                                if 'text' not in chunk:
                                    for key in ['content', 'document_content', 'chunk_text']:
                                        if key in chunk:
                                            chunk['text'] = chunk[key]
                                            break
                                           
                        # Add to all metadata store
                        print(f"Adding {len(chunks)} chunks from {proc_doc.document.filename}")
                        all_metadata_store.extend(chunks)
                       
                    except Exception as e:
                        print(f"Error loading metadata for {proc_doc.document.filename}: {str(e)}")
                        continue
               
                # Now search using the improved approach
                if all_metadata_store:
                    print(f"Searching in {len(all_metadata_store)} chunks across all documents")
                    # Get relevant content based on query - use modified_message if URLs detected
                    # Use updated search_similar_content that returns citation_sources
                    similar_contents, content_sources, chunk_citation_sources = self.search_similar_content(
                        modified_message if has_url_content else message,
                        processed_docs,
                        all_metadata_store
                    )
                   
                    # Store citation sources
                    citation_sources = chunk_citation_sources
                   
                    if similar_contents:
                        print(f"Found {len(similar_contents)} relevant content chunks")
                        all_chunks = [{'text': content, 'source': source} for content, source in zip(similar_contents, content_sources)]
                    else:
                        print("No relevant content found in documents")
                        all_chunks = []
                else:
                    print("No metadata found for any document")
                    all_chunks = []
           
            # Check if we should auto-detect the response format
            if response_format == 'auto_detect':
                context_snippets = [chunk.get('text', '') for chunk in all_chunks[:5]] if all_chunks else []
                detected_format, secondary_format, confidence = self.detect_question_format(modified_message if has_url_content else message, context_snippets)
               
                # Use the detected format if confidence is reasonable, otherwise default to 'natural'
                if confidence >= 5.0:
                    response_format = detected_format
                    print(f"Auto-detected format: {response_format} (confidence: {confidence})")
                else:
                    response_format = 'natural'
                    print(f"Low confidence in format detection ({confidence}), defaulting to 'natural'")
           
            # Get answer based on mode
            web_knowledge_response = None
            web_sources = []
            doc_citation_sources = {}
           
            # Get answer based on mode
            if general_chat_mode:
                # For general chat mode, consider URL content if available
                if has_url_content:
                    # Create a hybrid prompt that includes the URL content
                    augmented_query = f"""
                    Please answer this question: {modified_message}
                   
                    The question references content from URLs, which is provided here:
                   
                    {url_context}
                   
                    Based on this URL content, answer the user's question.
                    """
                   
                    # Use the existing general chat function but with our augmented query
                    answer = self.get_general_chat_answer(augmented_query, use_web_knowledge, response_length, response_format, user=user)
                    print("Generated response using general chat mode with URL content")
                else:
                    # Standard general chat response without URL content
                    answer = self.get_general_chat_answer(message, use_web_knowledge, response_length, response_format, user=user)
                    print("Generated response using general chat mode")
            else:
                # Document chat mode - now handles URL content too
                document_content_exists = bool(all_chunks)
               
                if document_content_exists or has_url_content:
                    # If we have URL content, add it to the context
                    if has_url_content:
                        # Add URL context to similar_contents and content_sources
                        combined_contents = similar_contents.copy() if similar_contents else []
                        combined_contents.append(url_context)
                       
                        combined_sources = content_sources.copy() if content_sources else []
                        combined_sources.append("URL Content")
                       
                        # Add URL to citation sources
                        if citation_sources:
                            url_citation_key = max(citation_sources.keys()) + 1 if citation_sources else 1
                            citation_sources[url_citation_key] = {
                                'source_file': 'URL Content',
                                'text': url_context[:500] + "..." if len(url_context) > 500 else url_context,
                                'document_id': 'url_content',
                                'score': 0.9,  # High score for URL content
                                'display_num': url_citation_key
                            }
                       
                        # Generate response using combined context with citation sources
                        if response_length == 'short':
                            document_answer, doc_citation_sources = self.generate_short_response(
                                modified_message,
                                combined_contents,
                                combined_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                        else:  # Default to comprehensive
                            document_answer, doc_citation_sources = self.generate_response(
                                modified_message,
                                combined_contents,
                                combined_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                       
                        print(f"Generated document-based answer with URL content using {response_length} response length")
                    else:
                        # Standard document-based response without URL content
                        if response_length == 'short':
                            document_answer, doc_citation_sources = self.generate_short_response(
                                message,
                                similar_contents,
                                content_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                        else:  # Default to comprehensive
                            document_answer, doc_citation_sources = self.generate_response(
                                message,
                                similar_contents,
                                content_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                       
                        print(f"Generated document-based answer using {response_length} response length")
                   
                    # If web knowledge is requested, get web response using document context to enhance the query
                    if use_web_knowledge:
                        web_response_data = self.get_web_knowledge_response(
                            modified_message if has_url_content else message,
                            user=user,
                            document_context=similar_contents  # Pass document context to enhance web search  
                        )
 
                        # Extract the components from the JSON response
                        if isinstance(web_response_data, dict):
                            web_knowledge_response = web_response_data.get("content", "No web response received")
                            web_sources = web_response_data.get("web_sources", [])
                        else:
                            # Fallback for backward compatibility
                            web_knowledge_response = str(web_response_data)
                            web_sources = []
                        print(f"Web knowledge response received with document context, source count: {len(web_sources)}")
                       
                        # Combine document and web responses
                   
                        combined_response_data, combined_citation_sources = self.combine_document_and_web_responses(
                            modified_message if has_url_content else message,
                            document_answer,
                            web_knowledge_response,
                            combined_sources if has_url_content else content_sources,
                            web_sources,
                            response_format,
                            conversation_context,
                            original_doc_context=similar_contents,
                        )
 
                        # Extract the answer from the JSON response
                        if isinstance(combined_response_data, dict):
                            answer = combined_response_data.get("content", "Error combining responses")
                        else:
                            answer = str(combined_response_data)
                        print("Combined document and web responses")
                       
                        # Update citation sources with combined sources
                        doc_citation_sources = combined_citation_sources
                    else:
                        # Just use document answer
                        answer = document_answer
                else:
                    # No document content found
                    if use_web_knowledge:
                        # Only web knowledge available
                        web_response_data = self.get_web_knowledge_response(
                            modified_message if has_url_content else message,
                            user=user
                        )
 
                        if isinstance(web_response_data, dict):
                            answer = web_response_data.get("content", "No web response received")
                            web_sources = web_response_data.get("web_sources", [])
                        else:
                            answer = str(web_response_data)
                            web_sources = []
                        print("Using web knowledge response as document search returned no results")
                    else:
                        print("No document content found and web knowledge not enabled")
                        answer = "I couldn't find any relevant information in the documents."
           
            # Extract main response and citation info (if any)
            if "\n\n*Sources:" in answer:
                parts = answer.split("\n\n*Sources:")
                clean_response = parts[0]
                source_info = parts[1].split('*\n')[0] if len(parts) > 1 else ""
            else:
                clean_response = answer
                source_info = "Web search results" if use_web_knowledge else ""
               
                # Add URL sources if applicable
                if has_url_content:
                    urls_domains = [self.extract_domain(url) for url in extracted_urls]
                    url_sources = ", ".join(urls_domains)
                    if source_info:
                        source_info += f", URLs: {url_sources}"
                    else:
                        source_info = f"URLs: {url_sources}"

            clean_response = self.normalize_citation_markers(clean_response)
           
            # Generate follow-up questions (either from documents or general chat)
            if general_chat_mode:
                follow_up_questions = self.generate_general_follow_up_questions(
                    modified_message if has_url_content else message,
                    clean_response,
                    user=user
                )
            else:
                if all_chunks:
                    context_texts = [chunk.get('text', '') for chunk in all_chunks[:3]]
                    follow_up_questions = self.generate_follow_up_questions(context_texts, user=user)
                else:
                    follow_up_questions = [
                        "What else would you like to know about this content?",
                        "Would you like me to elaborate on any specific point?",
                        "Do you have any other questions I can help with?"
                    ]
           
            # Create citations from chunks and citation sources
            citations = []
            if not general_chat_mode:
                # First add citations from doc_citation_sources (our enhanced citations)
                for citation_num, citation_info in doc_citation_sources.items():
                    # Format the snippet to properly handle tables
                    formatted_snippet = self.format_citation_content(citation_info.get('snippet', ''))
                    citations.append({
                        'source_file': citation_info.get('source_file', 'Unknown'),
                        'page_number': citation_info.get('page_number', 'Unknown'),
                        'section_title': citation_info.get('section_title', 'Unknown'),
                        'snippet': formatted_snippet,
                        'document_id': citation_info.get('document_id', 'Unknown')
                    })
               
                # If no enhanced citations, fall back to chunk-based citations (original method)
                if not citations:
                    for chunk in all_chunks:
                        chunk_text = chunk.get('text', '')
                        formatted_snippet = self.format_citation_content(chunk_text)
                        citations.append({
                            'source_file': chunk.get('source', 'Unknown'),
                            'page_number': chunk.get('chunk_id', 'Unknown'),
                            'section_title': 'Unknown',
                            'snippet': formatted_snippet[:200] + "..." if len(formatted_snippet) > 200 else formatted_snippet,
                            'document_id': next((str(doc.document.id) for doc in processed_docs if doc.document.filename == chunk.get('source')), 'Unknown')
                        })
               
                # Add URL citations if applicable
                if has_url_content:
                    for idx, url in enumerate(extracted_urls):
                        domain = self.extract_domain(url)
                        citations.append({
                            'source_file': domain,
                            'page_number': 'URL',
                            'section_title': 'URL Content',
                            'snippet': f"Content from {domain}",
                            'document_id': f"url_{idx}",
                            'url': url
                        })
               
                # Add web citations if applicable
                if use_web_knowledge and web_sources:
                    for idx, source in enumerate(web_sources):
                        citations.append({
                            'source_file': source.get('title', 'Web Source'),
                            'page_number': 'Web',
                            'section_title': 'Web Search Result',
                            'snippet': source.get('snippet', '')[:200] + "..." if source.get('snippet', '') else "Web content",
                            'document_id': f"web_{idx}",
                            'url': source.get('url', '')
                        })
 
            # Prepare conversation details
            conversation_id = conversation_id or str(uuid.uuid4())
            title = f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"
 
            # Create or retrieve conversation
            if conversation_id:
                conversation, created = ChatHistory.objects.get_or_create(
                    user=user,
                    conversation_id=conversation_id,
                    main_project_id=main_project_id,
                    defaults={
                        'title': f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                        'is_active': True,
                        'follow_up_questions': follow_up_questions,
                    }
                )
 
            if created:
                # Create conversation transaction
                self.create_conversation_transaction(
                    user, conversation, main_project_id,
                    use_web_knowledge, response_format, response_length
                )
            else:
                # Update existing conversation transaction
                self.update_conversation_transaction(conversation, use_web_knowledge)
 
            # Create user message
            user_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='user',
                content=message
            )
 
            # Create AI response message
            ai_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='assistant',
                content=clean_response,
                sources=source_info,
                citations=citations,
                use_web_knowledge=use_web_knowledge,
                response_length=response_length,
                response_format=response_format,
            )
 
            # Update or create memory buffer
            memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                conversation=conversation
            )
           
            # Get all messages for this conversation
            all_messages = ChatMessage.objects.filter(
                chat_history=conversation
            ).order_by('created_at')
           
            # Update memory with all messages
            memory_buffer.update_memory(all_messages)
 
            # Add selected documents to the conversation (skip if general chat mode)
            if selected_documents and not general_chat_mode:
                documents = Document.objects.filter(
                    id__in=selected_documents,
                    user=user
                )
                conversation.documents.set(documents)
 
            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
 
            # Update conversation details
            conversation.title = title
            conversation.follow_up_questions = follow_up_questions
            conversation.save()
 
            # Use citation_sources in the response data if available
            response_data = {
                'response': clean_response,
                'follow_up_questions': follow_up_questions,
                'conversation_id': str(conversation.conversation_id),
                'citations': citations,
                'active_document_id': request.session.get('active_document_id') if not general_chat_mode else None,
                'sources_info': source_info,
                'use_web_knowledge': use_web_knowledge,
                'general_chat_mode': general_chat_mode,
                'response_length': response_length,  # Include in response for tracking
                'response_format': response_format,   # Include detected/used format in response
                'url_content_used': has_url_content,  # Flag if URL content was used
                'extracted_urls': extracted_urls if has_url_content else [],  # List of extracted URLs
                'storage_types_used': self.get_storage_types_used(processed_docs) if not general_chat_mode else []  # Debug info
            }
 
            # Print detailed chat response information
            print("\n--- Chat Interaction Logged ---")
            print(f"User Question: {message}")
            print(f"Mode: {'Web Knowledge' if use_web_knowledge else 'General Chat' if general_chat_mode else 'Document Chat'}")
            print(f"Format: {response_format}")
            print(f"Length: {response_length}")
            print(f"URL content used: {has_url_content}")
            if has_url_content:
                print(f"Extracted URLs: {extracted_urls}")
            if not general_chat_mode:
                    storage_types = self.get_storage_types_used(processed_docs)
            print(f"Storage types used: {storage_types}")
            print(f"Assistant Response: {clean_response[:500]}...")  # First 500 chars
            print(f"Citation count: {len(citations)}")
            print(f"Follow-up Questions: {len(follow_up_questions)}")
            print("-----------------------------\n")

            # Get all messages for this conversation (with IDs)
            all_messages = ChatMessage.objects.filter(
                chat_history=conversation
            ).order_by('created_at')

            message_list = []
            for message in all_messages:
                message_data = {
                    'id': message.id,
                    'role': message.role,
                    'content': message.content,
                    'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                    'citations': message.citations or [],
                    'sources_info': getattr(message, 'sources', ''),  # Add this line
                    'extracted_urls': getattr(message, 'extracted_urls', []),  # Add this line
                   
                

                }
                # Add badge properties for assistant messages
                if message.role == 'assistant':
                    message_data['use_web_knowledge'] = getattr(message, 'use_web_knowledge', False)
                    message_data['response_length'] = getattr(message, 'response_length', 'comprehensive')
                    message_data['response_format'] = getattr(message, 'response_format', 'natural')
                message_list.append(message_data)

            response_data['messages'] = message_list
 
            return Response(response_data, status=status.HTTP_200_OK)
 
        except Exception as e:
            logger.error(f"Unexpected error in ChatView: {str(e)}", exc_info=True)
            return Response(
                {'error': f'An unexpected error occurred: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
 
    # ===================================================================
    # Helper methods for pgvector integration
    # ===================================================================
   
    def get_storage_types_used(self, processed_docs):
        """Get list of storage types used in the current query for debugging"""
        storage_types = []
        for doc in processed_docs:
            if doc.storage_type not in storage_types:
                storage_types.append(doc.storage_type)
        return storage_types
   
    def fallback_text_search_pgvector(self, query, pgvector_docs, limit=10):
        """
        Fallback text search when pgvector similarity search fails
        """
        try:
            from django.db.models import Q
           
            query_lower = query.lower()
            all_results = []
            all_sources = []
            citation_sources = {}
           
            for proc_doc in pgvector_docs:
                # Search for text matches in document embeddings
                text_matches = DocumentEmbedding.objects.filter(
                    document=proc_doc.document,
                    content__icontains=query_lower
                )[:limit]
               
                for idx, embedding_obj in enumerate(text_matches):
                    citation_key = len(all_results) + 1
                   
                    citation_sources[citation_key] = {
                        'source': proc_doc.document.filename,
                        'text': embedding_obj.content,
                        'relevance_score': 0.5,  # Default score for text match
                        'document_id': proc_doc.document.id,
                        'chunk_idx': embedding_obj.chunk_id,
                        'display_num': citation_key
                    }
                   
                    all_results.append(embedding_obj.content)
                    all_sources.append(proc_doc.document.filename)
           
            if all_results:
                print(f"Fallback text search found {len(all_results)} results")
                return all_results, all_sources, citation_sources
            else:
                return None
               
        except Exception as e:
            print(f"Error in fallback text search: {str(e)}")
            return None
    def create_conversation_transaction(self, user, conversation, main_project_id, use_web_knowledge, response_format, response_length):
        """Create transaction record for conversation"""
        try:
            # Get the project
            main_project = Project.objects.get(id=main_project_id, user=user)
           
            # Create main transaction record
            user_transaction = UserTransaction.objects.create(
                user=user,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                conversation_title=conversation.title,
                conversation_id=conversation.conversation_id,
                main_project=main_project,
                metadata={
                    'creation_timestamp': timezone.now().isoformat(),
                    'web_knowledge_used': use_web_knowledge,
                    'response_format': response_format,
                    'response_length': response_length
                }
            )
           
            # Create detailed conversation transaction
            ConversationTransaction.objects.create(
                user_transaction=user_transaction,
                message_count=1,  # Initial message
                web_knowledge_used=use_web_knowledge,
                response_format=response_format,
                response_length=response_length
            )
           
            print(f"Transaction recorded for conversation: {conversation.title}")
           
        except Exception as e:
            print(f"Error creating conversation transaction: {str(e)}")
 
    def update_conversation_transaction(self, conversation, use_web_knowledge):
        """Update existing conversation transaction when new messages are added"""
        try:
            # Find existing transaction for this conversation
            user_transaction = UserTransaction.objects.filter(
                conversation_id=conversation.conversation_id,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                is_active=True
            ).first()
           
            if user_transaction and hasattr(user_transaction, 'conversation_details'):
                # Update message count
                user_transaction.conversation_details.message_count += 1
                if use_web_knowledge:
                    user_transaction.conversation_details.web_knowledge_used = True
                user_transaction.conversation_details.save()
               
                # Update metadata
                user_transaction.metadata['last_message_timestamp'] = timezone.now().isoformat()
                user_transaction.save()
               
        except Exception as e:
            print(f"Error updating conversation transaction: {str(e)}")
 

    def _parse_json_response(self, response_content, fallback_message="Error processing response"):
        """Helper method to safely parse JSON responses from LLM"""
        try:
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing failed: {str(e)}")
            return {
                "content": fallback_message,
                "error": True,
                "format_type": "error"
            }
   

    def format_citation_content(self, snippet):
        """
        Format citation snippets to optimize display of complex tables with long content
        Designed to work with any table structure from any file type
        """
        if not snippet:
            return snippet
            
        # Generic table pattern detection - using structural indicators, not content-specific terms
        has_table_pattern = (
            # Standard markdown table patterns
            ('|' in snippet and any(separator in snippet for separator in ['|---', '| ---', '|----', '------'])) or
            # Table rows with pipes at beginning and end
            ('|' in snippet and any(row.strip().startswith('|') and row.strip().endswith('|') for row in snippet.split('\n'))) or
            # Any structured data that resembles a table (multiple lines with consistent pipe counts)
            self._detect_tabular_structure(snippet)
        )
        
        if has_table_pattern:
            try:
                # Start with basic cleanup
                content = snippet
                
                # Replace emoji arrows with text versions for better compatibility
                content = content.replace('🔻', '↓').replace('🔼', '↑')
                
                # Step 1: Process section headers to ensure they're properly formatted
                # Match any line starting with # and make sure there's a newline after it
                content = re.sub(r'(#+\s*[^\n]+)(?![\n])', r'\1\n', content)
                
                # Step 2: Make sure all headers have proper spacing
                content = re.sub(r'(#+)(\S)', r'\1 \2', content)
                
                # Step 3: Detect and format tabular structures
                lines = content.split('\n')
                formatted_lines = []
                
                in_table = False
                table_lines = []
                table_header_detected = False
                pending_table_header = None
                
                for i, line in enumerate(lines):
                    stripped_line = line.strip()
                    
                    # Check if this is a header followed by a potential table
                    is_header = stripped_line.startswith('#')
                    next_line_is_table = (i < len(lines) - 1 and 
                                        '|' in lines[i+1] and 
                                        lines[i+1].strip().startswith('|') and 
                                        lines[i+1].strip().endswith('|'))
                    
                    # Handle headers that might be associated with tables
                    if is_header and next_line_is_table:
                        formatted_lines.append(line)
                        continue
                    
                    # Check if this is a table line using structural patterns
                    is_table_line = '|' in stripped_line and stripped_line.startswith('|') and stripped_line.endswith('|')
                    
                    # Check if this is a separator line
                    is_separator = is_table_line and any(c in stripped_line for c in ['-', '='])
                    
                    # Check for potential table headers based on structure (first row of pipe-delimited content)
                    is_potential_header = (is_table_line and 
                                        not is_separator and 
                                        not in_table and
                                        i < len(lines) - 1 and
                                        '|' in lines[i+1])
                    
                    # Fix incomplete tables missing separators
                    if is_potential_header and i < len(lines) - 1:
                        next_line = lines[i+1].strip()
                        is_next_separator = '|' in next_line and any(c in next_line for c in ['-', '='])
                        
                        if '|' in next_line and not is_next_separator:
                            # This looks like a header row without a separator, let's add one
                            pending_table_header = line
                            table_header_detected = True
                            continue
                    
                    if table_header_detected and is_table_line and not is_separator:
                        # We found a data row after a header without separator, add the header and create a separator
                        if pending_table_header:
                            formatted_lines.append(pending_table_header)
                            # Create a separator row based on the header
                            parts = [p.strip() for p in pending_table_header.split('|')]
                            separator = '|'
                            for part in parts[1:-1]:  # Skip first and last empty parts
                                separator += ' ' + '-' * len(part) + ' |'
                            formatted_lines.append(separator)
                            pending_table_header = None
                            table_header_detected = False
                            in_table = True
                        formatted_lines.append(line)
                        continue
                    
                    if is_table_line:
                        if not in_table:
                            # Start of a new table
                            in_table = True
                            # If we have pending lines, add them
                            if formatted_lines and not formatted_lines[-1].strip() == '':
                                formatted_lines.append('')  # Empty line before table
                        
                        # Add to current table
                        table_lines.append(line)
                    else:
                        if in_table:
                            # End of table, process it
                            in_table = False
                            
                            # Process and optimize the table
                            optimized_table = self.optimize_table_structure(table_lines)
                            
                            # Add the optimized table
                            formatted_lines.extend(optimized_table)
                            
                            if not stripped_line:  # If this is an empty line
                                formatted_lines.append('')  # Empty line after table
                            else:
                                # If next line has content, add a separator
                                formatted_lines.append('')  # Empty line after table
                                formatted_lines.append(line)  # Add the current non-table line
                            
                            table_lines = []
                            continue
                        
                        # Add non-table line
                        formatted_lines.append(line)
                
                # Don't forget to process any table that's at the end
                if in_table and table_lines:
                    optimized_table = self.optimize_table_structure(table_lines)
                    formatted_lines.extend(optimized_table)
                
                # Join everything back together
                content = '\n'.join(formatted_lines)
                
                # Final cleanup, remove excessive newlines
                content = re.sub(r'\n{3,}', '\n\n', content)
                
                return content
                
            except Exception as e:
                print(f"Error formatting complex table content: {str(e)}")
                # Return original content on error
                return snippet
        
        # If not a table or processing failed, return the original content
        return snippet

    def _detect_tabular_structure(self, text):
        """
        Generic detection of tabular structures in text.
        Looks for consistent patterns of delimiter characters that indicate a table.
        """
        if not text or '|' not in text:
            return False
            
        lines = text.split('\n')
        pipe_counts = []
        
        # Count pipes in consecutive non-empty lines
        consecutive_lines = 0
        for line in lines:
            if not line.strip():
                consecutive_lines = 0
                continue
                
            if '|' in line:
                pipe_count = line.count('|')
                pipe_counts.append(pipe_count)
                consecutive_lines += 1
                
                # If we have at least 3 consecutive lines with pipes, check pattern
                if consecutive_lines >= 3:
                    # Check if pipe counts are consistent (allowing small variations)
                    unique_counts = set(pipe_counts[-3:])
                    if len(unique_counts) <= 2 and max(unique_counts) - min(unique_counts) <= 2:
                        return True
        
        # Alternative detection: Check if we have a header-like row followed by a separator-like row
        for i in range(len(lines) - 1):
            current_line = lines[i].strip()
            next_line = lines[i+1].strip()
            
            if (current_line.startswith('|') and current_line.endswith('|') and
                next_line.startswith('|') and next_line.endswith('|') and
                '-' in next_line):
                return True
        
        return False

    def optimize_table_structure(self, table_lines):
        """
        Process and optimize a table's structure for better display
        Works with any table format from any file type
        """
        if not table_lines:
            return []
        
        # Find header and separator rows
        header_idx = -1
        separator_idx = -1
        
        for i, line in enumerate(table_lines):
            if '-' in line or '=' in line:
                separator_idx = i
                header_idx = max(0, i - 1)
                break
        
        # If no separator found, try to identify the header row and create a separator
        if separator_idx == -1 and len(table_lines) >= 2:
            # Assume first row is header if we have at least two rows with same number of columns
            header_row = table_lines[0]
            header_parts = header_row.split('|')
            
            first_data_row = table_lines[1]
            data_parts = first_data_row.split('|')
            
            # If first and second row have similar structure, assume first is header
            if abs(len(header_parts) - len(data_parts)) <= 1:
                header_idx = 0
                # Create a separator row
                parts = header_row.split('|')
                separator = '|'
                for part in parts[1:-1]:  # Skip first and last empty parts
                    separator += ' ' + '-' * len(part.strip()) + ' |'
                
                # Insert the separator row
                table_lines.insert(1, separator)
                separator_idx = 1
        
        # If we still don't have a separator, check if we have consistent columns
        if separator_idx == -1:
            # Try to convert any table-like structure to proper format
            has_consistent_columns = True
            if len(table_lines) >= 2:
                col_count = len(table_lines[0].split('|'))
                for line in table_lines[1:]:
                    if len(line.split('|')) != col_count:
                        has_consistent_columns = False
                        break
                
                if has_consistent_columns:
                    # Create a proper table with first row as header
                    header_idx = 0
                    separator = '|'
                    header_parts = table_lines[0].split('|')
                    for part in header_parts[1:-1]:
                        separator += ' ' + '-' * len(part.strip()) + ' |'
                    
                    table_lines.insert(1, separator)
                    separator_idx = 1
        
        # If still no structure found, create a default structure
        if header_idx == -1 or separator_idx == -1:
            # The simplest approach: consider the first row as header
            if len(table_lines) >= 1:
                header_idx = 0
                
                # Create and insert a separator row
                header_parts = table_lines[0].split('|')
                separator = '|'
                for part in header_parts[1:-1]:  # Skip first and last empty parts
                    separator += ' ' + '-' * len(part.strip()) + ' |'
                
                if len(table_lines) > 1:
                    table_lines.insert(1, separator)
                    separator_idx = 1
                else:
                    table_lines.append(separator)
                    separator_idx = 1
        
        # Enhancement: Format cells with indicators, preserving the original indicator characters
        for i in range(len(table_lines)):
            # Skip header and separator
            if i != header_idx and i != separator_idx:
                line = table_lines[i]
                parts = line.split('|')
                formatted_parts = []
                
                for part in parts:
                    trimmed = part.strip()
                    
                    # Detect any trend indicators without hardcoding specific symbols
                    # This looks for common directional symbols in a generic way
                    has_up_indicator = any(up_symbol in trimmed for up_symbol in ['↑', '▲', '🔼', '⬆'])
                    has_down_indicator = any(down_symbol in trimmed for down_symbol in ['↓', '▼', '🔻', '⬇'])
                    
                    if has_up_indicator or has_down_indicator:
                        # Keep original format but ensure consistent spacing
                        formatted_parts.append(' ' + trimmed + ' ')
                    else:
                        formatted_parts.append(part)
                
                table_lines[i] = '|'.join(formatted_parts)
        
        # Analyze table for column consistency
        column_counts = [len(line.split('|')) - 1 for line in table_lines]  # -1 because n+1 pipes make n columns
        
        if len(set(column_counts)) > 1:
            # Inconsistent column counts - normalize
            max_columns = max(column_counts)
            
            for i in range(len(table_lines)):
                parts = table_lines[i].split('|')
                current_columns = len(parts) - 1
                
                if current_columns < max_columns:
                    # Add missing columns
                    if table_lines[i].endswith('|'):
                        # Add columns before the trailing |
                        table_lines[i] = table_lines[i][:-1] + ('| ' * (max_columns - current_columns)) + '|'
                    else:
                        # Add columns including trailing |
                        table_lines[i] = table_lines[i] + ('| ' * (max_columns - current_columns)) + '|'
        
        # Check for long content cells and analyze column widths
        has_long_cells = False
        max_col_lengths = []
        
        # Initialize with header column lengths if header exists
        if header_idx >= 0 and header_idx < len(table_lines):
            header_parts = table_lines[header_idx].split('|')[1:-1]  # Skip first and last
            max_col_lengths = [len(part.strip()) for part in header_parts]
        else:
            # Initialize with zeros if no header
            max_columns = max(len(line.split('|')) - 1 for line in table_lines)
            max_col_lengths = [0] * max_columns
        
        # Update with actual content lengths
        for line in table_lines:
            if separator_idx >= 0 and line == table_lines[separator_idx]:
                continue
                
            parts = line.split('|')[1:-1]  # Skip first and last
            
            for i, part in enumerate(parts):
                if i < len(max_col_lengths):
                    max_col_lengths[i] = max(max_col_lengths[i], len(part.strip()))
                    if len(part.strip()) > 20:
                        has_long_cells = True
                elif i >= len(max_col_lengths) and len(parts) > len(max_col_lengths):
                    # Extend max_col_lengths if needed
                    max_col_lengths.extend([0] * (len(parts) - len(max_col_lengths)))
                    max_col_lengths[i] = len(part.strip())
                    if len(part.strip()) > 20:
                        has_long_cells = True
        
        # Optimize the formatting if we have long cells or inconsistent structure
        if has_long_cells or len(set(column_counts)) > 1:
            # Format the table with consistent column widths
            formatted_table = []
            
            for i, line in enumerate(table_lines):
                if i == separator_idx:
                    # Special handling for separator
                    parts = ['']  # Start with empty part for leading pipe
                    for j, max_length in enumerate(max_col_lengths):
                        parts.append(' ' + '-' * max_length + ' ')
                    parts.append('')  # End with empty part for trailing pipe
                    formatted_table.append('|'.join(parts))
                else:
                    parts = line.split('|')
                    formatted_parts = ['']  # Start with empty part for leading pipe
                    
                    for j in range(len(max_col_lengths)):
                        if j + 1 < len(parts):
                            # Part exists in current row
                            part = parts[j+1]
                            # Clean up whitespace but maintain padding
                            cleaned = part.strip()
                            # Ensure consistent width formatting
                            formatted_parts.append(' ' + cleaned.ljust(max_col_lengths[j]) + ' ')
                        else:
                            # Missing part in current row
                            formatted_parts.append(' ' + ' ' * max_col_lengths[j] + ' ')
                    
                    formatted_parts.append('')  # End with empty part for trailing pipe
                    formatted_table.append('|'.join(formatted_parts))
            
            return formatted_table
        
        # If no special formatting needed, just return the original table lines
        return table_lines


    
    def scrape_webpage(self, url):
        """Scrape content from a webpage"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36 Edg/91.0.864.59'
            }

            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Remove script and style elements
            for script_or_style in soup(['script', 'style', 'header', 'footer', 'nav']):
                script_or_style.decompose()
                
            # Extract title
            title = soup.title.string if soup.title else "No title found"
            
            # Extract main content - focus on article, main, or div with content
            content_elements = soup.select('article, main, .content, #content, .post, .entry')
            
            if not content_elements:
                # If no specific content elements, get all paragraphs
                paragraphs = soup.find_all('p')
                content = ' '.join([p.get_text() for p in paragraphs])
            else:
                # Use the first content element found
                content = content_elements[0].get_text()
            
            # Clean the content
            content = re.sub(r'\s+', ' ', content).strip()
            
            return {
                'title': title,
                'content': content[:15000],  # Limit content length
                'url': url
            }
        except Exception as e:
            logger.error(f"Error scraping webpage {url}: {str(e)}")
            return {
                'title': 'Error scraping page',
                'content': f"Could not retrieve content: {str(e)}",
                'url': url
            }

    def extract_urls_from_text(self, text):
        """
        Extract URLs from text using regex pattern matching.
        
        Args:
            text (str): The text to extract URLs from
            
        Returns:
            list: List of extracted URLs
        """
        import re
        
        # Pattern to match URLs - handles http, https, www formats
        url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!./?=&#+]*)*'
        www_pattern = r'www\.(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!./?=&#+]*)*'
        
        # Find all matches
        http_urls = re.findall(url_pattern, text)
        www_urls = re.findall(www_pattern, text)
        
        # Add http:// prefix to www URLs if they aren't already prefixed
        processed_www_urls = []
        for url in www_urls:
            if not any(url in http_url for http_url in http_urls):
                processed_www_urls.append(f"http://{url}")
        
        # Combine and return all URLs
        return http_urls + processed_www_urls

    def process_urls_in_query(self, query):
        """
        Process query to extract URLs, scrape their content, and return enhanced context.
        
        Args:
            query (str): The user's query that may contain URLs
            
        Returns:
            tuple: (modified_query, url_context, extracted_urls)
        """
        # Extract URLs from the query
        extracted_urls = self.extract_urls_from_text(query)
        
        if not extracted_urls:
            return query, "", []
        
        # Log found URLs
        logger.info(f"Found {len(extracted_urls)} URLs in query: {extracted_urls}")
        
        # Scrape content from each URL
        url_contents = []
        successful_urls = []
        
        for url in extracted_urls:
            try:
                scraped_content = self.scrape_webpage(url)
                if scraped_content and scraped_content.get('content'):
                    url_contents.append({
                        'url': url,
                        'title': scraped_content.get('title', 'No title'),
                        'content': scraped_content.get('content')
                    })
                    successful_urls.append(url)
                    logger.info(f"Successfully scraped content from {url}")
                else:
                    logger.warning(f"No content retrieved from {url}")
            except Exception as e:
                logger.error(f"Error scraping URL {url}: {str(e)}")
        
        if not url_contents:
            return query, "", extracted_urls
        
        # Build context from scraped content
        url_context = "CONTENT FROM URLS IN QUERY:\n\n"
        
        for content in url_contents:
            domain = self.extract_domain(content['url'])
            url_context += f"Source: {domain}\n"
            url_context += f"Title: {content['title']}\n"
            # Limit content length to avoid exceeding token limits
            url_context += f"Content: {content['content'][:5000]}...\n\n"
        
        # Create a modified query that references the URLs without including them
        modified_query = query
        for url in extracted_urls:
            # Replace URLs with a more readable reference
            domain = self.extract_domain(url)
            modified_query = modified_query.replace(url, f"[content from {domain}]")
        
        return modified_query, url_context, extracted_urls
    
    def extract_domain(self, url):
        """Extract domain from URL safely"""
        try:
            parsed = urlparse(url)
            domain = parsed.netloc
            if domain.startswith("www."):
                domain = domain[4:]
            return domain if domain else "unknown.com"
        except Exception:
            return "unknown.com"
 
    def get_web_knowledge_response(self, query, user=None, document_context=None):
        """
        Search the web for information and return direct search results without GPT reformatting.
        PRIMARY METHOD: Uses search_web function to get structured search results and returns them directly.
       
        Args:
            query (str): The user's query
            document_context (list, optional): List of context chunks from document search
           
        Returns:
            tuple: (response, sources)
        """
        try:
            # Step 1: Extract keywords and context from document context if available
            enhanced_query = query
            doc_context_text = ""
           
            if document_context and len(document_context) > 0:
                # Combine all document context into a single text for analysis
                combined_context = " ".join(document_context[:3])  # Use top 3 context chunks
               
                # Extract keywords using a smarter prompt that identifies entity types
                keyword_prompt = f"""
                Analyze the following document context and the user query to extract key information for a web search.
 
                USER QUERY: "{query}"
               
                DOCUMENT CONTEXT:
                {combined_context[:2500]}
               
                Follow these steps:
                1. Identify the main entities (brands, products, organizations, etc.) in the document
                2. For each entity, determine its type (e.g., "cosmetic brand", "financial service", "software company")
                3. Extract distinctive attributes or features of these entities
                4. Identify industry-specific terminology that might help with search precision
                5. Include geographical context if relevant
               
                Return ONLY a comma-separated list of 5-8 search terms that would help find the most relevant information online.
               
                Note: Focus on terms that will disambiguate search results and prevent confusion with similarly-named but unrelated entities.
                """
               
                try:
                    # Extract context-aware keywords using OpenAI
                    keyword_response = client.chat.completions.create(
                        model="gpt-4o",  # Using a smaller model for efficiency
                        messages=[
                            {"role": "system", "content": "You are a search optimization specialist who identifies the most effective search terms based on document context."},
                            {"role": "user", "content": keyword_prompt}
                        ],
                        temperature=0.3,
                        max_tokens=150
                    )
                   
                    search_terms = keyword_response.choices[0].message.content.strip()
                    logger.info(f"Extracted search terms: {search_terms}")
                   
                    # Create a more precise enhanced query that includes both the original query and the extracted search terms
                    enhanced_query = f"{query} {search_terms}"
                   
                    # Also extract key passages for the context section of our prompt
                    context_prompt = f"""
                    Analyze the following document context and select 2-3 short passages (2-3 sentences each) that contain the most important information related to this query: "{query}"
                   
                    DOCUMENT CONTEXT:
                    {combined_context[:3000]}
                   
                    Return ONLY the extracted passages with no additional commentary.
                    """
                   
                    context_response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You extract the most relevant passages from documents."},
                            {"role": "user", "content": context_prompt}
                        ],
                        temperature=0.3,
                        max_tokens=300
                    )
                   
                    # Get key passages for later use in the main prompt
                    doc_context_text = context_response.choices[0].message.content.strip()
                    logger.info(f"Extracted key passages for context")
                   
                except Exception as ke:
                    logger.warning(f"Error extracting keywords: {str(ke)}, proceeding with original query")
                    enhanced_query = query
           
            # Log the queries for comparison
            logger.info(f"Original query: {query}")
            logger.info(f"Enhanced query: {enhanced_query}")
           
            # Step 2: PRIMARY METHOD - Use search_web function to get structured search results
            logger.info("PRIMARY: Using search_web function to get search results")
            web_results = self.search_web(enhanced_query, max_results=8, user=user)
           
            if not web_results:
                logger.warning("No search results found")
                return "I couldn't find relevant information on the web for your query.", []
           
            # Step 3: Process search results and format them directly without GPT
            web_sources = []
            domains_seen = set()
           
            for i, result in enumerate(web_results, 1):
                title = result.get('title', 'Unknown Source')
                url = result.get('href', '')
                snippet = result.get('body', 'No description available')
               
                domain = self.extract_domain(url)
                if domain not in domains_seen and domain != "unknown.com":
                    domains_seen.add(domain)
                   
                    # Add to sources for return
                    web_sources.append({
                        'title': title,
                        'url': url,
                        'snippet': snippet[:300]
                    })
               
                # Format source information for display
                source_domains = list(domains_seen) if domains_seen else ["Google Search"]
                source_info = ", ".join(source_domains)
               
                logger.info("Successfully formatted direct search results")
                print("###########################################", web_results)
                return f"{web_results}\n\n*Sources: {source_info}*", web_sources
            else:
                logger.warning("No valid search results to format")
                return "I couldn't find relevant information on the web for your query.", []
           
        except Exception as e:
            logger.error(f"Error in web knowledge search: {str(e)}", exc_info=True)
            return f"An error occurred while searching the web: {str(e)}", []
 
    def search_web(self, query, max_results=5, user=None):
        """Search the web using Google Genai and return raw results with extracted sources"""
        try:
            # Initialize Google Genai client and tools
            import os
            import requests
            from google.genai.types import Tool, GoogleSearch, GenerateContentConfig
            from google import genai
           
            # Get API key from environment
            user_api_tokens = UserAPITokens.objects.get(user=user)
            GOOGLE_API_KEY = user_api_tokens.gemini_token
            if not GOOGLE_API_KEY:
                logger.error("GOOGLE_API_KEY not found in environment variables.")
                return []
           
            # Initialize Google Genai client
            genai_client = genai.Client(api_key=GOOGLE_API_KEY)
            model_id = "gemini-2.0-flash-exp"
           
            # Configure search tool
            search_tool = Tool(google_search=GoogleSearch())
           
            # System instruction to get raw search results
            system_instruction = (
                f"Search the web for: {query}. "
                f"Return the search results exactly as you find them with all original content, titles, URLs, and snippets. "
                f"CRITICAL: I need the ACTUAL source website URLs, not redirect URLs. "
                f"When you find search results, provide the original website URLs from where you found them (like moneycontrol.com, economictimes.com, etc.) NOT redirect URLs from vertexaisearch.cloud.google.com. "
                f"VERY IMPORTANT: Only include links that are VALID and ACCESSIBLE. Before including any link, verify that it leads to a real and working webpage (i.e., not a 404 or unreachable site). "
                f"Do NOT include placeholder links, broken URLs, or websites that fail to load. "
                f"Preserve all numerical data, prices, and specific information exactly as found."
            )
           
            # Configure generation parameters for raw results
            config = GenerateContentConfig(
                system_instruction=system_instruction,
                tools=[search_tool],
                response_modalities=["TEXT"],
                temperature=0.1,
                candidate_count=1
            )
           
            # Generate search results using Gemini
            logger.info(f"Searching for: {query}")
            response = genai_client.models.generate_content(
                model=model_id,
                contents=f"Search for: {query}. Return complete search results with titles, URLs, and content.",
                config=config
            )
           
            # Get the response text
            result_text = response.text
            logger.info("Received search response from Google Genai")
           
            results = []
           
            # Extract URLs from the response for source tracking
            url_pattern = r'https?://[^\s\)\]\"\'>]+'
            all_urls = re.findall(url_pattern, result_text)
           
            # Clean and filter URLs
            domains_seen = set()
           
            for url in all_urls:
                clean_url = url.rstrip('.,);\'\"')
               
                # Skip Google redirect URLs but try to resolve them
                if 'vertexaisearch.cloud.google.com' in clean_url:
                    try:
                        resolved_url = self.resolve_redirect_url(clean_url)
                        if resolved_url and resolved_url != clean_url:
                            clean_url = resolved_url
                        else:
                            continue  # Skip if can't resolve
                    except:
                        continue
               
                domain = self.extract_domain(clean_url)
                if domain != "unknown.com" and domain not in domains_seen:
                    domains_seen.add(domain)
                   
                    # Create individual result for each unique source
                    results.append({
                        'title': f"Information from {domain}",
                        'href': clean_url,
                        'body': result_text  # Use the raw response text directly
                    })
                   
                    if len(results) >= max_results:
                        break
           
            # If no sources found, return the raw content with generic source
            if not results:
                results = [{
                    'title': f"Search Results for: {query}",
                    'href': 'https://www.google.com',
                    'body': result_text  # Raw unprocessed content
                }]
           
            logger.info(f"Found {len(results)} search results with sources")
            return results
           
        except Exception as e:
            logger.error(f"Error searching the web via Google Genai: {str(e)}", exc_info=True)
            return []
 
    def resolve_redirect_url(self, redirect_url):
        """Resolve a redirect URL to get the actual destination URL"""
        try:
            import requests
           
            # Set a reasonable timeout and don't follow too many redirects
            response = requests.head(redirect_url, allow_redirects=True, timeout=5,
                                headers={'User-Agent': 'Mozilla/5.0 (compatible; SearchBot)'})
           
            # Return the final URL after following redirects
            final_url = response.url
            logger.info(f"Resolved redirect: {redirect_url} -> {final_url}")
            return final_url
           
        except Exception as e:
            logger.warning(f"Failed to resolve redirect URL {redirect_url}: {e}")
            return None
 
    def get_general_chat_answer(self, query, use_web_knowledge=False, response_length='comprehensive', response_format='natural', user=None):
        """
        Generate an answer for general chat mode, optionally using web knowledge.
       
        Args:
            query (str): The user's query
            use_web_knowledge (bool): Whether to use web search
            response_length (str): 'short' or 'comprehensive'
            response_format (str): The format to use for the response
           
        Returns:
            str: Generated response
        """
        # Check if the query is a simple greeting
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
           
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
           
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya",
           
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
           
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
           
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
           
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
       
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
       
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
       
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
       
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
       
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index]
       
        # If not a greeting, proceed with the regular flow
        # If web knowledge is enabled, implement the full web search flow
        if use_web_knowledge:
            try:
                # Step 1: Search the web using DuckDuckGo
                print(f"Searching web for: {query}")
                web_response, web_sources = self.get_web_knowledge_response(query, user=user, document_context=None)
 
                if not web_response or not web_sources:
                    return "I couldn't find relevant information on the web for your query."
                   
                # Extract web source domains for display
                web_source_domains = []
                for source in web_sources:
                    domain = self.extract_domain(source.get('url', ''))
                    if domain and domain != "unknown.com":
                        web_source_domains.append(domain)
               
                # Remove existing sources from web_response to avoid duplication in GPT processing
                if "*Sources:" in web_response:
                    web_response_clean = web_response.split("*Sources:")[0].strip()
                else:
                    web_response_clean = web_response
 
                # Get format-specific guidance
                format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
               
                # Create the prompt for OpenAI
                # Replace your existing prompt section with this:
 
                prompt = f"""
                You are a web research assistant. Based ONLY on the following information from multiple web sources, answer the user question.
 
                If relevant details are missing or contradictory, clearly acknowledge this and summarize available related content. Be helpful by highlighting potentially useful information even if it doesn't fully resolve the question. Provide quantitative details where available.
 
                RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
 
                {format_guidance}
 
                RESPONSE GENERATION GUIDELINES:
                - Provide a DETAILED, COMPREHENSIVE, and THOROUGH answer.
                - Include ALL relevant information from the web sources below.
                - Prioritize completeness over brevity — include important details.
                - If multiple perspectives or data points exist, include all of them.
                - When appropriate, structure the information clearly with headings.
                - Maintain a natural, conversational tone.
                - If the web sources do NOT contain relevant information, clearly state that and summarize any related or tangential insights.
 
                QUESTION: {query}
 
                INFORMATION FROM WEB SOURCES:
                {web_response_clean}
 
                RESPONSE FORMAT REQUIREMENTS:
                1. Follow with detailed elaboration on each relevant point.
                2. Include specific details, examples, and support from the source content.
                3. Use clear, logical sections when needed.
                4. Conclude with any other contextually relevant insights.
                5. Use <b> for section headers.
                6. Use <p> for body text.
                7. Use <ul> / <li> for lists.
                8. Use <table>, <tr>, <td> for tables, if applicable.
                9. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.
 
                CRITICAL INSTRUCTION - MUST BE FOLLOWED:
                At the very end of your response, after all content, add a new line and provide ONLY this exact format:
                SOURCES_JSON: {{"sources": ["domain1.com", "domain2.com", "domain3.com"]}}
 
                Extract the actual website domains mentioned in the web sources content. Do not include generic domains. Look for domain names that appear in the source content.
 
                RESPONSE LENGTH: {'Provide a focused, concise response prioritizing the most important information.' if response_length == 'short' else 'Provide a comprehensive, detailed response that thoroughly covers the topic.'}
                """
 
                system_message = "You are a web search analysis expert. Provide an EXTREMELY DETAILED and COMPREHENSIVE response and ALWAYS end with the SOURCES_JSON format."
 
                temperature = 0.3 if response_format in ['factual_brief', 'technical_deep_dive'] else 0.5
                max_tokens = 2000 if response_length == 'short' else 2000
 
                print(f"Calling OpenAI API with temperature={temperature}, max_tokens={max_tokens}")
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=temperature,
                    max_tokens=max_tokens
                )
 
                formatted_web_response = response.choices[0].message.content
 
                # Debug: Print the full response to see what GPT returned
                print("=== FULL GPT RESPONSE ===")
                print(formatted_web_response)
                print("=== END FULL RESPONSE ===")
 
                # Extract JSON sources from the response
                import json
                import re
                sources_extracted = []
 
                if "SOURCES_JSON:" in formatted_web_response:
                    try:
                        # Split the response to get the JSON part
                        parts = formatted_web_response.split("SOURCES_JSON:")
                        main_response = parts[0].strip()
                        json_part = parts[1].strip()
                       
                        print(f"=== JSON PART EXTRACTED ===")
                        print(f"JSON part: '{json_part}'")
                        print("=== END JSON PART ===")
                       
                        # Parse the JSON
                        sources_data = json.loads(json_part)
                        sources_extracted = sources_data.get("sources", [])
                       
                        print(f"=== EXTRACTED SOURCES ===")
                        print(f"Sources: {sources_extracted}")
                        print("=== END EXTRACTED SOURCES ===")
                       
                        # Use the main response without the JSON part
                        formatted_web_response = main_response
                       
                    except (json.JSONDecodeError, IndexError) as e:
                        logger.warning(f"Failed to parse sources JSON: {e}")
                        print(f"=== JSON PARSING ERROR ===")
                        print(f"Error: {e}")
                        print(f"Raw JSON part: '{json_part if 'json_part' in locals() else 'Not extracted'}'")
                        print("=== END JSON ERROR ===")
                       
                        # Manual extraction as fallback - look for patterns in the response
                        sources_extracted = []
                       
                        # Look for domain patterns in the response using regex
                        domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                        domains_found = re.findall(domain_pattern, formatted_web_response.lower())
                       
                        # Filter out common generic domains
                        generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                        for domain in domains_found:
                            if domain not in generic_domains and domain not in sources_extracted:
                                sources_extracted.append(domain)
                       
                        print(f"=== MANUAL EXTRACTION RESULTS ===")
                        print(f"Manually extracted sources: {sources_extracted}")
                        print("=== END MANUAL EXTRACTION ===")
                       
                else:
                    print("=== NO SOURCES_JSON FOUND ===")
                    print("SOURCES_JSON marker not found in response")
                   
                    # Manual extraction as fallback - look for domain patterns
                    sources_extracted = []
                   
                    # Look for domain patterns in the response using regex
                    domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                    domains_found = re.findall(domain_pattern, formatted_web_response.lower())
                   
                    # Filter out common generic domains
                    generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                    for domain in domains_found:
                        if domain not in generic_domains and domain not in sources_extracted:
                            sources_extracted.append(domain)
                   
                    print(f"=== FALLBACK EXTRACTION RESULTS ===")
                    print(f"Fallback extracted sources: {sources_extracted}")
                    print("=== END FALLBACK EXTRACTION ===")
 
                # Use extracted sources if available, otherwise fall back to web_source_domains
                if sources_extracted:
                    source_info = ", ".join(sources_extracted)
                    final_response = f"{formatted_web_response}\n\n*Sources: {source_info}*"
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", final_response)
                    return final_response
                else:
                    # Fallback to existing method
                    if web_source_domains:
                        source_info = ", ".join(web_source_domains)
                        final_response = f"{formatted_web_response}\n\n*Sources: {source_info}*"
                        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", final_response)
                        return final_response
                    else:
                        return formatted_web_response
 
               
            except Exception as e:
                logger.error(f"Error in web knowledge general chat: {str(e)}", exc_info=True)
                return f"I encountered an error while searching the web: {str(e)}. Please try a different question or try again later."
       
        # If no web knowledge requested, use standard chat completion
        else:
            # Format-specific guidance
            format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
           
            # Configure conciseness based on response length
            conciseness = "Be concise and to-the-point." if response_length == 'short' else "Provide a comprehensive and detailed response."
           
            # Create a prompt for the general chat
            prompt = f"""
            Please answer the following question in a helpful, informative way.
           
            RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
           
            {format_guidance}
           
            {conciseness}
           
            Use HTML tags for structure (<b>, <p>, <ul>, <li>) to enhance readability.
           
            USER QUERY: {query}
            """
           
            try:
                # Use the OpenAI API
                completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are a helpful, friendly AI assistant providing informative and thoughtful responses."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.5,
                    max_tokens=2000 if response_length == 'comprehensive' else 800
                )
               
                return completion.choices[0].message.content
               
            except Exception as e:
                logger.error(f"Error in general chat: {str(e)}", exc_info=True)
                return f"I'm sorry, I encountered an error while processing your request. Please try again or rephrase your question."
 
 
    
    def combine_document_and_web_responses(self, query, document_response, web_response, doc_sources, web_sources, response_format, conversation_context, original_doc_context=None, doc_citation_sources=None):
        """
        Combine document-based response and web-based response into a single coherent answer.
        Returns JSON format with enhanced citation handling.
        
        Args:
            query (str): The original user query
            document_response (str): Response generated from document context
            web_response (str or dict): Response generated from web search (can be JSON)
            doc_sources (list): Document sources
            web_sources (list): Web sources
            response_format (str): Desired format for the response
            conversation_context (str): Previous conversation context
            original_doc_context (list, optional): The original document context chunks
            doc_citation_sources (dict, optional): Document citation sources
            
        Returns:
            tuple: (combined_response_dict, combined_citation_sources)
        """
        # Check if the query is a simple greeting
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya",
            
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
        
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
        
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
        
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
        
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
        
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            
            return {
                "content": greeting_responses[response_index],
                "sources": "",
                "citation_sources": {},
                "response_type": "greeting",
                "json_response": True
            }, {}

        # Initialize combined citation sources with document citations
        combined_citation_sources = doc_citation_sources.copy() if doc_citation_sources else {}
        
        # Handle web response format (could be string or dict)
        if isinstance(web_response, dict):
            web_content = web_response.get("content", str(web_response))
            web_sources_from_response = web_response.get("sources", [])
        else:
            web_content = str(web_response)
            web_sources_from_response = []
        
        # If not a greeting, proceed with the regular flow
        # Clean up responses by removing source sections
        if "\n\n*Sources:" in document_response:
            document_response = document_response.split("\n\n*Sources:")[0]
            
        if "\n\n*Sources:" in web_content:
            web_response_clean = web_content.split("\n\n*Sources:")[0]
        else:
            web_response_clean = web_content
        
        # Extract web source domains for final source list
        web_source_domains = []
        for source in web_sources:
            domain = self.extract_domain(source.get('url', ''))
            if domain and domain != "unknown.com":
                web_source_domains.append(domain)
        
        # Detect if this is a comparison query that needs special handling
        comparison_words = ["compare", "comparison", "versus", "vs", "difference", "similarities",
                            "better than", "worse than", "alternative", "competitors", "competition",
                            "similar to", "different from", "compared to"]
        
        is_comparison_query = any(word in query_lower for word in comparison_words)
        
        # Analyze the document response and original context to extract key entities for comparison
        # Prepare document context if available
        doc_context_text = ""
        if original_doc_context and len(original_doc_context) > 0:
            # Use up to 5 most relevant chunks from the original document context
            doc_context_text = "\n\n".join(original_doc_context[:5])
        
        # Create analysis prompt for JSON response
        analyze_prompt = f"""
        Analyze the following information to identify key details about the topic being discussed.
        
        USER QUERY: {query}
        
        DOCUMENT RESPONSE:
        {document_response}
        
        {f"ORIGINAL DOCUMENT CONTEXT:{doc_context_text}" if doc_context_text else ""}
        
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "main_entities": ["list", "of", "exact", "entity", "names"],
            "domain": "specific domain or category",
            "key_attributes": ["important", "attributes", "mentioned"],
            "mentioned_competitors": ["any", "competitors", "mentioned"],
            "is_comparison": true/false,
            "comparison_focus": "what aspects are being compared if applicable",
            "industry_context": "specific industry or market context",
            "missing_information": ["what", "important", "info", "seems", "missing"]
        }}
        """
        
        try:
            # Extract key information from document response to enhance web results integration
            analysis_response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert analyst who extracts structured information from text. Respond only in JSON format."},
                    {"role": "user", "content": analyze_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON
                temperature=0.2,
                max_tokens=500
            )
            
            # Parse JSON analysis
            analysis_json = self._parse_json_response(analysis_response.choices[0].message.content)
            logger.info(f"Document analysis: {analysis_json}")
            
            # Check if web response seems insufficient based on analysis
            web_info_lacking = False
            if "Unfortunately, there was no additional relevant information found on the web" in web_response_clean or "no relevant information" in web_response_clean.lower():
                web_info_lacking = True
                logger.info("Web information appears to be lacking")
            
            # Determine if we need to use model knowledge for a better response
            use_model_knowledge = False
            
            # If this is a comparison query and web info is lacking, use model knowledge
            if is_comparison_query and (web_info_lacking or not web_sources):
                use_model_knowledge = True
                logger.info("Using model knowledge to supplement missing web information for comparison")
            
            # If web is lacking information about mentioned entities, use model knowledge
            if web_info_lacking and analysis_json.get('main_entities'):
                use_model_knowledge = True
                logger.info("Using model knowledge to supplement information about entities")
            
            # Create system message for JSON response combination
            system_message = """You are an expert at combining information from multiple sources and using your knowledge to provide comprehensive answers.
            You must respond in JSON format with the following structure:
            {
                "content": "Your comprehensive response with proper HTML formatting and citations",
                "sources_combined": ["list", "of", "all", "sources", "used"],
                "response_sections": {
                    "document_info": "key points from documents",
                    "web_info": "key points from web",
                    "additional_knowledge": "any additional knowledge provided"
                },
                "combination_type": "standard/knowledge_enhanced",
                "citations_included": true/false,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }
            
            CRITICAL INSTRUCTION FOR WEB SOURCES - MUST BE FOLLOWED:
            In your JSON response, include a "web_sources_extracted" field that contains an array of actual website domains mentioned in the web response content. Extract the actual website domains mentioned in the web response content. Do not include generic domains. Look for domain names that appear in the web source content."""
            
            # Create enhanced prompt based on the analysis and whether to use model knowledge
            if use_model_knowledge:
                # Extract key information from the analysis
                main_entities = analysis_json.get('main_entities', [])
                domain = analysis_json.get('domain', '')
                mentioned_competitors = analysis_json.get('mentioned_competitors', [])
                industry_context = analysis_json.get('industry_context', '')
                
                # Enhanced prompt without any explicit brand examples
                enhanced_prompt = f"""
                You have a document response and a web response to the user's query. Your task is to combine these into a coherent response.

                DETAILED CONTEXT:
                - Domain: {domain if domain else "Not clearly identified"}
                - Industry: {industry_context if industry_context else "Not specified"}
                - Main entities: {', '.join(main_entities) if main_entities else "Not clearly identified"}
                - Mentioned competitors: {', '.join(mentioned_competitors) if mentioned_competitors else "None explicitly mentioned"}

                USER QUERY: {query}

                DOCUMENT RESPONSE:
                {document_response}

                WEB RESPONSE:
                {web_response_clean}

                SPECIAL INSTRUCTIONS:
                1. The web search did not return sufficiently helpful information, especially about additional competitors or alternatives not mentioned in the document.

                2. Use your built-in knowledge to SIGNIFICANTLY EXPAND the information, focusing on:
                - ADDITIONAL competitors or alternatives NOT mentioned in the document
                - The broader competitive landscape beyond just the entities mentioned in the document
                - Different market segments and their leading brands relevant to this industry
                - Comparing the main entities with other market players not covered in the document or web results
                
                3. When using your built-in knowledge:
                - Clearly label it as "Additional Market Information"
                - Focus on providing NEW information not found in either the document or web results
                - Include major competing products/services relevant to this specific industry
                - Include international alternatives if relevant
                - Mention different positioning of these alternatives (by demographic, price point, features, etc.)

                4. Format your response with these sections:
                - <h3>Information from Documents</h3> (prioritize this information)
                - <h3>Information from the Web</h3> (include relevant web information)
                - <h3>Additional Market Information</h3> (for supplementary information from your knowledge, focusing on NEW competitors and alternatives)
                
                5. In cases where you are adding information about competitors or alternatives:
                - Focus on major players NOT mentioned in either the document or web response
                - Include at least 3-5 additional competing entities if relevant to the domain
                - Highlight how they differ from the main entities in the document
                - Include details about:
                    * Target audience or use case
                    * Distinctive features or characteristics
                    * Market positioning (premium, mid-range, budget)
                    * Distribution approach
                    * Unique selling propositions
                - Organize the information in a clear, structured way (tables or bullet points if appropriate)

                6. The final response should:
                - Begin with document information (it's specific to the user's context)
                - Include relevant web information when available
                - Add substantial new knowledge about additional competitors and market alternatives
                - Be clearly structured and easy to follow
                - Use proper HTML formatting (<h3>, <b>, <p>, <ul>, <li>, <table> if appropriate)

                The response format should be: {response_format.replace('_', ' ').title()}
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "content": "Your comprehensive response with HTML formatting",
                    "sources_combined": ["document sources", "web sources", "knowledge base"],
                    "response_sections": {{
                        "document_info": "summary of document insights",
                        "web_info": "summary of web insights", 
                        "additional_knowledge": "summary of additional knowledge provided"
                    }},
                    "combination_type": "knowledge_enhanced",
                    "citations_included": false,
                    "web_sources_extracted": ["domain1.com", "domain2.com"]
                }}
                """
                
            else:
                # Update the prompt to include citation information and source JSON extraction
                citation_instructions = """
                CITATION GUIDELINES:
                When referencing information from the document sources, cite the source using [n] format.
                When referencing information from web sources, cite the source using [Wn] format (e.g., [W1], [W2]).
                Every statement that comes directly from a source should include a citation.
                Place citations immediately after the information they support.
                """
                
                # Standard combination prompt if we don't need model knowledge
                enhanced_prompt = f"""
                You have two responses to the user's query: one based on their documents and one based on web search.
                Your task is to combine these into a single coherent response that prioritizes the document information
                (since that is more specific to the user) but supplements with web information when valuable.

                {citation_instructions}

                The response format should be: {response_format.replace('_', ' ').title()}

                Previous conversation context:
                {conversation_context}

                User query: {query}

                RESPONSE FROM DOCUMENTS:
                {document_response}

                RESPONSE FROM WEB SEARCH:
                {web_response_clean}

                GUIDELINES FOR COMBINATION:
                1. When information appears in both sources, prioritize the document information.

                2. Clearly distinguish when you're providing information from the web vs. from documents.

                3. The web information should EXPAND BEYOND what's in the documents - focus on adding NEW information:
                - Additional competitors or alternatives NOT mentioned in the documents
                - Information about different market segments not covered in the documents
                - Broader competitive landscape and market positioning
                - If the query is about competitors, aim to mention additional major players not in the document

                4. Maintain a natural tone and logical flow between document and web information.

                5. If there are contradictions, note them and explain the different perspectives.

                6. Use HTML tags for structure: <p>, <ul>, <li>, <b>, etc.

                7. Structure your response with proper headings and sections for readability.

                8. Use <h3> tags for different sources:
                - <h3>Information from Documents</h3> for document-based information.
                - <h3>Information from the Web</h3> for web-based information.

                9. For document information, maintain the existing citation numbers [n].
                10. For web information, use [W1], [W2], etc. format.

                11. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.

                Create a single, coherent, well-structured response that combines both sources of information.
                Make sure your response is contextually relevant to the ongoing conversation.

                IMPORTANT: If the web search doesn't provide much additional information beyond what's in the documents,
                use your background knowledge to expand the response with additional relevant information about competitors,
                alternatives, or market players that aren't mentioned in either source.
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "content": "Your comprehensive response with HTML formatting",
                    "sources_combined": ["list", "of", "all", "sources"],
                    "response_sections": {{
                        "document_info": "key insights from documents",
                        "web_info": "key insights from web"
                    }},
                    "combination_type": "standard",
                    "citations_included": true,
                    "web_sources_extracted": ["domain1.com", "domain2.com"]
                }}
                """
            
            # Generate combined response with JSON format
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON
                temperature=0.3,
                max_tokens=2500
            )
            
            # Parse the JSON response
            combined_json = self._parse_json_response(completion.choices[0].message.content)
            combined_response = combined_json.get("content", "Error generating combined response")
            
            # Debug: Print the full response to see what GPT returned
            print("=== FULL COMBINED RESPONSE ===")
            print(combined_response)
            print("=== END FULL RESPONSE ===")
            
            # Extract web sources from JSON response
            web_sources_extracted = combined_json.get("web_sources_extracted", [])
            
            print(f"=== EXTRACTED WEB SOURCES FROM JSON ===")
            print(f"Web Sources: {web_sources_extracted}")
            print("=== END EXTRACTED WEB SOURCES ===")
            
            # If no web sources extracted from JSON, fall back to manual extraction
            if not web_sources_extracted:
                import re
                # Look for domain patterns in the response using regex
                domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                domains_found = re.findall(domain_pattern, combined_response.lower())
                
                # Filter out common generic domains
                generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                for domain in domains_found:
                    if domain not in generic_domains and domain not in web_sources_extracted:
                        web_sources_extracted.append(domain)
                
                print(f"=== FALLBACK EXTRACTION RESULTS ===")
                print(f"Fallback extracted web sources: {web_sources_extracted}")
                print("=== END FALLBACK EXTRACTION ===")
            
            # Process citations in the combined response to include web sources
            web_citation_pattern = r'\[W(\d+)\]'
            web_citations = re.findall(web_citation_pattern, combined_response)

            # Add web sources to combined_citation_sources
            if web_citations and web_sources:
                web_citation_base = max(combined_citation_sources.keys()) + 1 if combined_citation_sources else 1
                
                for i, web_citation in enumerate(set(web_citations)):
                    web_idx = int(web_citation) - 1
                    if web_idx < len(web_sources):
                        web_source = web_sources[web_idx]
                        citation_num = web_citation_base + i
                        
                        # Add web source to citation sources
                        combined_citation_sources[citation_num] = {
                            'source_file': web_source.get('title', 'Web Source'),
                            'page_number': 'Web',
                            'section_title': 'Web Search Result',
                            'snippet': web_source.get('snippet', '')[:200] + "..." if web_source.get('snippet', '') else "Web content",
                            'document_id': f"web_{web_idx}",
                            'url': web_source.get('url', '')
                        }
                        
                        # Replace [Wn] with [citation_num]
                        combined_response = combined_response.replace(f"[W{web_citation}]", f"[{citation_num}]")

            # Create combined sources list for final display
            all_doc_sources = list(set(doc_sources))
            
            # Use extracted web sources if available, otherwise fall back to web_source_domains
            if web_sources_extracted:
                final_web_sources = web_sources_extracted
            elif web_source_domains:
                final_web_sources = web_source_domains
            else:
                final_web_sources = []
            
            # Include web sources from JSON response if available
            if isinstance(web_response, dict) and web_response.get("sources"):
                final_web_sources.extend(web_response["sources"])
            
            all_sources = all_doc_sources + list(set(final_web_sources))
            
            # If we used model knowledge, add a note
            if use_model_knowledge:
                if all_sources:
                    all_sources_str = ", ".join(all_sources) + ", and general knowledge"
                else:
                    all_sources_str = "Document context and general knowledge"
            else:
                all_sources_str = ", ".join(all_sources) if all_sources else "Document context only"

            # Clean up and process citations in the final response
            processed_response, processed_citation_sources = self._format_citations_for_response(combined_response, combined_citation_sources)
            
            # Return JSON structure instead of plain text
            print("$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$ FINAL COMBINED RESPONSE", f"{processed_response}\n\n*Sources: {all_sources_str}*")
            return {
                "content": f"{processed_response}\n\n*Sources: {all_sources_str}*",
                "sources": all_sources_str,
                "citation_sources": processed_citation_sources,
                "document_sources": all_doc_sources,
                "web_sources": final_web_sources,
                "combination_type": combined_json.get("combination_type", "standard"),
                "response_sections": combined_json.get("response_sections", {}),
                "used_model_knowledge": use_model_knowledge,
                "is_comparison_query": is_comparison_query,
                "json_response": True,
                "success": True
            }, processed_citation_sources
            
        except Exception as e:
            logger.error(f"Error combining responses: {str(e)}", exc_info=True)
            
            # Fallback: Simply concatenate the responses with a divider in JSON format
            fallback_system_message = """You are an assistant that combines information from multiple sources.
            You must respond in JSON format with this structure:
            {
                "content": "Your combined response with HTML formatting",
                "sources": "list of sources used",
                "combination_type": "fallback",
                "error_handled": true,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }"""
            
            fallback_prompt = f"""
            Combine the following information sources for the user's query: "{query}"
            
            Document Information:
            {document_response}
            
            Web Information:
            {web_response_clean}
            
            Create a coherent response that includes both sources of information.
            Use <h3>Information from Your Documents:</h3> and <h3>Additional Information from Web Search:</h3> as section headers.
            
            CRITICAL: Your response MUST be a valid JSON object with this structure:
            {{
                "content": "Your combined response with HTML formatting",
                "sources": "list of sources",
                "combination_type": "fallback",
                "error_handled": true,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }}
            """
            
            try:
                fallback_response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": fallback_system_message},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON
                    temperature=0.3,
                    max_tokens=1500
                )
                
                fallback_json = self._parse_json_response(fallback_response.choices[0].message.content)
                fallback_content = fallback_json.get("content", "Error generating fallback response")
                fallback_web_sources = fallback_json.get("web_sources_extracted", [])
                
            except Exception as fallback_error:
                logger.error(f"Fallback response generation failed: {str(fallback_error)}")
                # Ultimate fallback - create JSON manually
                fallback_content = f"""
                <h3>Information from Your Documents:</h3>
                <div>
                {document_response}
                </div>
                
                <h3>Additional Information from Web Search:</h3>
                <div>
                {web_response_clean}
                </div>
                """
                fallback_web_sources = []
            
            # Add sources for fallback
            all_doc_sources = list(set(doc_sources))
            
            # Use extracted web sources if available from fallback, otherwise fall back to web_source_domains
            if fallback_web_sources:
                final_web_sources = fallback_web_sources
            elif web_source_domains:
                final_web_sources = web_source_domains
            else:
                final_web_sources = []
                
            # Include web sources from JSON response if available
            if isinstance(web_response, dict) and web_response.get("sources"):
                final_web_sources.extend(web_response["sources"])
                
            all_sources = all_doc_sources + list(set(final_web_sources))
            all_sources_str = ", ".join(all_sources) if all_sources else "Document and web sources"
            
            return {
                "content": f"{fallback_content}\n\n*Sources: {all_sources_str}*",
                "sources": all_sources_str,
                "citation_sources": doc_citation_sources if doc_citation_sources else {},
                "document_sources": all_doc_sources,
                "web_sources": final_web_sources,
                "combination_type": "fallback",
                "error": str(e),
                "error_handled": True,
                "json_response": True,
                "success": False
            }, doc_citation_sources if doc_citation_sources else {}

    

    def detect_question_format(self, query, context_snippets=None):
        """
        Analyzes the question and detects the most appropriate response format
        based on the query content and any available document context.
        
        Args:
            query (str): The user's question
            context_snippets (list): Optional list of document context snippets
            
        Returns:
            tuple: (primary_format, secondary_format, confidence)
        """
        # Implementation remains the same as in your original code
        print("###########################################################")
        
        # Normalize query for keyword matching
        query_lower = query.lower()
        
        # Format detection keywords dictionary
        format_keywords = {
            "executive_summary": [
                "summarize", "summary", "overview", "high level", "highlights", "key points", 
                "brief overview", "executive summary", "main points", "takeaways", "tldr",
                "in a nutshell", "brief summary", "key highlights", "summarise", "recap",
                "give me the gist", "summarization", "outline the main", "key takeaways",
                "bullet points of", "boil down", "distill", "condense"
            ],
            "comparative_analysis": [
                "compare", "comparison", "versus", "vs", "difference", "differences", 
                "similarities", "similarity", "better", "worse", "advantage", "disadvantage",
                "pros and cons", "strengths and weaknesses", "contrast", "choose between",
                "which is better", "how does", "stack up", "relative to", "compared to",
                "weighing", "juxtapose", "differentiate between", "distinction", "compare and contrast",
                "side by side", "match up against", "measure against"
            ],
            "detailed_analysis": [
                "analyze", "analysis", "deep dive", "in-depth", "detailed", "breakdown",
                "examine", "investigate", "elaborate on", "explain in detail", "thorough",
                "comprehensive", "drill down", "what factors", "analyze in detail",
                "dissect", "unpack", "scrutinize", "evaluate", "assessment", "diagnosis",
                "root cause", "deconstructing", "implications of", "examine carefully", 
                "look closely at", "study"
            ],
            "strategic_recommendation": [
                "recommend", "recommendation", "suggest", "advice", "should", "best course",
                "action plan", "strategic", "next steps", "how should", "what would you advise",
                "strategy", "approach", "what should", "best way to", "how to best",
                "action items", "propose", "guidance", "direction", "advise on",
                "what's the right approach", "strategic guidance", "roadmap", "plan for",
                "steps to take", "best practice", "optimal strategy"
            ],
            "market_insights": [
                "market trends", "industry", "segment", "consumer behavior", "market data", 
                "competitive landscape", "market share", "industry trend", "consumer preference",
                "market analysis", "market research", "market performance", "market outlook",
                "sector analysis", "competitors", "customer demographics", "forecast",
                "target audience", "customer segment", "market growth", "consumer insights",
                "purchase behavior", "audience analysis", "industry position"
            ],
            "factual_brief": [
                "what is", "when was", "where is", "who is", "facts about", "define", 
                "information on", "tell me about", "data on", "evidence of", "statistics on",
                "key facts", "figures on", "numbers for", "metrics on", "how many",
                "percentage of", "rate of", "frequency of", "occurrence of", "incidents of",
                "quantity of", "value of", "amount of", "total number", "statistics for",
                "list the", "enumerate", "identify", "reference"
            ],
            "technical_deep_dive": [
                "how does it work", "technical", "technology", "mechanism", "process", "methodology",
                "framework", "system", "architecture", "technical details", "underlying", "step by step",
                "explain the process", "technical explanation", "in technical terms",
                "inner workings", "components", "operation", "workflow", "algorithm",
                "explain technically", "implementation", "procedure", "specification",
                "design", "function", "operation of", "engineering behind"
            ]
        }
        
        # Check for format-specific keywords
        matched_formats = []
        for format_key, keywords in format_keywords.items():
            # Count matches and weight longer phrases higher
            match_score = 0
            for keyword in keywords:
                if keyword in query_lower:
                    # Weight multi-word phrases higher
                    words_in_keyword = len(keyword.split())
                    match_score += words_in_keyword  
            
            if match_score > 0:
                matched_formats.append((format_key, match_score))
        
        # Sort by match score (higher score = stronger signal)
        matched_formats.sort(key=lambda x: x[1], reverse=True)
        
        # If we have keyword matches, use the best match
        if matched_formats and matched_formats[0][1] >= 1:
            primary_format = matched_formats[0][0]
            # Set confidence based on match score
            base_confidence = 5.0
            # Boost confidence based on match score
            match_bonus = min(matched_formats[0][1], 4.0)  # Cap the bonus
            confidence = base_confidence + match_bonus
            
            # Get secondary format if available
            secondary_format = matched_formats[1][0] if len(matched_formats) > 1 else None
            
            return primary_format, secondary_format, confidence
        
        # If no keyword matches, we can try using LLM for classification
        # This would integrate with your existing classify_question_intent function
        # but we'll keep it simple for now
        
        # Default to natural if no clear classification
        return "natural", None, 4.0


    # Add this helper method to get format-specific guidance
    def _get_format_guidance(self, response_format, length_preference):
        """
        Get format-specific guidance based on response format and length preference.
        
        Args:
            response_format (str): The requested response format
            length_preference (str): 'short' or 'comprehensive'
            
        Returns:
            str: Format-specific guidance text
        """
        # Format guidance dictionary
        format_guidance = {
            "executive_summary": {

                 """
    This is an EXECUTIVE SUMMARY request. Provide a concise, high-level overview of the key points and implications. Specifically:

    1. Begin with a brief 1-2 sentence overview of the main finding or takeaway
    2. Structure your summary with clear sections for different key points
    3. Prioritize only the most important findings and insights (quality over quantity)
    4. Include critical data points and evidence that support the key messages
    5. Ensure every statement adds value and nothing is redundant
    6. Conclude with the most important implications or next steps
    7. Keep the entire summary focused and concise - aim for brevity while maintaining completeness

    Your summary should give the reader a clear understanding of the essential points without needing to read a longer document.
    """
            },
            "detailed_analysis": {

                 """
    This is a DETAILED ANALYSIS request. Provide a comprehensive analysis with in-depth examination of the information and its implications. Specifically:

    1. Begin with an executive summary of your key findings
    2. Include a rigorous analysis of the data, breaking it down into logical sections
    3. Examine patterns, trends, anomalies, and relationships in the data
    4. Support your analysis with specific evidence and examples from the documents
    5. Consider multiple perspectives and interpretations of the information
    6. Discuss implications of your findings for decision-making
    7. Include relevant charts or tables to visualize key data points
    8. Conclude with synthesis of the most important insights

    Your analysis should be thorough and demonstrate critical thinking throughout.
    """
            },
            "strategic_recommendation": {

               """
    This is a STRATEGIC RECOMMENDATION request. Provide actionable, prioritized recommendations with clear strategic direction. Specifically:

    1. Begin with a brief overview of the current situation or challenge
    2. Present 3-5 specific, actionable recommendations in priority order
    3. For each recommendation:
    - Clearly state what should be done
    - Explain why this approach is recommended (with evidence)
    - Outline expected outcomes or benefits
    - Address potential implementation challenges
    4. Include a high-level implementation timeline or roadmap if appropriate
    5. Conclude with the critical success factors for implementation

    Your recommendations should be practical, evidence-based, and strategically sound.
    """
            },
            "comparative_analysis": {
                """
    This is a COMPARATIVE question. Organize your response around comparing different options, approaches, or entities mentioned in the query. Specifically:

    1. Begin with a brief introduction of the items being compared
    2. Use a structure that clearly contrasts the different options (side-by-side format where possible)
    3. Create comparison tables that directly align comparable attributes
    4. Identify key differentiating factors between the options
    5. Discuss relative advantages and disadvantages of each
    6. Include direct quantitative comparisons where data is available
    7. Conclude with insights about the relative positioning of the compared items

    When presenting comparative data, use tables with this format when appropriate:

    | Attribute | Option A | Option B | Key Difference |
    |-----------|----------|----------|---------------|
    | Feature 1 | Value    | Value    | Insight       |
    """
            },
            "market_insights": {

               """
    This is a MARKET INSIGHTS request. Focus on market trends, consumer behavior, and competitive dynamics. Specifically:

    1. Begin with an overview of the key market dynamics or trends
    2. Analyze market segments, sizes, and growth rates with supporting data
    3. Examine competitive landscape and relative positioning of key players
    4. Discuss consumer behavior patterns, preferences, and changing demands
    5. Identify market opportunities and potential threats
    6. Present relevant market statistics in tabular format where appropriate
    7. Conclude with strategic implications of these market insights

    Your response should combine data-driven market analysis with strategic interpretation.
    """
            },
            "factual_brief": {

             """
    This is a FACTUAL INQUIRY. Provide direct, fact-based answers with supporting evidence. Specifically:

    1. Begin with a clear, direct answer to the question asked
    2. Present the key facts in a structured, logical order
    3. Support each fact with specific evidence from the documents
    4. Include relevant numbers, statistics, or data points
    5. Use bullet points or numbered lists where appropriate to improve readability
    6. Avoid speculation or unsupported assertions
    7. Clearly indicate if certain requested information is not available in the documents

    Your response should be factual, concise, and directly address the query.
    """
            },
            "technical_deep_dive": {

             """
    This is a TECHNICAL question. Include detailed explanations of processes, methodologies, or procedures. Specifically:

    1. Begin with a high-level overview of the technical concept or process
    2. Break down the technical components or steps in a logical sequence
    3. Explain how each component works or contributes to the whole
    4. Include relevant technical specifications, parameters, or metrics
    5. Use clear examples to illustrate technical concepts
    6. Consider including diagrams or technical illustrations where helpful
    7. Address limitations, constraints, or technical considerations
    8. Conclude with practical applications or implications

    Your explanation should be technically accurate while remaining accessible to a business audience.
    """
            },
            "natural": {
                "short": "Structure your response in a natural way that best addresses the query, while keeping it concise and to-the-point.",
                "comprehensive": "Structure your response in the most natural way that best addresses the query, with appropriate depth and detail."
            }
        }
        
        # Return the appropriate guidance based on format and length preference
        if response_format in format_guidance and length_preference in format_guidance[response_format]:
            return format_guidance[response_format][length_preference]
        
        # Default fallback
        return "Provide a helpful response to the query based on the document context."


    # Add these citation-specific prompt templates to the ChatView class

    CITATION_QA_TEMPLATE = """
    Please provide an answer based solely on the provided sources. 
    When referencing information from a source, cite the appropriate source(s) using their corresponding numbers within square brackets.
    Every statement in your answer should include at least one source citation.
    Only cite a source when you are explicitly referencing it.
    If none of the sources are helpful, you should indicate that.

    Examples of good citations:
    1. 'The sky is red in the evening [1]'
    2. 'Water is wet when the sky is red [2], which occurs in the evening [1]'
    3. 'According to the research [3], both findings [1][2] support this conclusion'

    Your answer should be clear, accurate, and directly tied to the source material.
    Below are several numbered sources of information:

    {context_str}

    Query: {query_str}
    Answer: 
    """

    CITATION_REFINE_TEMPLATE = """
    You are refining an existing answer to make it more accurate and well-cited.
    When referencing information from a source, cite the appropriate source(s) using their corresponding numbers within square brackets.
    Every statement in your answer should include at least one source citation.
    Only cite a source when you are explicitly referencing it.

    We have provided an existing answer: {existing_answer}

    Below are numbered sources of information.
    Use them to refine the existing answer, ensuring all statements have proper citations.

    {context_msg}

    Query: {query_str}
    Refined Answer:
    """


    def _clean_citations(self, text, citation_sources):

        """Clean up citation format to ensure consistent [n][m] format and remove duplicates."""

        # Fix citations that might not have brackets

        text = re.sub(r'(\s)(\d+)(\s)', lambda m: m.group(1) + '[' + m.group(2) + ']' + m.group(3) 

                    if self._is_likely_citation(int(m.group(2)), citation_sources) else m.group(0), text)

        # Fix spaces in citations and normalize format

        text = re.sub(r'\[\s*(\d+)\s*\]', r'[\1]', text)

        # Handle multiple citations in various formats and convert to separate brackets

        # Match patterns like [1, 2], [1,2], [1 2], etc.

        def normalize_multiple_citations(match):

            citation_content = match.group(1)

            # Extract all numbers

            numbers = re.findall(r'\d+', citation_content)

            # Remove duplicates while preserving order

            seen = set()

            unique_numbers = []

            for num in numbers:

                if num not in seen and self._is_likely_citation(int(num), citation_sources):

                    seen.add(num)

                    unique_numbers.append(num)

            # Return as separate brackets

            return ''.join(f'[{num}]' for num in unique_numbers)

        # Apply normalization for multiple citation patterns

        text = re.sub(r'\[([0-9,\s]+)\]', normalize_multiple_citations, text)

        # Remove duplicate citations and clean up

        text = self._remove_consecutive_duplicate_citations(text)

        return text
    


    def _ensure_separate_citation_brackets(self, text):
        """Ensure all citations are in separate brackets [1][2] not [1,2] or [1, 2]."""
        # Handle any remaining grouped citations
        def separate_brackets(match):
            numbers = re.findall(r'\d+', match.group(1))
            # Remove duplicates while preserving order
            seen = set()
            unique_numbers = []
            for num in numbers:
                if num not in seen:
                    seen.add(num)
                    unique_numbers.append(num)
            return ''.join(f'[{num}]' for num in unique_numbers)
        # Match any remaining grouped citations
        text = re.sub(r'\[([0-9,\s]+)\]', separate_brackets, text)
        # Final cleanup - remove any remaining duplicates in the same vicinity
        # Look for patterns like [1][1] within a short span
        words = text.split()
        cleaned_words = []
        for word in words:
            # If word contains multiple identical citations, deduplicate
            if '[' in word and ']' in word:
                citations = re.findall(r'\[(\d+)\]', word)
                if len(citations) != len(set(citations)):
                    # Remove duplicates
                    unique_citations = []
                    seen = set()
                    for cite in citations:
                        if cite not in seen:
                            unique_citations.append(cite)
                            seen.add(cite)
                    # Rebuild the word
                    citation_pattern = r'\[\d+\]'
                    parts = re.split(citation_pattern, word)
                    # Reconstruct with unique citations
                    result = parts[0]
                    for i, cite in enumerate(unique_citations):
                        result += f'[{cite}]'
                        if i + 1 < len(parts):
                            result += parts[i + 1]
                    cleaned_words.append(result)
                else:
                    cleaned_words.append(word)
            else:
                cleaned_words.append(word)
        return ' '.join(cleaned_words)
    
    def _is_likely_citation(self, num, citation_sources):
        """Check if a number is likely to be a citation based on available sources."""
        return num in citation_sources and 1 <= num <= len(citation_sources)

    def _format_citations_for_response(self, response_text, citation_sources):
        """Format the citations for response and ensure sequential numbering"""
        # First extract ALL citations from the text
        citation_pattern = r'\[(\d+)\]'
        all_citations = re.findall(citation_pattern, response_text)
        
        # Get unique citations while preserving order of first appearance
        unique_citations = []
        seen = set()
        for cite in all_citations:
            if cite.isdigit() and int(cite) in citation_sources and cite not in seen:
                seen.add(cite)
                unique_citations.append(cite)
        
        # Create mapping from original to sequential numbers (1,2,3...)
        citation_mapping = {int(old): new for new, old in enumerate(unique_citations, 1)}
        
        # Replace ALL citations in the text with sequential numbers
        def replace_citation(match):
            old_num = match.group(1)
            if old_num.isdigit() and int(old_num) in citation_mapping:
                return f'[{citation_mapping[int(old_num)]}]'
            return match.group(0)
        
        processed_text = re.sub(citation_pattern, replace_citation, response_text)
        
        # Create new citation sources with sequential numbers
        display_citation_sources = {}
        for display_num, original_num in enumerate(unique_citations, 1):
            original_num = int(original_num)
            if original_num in citation_sources:
                source_info = citation_sources[original_num]
                snippet = source_info.get('text', '')
                snippet_length = 1500 
                display_citation_sources[display_num] = {
                    'source_file': source_info.get('source_file', 'Unknown'),
                    'page_number': 'Unknown',
                    'section_title': 'Unknown',
                    'snippet': snippet[:snippet_length] + "..." if len(snippet) > snippet_length else snippet,
                    'document_id': source_info.get('document_id', 'Unknown')
                }
        
        return processed_text, display_citation_sources

    def _remove_consecutive_duplicate_citations(self, text):
        """Remove consecutive duplicate citations and ensure proper bracket separation."""
        # First, handle grouped citations like [1, 11] and convert to [1][11]
        def separate_grouped_citations(match):
            citation_content = match.group(1)
            # Split by comma and clean up
            citation_nums = [num.strip() for num in citation_content.split(',') if num.strip().isdigit()]
            # Remove duplicates while preserving order
            seen = set()
            unique_citations = []
            for num in citation_nums:
                if num not in seen:
                    seen.add(num)
                    unique_citations.append(num)
            # Return as separate brackets
            return ''.join(f'[{num}]' for num in unique_citations)
        # Pattern to match grouped citations like [1, 11] or [1,11] or [1 11]
        grouped_pattern = r'\[([0-9,\s]+)\]'
        text = re.sub(grouped_pattern, separate_grouped_citations, text)
        # Remove consecutive identical citations like [1][1] -> [1]
        consecutive_pattern = r'\[(\d+)\](?:\s*\[\1\])+' 
        def replace_consecutive_duplicates(match):
            citation_num = match.group(1)
            return f'[{citation_num}]'
        text = re.sub(consecutive_pattern, replace_consecutive_duplicates, text)
        # Remove citations separated by minimal punctuation like [1], [1] -> [1]
        punctuation_pattern = r'\[(\d+)\]\s*[,.;:!?]*\s*\[\1\]'
        text = re.sub(punctuation_pattern, r'[\1]', text)
        # Remove duplicate citations within the same sentence (more aggressive)
        # Find all citations in each sentence and deduplicate
        sentences = text.split('.')
        processed_sentences = []
        for sentence in sentences:
            # Find all citations in this sentence
            citations_in_sentence = re.findall(r'\[(\d+)\]', sentence)
            if len(citations_in_sentence) > len(set(citations_in_sentence)):
                # There are duplicates, remove them
                seen_in_sentence = set()
                def replace_if_seen(match):
                    num = match.group(1)
                    if num in seen_in_sentence:
                        return ''  # Remove duplicate
                    else:
                        seen_in_sentence.add(num)
                        return match.group(0)  # Keep first occurrence
                sentence = re.sub(r'\[(\d+)\]', replace_if_seen, sentence)
            processed_sentences.append(sentence)
        return '.'.join(processed_sentences)

    def generate_response(self, query, context, sources, use_web_knowledge=False, response_format='natural', conversation_context=""):
        """
        Generate a response using the provided context and sources with web search capability.
        Enhanced to handle URL content in context. Also Enhanced with accurate citation
        
        Args:
            query (str): User's original query
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            use_web_knowledge (bool): Whether to use web search in addition to documents
            response_format (str): Format style for the response
            conversation_context (str): Previous conversation history context
            
        Returns:
            str: Generated response based on the context and/or web search
        """
        # Check if the query is a simple greeting
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning", 
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya", 
            
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
        
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
        
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
        
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
        
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
        
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index], {}
        
        # If not a greeting, proceed with the regular flow
        if not context and not use_web_knowledge:
            return "I cannot answer this question based on the provided documents.", {}
        
        # Check for URL content in the context
        url_context = ""
        document_context = []
        document_sources = []
        url_sources = []
        
        # Separate URL content from document content
        for i, content_item in enumerate(context):
            if i < len(sources) and sources[i] == "URL Content" and "CONTENT FROM URLS IN QUERY:" in content_item:
                url_context = content_item
                url_sources.append(sources[i])
            else:
                document_context.append(content_item)
                if i < len(sources):
                    document_sources.append(sources[i])
        
        # Updated context preparation with citation mapping
        citation_sources = {}
        selected_context, selected_sources = self._prepare_context(document_context, document_sources)
        
        # Create context with source information
        contextualized_content = []
        for i, (content, source) in enumerate(zip(selected_context, selected_sources), 1):
            # Add citation number to each chunk
            citation_sources[i] = {
                "source_file": source,
                "text": content,
                "document_id": "Unknown",  # You might need to add logic to get the document ID
                "score": 1.0 - (i * 0.05),  # Simple scoring based on position (higher for earlier chunks)
                "display_num": i
            }
            contextualized_content.append(f"Source {i}:\n{content}")
            
        if url_context:
            citation_count = len(contextualized_content) + 1
            contextualized_content.append(f"Source {citation_count}:\n{url_context}")
            selected_sources.append("URL Content")
            citation_sources[citation_count] = {
                "source_file": "URL Content",
                "text": url_context,
                "document_id": "url_content",
                "score": 0.8,  # Fixed score for URL content
                "display_num": citation_count
            }

        full_context = "\n\n".join(contextualized_content)

        # Get format-specific guidance
        format_guidance = self._get_format_guidance(response_format, 'comprehensive')
        
        # Include conversation context if available
        conversation_prompt = ""
        if conversation_context:
            conversation_prompt = f"""
            RECENT CONVERSATION HISTORY:
            {conversation_context}
            
            The above messages are the previous parts of the same ongoing conversation with the user. 
            When generating your response:
            - Assume that the user may be referring back to earlier topics or questions from this history.
            - Resolve any references, clarifications, or follow-up questions based on this conversation history.
            - Maintain consistency with the user's earlier context, assumptions, and preferred response style if any.
            - If something from the conversation history clearly answers or helps answer the current question, incorporate it thoughtfully.
            - If there is ambiguity, prefer interpretations that align with prior conversation context.
            
            Please use this conversation history to maintain continuity, relevance, and coherence in your response.
            """

        # Get project description if available
        project_description = self._get_project_description(query)

        project_guidance = ""
        if project_description and len(project_description.strip()) > 5:
            project_guidance = f"""
            PROJECT CONTEXT:
            {project_description}

            IMPORTANT: This project context is only to help you understand the domain and purpose of this project. Your responses MUST be derived exclusively from the uploaded documents provided in the document context.

            When formulating your response:
            1. Use this project context solely to understand the general domain and focus of the project
            2. Draw ALL factual information, data, and specific content ONLY from the uploaded documents
            3. Shape your response style and tone to align with this project's domain, but never invent content
            4. DO NOT add information from your general knowledge even if relevant to the project description
            5. If the documents don't contain information relevant to the query, clearly state this limitation

            Remember: The project description helps you understand the context, but all response content must come from the uploaded documents.
            """

        # Updated system message for JSON response
        system_message = """You are a document analysis expert with conversation memory. 
        You must respond in JSON format with the following structure:
        {
            "content": "Your detailed response content here with proper HTML formatting and citations",
            "citations_used": [list of citation numbers used in the response],
            "format_type": "response_format_used",
            "sources_referenced": ["list", "of", "source", "names"]
        }
        
        In the content field, use HTML tags for formatting (<b>, <p>, <ul>, <li>, <table>, etc.).
        Use source citations [n] for all information from the documents.
        Provide a comprehensive and detailed answer while maintaining conversation continuity."""

        # Define the user prompt for JSON response
        user_prompt = f"""
        Based ONLY on the following context from multiple documents and URLs, answer the question. If relevant details are not fully available, provide the information that is present and kindly note any specific information that is missing. Be helpful by mentioning related information that may assist in answering the question, and offer to expand on available details if useful. Additionally, provide quantitative details where needed.

        RESPONSE FORMAT: {response_format.replace('_', ' ').title()}

        {format_guidance}

        {conversation_prompt}

        {project_guidance}

        RESPONSE GENERATION GUIDELINES:
        - PRIORITIZE the current user query over previous conversation context - focus primarily on answering the current question.
        - Previous conversation context should only inform your response when directly relevant to the current query.
        - Provide a DETAILED, COMPREHENSIVE, and THOROUGH answer.
        - Include ALL relevant information from the context in your response.
        - Prioritize completeness over brevity - make sure to include details.
        - If multiple perspectives or data points exist, include all of them.
        - When appropriate, organize information with clear sections and structure.
        - Maintain a natural, conversational tone.
        - Ensure the response is directly derived from the provided document context.
        - If URL content is present, specifically address content from URLs that is relevant to the query.
        - If the document does NOT contain relevant information, explicitly state that and summarize related available content instead.
        - Consider the conversation history to maintain continuity only when it adds valuable context to the current query.

        DOCUMENT AND URL CONTEXT:
        {full_context}

        USER QUERY: {query}

        CRITICAL: Your response MUST be a valid JSON object with this exact structure:
        {{
            "content": "Your comprehensive response with HTML formatting and [n] citations",
            "citations_used": [list of citation numbers used],
            "format_type": "{response_format}",
            "sources_referenced": ["list of source names referenced"]
        }}

        CRITICAL CONSTRAINTS:
        - Use ONLY information from the provided document and URL context.
        - DO NOT generate speculative or external information.
        - Your response must be based EXCLUSIVELY on the document and URL context provided.
        - Ensure clarity, accuracy, and comprehensive coverage of all relevant details.
        - Avoid summarizing or condensing important information - include all relevant details.
        - Use <b> tags for key section headings.
        - If the context contains URL content, clearly attribute information from URLs in your response.
        """
        
        try:
            # Call the OpenAI chat completion API with JSON format
            completion = client.chat.completions.create(
                model="o3-mini",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
              
               
            )
             # ===== LOG RAW LLM RESPONSE =====
            raw_llm_response = completion.choices[0].message.content
            print("\n" + "="*100)
            print("🤖 RAW LLM RESPONSE FROM OPENAI (generate_short_response)")
            print("="*100)
            print("Model used:", completion.model)
            print("Response ID:", completion.id)
            print("Raw response content:")
            print(raw_llm_response)
            print("Response length:", len(raw_llm_response))
            print("Response type:", type(raw_llm_response))
            print("="*100)
            # ===== END LOG RAW LLM RESPONSE =====
            # Parse the JSON response
            json_response = self._parse_json_response(completion.choices[0].message.content)
            
            if json_response.get("error"):
                # If JSON parsing failed, return the fallback message
                answer = json_response.get("content", "An error occurred while generating the response.")
            else:
                # Extract the content from JSON response
                answer = json_response.get("content", "No response content found.")

            # Clean and process citations in the answer
            answer = self._clean_citations(answer, citation_sources)
            answer = self._ensure_separate_citation_brackets(answer)
            processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
            
            # If answer seems too short, try to generate a more detailed one
            if len(answer.split()) < 100:  # If less than ~100 words
                enhanced_system_message = """You are a document analysis expert with conversation memory. 
                You must respond in JSON format with the following structure:
                {
                    "content": "Your EXTREMELY DETAILED and COMPREHENSIVE response with proper HTML formatting and citations",
                    "citations_used": [list of citation numbers used],
                    "format_type": "response_format_used",
                    "sources_referenced": ["list", "of", "source", "names"]
                }
                
                Provide an EXTREMELY DETAILED and COMPREHENSIVE response including ALL information from the context while maintaining conversational continuity. Use source citations [n] for all information from the documents."""
                
                completion = client.chat.completions.create(
                    model="o3-mini",
                    messages=[
                        {"role": "system", "content": enhanced_system_message},
                        {"role": "user", "content": user_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                   
                   
                )
                
                # Parse the enhanced JSON response
                json_response = self._parse_json_response(completion.choices[0].message.content)
                answer = json_response.get("content", answer)  # Use original if parsing fails
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
                
        except Exception as e:
            # If there's an error (e.g., token limit), try with fewer context chunks
            reduced_context = "\n\n".join(contextualized_content[:8])
            # Create a shorter version of conversation context if needed
            short_conv_context = ""
            if conversation_context:
                short_conv_context = f"Previous conversation context: {conversation_context[:300]}..."
                
            fallback_prompt = f"""
            Based ONLY on the following context, provide a comprehensive answer to the question: {query}
            
            {short_conv_context}
            
            CONTEXT:
            {reduced_context}
            
            CRITICAL: Respond in JSON format:
            {{
                "content": "Your detailed response with HTML formatting",
                "citations_used": [],
                "format_type": "fallback",
                "sources_referenced": []
            }}
            """
            
            try:
                fallback_completion = client.chat.completions.create(
                    model="o3-mini",
                    messages=[
                        {"role": "system", "content": "You are a document analysis expert. Respond in JSON format."},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    
                  
                )
                
                json_response = self._parse_json_response(fallback_completion.choices[0].message.content)
                answer = json_response.get("content", "An error occurred while generating the response.")
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
            except Exception as nested_e:
                # Last resort fallback
                answer = f"An error occurred while generating the response: {str(e)}. Please try a more specific question or with fewer documents."
                processed_answer = answer
                processed_citations = {}

        # Add source information
        source_list = list(set(selected_sources))
        source_info = ", ".join(source_list)
        return f"{processed_answer}\n\n*Sources: {source_info}*", processed_citations

    def generate_short_response(self, query, context, sources, use_web_knowledge=False, response_format='natural', conversation_context=""):
        """
        Generate a shorter, concise response using the provided context with web search capability.
        Enhanced to handle URL content in context.
        
        Args:
            query (str): User's original query
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            use_web_knowledge (bool): Whether to use web search in addition to documents
            response_format (str): Format style for the response
            conversation_context (str): Previous conversation history context
            
        Returns:
            str: Generated response based on the context and/or web search
        """
        # Check if the query is a simple greeting (same logic as above)
        greeting_keywords = [
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning", 
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            "hlw", "g'day", "yo", "heya", 
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
        
        query_lower = query.lower().strip()
        is_greeting = False
        
        if query_lower in greeting_keywords:
            is_greeting = True
        
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:
                is_greeting = True
                break
        
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index], {}
        
        if not context and not use_web_knowledge:
            return "I cannot answer this question based on the provided documents."
        
        # [Same URL and context processing logic as in generate_response]
        url_context = ""
        document_context = []
        document_sources = []
        url_sources = []
        
        for i, content_item in enumerate(context):
            if i < len(sources) and sources[i] == "URL Content" and "CONTENT FROM URLS IN QUERY:" in content_item:
                url_context = content_item
                url_sources.append(sources[i])
            else:
                document_context.append(content_item)
                if i < len(sources):
                    document_sources.append(sources[i])
        
        citation_sources = {}
        selected_context, selected_sources = self._prepare_context(document_context, document_sources)
        
        contextualized_content = []
        for i, (content, source) in enumerate(zip(selected_context, selected_sources), 1):
            citation_sources[i] = {
                "source_file": source,
                "text": content,
                "document_id": "Unknown",
                "score": 1.0 - (i * 0.05),
                "display_num": i
            }
            contextualized_content.append(f"Source {i}:\n{content}")

        if url_context:
            citation_count = len(contextualized_content) + 1
            contextualized_content.append(f"Source {citation_count}:\n{url_context}")
            selected_sources.append("URL Content")
            citation_sources[citation_count] = {
                "source_file": "URL Content",
                "text": url_context,
                "document_id": "url_content",
                "score": 0.8,
                "display_num": citation_count
            }

        full_context = "\n\n".join(contextualized_content)

        # Get format-specific guidance for short responses
        format_guidance = self._get_format_guidance(response_format, 'short')
        
        # Include conversation context if available (same as above)
        conversation_prompt = ""
        if conversation_context:
            conversation_prompt = f"""
            RECENT CONVERSATION HISTORY:
            {conversation_context}
            
            The above messages are the previous parts of the same ongoing conversation with the user. 
            When generating your response:
            - Assume that the user may be referring back to earlier topics or questions from this history.
            - Resolve any references, clarifications, or follow-up questions based on this conversation history.
            - Maintain consistency with the user's earlier context, assumptions, and preferred response style if any.
            - If something from the conversation history clearly answers or helps answer the current question, incorporate it thoughtfully.
            - If there is ambiguity, prefer interpretations that align with prior conversation context.
            
            Please use this conversation history to maintain continuity, relevance, and coherence in your response.
            """

        # Get project description if available
        project_description = self._get_project_description(query)
        
        project_guidance = ""
        if project_description and len(project_description.strip()) > 5:
            project_guidance = f"""
            PROJECT CONTEXT:
            {project_description}
            
            IMPORTANT: This project context is only to help you understand the domain and purpose of this project. Your responses MUST be derived exclusively from the uploaded documents provided in the document context.
            
            When formulating your response:
            1. Use this project context solely to understand the general domain and focus of the project
            2. Draw ALL factual information, data, and specific content ONLY from the uploaded documents
            3. Shape your response style and tone to align with this project's domain, but never invent content
            4. DO NOT add information from your general knowledge even if relevant to the project description
            5. If the documents don't contain information relevant to the query, clearly state this limitation
            
            Remember: The project description helps you understand the context, but all response content must come from the uploaded documents.
            """

        # Updated system message for JSON response
        system_message = """You are a document analysis expert with conversation memory. 
        You must respond in JSON format with the following structure:
        {
            "content": "Your concise yet informative response with proper HTML formatting and citations",
            "citations_used": [list of citation numbers used],
            "format_type": "response_format_used",
            "sources_referenced": ["list", "of", "source", "names"]
        }
        
        Provide a concise yet informative response that maintains conversation continuity. Use source citations [n] for information from the documents."""

        # Define the user prompt for short response with conversation context
        user_prompt = f"""
        Based ONLY on the following context from multiple documents and URLs, answer the question. If relevant details are not fully available, provide the information that is present and kindly note any specific information that is missing. Be helpful by mentioning related information that may assist in answering the question, and offer to expand on available details if useful. Additionally, provide quantitative details where needed.

        RESPONSE FORMAT: {response_format.replace('_', ' ').title()}

        {format_guidance}

        {conversation_prompt}

        {project_guidance}

        RESPONSE GENERATION GUIDELINES:
        - PRIORITIZE the current user query over previous conversation context - focus primarily on answering the current question.
        - Previous conversation context should only inform your response when directly relevant to the current query.
        - Provide a CONCISE yet THOROUGH answer.
        - Focus on the most relevant information from the context.
        - Prioritize clarity and directness over excessive detail.
        - When appropriate, organize information with clear sections and structure.
        - Maintain a natural, conversational tone.
        - Ensure the response is directly derived from the provided document and URL context.
        - If URL content is present, specifically address content from URLs that is relevant to the query.
        - If the document does NOT contain relevant information, explicitly state that.
        - Consider the conversation history to maintain continuity only when it adds valuable context to the current query.

        DOCUMENT AND URL CONTEXT:
        {full_context}

        USER QUERY: {query}

        CRITICAL: Your response MUST be a valid JSON object with this exact structure:
        {{
            "content": "Your concise response with HTML formatting and [n] citations",
            "citations_used": [list of citation numbers used],
            "format_type": "{response_format}",
            "sources_referenced": ["list of source names referenced"]
        }}

        CRITICAL CONSTRAINTS:
        - Use ONLY information from the provided document and URL context.
        - DO NOT generate speculative or external information.
        - Your response must be based EXCLUSIVELY on the document and URL context provided.
        - Ensure clarity, accuracy, and comprehensive coverage of key relevant details.
        - If the context contains URL content, clearly attribute information from URLs in your response.
        """
        
        try:
            # Call the OpenAI chat completion API for short response with JSON format
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
                temperature=0.4,
                max_tokens=2000
            )
                # ===== LOG RAW LLM RESPONSE =====
            raw_llm_response = completion.choices[0].message.content
            print("\n" + "="*100)
            print("🤖 RAW LLM RESPONSE FROM OPENAI (generate_short_response)")
            print("="*100)
            print("Model used:", completion.model)
            print("Response ID:", completion.id)
            print("Raw response content:")
            print(raw_llm_response)
            print("Response length:", len(raw_llm_response))
            print("Response type:", type(raw_llm_response))
            print("="*100)
            # ===== END LOG RAW LLM RESPONSE =====
            # Parse the JSON response
            json_response = self._parse_json_response(completion.choices[0].message.content)
            answer = json_response.get("content", "No response content found.")
            answer = self._clean_citations(answer, citation_sources)
            answer = self._ensure_separate_citation_brackets(answer)
            processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
                
        except Exception as e:
            # If there's an error, try with even fewer context chunks
            reduced_context = "\n\n".join(contextualized_content[:5])
            fallback_prompt = f"""
            Based ONLY on the following context, provide a brief answer to the question: {query}
            
            CONTEXT:
            {reduced_context}
            
            CRITICAL: Respond in JSON format:
            {{
                "content": "Your brief response with HTML formatting",
                "citations_used": [],
                "format_type": "fallback",
                "sources_referenced": []
            }}
            """
            
            try:
                fallback_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are a document analysis expert. Be brief. Respond in JSON format."},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.4,
                    max_tokens=1024
                )
                
                json_response = self._parse_json_response(fallback_completion.choices[0].message.content)
                answer = json_response.get("content", "An error occurred while generating the response.")
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
            except Exception as nested_e:
                # Last resort fallback
                answer = f"An error occurred while generating the response. Please try a more specific question."
                processed_answer = answer
                processed_citations = {}

        # Add source information
        source_list = list(set(selected_sources))
        source_info = ", ".join(source_list)
        return f"{processed_answer}\n\n*Sources: {source_info}*", processed_citations

    def _get_project_description(self, query):
        """
        Get the project description for the current main project
        
        Returns:
            str: Project description or empty string if not found
        """
        try:
            # Get the main_project_id from the current request
            main_project_id = self.request.data.get('main_project_id')
            
            if not main_project_id:
                return ""
                
            # Get the project from the database
            project = Project.objects.filter(id=main_project_id).first()
            
            if project and project.description:
                return project.description
                
            return ""
        except Exception as e:
            logger.error(f"Error getting project description: {str(e)}")
            return ""
    

    def get_general_chat_answer(self, query, use_web_knowledge=False, response_length='comprehensive', response_format='natural', user=None):
        """
        Generate an answer for general chat mode, optionally using web knowledge.
    
        Args:
            query (str): The user's query
            use_web_knowledge (bool): Whether to use web search
            response_length (str): 'short' or 'comprehensive'
            response_format (str): The format to use for the response
        
        Returns:
            str: Generated response
        """
        # Check if the query is a simple greeting
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
        
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
        
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya",
        
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
        
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
        
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
        
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
    
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
    
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
    
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
    
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
    
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index]
    
        # If not a greeting, proceed with the regular flow
        # If web knowledge is enabled, implement the full web search flow
        if use_web_knowledge:
            try:
                # Step 1: Search the web using DuckDuckGo
                print(f"Searching web for: {query}")
                web_response, web_sources = self.get_web_knowledge_response(query, user=user, document_context=None)

                if not web_response or not web_sources:
                    return "I couldn't find relevant information on the web for your query."
                
                # Extract web source domains for display
                web_source_domains = []
                for source in web_sources:
                    domain = self.extract_domain(source.get('url', ''))
                    if domain and domain != "unknown.com":
                        web_source_domains.append(domain)
            
                # Remove existing sources from web_response to avoid duplication in GPT processing
                if "*Sources:" in web_response:
                    web_response_clean = web_response.split("*Sources:")[0].strip()
                else:
                    web_response_clean = web_response

                # Get format-specific guidance
                format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
            
                # Updated system message for JSON response
                system_message = """You are a web research assistant. 
                You must respond in JSON format with the following structure:
                {
                    "content": "Your response with proper HTML formatting",
                    "web_sources_used": true,
                    "response_type": "web_enhanced",
                    "format_type": "response_format_used"
                }
                
                Provide detailed, comprehensive responses based on web sources."""

                # Create the prompt for OpenAI with JSON requirement
                prompt = f"""
                You are a web research assistant. Based ONLY on the following information from multiple web sources, answer the user question.

                If relevant details are missing or contradictory, clearly acknowledge this and summarize available related content. Be helpful by highlighting potentially useful information even if it doesn't fully resolve the question. Provide quantitative details where available.

                RESPONSE FORMAT: {response_format.replace('_', ' ').title()}

                {format_guidance}

                RESPONSE GENERATION GUIDELINES:
                - Provide a DETAILED, COMPREHENSIVE, and THOROUGH answer.
                - Include ALL relevant information from the web sources below.
                - Prioritize completeness over brevity — include important details.
                - If multiple perspectives or data points exist, include all of them.
                - When appropriate, structure the information clearly with headings.
                - Maintain a natural, conversational tone.
                - If the web sources do NOT contain relevant information, clearly state that and summarize any related or tangential insights.

                QUESTION: {query}

                INFORMATION FROM WEB SOURCES:
                {web_response_clean}

                CRITICAL: Your response MUST be a valid JSON object with this exact structure:
                {{
                    "content": "Your response with HTML tags for structure and formatting",
                    "web_sources_used": true,
                    "response_type": "web_enhanced",
                    "format_type": "{response_format}"
                }}

                RESPONSE FORMAT REQUIREMENTS:
                1. Follow with detailed elaboration on each relevant point.
                2. Include specific details, examples, and support from the source content.
                3. Use clear, logical sections when needed.
                4. Conclude with any other contextually relevant insights.
                5. Use <b> for section headers.
                6. Use <p> for body text.
                7. Use <ul> / <li> for lists.
                8. Use <table>, <tr>, <td> for tables, if applicable.
                9. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.

                CRITICAL INSTRUCTION - MUST BE FOLLOWED:
                At the very end of your content (within the JSON content field), after all content, add a new line and provide ONLY this exact format:
                SOURCES_JSON: {{"sources": ["domain1.com", "domain2.com", "domain3.com"]}}

                Extract the actual website domains mentioned in the web sources content. Do not include generic domains. Look for domain names that appear in the source content.

                CRITICAL CONSTRAINTS:
                - Ensure clarity, accuracy, and full coverage of the relevant content.

                RESPONSE LENGTH: {'Provide a focused, concise response prioritizing the most important information.' if response_length == 'short' else 'Provide a comprehensive, detailed response that thoroughly covers the topic.'}
                """

                temperature = 0.3 if response_format in ['factual_brief', 'technical_deep_dive'] else 0.5
                max_tokens = 2000 if response_length == 'short' else 2000

                print(f"Calling OpenAI API with temperature={temperature}, max_tokens={max_tokens}")
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=temperature,
                    max_tokens=max_tokens
                )

                # Parse the JSON response
                json_response = self._parse_json_response(response.choices[0].message.content)
                formatted_web_response = json_response.get("content", "No response content found.")

                # Debug: Print the full response to see what GPT returned
                print("=== FULL GPT RESPONSE ===")
                print(formatted_web_response)
                print("=== END FULL RESPONSE ===")

                # Extract JSON sources from the response
                import json
                import re
                sources_extracted = []

                if "SOURCES_JSON:" in formatted_web_response:
                    try:
                        # Split the response to get the JSON part
                        parts = formatted_web_response.split("SOURCES_JSON:")
                        main_response = parts[0].strip()
                        json_part = parts[1].strip()
                    
                        print(f"=== JSON PART EXTRACTED ===")
                        print(f"JSON part: '{json_part}'")
                        print("=== END JSON PART ===")
                    
                        # Parse the JSON
                        sources_data = json.loads(json_part)
                        sources_extracted = sources_data.get("sources", [])
                    
                        print(f"=== EXTRACTED SOURCES ===")
                        print(f"Sources: {sources_extracted}")
                        print("=== END EXTRACTED SOURCES ===")
                    
                        # Use the main response without the JSON part
                        formatted_web_response = main_response
                    
                    except (json.JSONDecodeError, IndexError) as e:
                        logger.warning(f"Failed to parse sources JSON: {e}")
                        print(f"=== JSON PARSING ERROR ===")
                        print(f"Error: {e}")
                        print(f"Raw JSON part: '{json_part if 'json_part' in locals() else 'Not extracted'}'")
                        print("=== END JSON ERROR ===")
                    
                        # Manual extraction as fallback - look for patterns in the response
                        sources_extracted = []
                    
                        # Look for domain patterns in the response using regex
                        domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                        domains_found = re.findall(domain_pattern, formatted_web_response.lower())
                    
                        # Filter out common generic domains
                        generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                        for domain in domains_found:
                            if domain not in generic_domains and domain not in sources_extracted:
                                sources_extracted.append(domain)
                    
                        print(f"=== MANUAL EXTRACTION RESULTS ===")
                        print(f"Manually extracted sources: {sources_extracted}")
                        print("=== END MANUAL EXTRACTION ===")
                    
                else:
                    print("=== NO SOURCES_JSON FOUND ===")
                    print("SOURCES_JSON marker not found in response")
                
                    # Manual extraction as fallback - look for domain patterns
                    sources_extracted = []
                
                    # Look for domain patterns in the response using regex
                    domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                    domains_found = re.findall(domain_pattern, formatted_web_response.lower())
                
                    # Filter out common generic domains
                    generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                    for domain in domains_found:
                        if domain not in generic_domains and domain not in sources_extracted:
                            sources_extracted.append(domain)
                
                    print(f"=== FALLBACK EXTRACTION RESULTS ===")
                    print(f"Fallback extracted sources: {sources_extracted}")
                    print("=== END FALLBACK EXTRACTION ===")

                # Use extracted sources if available, otherwise fall back to web_source_domains
                if sources_extracted:
                    source_info = ", ".join(sources_extracted)
                    final_response = f"{formatted_web_response}\n\n*Sources: {source_info}*"
                    print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", final_response)
                    return final_response
                else:
                    # Fallback to existing method
                    if web_source_domains:
                        source_info = ", ".join(web_source_domains)
                        final_response = f"{formatted_web_response}\n\n*Sources: {source_info}*"
                        print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", final_response)
                        return final_response
                    else:
                        return formatted_web_response
            
            except Exception as e:
                logger.error(f"Error in web knowledge general chat: {str(e)}", exc_info=True)
                return f"I encountered an error while searching the web: {str(e)}. Please try a different question or try again later."
    
        # If no web knowledge requested, use standard chat completion
        else:
            # Format-specific guidance
            format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
        
            # Configure conciseness based on response length
            conciseness = "Be concise and to-the-point." if response_length == 'short' else "Provide a comprehensive and detailed response."
        
            # Updated system message for JSON response
            system_message = """You are a helpful, friendly AI assistant. 
            You must respond in JSON format with the following structure:
            {
                "content": "Your response with proper HTML formatting",
                "response_type": "general_chat",
                "format_type": "response_format_used",
                "web_sources_used": false
            }
            
            Provide informative and thoughtful responses using HTML tags for structure."""

            # Create a prompt for the general chat with JSON requirement
            prompt = f"""
            Please answer the following question in a helpful, informative way.
            
            RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
            
            {format_guidance}
            
            {conciseness}
            
            Use HTML tags for structure (<b>, <p>, <ul>, <li>) to enhance readability.
            
            USER QUERY: {query}
            
            CRITICAL: Your response MUST be a valid JSON object with this exact structure:
            {{
                "content": "Your response with HTML formatting",
                "response_type": "general_chat",
                "format_type": "{response_format}",
                "web_sources_used": false
            }}
            """
        
            try:
                # Use the OpenAI API with JSON format
                completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.5,
                    max_tokens=2000 if response_length == 'comprehensive' else 800
                )
                
                # Parse the JSON response
                json_response = self._parse_json_response(completion.choices[0].message.content)
                return json_response.get("content", "No response content found.")
            
            except Exception as e:
                logger.error(f"Error in general chat: {str(e)}", exc_info=True)
                return f"I'm sorry, I encountered an error while processing your request. Please try again or rephrase your question."

    # Enhanced embeddings function from Streamlit version
    def get_embeddings(self, texts):
        """
        Gets embeddings using OpenAI's text-embedding-3-small model.
        """
        try:
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small"
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            logger.error(f"Error getting embeddings: {str(e)}")
            return []
    
    # # Improved search_similar_content with TF-IDF ranking from Streamlit
    # def search_similar_content(self, query, processed_docs, metadata_store, k=40):
    #     """
    #     Enhanced search function that handles both LlamaParse and local document parsing
    #     """
    #     # Get embeddings for the query
    #     query_embedding = self.get_embeddings([query])
    #     if not query_embedding:
    #         return [], [], {}  # Return empty citation mapping to
        
    #     # Search each document's FAISS index separately
    #     all_results = []
    #     all_distances = []
    #     all_sources = []
    #     citation_mapping = {}  # Add citation mapping
    #     citation_count = 0
        
    #     # Track content to avoid duplicates
    #     seen_content_hashes = set()
        
    #     # Check if we have any valid documents to search
    #     valid_docs_found = False
        
    #     for proc_doc in processed_docs:
    #         # Skip documents without FAISS indices
    #         if not proc_doc.faiss_index or not proc_doc.metadata:
    #             continue
            
    #         # Skip if files don't exist
    #         if not os.path.exists(proc_doc.faiss_index) or not os.path.exists(proc_doc.metadata):
    #             continue
            
    #         # First check if this is a LlamaParse document with markdown
    #         markdown_content = None
    #         if proc_doc.markdown_path and os.path.exists(proc_doc.markdown_path):
    #             try:
    #                 with open(proc_doc.markdown_path, 'r', encoding='utf-8') as f:
    #                     markdown_content = f.read()
    #                 # If we have markdown content, we'll use it later in our results
    #                 valid_docs_found = True
    #             except Exception as e:
    #                 logger.error(f"Error reading markdown file for {proc_doc.document.filename}: {str(e)}")
            
    #         # Load metadata for both document types
    #         try:
    #             with open(proc_doc.metadata, 'rb') as f:
    #                 chunks = pickle.load(f)
                
    #             # Make sure chunks is a list - sometimes it might be empty or None
    #             if not chunks:
    #                 logger.warning(f"No chunks found in metadata for {proc_doc.document.filename}")
    #                 chunks = []
    #             elif not isinstance(chunks, list):
    #                 logger.warning(f"Unexpected chunks format for {proc_doc.document.filename}: {type(chunks)}")
    #                 chunks = []
                
    #             # Add document source information to each chunk if missing
    #             for chunk in chunks:
    #                 if isinstance(chunk, dict) and 'source_file' not in chunk:
    #                     chunk['source_file'] = proc_doc.document.filename
    #                 # Ensure 'text' field exists in chunk
    #                 if isinstance(chunk, dict) and not chunk.get('text'):
    #                     # Try alternate field names
    #                     for field in ['content', 'document_content', 'chunk_text']:
    #                         if field in chunk:
    #                             chunk['text'] = chunk[field]
    #                             break
                
    #             # Try vector search with FAISS index
    #             valid_docs_found = True
    #             try:
    #                 index = faiss.read_index(proc_doc.faiss_index)
    #                 # Search this index with increased k for better diversity
    #                 distances, indices = index.search(np.array([query_embedding[0]]).astype('float32'), min(k, index.ntotal))
                    
    #                 # Extract results for this document
    #                 for i, idx in enumerate(indices[0]):
    #                     if idx < len(chunks):
    #                         content = chunks[idx].get('text', '')
    #                         # Only add non-empty content
    #                         if content and content.strip():
    #                             # Add numbered citation
    #                             citation_count += 1
    #                             citation_key = citation_count

    #                             # Store citation mapping with metadata
    #                             citation_mapping[citation_key] = {
    #                                 'source': proc_doc.document.filename,
    #                                 'text': content,
    #                                 'relevance_score': float(distances[0][i]),
    #                                 'document_id': proc_doc.document.id,
    #                                 'chunk_idx': idx
    #                             }
    #                             all_results.append(content)
    #                             all_distances.append(distances[0][i])
    #                             all_sources.append(proc_doc.document.filename)
    #             except Exception as e:
    #                 logger.error(f"Error searching FAISS index for {proc_doc.document.filename}: {str(e)}")
                    
    #                 # Fallback to basic text search for simple documents if vector search fails
    #                 if chunks:
    #                     # Use TF-IDF for matching if FAISS fails
    #                     query_lower = query.lower()
    #                     matched_chunks = []
                        
    #                     for chunk in chunks:
    #                         if isinstance(chunk, dict):
    #                             chunk_text = chunk.get('text', '')
    #                             if chunk_text and query_lower in chunk_text.lower():
    #                                 matched_chunks.append(chunk)
                        
    #                     # Sort by simple text similarity
    #                     if matched_chunks:
    #                         logger.info(f"Found {len(matched_chunks)} text matches for {proc_doc.document.filename}")
    #                         for chunk in matched_chunks[:5]:  # Limit to top 5 matches
    #                             all_results.append(chunk.get('text', ''))
    #                             all_distances.append(0.5)  # Middle value since we don't have real distances
    #                             all_sources.append(proc_doc.document.filename)
                
    #             # If we have markdown content but no vector results, do fallback text search on whole content
    #             if markdown_content and not all_results:
    #                 query_lower = query.lower()
    #                 if query_lower in markdown_content.lower():
    #                     # Split into paragraphs and find matching ones
    #                     paragraphs = markdown_content.split('\n\n')
    #                     for para in paragraphs:
    #                         if query_lower in para.lower():
    #                             all_results.append(para)
    #                             all_distances.append(0.5)  # Middle value since we don't have real distances
    #                             all_sources.append(proc_doc.document.filename)
                    
    #         except Exception as e:
    #             logger.error(f"Error processing metadata for {proc_doc.document.filename}: {str(e)}")
    #             continue
        
    #     # Combine and sort results by similarity score
    #     if not all_results:
    #         if valid_docs_found:
    #             logger.warning(f"No matching content found in documents for query: {query}")
    #         else:
    #             logger.warning(f"No valid documents found to search")
    #         return [], []
            
    #     # Apply TF-IDF ranking using the Streamlit approach
    #     try:
    #         from sklearn.feature_extraction.text import TfidfVectorizer
            
    #         # Create a TF-IDF vectorizer with increased max_features
    #         vectorizer = TfidfVectorizer(stop_words='english', max_features=100)
    #         try:
    #             tfidf_matrix = vectorizer.fit_transform([query] + all_results)
                
    #             # Calculate similarity scores
    #             tfidf_scores = tfidf_matrix[0].dot(tfidf_matrix.T).toarray().flatten()[1:]
                
    #             # Convert embedding distances to similarity scores
    #             similarity_scores = [1 / (1 + dist) for dist in all_distances]
                
    #             # Combine scores with same weighting (70% TF-IDF, 30% embedding)
    #             combined_scores = [0.7 * tf + 0.3 * sim for tf, sim in zip(tfidf_scores, similarity_scores)]
                
    #             # Re-sort results by combined score
    #             combined_results = list(zip(all_results, all_sources, combined_scores))
    #             sorted_results = sorted(combined_results, key=lambda x: x[2], reverse=True)
                
    #             # Extract sorted content and sources
    #             results = [res[0] for res in sorted_results]
    #             sources = [res[1] for res in sorted_results]

    #              # Update citation mapping with new ranking
    #             reordered_citation_mapping = {}
    #             for new_idx, (_, _, _, old_key) in enumerate(sorted_results, 1):
    #                 if old_key in citation_mapping:
    #                     reordered_citation_mapping[new_idx] = citation_mapping[old_key]
    #                     reordered_citation_mapping[new_idx]['display_num'] = new_idx
                
    #             citation_mapping = reordered_citation_mapping
    #         except Exception as e:
    #             logger.error(f"Error in TF-IDF processing: {str(e)}")
    #             # Fallback to original results if TF-IDF fails
    #             results = all_results
    #             sources = all_sources
    #     except Exception as e:
    #         logger.error(f"Error applying TF-IDF ranking: {str(e)}")
    #         # Fall back to basic distance-based ranking
    #         combined_results = list(zip(all_results, all_sources, all_distances))
    #         # Note: In the fallback, we sort by distance (smaller is better) so no reverse=True
    #         sorted_results = sorted(combined_results, key=lambda x: x[2])
            
    #         results = [res[0] for res in sorted_results]
    #         sources = [res[1] for res in sorted_results]

    #         # Update citation mapping with new ranking
    #         reordered_citation_mapping = {}
    #         for new_idx, (_, _, _, old_key) in enumerate(sorted_results, 1):
    #             if old_key in citation_mapping:
    #                 reordered_citation_mapping[new_idx] = citation_mapping[old_key]
    #                 reordered_citation_mapping[new_idx]['display_num'] = new_idx
            
    #         citation_mapping = reordered_citation_mapping
        
    #     # Remove duplicates while preserving order
    #     seen_content = set()
    #     filtered_results = []
    #     filtered_sources = []
    #     filtered_citation_mapping = {}
    #     new_citation_count = 0
        
    #     for content, source in zip(results, sources):
    #         # Use first 100 chars as a content signature
    #         content_hash = content[:100] if content else ""
    #         if content_hash and content_hash not in seen_content:
    #             seen_content.add(content_hash)
    #             filtered_results.append(content)
    #             filtered_sources.append(source)

    #             # Update citation mapping
    #             new_citation_count += 1
    #             old_idx = i + 1
    #             if old_idx in citation_mapping:
    #                 filtered_citation_mapping[new_citation_count] = citation_mapping[old_idx]
    #                 filtered_citation_mapping[new_citation_count]['display_num'] = new_citation_count
        
    #     # Return top matches (limit to 15 most relevant)
    #     max_results = min(len(filtered_results), 15)
    #     return filtered_results[:max_results], filtered_sources[:max_results], filtered_citation_mapping
    
    def search_similar_content(self, query, processed_docs, metadata_store, k=40):
        """
        Enhanced search function that handles both LlamaParse and local document parsing
        with improved citation deduplication
        """
        # Get embeddings for the query
        query_embedding = self.get_embeddings([query])
        if not query_embedding:
            return [], [], {}
        
        # Search each document's FAISS index separately
        all_results = []
        all_distances = []
        all_sources = []
        citation_mapping = {}
        citation_count = 0
        
        # Track content to avoid duplicates
        seen_content_hashes = set()
        
        # Check if we have any valid documents to search
        valid_docs_found = False
        
        for proc_doc in processed_docs:
            # Skip documents without FAISS indices
            if not proc_doc.faiss_index or not proc_doc.metadata:
                continue
            
            # Skip if files don't exist
            if not os.path.exists(proc_doc.faiss_index) or not os.path.exists(proc_doc.metadata):
                continue
            
            # First check if this is a LlamaParse document with markdown
            markdown_content = None
            if proc_doc.markdown_path and os.path.exists(proc_doc.markdown_path):
                try:
                    with open(proc_doc.markdown_path, 'r', encoding='utf-8') as f:
                        markdown_content = f.read()
                    valid_docs_found = True
                except Exception as e:
                    logger.error(f"Error reading markdown file for {proc_doc.document.filename}: {str(e)}")
            
            # Load metadata for both document types
            try:
                with open(proc_doc.metadata, 'rb') as f:
                    chunks = pickle.load(f)
                
                if not chunks:
                    logger.warning(f"No chunks found in metadata for {proc_doc.document.filename}")
                    chunks = []
                elif not isinstance(chunks, list):
                    logger.warning(f"Unexpected chunks format for {proc_doc.document.filename}: {type(chunks)}")
                    chunks = []
                
                # Add document source information to each chunk if missing
                for chunk in chunks:
                    if isinstance(chunk, dict) and 'source_file' not in chunk:
                        chunk['source_file'] = proc_doc.document.filename
                    if isinstance(chunk, dict) and not chunk.get('text'):
                        for field in ['content', 'document_content', 'chunk_text']:
                            if field in chunk:
                                chunk['text'] = chunk[field]
                                break
                
                # Try vector search with FAISS index
                valid_docs_found = True
                try:
                    index = faiss.read_index(proc_doc.faiss_index)
                    distances, indices = index.search(np.array([query_embedding[0]]).astype('float32'), min(k, index.ntotal))
                    
                    # Extract results for this document
                    for i, idx in enumerate(indices[0]):
                        if idx < len(chunks):
                            content = chunks[idx].get('text', '')
                            if content and content.strip():
                                # Create content hash to check for duplicates
                                content_hash = hash(content[:200])  # Use first 200 chars for hash
                                
                                if content_hash not in seen_content_hashes:
                                    seen_content_hashes.add(content_hash)
                                    
                                    # Add numbered citation
                                    citation_count += 1
                                    citation_key = citation_count

                                    # Store citation mapping with metadata
                                    citation_mapping[citation_key] = {
                                        'source': proc_doc.document.filename,
                                        'text': content,
                                        'relevance_score': float(distances[0][i]),
                                        'document_id': proc_doc.document.id,
                                        'chunk_idx': idx
                                    }
                                    all_results.append(content)
                                    all_distances.append(distances[0][i])
                                    all_sources.append(proc_doc.document.filename)
                                    
                except Exception as e:
                    logger.error(f"Error searching FAISS index for {proc_doc.document.filename}: {str(e)}")
                    # Fallback logic here...
                    # Fallback to basic text search for simple documents if vector search fails
                    if chunks:
                        # Use TF-IDF for matching if FAISS fails
                        query_lower = query.lower()
                        matched_chunks = []
                        
                        for chunk in chunks:
                            if isinstance(chunk, dict):
                                chunk_text = chunk.get('text', '')
                                if chunk_text and query_lower in chunk_text.lower():
                                    matched_chunks.append(chunk)
                        
                        # Sort by simple text similarity
                        if matched_chunks:
                            logger.info(f"Found {len(matched_chunks)} text matches for {proc_doc.document.filename}")
                            for chunk in matched_chunks[:5]:  # Limit to top 5 matches
                                all_results.append(chunk.get('text', ''))
                                all_distances.append(0.5)  # Middle value since we don't have real distances
                                all_sources.append(proc_doc.document.filename)
                
                # If we have markdown content but no vector results, do fallback text search on whole content
                if markdown_content and not all_results:
                    query_lower = query.lower()
                    if query_lower in markdown_content.lower():
                        # Split into paragraphs and find matching ones
                        paragraphs = markdown_content.split('\n\n')
                        for para in paragraphs:
                            if query_lower in para.lower():
                                all_results.append(para)
                                all_distances.append(0.5)  # Middle value since we don't have real distances
                                all_sources.append(proc_doc.document.filename)
                    
            except Exception as e:
                logger.error(f"Error processing metadata for {proc_doc.document.filename}: {str(e)}")
                continue
        
        # Rest of the method remains the same but with improved deduplication
        if not all_results:
            if valid_docs_found:
                logger.warning(f"No matching content found in documents for query: {query}")
            else:
                logger.warning(f"No valid documents found to search")
            return [], [], {}
        
        # Apply TF-IDF ranking and final deduplication
        try:
            from sklearn.feature_extraction.text import TfidfVectorizer
            
            vectorizer = TfidfVectorizer(stop_words='english', max_features=100)
            try:
                tfidf_matrix = vectorizer.fit_transform([query] + all_results)
                tfidf_scores = tfidf_matrix[0].dot(tfidf_matrix.T).toarray().flatten()[1:]
                similarity_scores = [1 / (1 + dist) for dist in all_distances]
                combined_scores = [0.7 * tf + 0.3 * sim for tf, sim in zip(tfidf_scores, similarity_scores)]
                
                # Combine results with citation keys for proper tracking
                combined_results = list(zip(all_results, all_sources, combined_scores, range(1, len(all_results) + 1)))
                sorted_results = sorted(combined_results, key=lambda x: x[2], reverse=True)
                
                results = [res[0] for res in sorted_results]
                sources = [res[1] for res in sorted_results]
                
                # Update citation mapping with new ranking
                reordered_citation_mapping = {}
                for new_idx, (_, _, _, old_key) in enumerate(sorted_results, 1):
                    if old_key in citation_mapping:
                        reordered_citation_mapping[new_idx] = citation_mapping[old_key]
                        reordered_citation_mapping[new_idx]['display_num'] = new_idx
                
                citation_mapping = reordered_citation_mapping
                
            except Exception as e:
                logger.error(f"Error in TF-IDF processing: {str(e)}")
                results = all_results
                sources = all_sources
                
        except Exception as e:
            logger.error(f"Error applying TF-IDF ranking: {str(e)}")
            # Fallback processing...
            results = all_results
            sources = all_sources
        
        # Final deduplication while preserving order
        seen_content = set()
        filtered_results = []
        filtered_sources = []
        filtered_citation_mapping = {}
        new_citation_count = 0
        
        for i, (content, source) in enumerate(zip(results, sources)):
            content_hash = content[:100] if content else ""
            if content_hash and content_hash not in seen_content:
                seen_content.add(content_hash)
                filtered_results.append(content)
                filtered_sources.append(source)

                new_citation_count += 1
                old_idx = i + 1
                if old_idx in citation_mapping:
                    filtered_citation_mapping[new_citation_count] = citation_mapping[old_idx]
                    filtered_citation_mapping[new_citation_count]['display_num'] = new_citation_count
        
        # Return top matches (limit to 15 most relevant)
        max_results = min(len(filtered_results), 40)
        return filtered_results[:max_results], filtered_sources[:max_results], filtered_citation_mapping
    
    # Helper method to prepare context for response generation
    def _prepare_context(self, context, sources):
        """
        Helper method to prepare context for response generation
        
        Args:
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            
        Returns:
            tuple: (selected_context, selected_sources)
        """
        # Limit number of tokens by selecting most relevant chunks
        selected_context = []
        selected_sources = []
        seen_content_hashes = set()
        
        # Take first 8 chunks (most relevant ones)
        initial_count = min(8, len(context))
        for i in range(initial_count):
            selected_context.append(context[i])
            selected_sources.append(sources[i])
            # Add a content hash to avoid repetition
            seen_content_hashes.add(context[i][:100])
        
        # Add more diverse chunks from the remaining context
        for i in range(initial_count, len(context)):
            content_hash = context[i][:100]
            # Check if content is sufficiently different from what we've seen
            if content_hash not in seen_content_hashes:
                selected_context.append(context[i])
                selected_sources.append(sources[i])
                seen_content_hashes.add(content_hash)
                # Limit to 15 total pieces of context for token limit reasons
                if len(selected_context) >= 15:
                    break
        
        return selected_context, selected_sources
            
    # Keep general follow-up question generation as is
    def generate_general_follow_up_questions(self, question, answer, user=None):
        system_message = "You are an expert at generating relevant follow-up questions. Always respond with valid JSON format."
       
        user_prompt = f"""
        Based on this user question and your answer, suggest 3 relevant follow-up questions that the user might want to ask next.
        The questions should be short, interesting, and directly related to the topic.
       
        User Question: {question}
        Your Answer (abbreviated): {answer[:500]}...
       
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "questions": ["question1", "question2", "question3"],
            "topic": "main_topic_discussed"
        }}
        """
       
        try:            
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.3,
                max_tokens=500
            )
           
            # Parse JSON response
            json_response = json.loads(completion.choices[0].message.content)
            questions = json_response.get("questions", [])
            return questions[:3]
           
        except Exception as e:
            logger.error(f"Error generating follow-up questions: {str(e)}")
            return [
                "What else would you like to know about this topic?",
                "Would you like me to elaborate on any specific point?",
                "Do you have any other questions I can help with?"
            ]
 
    def generate_follow_up_questions(self, context, user=None):
        context_sample = "\n".join(context[:3]) if context else ""
       
        system_message = "You are an expert at generating relevant follow-up questions based on context. Always respond with valid JSON format."
       
        user_prompt = f"""
        Based on this context, suggest 3 relevant follow-up questions. The questions should be short and directly related to the content.
       
        Context: {context_sample}
       
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "questions": ["question1", "question2", "question3"],
            "context_topic": "main_topic_from_context"
        }}
        """
       
        try:
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.3,
                max_tokens=500
            )
           
            # Parse JSON response
            json_response = json.loads(completion.choices[0].message.content)
            questions = json_response.get("questions", [])
            return questions[:3]
           
        except Exception as e:
            logger.error(f"Error generating follow-up questions: {str(e)}")
            return [
                "What additional information would you like to know?",
                "Would you like me to elaborate on any specific point?",
                "How can I help clarify this information further?"
            ]

class DeleteMessagePairView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        print("DeleteMessagePairView POST data:", request.data)
        user = request.user
        user_message_id = request.data.get('user_message_id')
        assistant_message_id = request.data.get('assistant_message_id')
        conversation_id = request.data.get('conversation_id')


        if not (user_message_id and assistant_message_id and conversation_id):
            return Response({"error": "Missing required parameters."}, status=status.HTTP_400_BAD_REQUEST)

        try:
            with transaction.atomic():
                # Ensure both messages belong to the user and conversation
                chat_history = ChatHistory.objects.get(conversation_id=conversation_id, user=user)
                user_msg = ChatMessage.objects.get(id=user_message_id, chat_history=chat_history, role="user")
                assistant_msg = ChatMessage.objects.get(id=assistant_message_id, chat_history=chat_history, role="assistant")

                user_msg.delete()
                assistant_msg.delete()

            return Response({"success": True}, status=status.HTTP_200_OK)
        except ChatHistory.DoesNotExist:
            return Response({"error": "Conversation not found."}, status=status.HTTP_404_NOT_FOUND)
        except ChatMessage.DoesNotExist:
            return Response({"error": "Message not found."}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR) 

class GetConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request, conversation_id=None):
        try:
            user = request.user
            
            if conversation_id:
                # Fetch specific conversation
                try:
                    conversation = ChatHistory.objects.get(
                        conversation_id=conversation_id, 
                        user=user
                    )
                    
                    # Retrieve all messages for this conversation
                    messages = conversation.messages.all().order_by('created_at')
                    
                    # Prepare message list - UPDATED to include all badge properties
                   

                    message_list = []
                    for message in messages:
                        message_data = {
                            'id': message.id,
                            'role': message.role,
                            'content': message.content,
                            'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': message.citations or []
                        }
                       
                        if message.role == 'assistant':
                            message_data['use_web_knowledge'] = getattr(message, 'use_web_knowledge', False)
                            message_data['response_length'] = getattr(message, 'response_length', 'comprehensive')
                            message_data['response_format'] = getattr(message, 'response_format', 'natural')
                            message_data['sources_info'] = getattr(message, 'sources', "")  # ✅ NEW LINE
                       
                        message_list.append(message_data)
 
                    
                    return Response({
                        'conversation_id': str(conversation.conversation_id),
                        'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                        'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                        'messages': message_list,
                        'follow_up_questions': conversation.follow_up_questions or [],
                        'selected_documents': [
                            str(doc.id) for doc in conversation.documents.all()
                        ]
                    }, status=status.HTTP_200_OK)
                
                except ChatHistory.DoesNotExist:
                    return Response(
                        {'error': 'Conversation not found'}, 
                        status=status.HTTP_404_NOT_FOUND
                    )
            
            else:
                # Fetch all conversations for the user
                conversations = ChatHistory.objects.filter(user=user) \
                    .filter(is_active=True) \
                    .order_by('-created_at')
                
                conversation_list = []
                for conversation in conversations:
                    # Get the first and last messages
                    messages = conversation.messages.all().order_by('created_at')
                    
                    if messages:
                        first_message = messages.first()
                        last_message = messages.last()
                        
                        conversation_list.append({
                            'conversation_id': str(conversation.conversation_id),
                            'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                            'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                            'first_message': first_message.content,
                            'last_message': last_message.content,
                            'message_count': messages.count(),
                            'follow_up_questions': conversation.follow_up_questions or []
                        })
                
                return Response(conversation_list, status=status.HTTP_200_OK)
        
        except Exception as e:
            return Response(
                {'error': f'Failed to fetch conversations: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    # Add PATCH method to handle title updates
    def patch(self, request, conversation_id):
        try:
            # Log incoming request data for debugging
            print(f"PATCH request for conversation {conversation_id}")
            print(f"Request data: {request.data}")

            # Validate input
            new_title = request.data.get('title')
            is_active = request.data.get('is_active', True)

            if not new_title or not new_title.strip():
                return Response(
                    {'error': 'Title cannot be empty'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Find the conversation
            try:
                conversation = ChatHistory.objects.get(
                    conversation_id=conversation_id, 
                    user=request.user
                )
            except ChatHistory.DoesNotExist:
                return Response(
                    {'error': 'Conversation not found'},
                    status=status.HTTP_404_NOT_FOUND
                )

            # Update the title and active status
            conversation.title = new_title.strip()
            conversation.is_active = is_active
            conversation.save()

            # Log successful update
            print(f"Conversation {conversation_id} updated successfully")
            print(f"New title: {conversation.title}")

            return Response({
                'message': 'Conversation title updated successfully',
                'conversation_id': str(conversation.conversation_id),
                'new_title': conversation.title,
                'is_active': conversation.is_active
            }, status=status.HTTP_200_OK)

        except Exception as e:
            # Comprehensive error logging
            print(f"Error updating conversation title: {str(e)}")
            logger.error(f"Conversation title update error: {str(e)}")

            return Response(
                {'error': 'Failed to update conversation title', 'details': str(e)},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    # Also implement PUT for completeness
    def put(self, request, conversation_id):
        return self.patch(request, conversation_id)

# Optional: Add a view to delete a conversation
class DeleteConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def delete(self, request, conversation_id):
        try:
            conversation = ChatHistory.objects.get(
                conversation_id=conversation_id, 
                user=request.user
            )
            
            conversation.delete()
            return Response(
                {'message': 'Conversation deleted successfully'}, 
                status=status.HTTP_200_OK
            )
        
        except ChatHistory.DoesNotExist:
            return Response(
                {'error': 'Conversation not found'}, 
                status=status.HTTP_404_NOT_FOUND
            )
        except Exception as e:
            return Response(
                {'error': f'Failed to delete conversation: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


class GenerateIdeaContextView(APIView):
    """
    API endpoint to extract structured idea generation parameters
    from documents or query results.
    """
   
    def post(self, request):
        user = request.user
        document_id = request.data.get('document_id')
        query = request.data.get('query')
        main_project_id = request.data.get('main_project_id')
       
        # Validate input - need either document_id or query
        if not document_id and not query:
            return Response({
                'error': 'Either document_id or query parameter is required'
            }, status=status.HTTP_400_BAD_REQUEST)
           
        try:
            # Case 1: Using Document ID - fetch existing parameters or extract new ones
            if document_id:
                document = get_object_or_404(Document, id=document_id, user=user)
                
                # Get document name without extension
                document_name_no_ext = self.remove_file_extension(document.filename)
                
                # Generate a unique project name - handle the case when main_project_id is None
                suggested_project_name = f"Ideas from {document_name_no_ext}"
                if main_project_id:
                    try:
                        suggested_project_name = self.generate_unique_project_name(document_name_no_ext, main_project_id)
                    except Exception as e:
                        # Log the error but continue with the default name
                        print(f"Error generating unique project name: {str(e)}")
               
                try:
                    # Check if we already have parameters stored
                    processed_index = ProcessedIndex.objects.get(document=document)
                   
                    # If parameters exist, return them
                    if processed_index.idea_parameters:
                        return Response({
                            'document_id': document_id,
                            'document_name': document.filename,
                            'document_name_no_ext': document_name_no_ext,
                            'idea_parameters': processed_index.idea_parameters,
                            'suggested_project_name': suggested_project_name
                        })
                   
                    # If no parameters yet, extract them from the document
                    index_file = processed_index.faiss_index
                    metadata_file = processed_index.metadata
                    
                    # First check if the document has a markdown path (LlamaParse document)
                    if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                        # This is a LlamaParse document, read the markdown content directly
                        try:
                            with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                                full_text = f.read()
                                
                            # Extract parameters from markdown content
                            idea_params = self.extract_idea_parameters(full_text)
                            
                            # Save the parameters for future use
                            processed_index.idea_parameters = idea_params
                            processed_index.save()
                            
                            return Response({
                                'document_id': document_id,
                                'document_name': document.filename,
                                'document_name_no_ext': document_name_no_ext,
                                'idea_parameters': idea_params,
                                'suggested_project_name': suggested_project_name
                            })
                        except Exception as e:
                            print(f"Error reading markdown file: {str(e)}")
                            # Continue with FAISS approach as fallback
                   
                    # Load index and metadata
                    index, chunks = self.load_faiss_index_from_paths(index_file, metadata_file)
                    
                    if not chunks:
                        return Response({
                            'error': 'No content found in the document'
                        }, status=status.HTTP_400_BAD_REQUEST)
                   
                    # Extract parameters from document content
                    full_text = " ".join([chunk.get('text', '') for chunk in chunks])
                    idea_params = self.extract_idea_parameters(full_text, chunks)
                   
                    # Save the parameters for future use
                    processed_index.idea_parameters = idea_params
                    processed_index.save()

                    if main_project_id:
                        update_project_timestamp(main_project_id, user)
            
                   
                    return Response({
                        'document_id': document_id,
                        'document_name': document.filename,
                        'document_name_no_ext': document_name_no_ext,
                        'idea_parameters': idea_params,
                        'suggested_project_name': suggested_project_name
                    })
                   
                except ProcessedIndex.DoesNotExist:
                    return Response({
                        'error': 'Document has not been processed yet'
                    }, status=status.HTTP_404_NOT_FOUND)
           
            # Case 2: Using Query - search across documents and extract from relevant chunks
            else:
                # Get active/selected document
                active_doc_id = request.session.get('active_document_id')
                if not active_doc_id:
                    return Response({
                        'error': 'No active document selected'
                    }, status=status.HTTP_400_BAD_REQUEST)
               
                document = get_object_or_404(Document, id=active_doc_id, user=user)
                processed_index = get_object_or_404(ProcessedIndex, document=document)
                
                # Get document name without extension
                document_name_no_ext = self.remove_file_extension(document.filename)
                
                # Generate a unique project name - handle the case when main_project_id is None
                suggested_project_name = f"Ideas from {document_name_no_ext}"
                if main_project_id:
                    try:
                        suggested_project_name = self.generate_unique_project_name(document_name_no_ext, main_project_id)
                    except Exception as e:
                        # Log the error but continue with the default name
                        print(f"Error generating unique project name: {str(e)}")
                
                # First check if the document has a markdown path (LlamaParse document)
                if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                    # This is a LlamaParse document, read the markdown content directly
                    try:
                        with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                            full_text = f.read()
                            
                        # Extract parameters from markdown content
                        idea_params = self.extract_idea_parameters(query, None, full_text)
                        
                        return Response({
                            'document_id': active_doc_id,
                            'document_name': document.filename,
                            'document_name_no_ext': document_name_no_ext,
                            'query': query,
                            'idea_parameters': idea_params,
                            'suggested_project_name': suggested_project_name
                        })
                    except Exception as e:
                        print(f"Error reading markdown file: {str(e)}")
                        # Continue with FAISS approach as fallback
               
                # Load index and metadata for FAISS approach
                index_file = processed_index.faiss_index
                metadata_file = processed_index.metadata
                index, chunks = self.load_faiss_index_from_paths(index_file, metadata_file)
               
                # Get embedding for query
                query_embedding = self.get_query_embedding(query)
               
                # Search for relevant chunks
                relevant_chunks = self.search_faiss_index(index, chunks, query_embedding, k=5)
               
                # Extract parameters from relevant chunks
                idea_params = self.extract_idea_parameters(query, relevant_chunks)

                if main_project_id:
                    update_project_timestamp(main_project_id, user)
            
               
                return Response({
                    'document_id': active_doc_id,
                    'document_name': document.filename,
                    'document_name_no_ext': document_name_no_ext,
                    'query': query,
                    'idea_parameters': idea_params,
                    'suggested_project_name': suggested_project_name
                })
               
        except Exception as e:
            print(f"Error generating idea context: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while generating idea context'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def remove_file_extension(self, filename):
        """
        Remove file extension from filename
        """
        import os
        return os.path.splitext(filename)[0]
    
    def generate_unique_project_name(self, document_name, main_project_id):
        """
        Generate a unique project name based on document name,
        adding (1), (2), etc. if needed to avoid duplicates
        """
        from ideaGen.models import Project
        
        base_name = f"Ideas from {document_name}"
        project_name = base_name
        counter = 1
        
        # Check for existing projects with this name in the main project
        while Project.objects.filter(
            name=project_name,
            main_project_id=main_project_id
        ).exists():
            # Increment counter and update name
            project_name = f"{base_name} ({counter})"
            counter += 1
            
        return project_name
   
    def load_faiss_index_from_paths(self, index_file, metadata_file):
        """Load FAISS index and metadata from file paths"""
        import faiss
        import pickle
       
        try:
            index = faiss.read_index(index_file)
            with open(metadata_file, "rb") as f:
                chunks = pickle.load(f)
            return index, chunks
        except Exception as e:
            print(f"Error loading FAISS index: {str(e)}")
            return None, []
   
    def get_query_embedding(self, query):
        """Get embedding for the query"""
        from openai import OpenAI
        client = OpenAI()  # Initialize the client
       
        try:
            response = client.embeddings.create(
                input=[query],
                model="text-embedding-3-small"
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"Error getting embedding for query: {str(e)}")
            raise
   
    def search_faiss_index(self, index, chunks, query_embedding, k=5):
        """Search FAISS index for relevant chunks"""
        import numpy as np
       
        # Check if index is None
        if index is None:
            return []
            
        # Convert embedding to numpy array
        query_vector = np.array([query_embedding]).astype('float32')
       
        # Search index
        distances, indices = index.search(query_vector, k=k)
       
        # Get relevant chunks
        results = []
        for i in indices[0]:
            if i < len(chunks):
                results.append(chunks[i])
       
        return results
   
    def extract_idea_parameters(self, context, relevant_chunks=None, full_text=None):
        """
        Extract structured parameters for the Idea Generator
       
        Args:
            context (str): Either full document content or query
            relevant_chunks (list, optional): List of relevant chunks from search
            full_text (str, optional): For LlamaParse documents, the full markdown text
           
        Returns:
            dict: Structured parameters for idea generation
        """
        from openai import OpenAI
        client = OpenAI()  # Initialize the client
        import json
       
        try:
            # Determine extraction context source
            if full_text:
                # If we have full text (e.g., from markdown), use that
                extraction_context = full_text
            elif relevant_chunks and len(relevant_chunks) > 0:
                # If relevant chunks are provided, use them for extraction context
                extraction_context = "\n\n".join([chunk.get('text', '') for chunk in relevant_chunks])
            else:
                # If no chunks available, use the context directly
                extraction_context = context
               
            # Create extraction prompt
            extraction_prompt = f"""
            Extract the following key parameters from the provided context.
            If a parameter isn't explicitly mentioned, infer it from context or leave it blank.
           
            Context:
            {extraction_context}
           
            Extract these parameters:
            - Brand_Name: The company or product brand mentioned
            - Category: The product category or market area
            - Concept: The main idea, campaign, or product concept
            - Benefits: Key benefits or value propositions (list up to 3)
            - RTB: Reason to Believe - evidence supporting the benefits (list up to 3)
            - Ingredients: Key ingredients or components (if applicable)
            - Features: Notable product/service features (list up to 3)
            - Theme: Overall theme or tone
            - Demographics: Target audience demographics
           
            Format the response as a valid JSON object with these fields.
            """
           
            # Call LLM to extract parameters
            response = client.chat.completions.create(
                model="gpt-4o",  
                messages=[
                    {"role": "system", "content": "You are a specialist in extracting structured information from documents."},
                    {"role": "user", "content": extraction_prompt}
                ],
                response_format={"type": "json_object"}
            )
           
            # Parse JSON response
            extracted_params = json.loads(response.choices[0].message.content)
           
            return extracted_params
       
        except Exception as e:
            print(f"Error extracting idea parameters: {str(e)}")
            # Return empty structure if extraction fails
            return {
                "Brand_Name": "",
                "Category": "",
                "Concept": "",
                "Benefits": "",
                "RTB": "",
                "Ingredients": "",
                "Features": "",
                "Theme": "",
                "Demographics": ""
            }

from django.db import transaction
from .models import UserAPITokens

class AdminUserStatsView(APIView):
    permission_classes = [IsAuthenticated]
 
    def get(self, request):
        # Only allow admin user
        if not request.user.username == 'admin':
            return Response({'error': 'Unauthorized'}, status=403)
 
        users = User.objects.all()
        stats = []
 
        for user in users:
            # Get all active documents for this user
            docs = Document.objects.filter(user=user)
            doc_stats = []
            for doc in docs:
                # Count user questions for this document
                # Find all conversations (ChatHistory) that reference this doc
                # (Assuming you have a ManyToMany from ChatHistory to Document as 'documents')
                from .models import ChatHistory
                conversations = ChatHistory.objects.filter(user=user, documents=doc)
                question_count = ChatMessage.objects.filter(
                    chat_history__in=conversations,
                    role='user'
                ).count()
                doc_stats.append({
                    'document_id': doc.id,
                    'filename': doc.filename,
                    'questions_asked': question_count,
                })
 
            # Count total document uploads
            doc_upload_count = docs.count()
 
            stats.append({
                'user_id': user.id,
                'username': user.username,
                'document_upload_count': doc_upload_count,
                'documents': doc_stats,
            })
 
        return Response({'user_stats': stats}, status=200)

 
from django.db import transaction
from .models import UserAPITokens

 
class AdminUserManagementView(APIView):
    def get(self, request):
        print("AdminUserManagementView GET request received")
        print(f"Request user: {request.user.username}")
        # Your existing GET method code stays the same
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can access this endpoint'
                }, status=status.HTTP_403_FORBIDDEN)
                    
            # Get all users with their API tokens
            users = User.objects.all()
                    
            user_data = []
            for user in users:
                # Get API tokens if they exist
                api_tokens = None
                try:
                    api_tokens = user.api_tokens
                except UserAPITokens.DoesNotExist:
                    pass
                  
                # Get disabled modules
                try:
                    module_permissions = user.module_permissions.disabled_modules
                except (AttributeError, UserModulePermissions.DoesNotExist):
                    module_permissions = {}
                    # Create module permissions if they don't exist
                    UserModulePermissions.objects.create(user=user, disabled_modules={})
                  
                # Get upload permissions
                try:
                    upload_permissions = UserUploadPermissions.objects.get(user=user)
                    can_upload = upload_permissions.can_upload
                except UserUploadPermissions.DoesNotExist:
                    # Default to allowed if permissions don't exist
                    can_upload = True
                    # Create default permissions
                    UserUploadPermissions.objects.create(user=user, can_upload=True)
                            
                user_info = {
                    'id': user.id,
                    'username': user.username,
                    'email': user.email,
                    'tokens_used': user.tokens_used if hasattr(user, 'tokens_used') else 0,                 
                    'token_limit': api_tokens.token_limit if api_tokens and api_tokens.token_limit is not None else None,
                    'page_limit': api_tokens.page_limit if api_tokens and api_tokens.page_limit is not None else None,
                    'disabled_modules': module_permissions,
                    'upload_permissions': {
                        'can_upload': can_upload
                    },
                    'api_tokens': {
                        'nebius_token': api_tokens.nebius_token if api_tokens else None,
                        'gemini_token': api_tokens.gemini_token if api_tokens else None,
                        'llama_token': api_tokens.llama_token if api_tokens else None,
                        'huggingface_token': api_tokens.huggingface_token if api_tokens else None,
                        'token_limit': api_tokens.token_limit if api_tokens and api_tokens.token_limit is not None else None,
                        'page_limit': int(api_tokens.page_limit) if api_tokens and api_tokens.page_limit is not None else None,
                    }
                }
                user_data.append(user_info)
                print("##############################:", user_data)  # Log user info for debugging
                    
            return Response(user_data, status=status.HTTP_200_OK)
                      
        except Exception as e:
            print(f"Error in AdminUserManagementView.get: {str(e)}")
            return Response(
                {'error': f'Failed to fetch user data: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def put(self, request):
        """Update user tokens - ADD THIS METHOD"""
        print(f"AdminUserManagementView PUT data: {request.data}")
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can access this endpoint'
                }, status=status.HTTP_403_FORBIDDEN)

            user_id = request.data.get('user_id')
            if not user_id:
                return Response({
                    'error': 'user_id is required'
                }, status=status.HTTP_400_BAD_REQUEST)

            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'User not found'
                }, status=status.HTTP_404_NOT_FOUND)

            # Get or create API tokens for the user
            api_tokens, created = UserAPITokens.objects.get_or_create(user=user)

            # Update tokens
            if 'nebius_token' in request.data:
                api_tokens.nebius_token = request.data['nebius_token']
            if 'gemini_token' in request.data:
                api_tokens.gemini_token = request.data['gemini_token']
            if 'llama_token' in request.data:
                api_tokens.llama_token = request.data['llama_token']
            # if 'token_limit' in request.data:
            #     api_tokens.token_limit = request.data['token_limit'] if request.data['token_limit'] else None

            if 'token_limit' in request.data:
                token_limit = request.data['token_limit']
                api_tokens.token_limit = int(token_limit) if token_limit not in [None, ""] else 1_000_000
            # FIX: Set default if page_limit is empty
            if 'page_limit' in request.data:
                page_limit = request.data['page_limit']
                api_tokens.page_limit = int(page_limit) if page_limit not in [None, ""] else 20

            api_tokens.save()

            return Response({
                'success': True,
                'message': 'User tokens updated successfully'
            }, status=status.HTTP_200_OK)

        except Exception as e:
            print(f"Error in AdminUserManagementView.put: {str(e)}")
            return Response(
                {'error': f'Failed to update user tokens: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def post(self, request):
        """Create new user - Admin only"""

        from dotenv import load_dotenv
        import os
 
        load_dotenv()
        
        try:
            # Restrict to admin only
            if request.user.username != 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can access this endpoint'
                }, status=status.HTTP_403_FORBIDDEN)
 
            username = request.data.get('username')
            email = request.data.get('email')
            password = request.data.get('password')
            nebius_token = request.data.get('nebius_token', '')
            gemini_token = request.data.get('gemini_token', '')
            llama_token = request.data.get('llama_token', '')
            huggingface_token = request.data.get('huggingface_token', '')
 
            if not all([username, email, password]):
                return Response({
                    'error': 'username, email, and password are required'
                }, status=status.HTTP_400_BAD_REQUEST)
 
            # Check duplicates
            if User.objects.filter(username=username).exists():
                return Response({'error': 'Username already exists'}, status=status.HTTP_400_BAD_REQUEST)
 
            if User.objects.filter(email=email).exists():
                return Response({'error': 'Email already exists'}, status=status.HTTP_400_BAD_REQUEST)
 
            # Create the user
            user = User.objects.create_user(username=username, email=email, password=password)
 
            # ✅ Update the existing UserAPITokens created by signal
            user.api_tokens.huggingface_token = huggingface_token or os.getenv("HUGGINGFACE_TOKEN", "")
            user.api_tokens.gemini_token = gemini_token or os.getenv("GOOGLE_API_KEY", "")
            user.api_tokens.nebius_token = nebius_token or os.getenv("NEBIUS_TOKEN", "")
            user.api_tokens.llama_token = llama_token or os.getenv("LLAMA_TOKEN", "")
            user.api_tokens.save()
 
            return Response({
                'success': True,
                'message': 'User created successfully',
                'user_id': user.id
            }, status=status.HTTP_201_CREATED)

        except Exception as e:
            print(f"Error in AdminUserManagementView.post: {str(e)}")
            return Response(
                {'error': f'Failed to create user: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def delete(self, request):
        """Delete user - ADD THIS METHOD"""
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can access this endpoint'
                }, status=status.HTTP_403_FORBIDDEN)

            user_id = request.GET.get('user_id')
            if not user_id:
                return Response({
                    'error': 'user_id is required'
                }, status=status.HTTP_400_BAD_REQUEST)

            try:
                user = User.objects.get(id=user_id)
                
                # Prevent deletion of admin user
                if user.username == 'admin':
                    return Response({
                        'error': 'Cannot delete admin user'
                    }, status=status.HTTP_400_BAD_REQUEST)
                
                user.delete()
                
                return Response({
                    'success': True,
                    'message': 'User deleted successfully'
                }, status=status.HTTP_200_OK)

            except User.DoesNotExist:
                return Response({
                    'error': 'User not found'
                }, status=status.HTTP_404_NOT_FOUND)

        except Exception as e:
            print(f"Error in AdminUserManagementView.delete: {str(e)}")
            return Response(
                {'error': f'Failed to delete user: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    # Token Usage Tracking Methods (keep these the same)
    @staticmethod
    def check_token_limit(user):
        """Check if user has exceeded their token limit"""
        try:
            api_tokens = user.api_tokens
            if api_tokens.token_limit:
                current_usage = user.tokens_used if hasattr(user, 'tokens_used') else 0
                if current_usage >= api_tokens.token_limit:
                    return False, "Token limit exceeded"
            return True, None
        except UserAPITokens.DoesNotExist:
            return True, None  # No limit set, allow usage

    @staticmethod
    def increment_token_usage(user, tokens_used):
        """Increment user's token usage"""
        try:
            current_usage = user.tokens_used if hasattr(user, 'tokens_used') else 0
            user.tokens_used = current_usage + tokens_used
            user.save()
            return True
        except Exception as e:
            print(f"Error incrementing token usage: {str(e)}")
            return False

    @staticmethod
    def get_remaining_tokens(user):
        """Get remaining tokens for a user"""
        try:
            api_tokens = user.api_tokens
            if api_tokens.token_limit:
                current_usage = user.tokens_used if hasattr(user, 'tokens_used') else 0
                return max(0, api_tokens.token_limit - current_usage)
            return None  # Unlimited
        except UserAPITokens.DoesNotExist:
            return None  # Unlimited


class AdminUserModuleView(APIView):
    def patch(self, request, user_id):
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can update module permissions'
                }, status=status.HTTP_403_FORBIDDEN)
            
            # Get the user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'User not found'
                }, status=status.HTTP_404_NOT_FOUND)
                
            # Extract disabled_modules from request data
            disabled_modules = request.data.get('disabled_modules', {})
            
            # Get or create module permissions
            module_permissions, created = UserModulePermissions.objects.get_or_create(user=user)
            module_permissions.disabled_modules = disabled_modules
            module_permissions.save()
            
            # Get API tokens if they exist
            api_tokens = None
            try:
                api_tokens = user.api_tokens  # Assuming this relation exists
            except AttributeError:
                pass
            
            # Return updated user data
            return Response({
                'id': user.id,
                'username': user.username,
                'email': user.email,
                'disabled_modules': module_permissions.disabled_modules,
                'api_tokens': {
                    'nebius_token': getattr(api_tokens, 'nebius_token', None) if api_tokens else None,
                    'gemini_token': getattr(api_tokens, 'gemini_token', None) if api_tokens else None,
		            'llama_token': getattr(api_tokens, 'llama_token', None) if api_tokens else None
                }
            }, status=status.HTTP_200_OK)
                
        except Exception as e:
            print(f"Error in AdminUserModuleView.patch: {str(e)}")
            return Response(
                {'error': f'Failed to update module permissions: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

# Add this to your views.py

# Add this to your views.py file
class OriginalDocumentView(APIView):
    permission_classes = [IsAuthenticated]
    
    # Fixed document viewing implementation for OriginalDocumentView.get method

    # Update OriginalDocumentView.get in views.py

    def get(self, request, document_id):
        try:
            # Get the document making sure it belongs to the user
            document = get_object_or_404(Document, id=document_id, user=request.user)
            
            # Get the file path
            file_path = document.file.path
            
            # Check if file exists
            if not os.path.exists(file_path):
                return Response(
                    {'error': 'Document file not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
                
            # Determine content type based on file extension
            content_type = self.get_content_type(file_path)
            
            # Create a file response
            response = FileResponse(
                open(file_path, 'rb'),
                content_type=content_type
            )
            
            # Set content disposition to attachment with the original filename
            response['Content-Disposition'] = f'inline; filename="{document.filename}"'
            
            # Log the file access
            self.log_document_view(request.user, document)
            
            return response
            
        except Exception as e:
            print(f"Error serving original document: {str(e)}")
            return Response(
                {'error': f'Failed to retrieve document: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    def get_content_type(self, file_path):
        """Determine content type based on file extension"""
        import os
        extension = os.path.splitext(file_path)[1].lower()
        
        content_types = {
            '.pdf': 'application/pdf',
            '.doc': 'application/msword',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xls': 'application/vnd.ms-excel',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.txt': 'text/plain',
            '.csv': 'text/csv',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.md': 'text/markdown'
        }
        
        return content_types.get(extension, 'application/octet-stream')
    
    def log_document_view(self, user, document):
        """Log document view for analytics (optional)"""
        try:
            # You can create a DocumentViewLog model to track this
            # For now, we'll just log it
            logger.info(f"User {user.username} viewed document {document.id}: {document.filename}")
            
        except Exception as e:
            logger.error(f"Error logging document view: {str(e)}")
            # Don't raise the exception - this is a non-critical feature

class DocumentViewLogView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request, document_id):
        """Log when a user views a document"""
        try:
            # Get the document ensuring it belongs to the user
            document = get_object_or_404(Document, id=document_id, user=request.user)
    
            
            # For simplicity, just update the document's view count
            document.view_count = F('view_count') + 1  # This requires adding view_count field to Document model
            document.last_viewed_at = timezone.now()  # This requires adding last_viewed_at field to Document model
            document.save(update_fields=['view_count', 'last_viewed_at'])

            # Get main_project_id from the document
            main_project_id = document.main_project_id
            
            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
            
            return Response({
                'success': True,
                'message': 'Document view logged successfully'
            }, status=status.HTTP_200_OK)
            
        except Document.DoesNotExist:
            return Response({
                'error': 'Document not found'
            }, status=status.HTTP_404_NOT_FOUND)
            
        except Exception as e:
            print(f"Error logging document view: {str(e)}")
            # Return success anyway - this is a non-critical feature
            return Response({
                'success': True,
                'message': 'Continued without logging view'
            }, status=status.HTTP_200_OK)




class DocumentContentSearchView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        try:
            user = request.user
            query = request.data.get('query')
            main_project_id = request.data.get('main_project_id')
            
            if not query or not main_project_id:
                return Response({
                    'error': 'Query and project ID are required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Get all documents for this user and project
            documents = Document.objects.filter(
                user=user, 
                main_project_id=main_project_id
            )
            
            search_results = []
            
            # For each document, search for query in the content
            for doc in documents:
                try:
                    # Get the processed index for this document
                    processed_index = ProcessedIndex.objects.get(document=doc)
                    
                    # First, check if this is a LlamaParse document
                    if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                        # This is a LlamaParse document, read the markdown content directly
                        try:
                            with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                                markdown_content = f.read()
                            
                            # Convert query to lowercase for case-insensitive search
                            query_lower = query.lower()
                            
                            # Check if query exists in the markdown content
                            if query_lower in markdown_content.lower():
                                # Found a match - get the surrounding context
                                index = markdown_content.lower().find(query_lower)
                                start = max(0, index - 100)
                                end = min(len(markdown_content), index + len(query) + 100)
                                match_context = markdown_content[start:end]
                                
                                # Count occurrences
                                match_count = markdown_content.lower().count(query_lower)
                                
                                search_results.append({
                                    'id': doc.id,
                                    'filename': doc.filename,
                                    'match_count': match_count,
                                    'match_context': match_context,
                                    'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M')
                                })
                                
                                # Continue to the next document
                                continue
                        except Exception as e:
                            logger.error(f"Error searching in markdown file for document {doc.id}: {str(e)}")
                            # Continue with the FAISS approach as fallback
                    
                    # Standard FAISS approach for non-LlamaParse documents
                    if processed_index.faiss_index and processed_index.metadata:
                        # Check if the files exist
                        if not os.path.exists(processed_index.faiss_index) or not os.path.exists(processed_index.metadata):
                            continue
                            
                        # Load the metadata (chunks)
                        with open(processed_index.metadata, 'rb') as f:
                            chunks = pickle.load(f)
                        
                        # Convert query to lowercase for case-insensitive search
                        query_lower = query.lower()
                        matches = []
                        
                        # Search in each chunk
                        for chunk in chunks:
                            chunk_text = chunk.get('text', '')
                            if chunk_text and query_lower in chunk_text.lower():
                                # Found a match, add the context
                                matches.append(chunk_text)
                        
                        # If we found matches, add to results
                        if matches:
                            # Use the first match for preview, but count all matches
                            search_results.append({
                                'id': doc.id,
                                'filename': doc.filename,
                                'match_count': len(matches),
                                'match_context': matches[0] if matches else "",
                                'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M')
                            })
                                
                except ProcessedIndex.DoesNotExist:
                    # Skip documents that haven't been processed
                    continue
                except Exception as e:
                    logger.error(f"Error searching document {doc.id}: {str(e)}")
                    continue
            
            # Sort results by number of matches (most relevant first)
            search_results.sort(key=lambda x: x['match_count'], reverse=True)
            
            return Response({
                'results': search_results,
                'total_count': len(search_results)
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            logger.error(f"Error in document content search: {str(e)}")
            return Response({
                'error': f'An error occurred: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        

class AdminUserProjectsView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request, user_id):
        """Get all projects for a specific user"""
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can access this endpoint'
                }, status=status.HTTP_403_FORBIDDEN)
            
            # Get the user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'User not found'
                }, status=status.HTTP_404_NOT_FOUND)
            
            # Get projects for the user
            projects = Project.objects.filter(user=user)
            
            # Format response
            project_list = []
            for project in projects:
                project_list.append({
                    'id': project.id,
                    'name': project.name,
                    'created_at': project.created_at.strftime('%Y-%m-%d %H:%M:%S')
                })
            
            return Response(project_list, status=status.HTTP_200_OK)
            
        except Exception as e:
            logger.error(f"Error in AdminUserProjectsView: {str(e)}")
            return Response(
                {'error': f'Failed to fetch user projects: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
        

class AdminUserUploadPermissionsView(APIView):
    permission_classes = [IsAuthenticated]
    
    def patch(self, request, user_id):
        """Update a user's upload permissions"""
        try:
            # Check if the requesting user is an admin
            if not request.user.username == 'admin':
                return Response({
                    'error': 'Unauthorized. Only admin users can update upload permissions',
                    'success': False
                }, status=status.HTTP_403_FORBIDDEN)
            
            # Get the user
            try:
                user = User.objects.get(id=user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found',
                    'success': False
                }, status=status.HTTP_404_NOT_FOUND)
                
            # Extract permissions data - with explicit conversion to boolean
            can_upload = request.data.get('can_upload', True)
            
            # Convert can_upload to boolean if it's a string
            if isinstance(can_upload, str):
                can_upload = can_upload.lower() in ('true', 't', 'yes', 'y', '1')
            
            # Get or create upload permissions
            permissions, created = UserUploadPermissions.objects.get_or_create(
                user=user,
                defaults={'can_upload': can_upload}
            )
            
            # Update permissions if not created
            if not created:
                permissions.can_upload = can_upload
                permissions.save()
            
            # Log the update for troubleshooting
            print(f"Admin user {request.user.username} set upload permissions for user {user.username} to {can_upload}")
            
            return Response({
                'id': user.id,
                'username': user.username,
                'upload_permissions': {
                    'can_upload': permissions.can_upload
                },
                'success': True
            }, status=status.HTTP_200_OK)
                
        except Exception as e:
            print(f"Error in AdminUserUploadPermissionsView: {str(e)}")
            return Response({
                'error': f'Failed to update upload permissions: {str(e)}',
                'success': False
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            

class UserUploadPermissionsMiddleware(MiddlewareMixin):
    def process_request(self, request):
        # Skip middleware for admin or non-upload endpoints
        if request.path_info.startswith('/admin/') or not request.path_info.endswith('/upload/'):
            return None
        
        # Skip for non-authenticated requests
        if not hasattr(request, 'user') or not request.user.is_authenticated:
            return None
        
        # Skip for admin users (they can always upload)
        if request.user.username == 'admin':
            return None
        
        # Check if user is allowed to upload
        try:
            permissions = UserUploadPermissions.objects.get(user=request.user)
            if not permissions.can_upload:
                return JsonResponse({
                    'error': 'You do not have permission to upload documents. Please contact your administrator.'
                }, status=403)
        except UserUploadPermissions.DoesNotExist:
            # Default to allowed if no permissions are explicitly set
            pass
        
        # Continue with the request
        return None
    
class CheckUploadPermissionsView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request):
        try:
            user = request.user
            
            # Skip permission check for admin users
            if user.username == 'admin':
                return Response({
                    'can_upload': True
                }, status=status.HTTP_200_OK)
            
            # Check if the user has upload permissions
            try:
                permissions = UserUploadPermissions.objects.get(user=user)
                can_upload = permissions.can_upload
            except UserUploadPermissions.DoesNotExist:
                # Default to allowed if no permissions are explicitly set
                can_upload = True
            
            return Response({
                'can_upload': can_upload
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            logger.error(f"Error in CheckUploadPermissionsView: {str(e)}")
            return Response({
                'error': f'Failed to check upload permissions: {str(e)}',
                'can_upload': False  # Default to disallowing upload on error
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class ProcessCitationsView(APIView):
    """
    API endpoint to process a response and map citations to relevant parts of the text.
    This can be called by the frontend when displaying a message.
    """
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        try:
            response_text = request.data.get('response_text')
            citations = request.data.get('citations')
            
            if not response_text or not citations:
                return Response({
                    'error': 'Both response_text and citations are required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Process the response text with citations
            processed_text, enhanced_citations = self.process_response_with_citations(response_text, citations)
            
            return Response({
                'processed_text': processed_text,
                'enhanced_citations': enhanced_citations
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            import traceback
            print(f"Citation processing error: {str(e)}")
            print(traceback.format_exc())
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def extract_citation_metadata(self, citation_text, source_file):
        """
        Extract metadata like page numbers and section titles from citation text.
        
        Args:
            citation_text (str): The citation text from the document
            source_file (str): The source file name
            
        Returns:
            dict: Enhanced citation metadata
        """
        # Default citation structure
        citation = {
            'source_file': source_file,
            'page_number': 'Unknown',
            'section_title': 'Unknown',
            'snippet': citation_text,
        }
        
        # Try to extract page numbers - common formats
        page_matches = re.findall(r'(page|p\.?)\s*(\d+)', citation_text, re.IGNORECASE)
        if page_matches:
            citation['page_number'] = page_matches[0][1]
        
        # Try to extract section titles - look for section headings
        section_matches = re.findall(r'(?:section|heading):\s*["\'"]([^\'"\']+)[\'"\']', citation_text, re.IGNORECASE)
        if section_matches:
            citation['section_title'] = section_matches[0]
        
        # Alternative: Look for text in quotes that might be section titles
        if citation['section_title'] == 'Unknown':
            quote_matches = re.findall(r'[\'"\']([^\'"\']{5,50})[\'"\']', citation_text)
            if quote_matches:
                citation['section_title'] = quote_matches[0]
        
        # Look for date information
        date_matches = re.findall(r'\b(\d{4}[-/]\d{1,2}[-/]\d{1,2})\b', citation_text)
        if date_matches:
            citation['date'] = date_matches[0]
        
        return citation
    
    def process_response_with_citations(self, response_text, citations):
        """
        Maps citations to specific parts of the response text based on content similarity.
        
        Args:
            response_text (str): The response text from the LLM
            citations (list): List of citation objects with source information
            
        Returns:
            tuple: (processed_text, enhanced_citations)
        """
        if not citations or not response_text:
            return response_text, citations
        
        # First, try to improve citation metadata
        enhanced_citations = []

        for idx, citation in enumerate(citations):
            # Get basic info
            source_file = citation.get('source_file', 'Unknown Document')
            snippet = citation.get('snippet', '')
            
            # If metadata is missing, try to extract it
            if citation.get('page_number') in [None, 'Unknown'] or citation.get('section_title') in [None, 'Unknown']:
                enhanced_citation = self.extract_citation_metadata(snippet, source_file)
                # Preserve original data where available
                for key, value in citation.items():
                    if value and value != 'Unknown':
                        enhanced_citation[key] = value
                enhanced_citations.append(enhanced_citation)
            else:
                enhanced_citations.append(citation)
        
        # Remove HTML tags for better text processing
        clean_text = strip_tags(response_text)
        
        # Step 1: Split the response into sentences
        sentences = sent_tokenize(clean_text)
        
        # Step 2: For each citation, find the most relevant sentence
        citation_mapping = {}  # Maps citation index to sentence index
        
        for citation_idx, citation in enumerate(enhanced_citations):
            snippet = citation.get('snippet', '').lower()
            if not snippet:
                continue
                
            # Generate keywords from the snippet
            keywords = set(re.findall(r'\b\w+\b', snippet.lower()))
            keywords = {k for k in keywords if len(k) > 3}  # Filter out short words
            
            # Score each sentence based on keyword overlap
            best_score = 0
            best_sentence_idx = -1
            
            for sent_idx, sentence in enumerate(sentences):
                sentence_lower = sentence.lower()
                sentence_words = set(re.findall(r'\b\w+\b', sentence_lower))
                
                # Calculate overlap score
                overlap = keywords.intersection(sentence_words)
                if overlap:
                    score = len(overlap) / max(len(keywords), 1)  # Avoid division by zero
                    
                    # Give higher score for exact phrase matches
                    for i in range(0, len(snippet), 10):
                        if i + 20 <= len(snippet):
                            phrase = snippet[i:i+20]
                            if phrase in sentence_lower:
                                score += 0.5
                    
                    if score > best_score:
                        best_score = score
                        best_sentence_idx = sent_idx
            
            # Map citation to sentence if we found a good match
            if best_score > 0.2 and best_sentence_idx >= 0:  # Threshold to ensure relevance
                if best_sentence_idx not in citation_mapping:
                    citation_mapping[best_sentence_idx] = []
                citation_mapping[best_sentence_idx].append(citation_idx)
        
        # Step 3: Reinsert the HTML and add citation markers
        # Start by marking positions in the original response_text where sentences end
        sentence_positions = []
        current_pos = 0
        
        for sentence in sentences:
            # Find this sentence in the original text
            sentence_pos = response_text.find(sentence, current_pos)
            if sentence_pos >= 0:
                sentence_end = sentence_pos + len(sentence)
                sentence_positions.append(sentence_end)
                current_pos = sentence_end
        
        # Now insert citation markers at the end of relevant sentences
        result = response_text
        offset = 0  # Track position offset as we insert markers
        
        for sent_idx, end_pos in enumerate(sentence_positions):
            if sent_idx in citation_mapping:
                # Create citation markers
                citation_markers = ""
                for citation_idx in citation_mapping[sent_idx]:
                    citation_markers += f'<citation id="{citation_idx}"></citation>'
                
                # Insert at the adjusted position
                insert_pos = end_pos + offset
                if insert_pos <= len(result):
                    result = result[:insert_pos] + citation_markers + result[insert_pos:]
                    offset += len(citation_markers)
        
        return result, enhanced_citations


# codes for chat download

class ChatExportView(APIView):
    permission_classes = [IsAuthenticated]
 
    def remove_citations_from_html(self, html):
        """
        Remove all [n], [n][m], [n, m] style citations from HTML content.
        """
        if not html:
            return html
        # Remove [number] and [number][number] patterns
        html = re.sub(r'(\[\d+(?:,\s*\d+)*\])+', '', html)
        # Optionally, remove citations inside <sup> or <span> tags (if any)
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup.find_all(['sup', 'span']):
            if tag.string and re.match(r'^\[\d+(?:,\s*\d+)*\]$', tag.string.strip()):
                tag.decompose()
        return str(soup)
   
    def post(self, request):
        try:
            # Extract parameters from request
            chats = request.data.get('chats', [])
            date_range = request.data.get('date_range')
            main_project_id = request.data.get('main_project_id')
            options = request.data.get('options', {})
           
            # Validate inputs
            if not chats and not (date_range and main_project_id):
                return Response({
                    'error': 'Either chat data or date_range with main_project_id is required'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Default options
            include_timestamps = options.get('includeTimestamps', True)
            include_metadata = options.get('includeChatMetadata', True)
            include_follow_ups = options.get('includeFollowUpQuestions', True)
            format_code = options.get('formatCode', True)
           
            # If we have date_range but no chats, fetch the chats from database
            if date_range and main_project_id and not chats:
                try:
                    start_date = datetime.fromisoformat(date_range.get('startDate').replace('Z', '+00:00'))
                    end_date = datetime.fromisoformat(date_range.get('endDate').replace('Z', '+00:00'))
                    end_date = end_date + timedelta(days=1)  # Include the entire end day
                   
                    # Fetch conversations from database
                    conversations = ChatHistory.objects.filter(
                        user=request.user,
                        main_project_id=main_project_id,
                        created_at__gte=start_date,
                        created_at__lt=end_date
                    ).order_by('-created_at')
                   
                    if not conversations:
                        return Response({
                            'error': 'No conversations found in the specified date range'
                        }, status=status.HTTP_404_NOT_FOUND)
                   
                    # Convert to the same format as frontend-processed chats
                    chats = [{
                        'conversation_id': conv.conversation_id,
                        'title': conv.title,
                        'created_at': conv.created_at.strftime('%Y-%m-%d %H:%M'),
                        'messages': [{
                            'role': msg.role,
                            'content': msg.content,
                            'created_at': msg.created_at.strftime('%Y-%m-%d %H:%M')
                        } for msg in conv.messages.all().order_by('created_at')],
                        'follow_up_questions': conv.follow_up_questions or []
                    } for conv in conversations]
                   
                except Exception as e:
                    return Response({
                        'error': f'Invalid date format: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
           
            # Create a DOCX file in memory
            doc = DocxDocument()
           
            # Set document styles
            styles = doc.styles
           
            # Heading styles
            heading1_style = styles['Heading 1']
            heading1_style.font.size = Pt(16)
            heading1_style.font.bold = True
            heading1_style.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue
           
            heading2_style = styles['Heading 2']
            heading2_style.font.size = Pt(14)
            heading2_style.font.bold = True
            heading2_style.font.color.rgb = RGBColor(0, 0, 100)  # Slightly darker blue
           
            # Normal style
            normal_style = styles['Normal']
            normal_style.font.size = Pt(11)
            normal_style.font.name = 'Calibri'
           
            # List styles
            if 'List Bullet' in styles:
                bullet_style = styles['List Bullet']
                bullet_style.paragraph_format.left_indent = Inches(0.25)
                bullet_style.paragraph_format.first_line_indent = Inches(-0.25)
               
            if 'List Number' in styles:
                number_style = styles['List Number']
                number_style.paragraph_format.left_indent = Inches(0.25)
                number_style.paragraph_format.first_line_indent = Inches(-0.25)
           
            # Code style
            if 'Code' not in styles:
                code_style = styles.add_style('Code', WD_STYLE_TYPE.PARAGRAPH)
                code_style.font.name = 'Consolas'
                code_style.font.size = Pt(10)
                code_style.paragraph_format.space_before = Pt(6)
                code_style.paragraph_format.space_after = Pt(6)
                code_style.paragraph_format.left_indent = Inches(0.25)
                code_style.font.color.rgb = RGBColor(50, 50, 50)
           
            # Set document margins
            sections = doc.sections
            for section in sections:
                section.left_margin = Inches(1.0)
                section.right_margin = Inches(1.0)
                section.top_margin = Inches(1.0)
                section.bottom_margin = Inches(1.0)
           
            # Document title
            if len(chats) == 1:
                doc.add_heading(f"Chat Conversation: {chats[0]['title']}", level=0)
            else:
                # Get the start date and add one day
                start_date_str = date_range.get('startDate').split('T')[0] if date_range else chats[-1]['created_at'].split(' ')[0]
                start_date = datetime.strptime(start_date_str, '%Y-%m-%d') + timedelta(days=1)
                start_date_str = start_date.strftime('%Y-%m-%d')
               
                # Get the end date and add one day
                end_date_str = date_range.get('endDate').split('T')[0] if date_range else chats[0]['created_at'].split(' ')[0]
                end_date = datetime.strptime(end_date_str, '%Y-%m-%d') + timedelta(days=1)
                end_date_str = end_date.strftime('%Y-%m-%d')
               
                doc.add_heading(f"Chat Conversations: {start_date_str} to {end_date_str}", level=0)
           
            # Process each conversation
            for i, chat in enumerate(chats):
                # Add conversation title
                section_heading = doc.add_heading(chat['title'], level=1)
               
                # Add metadata if requested
                if include_metadata:
                    metadata_para = doc.add_paragraph(style='Intense Quote')
                    metadata_para.add_run(f"Created: {chat['created_at']}")
               
                # Process each message
                for message in chat['messages']:
                    # Message role as heading
                    role_heading = "User Question:" if message['role'] == 'user' else "Assistant Answer:"
                    doc.add_heading(role_heading, level=2)
                   
                    # Add timestamp if requested
                    if include_timestamps:
                        timestamp = doc.add_paragraph(style='Subtitle')
                        timestamp.add_run(f"{message['created_at']}")
                   
                    content = message['content']
                    if message['role'] == 'assistant':
                        content = self.remove_citations_from_html(content)
 
                    # Add the content (already processed by frontend)
                    self.add_html_to_docx(doc, content, format_code)
               
                # Add follow-up questions if requested
                if include_follow_ups and chat.get('follow_up_questions'):
                    doc.add_heading("Suggested Follow-up Questions:", level=2)
                    for question in chat['follow_up_questions']:
                        p = doc.add_paragraph(style='List Bullet')
                        p.add_run(question)
               
                # Add page break between conversations if multiple
                if i < len(chats) - 1:
                    doc.add_page_break()
           
            # Save document to a BytesIO object
            docx_file = BytesIO()
            doc.save(docx_file)
            docx_file.seek(0)
 
            # Update project timestamp if needed
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
           
            # Create response with file
            if len(chats) == 1:
                safe_title = ''.join(c if c.isalnum() or c in ' _-' else '_' for c in chat['title'])
                safe_title = '_'.join(safe_title.split()[:3])
                filename = f"Chat_{safe_title}.docx"
            else:
                start_date = date_range.get('startDate').split('T')[0] if date_range else chats[-1]['created_at'].split(' ')[0]
                end_date = date_range.get('endDate').split('T')[0] if date_range else chats[0]['created_at'].split(' ')[0]
                filename = f"Chats_{start_date.replace('-', '')}_to_{end_date.replace('-', '')}.docx"
           
            response = HttpResponse(
                docx_file.getvalue(),
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
           
            return response
           
        except Exception as e:
            logger.error(f"Error exporting chat: {str(e)}", exc_info=True)
            return Response({
                'error': f'Failed to export chat: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    def add_html_to_docx(self, doc, html_content, format_code=True):
        """Parse HTML content and add it to the document with proper formatting."""
        soup = BeautifulSoup(html_content, "html.parser")
       
        # Process each top-level element
        for elem in soup.children:
            if isinstance(elem, NavigableString) and elem.strip():
                # Plain text
                doc.add_paragraph(elem.strip())
            elif elem.name:
                self.process_element(doc, elem, format_code, list_level=0)
   
    def process_element(self, doc, element, format_code=True, list_level=0):
        """Process a single HTML element and add it to the document."""
        # Handle different element types
        if element.name == 'p':
            p = doc.add_paragraph()
            self.process_inline_elements(p, element)
       
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])
            doc.add_heading(element.get_text(), level=min(level, 9))
       
        elif element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.left_indent = Inches(0.25 + (0.25 * list_level))
                p.paragraph_format.first_line_indent = Inches(-0.25)
               
                has_nested_list = False
                nested_list = None
               
                for child in li.children:
                    if child.name in ['ul', 'ol']:
                        has_nested_list = True
                        nested_list = child
                        break
               
                if has_nested_list and nested_list:
                    nested_list.extract()
               
                self.process_inline_elements(p, li)
               
                if has_nested_list and nested_list:
                    self.process_element(doc, nested_list, format_code, list_level + 1)
       
        elif element.name == 'ol':
            for i, li in enumerate(element.find_all('li', recursive=False), 1):
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.left_indent = Inches(0.25 + (0.25 * list_level))
                p.paragraph_format.first_line_indent = Inches(-0.25)
               
                has_nested_list = False
                nested_list = None
               
                for child in li.children:
                    if child.name in ['ul', 'ol']:
                        has_nested_list = True
                        nested_list = child
                        break
               
                if has_nested_list and nested_list:
                    nested_list.extract()
               
                self.process_inline_elements(p, li)
               
                if has_nested_list and nested_list:
                    self.process_element(doc, nested_list, format_code, list_level + 1)
       
        elif element.name == 'pre' and format_code:
            code_element = element.find('code')
            code_text = code_element.get_text() if code_element else element.get_text()
            doc.add_paragraph(code_text, style='Code')
       
        elif element.name == 'code' and element.parent.name != 'pre' and format_code:
            p = doc.add_paragraph()
            code_run = p.add_run(element.get_text())
            code_run.font.name = 'Consolas'
            code_run.font.size = Pt(10)
           
        elif element.name == 'table':
            rows = element.find_all('tr')
            if rows:
                max_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
                if max_cols > 0:
                    table = doc.add_table(rows=len(rows), cols=max_cols)
                    table.style = 'Table Grid'
                   
                    for i, row in enumerate(rows):
                        cells = row.find_all(['td', 'th'])
                        for j, cell in enumerate(cells):
                            if j < max_cols:
                                cell_para = table.cell(i, j).paragraphs[0]
                                if cell.name == 'th':
                                    run = cell_para.add_run(cell.get_text())
                                    run.bold = True
                                else:
                                    self.process_inline_elements(cell_para, cell)
       
        elif element.name == 'blockquote':
            p = doc.add_paragraph(style='Quote')
            self.process_inline_elements(p, element)
       
        elif element.name == 'div':
            for child in element.children:
                if isinstance(child, NavigableString) and child.strip():
                    doc.add_paragraph(child.strip())
                elif child.name:
                    self.process_element(doc, child, format_code)
       
        else:
            text = element.get_text(strip=True)
            if text:
                doc.add_paragraph(text)
   
    def process_inline_elements(self, paragraph, element):
        """Process inline elements like bold, italic, etc. within a paragraph."""
        if element is None:
            return
           
        for child in element.children:
            if isinstance(child, NavigableString):
                if str(child).strip():
                    paragraph.add_run(str(child))
            elif child.name == 'b' or child.name == 'strong':
                run = paragraph.add_run(child.get_text())
                run.bold = True
            elif child.name == 'i' or child.name == 'em':
                run = paragraph.add_run(child.get_text())
                run.italic = True
            elif child.name == 'u':
                run = paragraph.add_run(child.get_text())
                run.underline = True
            elif child.name == 'a':
                run = paragraph.add_run(child.get_text())
                run.underline = True
                run.font.color.rgb = RGBColor(0, 0, 255)
            elif child.name == 'code':
                run = paragraph.add_run(child.get_text())
                run.font.name = 'Consolas'
                run.font.size = Pt(10)
            elif child.name == 'br':
                paragraph.add_run("\n")
            elif child.name not in ['ul', 'ol']:
                self.process_inline_elements(paragraph, child)

  