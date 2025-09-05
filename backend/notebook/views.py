#views.py original
from django.contrib.auth.models import User
from django.contrib.auth import authenticate, login
from rest_framework import status
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.permissions import IsAuthenticated, AllowAny  
from chat.models import UserUploadPermissions
from rest_framework.parsers import JSONParser
from rest_framework.permissions import IsAuthenticated
from rest_framework.decorators import api_view, permission_classes
from moviepy import VideoFileClip, AudioFileClip
import time
import faiss
import numpy as np
import os
import pickle
import json
import re
from datetime import datetime
from django.core.files.storage import default_storage
from django.conf import settings
from django.utils.deprecation import MiddlewareMixin
# import google.generativeai as genai  
from google import genai
from google.genai import types
from core.utils import update_project_timestamp
from .models import (
    ChatHistory,
    ChatMessage,
    Document,
    ProcessedIndex,
    ConversationMemoryBuffer,
    UserModulePermissions,
    Note,
    ConversationTransaction,
    UserTransaction,
    DocumentTransaction,
    TransactionType,
    DocumentProcessingStatus,
    DocumentEmbedding,
)
import uuid
from rest_framework.authtoken.models import Token
from django.utils.safestring import mark_safe
import logging

logger = logging.getLogger(__name__)
from core.models import Project
from openai import OpenAI
from dotenv import load_dotenv
import os

from io import StringIO, BytesIO
from django.http import HttpResponse, FileResponse, JsonResponse
from django.utils import timezone
from django.utils.html import strip_tags
from django.shortcuts import get_object_or_404
from django.contrib.auth import update_session_auth_hash
from llama_parse import LlamaParse
import requests
from bs4 import BeautifulSoup
import re
from duckduckgo_search import DDGS
from urllib.parse import urlparse
import nltk
from nltk.tokenize import sent_tokenize
from django.utils.html import strip_tags
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

from moviepy import VideoFileClip

import tempfile
from django.core.files import File

from chat.models import UserAPITokens
from io import BytesIO
from datetime import datetime, timedelta
from docx import Document as DocxDocument
from docx.shared import Pt, Inches, RGBColor
from docx.enum.style import WD_STYLE_TYPE
from bs4 import BeautifulSoup, NavigableString
from django.http import HttpResponse
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated
from rest_framework import status
from rest_framework.response import Response
import logging
from django.db.models import F
import json
from typing import List, Dict, Tuple
import yt_dlp
from django.db import transaction
import tempfile
from google import genai


logger = logging.getLogger(__name__)




load_dotenv()

OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

if not OPENAI_API_KEY:
    raise ValueError("Missing required API keys in environment variables")



client = OpenAI(api_key=OPENAI_API_KEY)




class GetUserDocumentsView(APIView):
    def get(self, request):
        try:
            user = request.user
            main_project_id = request.query_params.get('main_project_id')
            print(f"Getting documents for user {user.username} and project {main_project_id}")
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            documents = Document.objects.filter(user=user, main_project_id=main_project_id).select_related('processedindex')
        
            document_list = []
            for doc in documents:
                try:
                    processed = doc.processedindex
                    document_data = {
                        'id': doc.id,
                        'filename': doc.filename,
                        'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M'),
                        'summary': processed.summary,
                        'key_points' : processed.key_points,
                        'follow_up_questions': [
                            processed.follow_up_question_1,
                            processed.follow_up_question_2,
                            processed.follow_up_question_3
                        ] if all([processed.follow_up_question_1,
                                processed.follow_up_question_2,
                                processed.follow_up_question_3]) else []
                    }
                    document_list.append(document_data)

                   
                except ProcessedIndex.DoesNotExist:
                    document_data = {
                        'id': doc.id,
                        'filename': doc.filename,
                        'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M'),
                        'summary': 'Document processing pending',
                        'follow_up_questions': []
                    }
                    document_list.append(document_data)
            
            # Update project timestamp to track activity
            update_project_timestamp(main_project_id, user)
            # print("########################", document_list)
            return Response(document_list, status=status.HTTP_200_OK)
        
        except Exception as e:
            print(f"Error in GetUserDocumentsView: {str(e)}")
            return Response(
                {'error': f'Failed to fetch documents: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class ManageConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def update_conversation(self, request, conversation_id):
        try:
            # Log incoming request data for debugging
            print(f"Incoming update request for conversation {conversation_id}")
            print(f"Request data: {request.data}")

            # Validate input
            new_title = request.data.get('title')
            is_active = request.data.get('is_active', True)

            if not new_title or not new_title.strip():
                return Response(
                    {'error': 'Title cannot be empty'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Find the conversation
            try:
                conversation = ChatHistory.objects.get(
                    conversation_id=conversation_id, 
                    user=request.user
                )
            except ChatHistory.DoesNotExist:
                return Response(
                    {'error': 'Conversation not found'},
                    status=status.HTTP_404_NOT_FOUND
                )

            # Update the title and active status
            conversation.title = new_title.strip()
            conversation.is_active = is_active
            conversation.save()

            # Update project timestamp if main_project_id exists
            if hasattr(conversation, 'main_project_id') and conversation.main_project_id:
                update_project_timestamp(conversation.main_project_id, request.user)

            # Log successful update
            print(f"Conversation {conversation_id} updated successfully")
            print(f"New title: {conversation.title}")

            return Response({
                'message': 'Conversation title updated successfully',
                'conversation_id': str(conversation.conversation_id),
                'new_title': conversation.title,
                'is_active': conversation.is_active
            }, status=status.HTTP_200_OK)

        except Exception as e:
            # Comprehensive error logging
            print(f"Error updating conversation title: {str(e)}")
            logger.error(f"Conversation title update error: {str(e)}")

            return Response(
                {'error': 'Failed to update conversation title', 'details': str(e)},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def put(self, request, conversation_id):
        return self.update_conversation(request, conversation_id)

    def patch(self, request, conversation_id):
        return self.update_conversation(request, conversation_id)


# class DocumentProcessingMixin:
#     from google import genai
#     def _parse_json_response(self, response_content, fallback_message="Error processing response"):
#         """Helper method to safely parse JSON responses from LLM"""
#         try:
#             return json.loads(response_content)
#         except json.JSONDecodeError as e:
#             logger.error(f"JSON parsing failed: {str(e)}")
#             return {
#                 "content": fallback_message,
#                 "error": True,
#                 "format_type": "error"
#             }
 
#     def init_gemini(self, user=None):
#         """Initialize Gemini 2.5 Flash client"""
#         try:
#             # Get the user's Gemini API token
#             gemini_api_key = None
#             if user:
#                 try:
#                     from .models import UserAPITokens  # Adjust import based on your models location
#                     user_api_tokens = UserAPITokens.objects.get(user=user)
#                     gemini_api_key = user_api_tokens.gemini_token
                   
#                     if not gemini_api_key:
#                         logger.warning(f"No Gemini API token found for user {user.username}")
#                         # Fallback to environment variable
#                         gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
#                         if not gemini_api_key:
#                             raise ValueError("Gemini API token is required for processing")
                       
#                 except Exception as e:
#                     logger.error(f"Error getting API tokens for user {user.username}: {str(e)}")
#                     # Fallback to environment variable
#                     gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
#                     if not gemini_api_key:
#                         raise ValueError("Gemini API token is required for processing")
#             else:
#                 # Fallback to environment variable if no user provided
#                 gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
#                 if not gemini_api_key:
#                     raise ValueError("Gemini API token is required for processing")
           
#             from google import genai
#             client = genai.Client(api_key=gemini_api_key)
#             return client
#         except Exception as e:
#             logger.error(f"Error initializing Gemini: {str(e)}")
#             return None
 
#     def detect_document_structure(self, text: str, client) -> str:
#         """Analyze document structure to determine if it has clear headings/titles"""
#         try:
#             system_message = "You are an expert at analyzing document structure. Determine if documents have clear headings and sections."
           
#             user_prompt = f"""
#             Analyze this document text and determine if it has clear structure with headings, titles, or sections:
           
#             Text: {text[:2000]}
           
#             Look for:
#             - Clear headings (numbered sections, bold titles, etc.)
#             - Section breaks
#             - Structured format
#             - Table of contents references
           
#             Respond with only: "STRUCTURED" or "UNSTRUCTURED"
#             """
           
#             completion = client.chat.completions.create(
#                 model="gpt-4o",
#                 messages=[
#                     {"role": "system", "content": system_message},
#                     {"role": "user", "content": user_prompt}
#                 ],
#                 temperature=0.1,
#                 max_tokens=10
#             )
           
#             response_text = completion.choices[0].message.content.strip().upper()
#             return response_text
#         except Exception as e:
#             logger.error(f"Error detecting document structure: {str(e)}")
#             return "UNSTRUCTURED"  # Default to unstructured if analysis fails
 
#     def generate_topics_from_content(self, text: str, client) -> List[str]:
#         """Generate key topics from document content when no clear headings exist"""
#         try:
#             system_message = "You are an expert at identifying main topics and themes from documents. Generate concise, specific topic phrases."
           
#             user_prompt = f"""
#             Analyze this document and identify 5-7 main topics or themes discussed:
           
#             Text: {text[:6000]}
           
#             Generate concise topic phrases (3-5 words each) that represent the main concepts, ideas, or subjects covered in this document.
           
#             IMPORTANT:
#             - Focus on the most important themes and concepts
#             - Use clear, descriptive phrases
#             - Make topics specific to the document content
#             - Avoid generic terms
#             - Each topic should be 3-5 words
           
#             Format as a simple numbered list:
#             1. Topic phrase
#             2. Topic phrase
#             ...
#             """
           
#             completion = client.chat.completions.create(
#                 model="gpt-4o",
#                 messages=[
#                     {"role": "system", "content": system_message},
#                     {"role": "user", "content": user_prompt}
#                 ],
#                 temperature=0.4,
#                 max_tokens=500
#             )
           
#             response_text = completion.choices[0].message.content.strip()
#             topics = []
           
#             # Parse the numbered list
#             lines = response_text.split('\n')
#             for line in lines:
#                 line = line.strip()
#                 if line and (line[0].isdigit() or line.startswith('•') or line.startswith('-')):
#                     # Extract topic after number/bullet
#                     topic = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
#                     if topic:
#                         topics.append(topic)
           
#             return topics[:7]  # Limit to 7 topics
#         except Exception as e:
#             logger.error(f"Error generating topics: {str(e)}")
#             return []
 
#     def extract_headings_from_structured_doc(self, text: str, client) -> List[str]:
#         """Extract actual headings from structured documents"""
#         try:
#             system_message = "You are an expert at extracting headings and section titles from structured documents. Extract only exact headings that appear in the document."
           
#             user_prompt = f"""
#             Extract the main headings and section titles from this structured document:
           
#             Text: {text[:8000]}
           
#             IMPORTANT:
#             - Extract ONLY actual headings, titles, or section names that appear in the document
#             - Use exact words from the document (3-6 words per heading)
#             - Look for numbered sections, bold text, capitalized headings
#             - Do not create new phrases - use exact text from document
           
#             Format as a simple numbered list:
#             1. Exact heading from document
#             2. Exact heading from document
#             ...
#             """
           
#             completion = client.chat.completions.create(
#                 model="gpt-4o",
#                 messages=[
#                     {"role": "system", "content": system_message},
#                     {"role": "user", "content": user_prompt}
#                 ],
#                 temperature=0.2,
#                 max_tokens=800
#             )
           
#             response_text = completion.choices[0].message.content.strip()
#             headings = []
           
#             # Parse the numbered list
#             lines = response_text.split('\n')
#             for line in lines:
#                 line = line.strip()
#                 if line and (line[0].isdigit() or line.startswith('•') or line.startswith('-')):
#                     # Extract heading after number/bullet
#                     heading = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
#                     if heading:
#                         headings.append(heading)
           
#             return headings[:7]  # Limit to 7 headings
#         except Exception as e:
#             logger.error(f"Error extracting headings: {str(e)}")
#             return []
 
#     def generate_summary(self, text, file_name, user=None):
#         """
#         REPLACED: Generate summary and key points using GPT-4 with smart topic detection
#         Returns: (summary_text, key_points_list)
#         """
#         if not text.strip():
#             return "No extractable text found in the document.", []
       
#         try:
           
#             # First, detect if document has structure
#             doc_structure = self.detect_document_structure(text, client)
           
#             # Generate summary
#             summary_system_message = """You are an expert document analyzer. Provide comprehensive document summaries with proper HTML-like formatting."""
           
#             summary_user_prompt = f"""
#             Please analyze the following document '{file_name}' and provide a comprehensive quick summary (8-10 lines):
           
#             Text: {text[:8000]}
           
#             IMPORTANT FOR SUMMARY:
#             - Must be 8-10 lines providing comprehensive coverage
#             - FOCUS ON EXECUTION AND VALIDATION: Emphasize practical implementation steps, validation methods, testing approaches, and measurable outcomes
#             - DESCRIBE MORE CLEARLY: Use specific language, avoid vague terms, provide clear explanations of processes and concepts
#             - MENTION SPECIFIC SECTORS: Identify and name the exact industries, market segments, business sectors, or domains discussed
#             - BE MORE CONCRETE: Include specific numbers, percentages, timeframes, locations, company names, product names, or measurable data points
#             - COMPREHENSIVE AND APPLIED PICTURE: Show how the analysis translates to real-world applications, practical benefits, and actionable insights
#             - Include key findings, main themes, important data/statistics if present
#             - Cover methodology, results, conclusions, and recommendations if applicable
#             - Explain the context, significance, and implications with specific examples
#             - Make it thorough but still readable and well-structured
           
#             ### Expected Response Format:
#             <b>Summary Overview</b>
#             <p>High-level introduction to the document's main theme</p>
 
#             <b>Key Highlights</b>
#             <ul>
#                 <li>First major insight</li>
#                 <li>Second major insight</li>
#                 <li>Third major insight</li>
#             </ul>
 
#             <b>Detailed Insights</b>
#             <p>Expanded explanation of the document's core content and significance</p>
#             """
           
#             summary_completion = client.chat.completions.create(
#                 model="gpt-4o",
#                 messages=[
#                     {"role": "system", "content": summary_system_message},
#                     {"role": "user", "content": summary_user_prompt}
#                 ],
#                 temperature=0.4,
#                 max_tokens=2000
#             )
           
#             summary = summary_completion.choices[0].message.content.strip()
           
#             # Format the summary response to ensure proper HTML-like formatting
#             summary = self.format_summary_response(summary)
           
#             # Generate key points based on document structure
#             if doc_structure == "STRUCTURED":
#                 key_points = self.extract_headings_from_structured_doc(text, client)
#             else:
#                 key_points = self.generate_topics_from_content(text, client)
           
#             # Fallback if no key points were generated
#             if not key_points:
#                 fallback_system_message = "You are an expert at identifying main topics and themes from documents."
               
#                 fallback_user_prompt = f"""
#                 From this document, identify 5 main topics or themes:
               
#                 Text: {text[:4000]}
               
#                 Provide 5 short phrases (3-4 words each) that capture the main ideas.
#                 Format as:
#                 1. Topic phrase
#                 2. Topic phrase
#                 ...
#                 """
               
#                 fallback_completion = client.chat.completions.create(
#                     model="gpt-4o",
#                     messages=[
#                         {"role": "system", "content": fallback_system_message},
#                         {"role": "user", "content": fallback_user_prompt}
#                     ],
#                     temperature=0.4,
#                     max_tokens=500
#                 )
               
#                 fallback_response = fallback_completion.choices[0].message.content.strip()
#                 lines = fallback_response.split('\n')
#                 key_points = []
#                 for line in lines:
#                     line = line.strip()
#                     if line and (line[0].isdigit() or line.startswith('•')):
#                         topic = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
#                         if topic:
#                             key_points.append(topic)
           
#             logger.info(f"Successfully generated summary and {len(key_points)} key points for {file_name}")
#             print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", key_points)
#             return summary, key_points
           
#         except Exception as e:
#             logger.error(f"Error generating content with GPT-4: {str(e)}")
#             return self.format_error_message(file_name), []
   
#     def format_summary_response(self, response_text):
#         """Ensures proper HTML-like formatting in the summary response."""
#         import re
       
#         # Wrap paragraphs in <p> tags if not already wrapped
#         response_text = re.sub(r'^(?!<[p|b|u|ul|li])(.*?)$', r'<p>\1</p>', response_text, flags=re.MULTILINE)
       
#         # Ensure bold tags for section headers
#         section_headers = ['Summary Overview', 'Key Highlights', 'Detailed Insights']
#         for header in section_headers:
#             response_text = response_text.replace(header, f'<b>{header}</b>')
       
#         return response_text
 
#     def format_error_message(self, file_name):
#         """Returns a formatted error message."""
#         return f"""
#         <b>Summary Generation Error</b>
#         <p>Unable to generate a comprehensive summary for {file_name}</p>
       
#         <b>Possible Reasons</b>
#         <ul>
#             <li>Document may be too complex</li>
#             <li>Parsing issues encountered</li>
#             <li>Insufficient context extracted</li>
#         </ul>
#         """
 
 
#     def convert_video_to_audio(self, video_path):
#         """
#         Convert video files to audio format using MoviePy.
#         Returns the path to the created audio file.
       
#         Args:
#             video_path: Path to the video file
#         """
       
#         from moviepy import VideoFileClip
#         import logging
       
#         logger = logging.getLogger(__name__)
       
#         # Set up the audio directory
#         AUDIO_DIR = "audio_files"
#         os.makedirs(AUDIO_DIR, exist_ok=True)
       
#         try:
#             # Generate output path for the audio file
#             base_name = os.path.splitext(os.path.basename(video_path))[0]
#             output_path = os.path.join(AUDIO_DIR, f"{base_name}.mp3")
           
#             logger.info(f"Loading video: {video_path}")
#             video = VideoFileClip(video_path)
           
#             logger.info("Extracting audio...")
#             audio = video.audio
           
#             if audio is None:
#                 logger.warning(f"No audio stream found in video: {video_path}")
#                 return None
           
#             logger.info(f"Converting to MP3: {output_path}")
#             audio.write_audiofile(output_path, codec='mp3')
           
#             # Close files
#             audio.close()
#             video.close()
           
#             logger.info(f"Video-to-audio conversion complete: {output_path}")
#             return output_path
       
#         except Exception as e:
#             logger.error(f"Error during video-to-audio conversion: {str(e)}")
#             return None
#     def extract_text_from_video(self, file_path, user=None):
#         """
#         Extract text from video files by first converting to audio, then transcribing.
#         Returns the transcribed text.
       
#         Args:
#             file_path: Path to the video file
#             user: Django user object to get API token
#         """
       
#         logger = logging.getLogger(__name__)
#         print(f"Starting video-to-text extraction for: {file_path}")
       
#         try:
#             # First convert video to audio
#             print("Attempting to convert video to audio...")
#             audio_path = self.convert_video_to_audio(file_path)
           
#             if not audio_path:
#                 print(f"Failed to convert video to audio: {file_path}")
#                 logger.error(f"Failed to convert video to audio: {file_path}")
#                 return f"Error: Failed to extract audio from video file."
           
#             print(f"Successfully converted video to audio: {audio_path}")
           
#             # Then extract text from the resulting audio file
#             print(f"Starting audio transcription...")
#             extracted_text = self.extract_text_from_audio(audio_path, user=user)
           
#             print(f"Audio transcription completed, text length: {len(extracted_text) if extracted_text else 0}")
           
#             #Clean up the temporary audio file if needed
#             #Uncomment if you want to delete the audio file after processing
#             if os.path.exists(audio_path):
#                 os.remove(audio_path)
           
#             return extracted_text
       
#         except Exception as e:
#             print(f"Exception in extract_text_from_video: {str(e)}")
#             logger.error(f"Error extracting text from video: {str(e)}")
#             return f"Error transcribing video: {str(e)}"
 
#     def process_document_from_text(self, text, filename):
#         """
#         Process document directly from provided text.
#         For use with transcripts that are already extracted.
#         """
       
#         try:
#             # Clean the text
#             cleaned_text = self.clean_text(text)
           
#             # Create document record for processing
#             doc = {
#                 'text': cleaned_text,
#                 'name': filename
#             }
           
#             all_chunks = []
#             # Split text into chunks, with smaller size for transcripts
#             chunks = self.split_text_into_chunks(doc['text'], chunk_size=800, chunk_overlap=200)
           
#             for i, chunk in enumerate(chunks):
#                 all_chunks.append({
#                     'text': chunk,
#                     'source': doc['name'],
#                     'source_file': doc['name'],
#                     'chunk_id': i,
#                     'is_transcript': True  # Mark this as a transcript
#                 })
           
#             # Get embeddings for all chunks
#             text_chunks = [chunk['text'] for chunk in all_chunks]
#             print(f"Getting embeddings for {len(text_chunks)} text chunks")
#             embeddings = self.get_embeddings(text_chunks)
           
#             if not embeddings:
#                 raise ValueError("Failed to generate embeddings for document chunks")
           
#             # Create FAISS index
#             index = self.create_faiss_index(embeddings)
           
#             # Generate a unique session ID for this document
#             session_id = uuid.uuid4().hex
           
#             # Save the index and chunks
#             index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
           
#             print(f"Transcript processed successfully: {len(all_chunks)} chunks created")
           
#             return {
#                 'index_path': index_file,
#                 'metadata_path': pickle_file,
#                 'full_text': doc['text']
#             }
           
#         except Exception as e:
#             print(f"Error in process_document_from_text: {str(e)}")
#             raise
 
 
#     def extract_text_from_audio(self, file_path, user=None):
#         """
#         Extract text from audio files using Google's Gemini API.
#         For large files, chunks the audio before transcription.
#         Returns the transcribed text.
       
#         Args:
#             file_path: Path to the audio file
#             user: Django user object to get API token
#         """
 
#         logger = logging.getLogger(__name__)
       
#         # Set up the transcripts directory
#         TRANSCRIPT_DIR = "transcripts"
#         os.makedirs(TRANSCRIPT_DIR, exist_ok=True)
       
#         # Audio chunking settings
#         CHUNK_LENGTH_MINUTES = 5
#         CHUNK_THRESHOLD_MB = 20
#         DURATION_THRESHOLD_MINUTES = 10
#         MAX_RETRIES = 3
       
#         try:
#             # Initialize Gemini client
#             client = self.init_gemini(user)
#             if not client:
#                 raise ValueError("Failed to initialize Gemini client")
           
#             # Generate a unique base filename for this transcription
#             base_filename = f"audio_transcript_{uuid.uuid4().hex}"
#             full_transcript = ""
#             temp_files = []
           
#             # Check if audio file needs chunking
#             file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
           
#             # Load audio file to check duration using MoviePy
#             audio_clip = AudioFileClip(file_path)
#             duration_minutes = audio_clip.duration / 60  # Convert seconds to minutes
           
#             logger.info(f"Audio file: {file_path}, Size: {file_size_mb:.2f}MB, Duration: {duration_minutes:.2f} minutes")
           
#             # Determine if chunking is needed
#             need_chunking = False
#             if file_size_mb > CHUNK_THRESHOLD_MB or duration_minutes > DURATION_THRESHOLD_MINUTES:
#                 need_chunking = True
#                 logger.info(f"File exceeds size/duration thresholds, will chunk for processing")
           
#             if need_chunking:
#                 # Calculate parameters for chunking
#                 chunk_length_seconds = CHUNK_LENGTH_MINUTES * 60
#                 overlap_seconds = 5  # 5 second overlap
#                 total_chunks = int(audio_clip.duration / chunk_length_seconds) + (1 if audio_clip.duration % chunk_length_seconds > 0 else 0)
               
#                 logger.info(f"Splitting audio into {total_chunks} chunks of {CHUNK_LENGTH_MINUTES} minutes each")
               
#                 for i in range(total_chunks):
#                     # Calculate start and end times with overlap
#                     start_time = max(0, i * chunk_length_seconds - overlap_seconds) if i > 0 else 0
#                     end_time = min(audio_clip.duration, (i + 1) * chunk_length_seconds + overlap_seconds)
                   
#                     # Create a subclip
#                     chunk = audio_clip.subclipped(start_time, end_time)
                   
#                     # Export the chunk to a temporary file
#                     chunk_filename = f"{os.path.splitext(file_path)[0]}_chunk_{i}.mp3"
#                     chunk.write_audiofile(chunk_filename, codec='mp3')
#                     temp_files.append(chunk_filename)
#                     chunk.close()  # Close the subclip to free resources
                   
#                     # Process the chunk
#                     logger.info(f"Processing chunk {i + 1}/{total_chunks}")
                   
#                     # Transcribe each chunk with retry logic
#                     retry_count = 0
#                     success = False
                   
#                     while retry_count < MAX_RETRIES and not success:
#                         try:
#                             myfile = client.files.upload(file=chunk_filename)
 
#                             prompt = """You are a transcription and translation assistant. Your job is to:
#                             1. Transcribe spoken content from audio.
#                             2. Translate everything fully into English.
#                             3. Output only the translated English version.
#                             4. If other languages are spoken, translate them fully into English.
#                             5. Identify intelligently the Moderator and Respondent(s) in the conversation.
#                             6. If a speaker is continously speaking, do not repeat the speaker label.
 
#                             **Formatting Instructions:**
#                             - Format each line like this:
#                             M: [Moderator's translated English speech]
#                             R: [Respondent's translated English speech]
#                             - Use "M:" for the moderator and "R:" for any respondent. Do not attempt to label or distinguish individual respondents.
#                             - Do not include names or speaker numbers.
#                             - Do not include the original language (e.g., Hindi).
#                             - Do not paraphrase or summarize. Translate every spoken sentence accurately.
#                             - Output should be clean, easy-to-read translated English dialogue.
#                             - Do not use JSON or structured formats.
 
#                             **Example Output:**
#                             M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
#                             R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
#                             R: I live with my parents and younger brother.
 
#                             Only output the above format. Nothing else."""
                               
#                             # Generate the transcription using new API
#                             response = client.models.generate_content(
#                                 model="gemini-2.5-flash",
#                                 contents=[prompt, myfile]
#                             )
#                             chunk_transcription = response.text
                           
#                             # Add a separator between chunks
#                             if i > 0:
#                                 full_transcript += "\n\n--- NEXT SEGMENT ---\n\n"
                           
#                             full_transcript += chunk_transcription
#                             success = True
                           
#                         except Exception as e:
#                             retry_count += 1
#                             logger.error(f"Error in chunk transcription (attempt {retry_count}): {str(e)}")
#                             if retry_count >= MAX_RETRIES:
#                                 full_transcript += f"\n\n[Error transcribing segment {i + 1}: {str(e)}]\n\n"
#                             else:
#                                 time.sleep(2)  # Wait before retry
               
#                 # Close the main audio clip
#                 audio_clip.close()
                   
#             else:
#                 # Process the entire file at once if it's small enough
#                 logger.info("Processing audio file in a single pass")
#                 audio_clip.close()  # Close the clip since we're not using it for chunking
               
#                 myfile = client.files.upload(file=file_path)
               
#                 # Define the prompt for transcription
#                 prompt = """You are a transcription and translation assistant. Your job is to:
#                             1. Transcribe spoken content from audio.
#                             2. Translate everything fully into English.
#                             3. Output only the translated English version.
#                             4. If other languages are spoken, translate them fully into English.
#                             5. Identify intelligently the Moderator and Respondent(s) in the conversation.
#                             6. If a speaker is continously speaking, do not repeat the speaker label.
 
#                             **Formatting Instructions:**
#                             - Format each line like this:
#                             M: [Moderator's translated English speech]
#                             R: [Respondent's translated English speech]
#                             - Use "M:" for the moderator and "R:" for any respondent. Do not attempt to label or distinguish individual respondents.
#                             - Do not include names or speaker numbers.
#                             - Do not include the original language (e.g., Hindi).
#                             - Do not paraphrase or summarize. Translate every spoken sentence accurately.
#                             - Output should be clean, easy-to-read translated English dialogue.
#                             - Do not use JSON or structured formats.
 
#                             **Example Output:**
#                             M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
#                             R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
#                             R: I live with my parents and younger brother.
 
#                             Only output the above format. Nothing else.
#                             """
               
#                 # Generate the transcription
#                 response = client.models.generate_content(
#                     model="gemini-2.5-flash",
#                     contents=[prompt, myfile]
#                 )
#                 full_transcript = response.text
           
#             # Save the full transcription as a text file
#             txt_path = os.path.join(TRANSCRIPT_DIR, f"{base_filename}.txt")
#             with open(txt_path, "w", encoding="utf-8") as f:
#                 f.write(full_transcript)
           
#             # Clean up temporary chunk files
#             for temp_file in temp_files:
#                 try:
#                     if os.path.exists(temp_file):
#                         os.remove(temp_file)
#                         logger.info(f"Removed temporary file: {temp_file}")
#                 except Exception as e:
#                     logger.error(f"Error removing file {temp_file}: {str(e)}")
           
#             # Extract the text from the saved text file using the existing method
#             extracted_text = self.extract_text_from_txt(txt_path)
           
#             return extracted_text
       
#         except Exception as e:
#             logger.error(f"Error transcribing audio file: {str(e)}")
#             return f"Error transcribing audio: {str(e)}"
           
#     def extract_text_from_file(self, file_path, user=None):
#         """Extract text from different file types."""
#         from pathlib import Path
#         ext = Path(file_path).suffix.lower()
#         if ext == '.pdf':
#             return self.extract_text_from_pdf(file_path)
#         elif ext == '.docx':
#             return self.extract_text_from_docx(file_path)
#         elif ext == '.pptx':
#             return self.extract_text_from_pptx(file_path)
#         elif ext == '.xlsx':
#             return self.extract_text_from_xlsx(file_path)
#         elif ext == '.txt':
#             return self.extract_text_from_txt(file_path)
       
#         elif ext in ['.mp3', '.wav', '.mpeg']:
#             return self.extract_text_from_audio(file_path, user)
#         else:
#             return ""
 
#     def extract_text_from_pdf(self, file_path):
#         import fitz
#         text = ""
#         with fitz.open(file_path) as doc:
#             for page in doc:
#                 text += page.get_text()
#         return text
 
#     def extract_text_from_docx(self, file_path):
#         import docx
#         doc = docx.Document(file_path)
#         text = [para.text for para in doc.paragraphs]
#         return "\n".join(text)
 
#     def extract_text_from_pptx(self, file_path):
#         import pptx
#         pres = pptx.Presentation(file_path)
#         text = []
#         for slide in pres.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text.append(shape.text)
#         return "\n".join(text)
 
#     def extract_text_from_xlsx(self, file_path):
#         import openpyxl
#         wb = openpyxl.load_workbook(file_path, data_only=True)
#         text = []
#         for sheet in wb.sheetnames:
#             ws = wb[sheet]
#             text.append(f"Sheet: {sheet}")
#             for row in ws.iter_rows(values_only=True):
#                 row_text = [str(cell) if cell is not None else "" for cell in row]
#                 text.append(" | ".join(row_text))
#         return "\n".join(text)
   
#     def extract_text_from_txt(self, file_path):
#         """Extract text from a plain text file."""
#         try:
#             # Try multiple encodings in case UTF-8 fails
#             encodings = ['utf-8', 'latin-1', 'cp1252']
           
#             for encoding in encodings:
#                 try:
#                     with open(file_path, 'r', encoding=encoding) as file:
#                         text = file.read()
#                     return text
#                 except UnicodeDecodeError:
#                     continue
           
#             # If all encodings fail, try binary mode and decode with errors='replace'
#             with open(file_path, 'rb') as file:
#                 binary_data = file.read()
#                 return binary_data.decode('utf-8', errors='replace')
               
#         except Exception as e:
#             print(f"Error extracting text from txt file: {str(e)}")
#             return f"Error extracting text: {str(e)}"
 
#     def clean_text(self, text):
#         """Clean and normalize extracted text."""
#         import re
#         import unicodedata
       
#         # Remove excessive whitespace while preserving paragraph breaks
#         text = re.sub(r'\s+', ' ', text)
#         text = re.sub(r'\n\s*\n', '\n\n', text)
#         text = text.strip()
       
#         # Remove any control characters
#         text = ''.join(char for char in text if unicodedata.category(char)[0] != 'C')
       
#         return text
 
 

class DocumentProcessingMixin:
    from google import genai
    def _parse_json_response(self, response_content, fallback_message="Error processing response"):
        """Helper method to safely parse JSON responses from LLM"""
        try:
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing failed: {str(e)}")
            return {
                "content": fallback_message,
                "error": True,
                "format_type": "error"
            }

    def init_gemini(self, user=None):
        """Initialize Gemini 2.5 Flash client"""
        try:
            # Get the user's Gemini API token
            gemini_api_key = None
            if user:
                try:
                    from .models import UserAPITokens  # Adjust import based on your models location
                    user_api_tokens = UserAPITokens.objects.get(user=user)
                    gemini_api_key = user_api_tokens.gemini_token
                    
                    if not gemini_api_key:
                        logger.warning(f"No Gemini API token found for user {user.username}")
                        # Fallback to environment variable
                        gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                        if not gemini_api_key:
                            raise ValueError("Gemini API token is required for processing")
                        
                except Exception as e:
                    logger.error(f"Error getting API tokens for user {user.username}: {str(e)}")
                    # Fallback to environment variable
                    gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                    if not gemini_api_key:
                        raise ValueError("Gemini API token is required for processing")
            else:
                # Fallback to environment variable if no user provided
                gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                if not gemini_api_key:
                    raise ValueError("Gemini API token is required for processing")
            
            from google import genai
            gemini_client = genai.Client(api_key=gemini_api_key)
            return gemini_client
        except Exception as e:
            logger.error(f"Error initializing Gemini: {str(e)}")
            return None

    def detect_document_structure(self, text: str, client) -> str:
        """Analyze document structure to determine if it has clear headings/titles"""
        try:
            system_message = "You are an expert at analyzing document structure. Determine if documents have clear headings and sections."
           
            user_prompt = f"""
            Analyze this document text and determine if it has clear structure with headings, titles, or sections:
           
            Text: {text[:2000]}
           
            Look for:
            - Clear headings (numbered sections, bold titles, etc.)
            - Section breaks
            - Structured format
            - Table of contents references
           
            Respond with only: "STRUCTURED" or "UNSTRUCTURED"
            """
           
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.1,
                max_tokens=10
            )
           
            response_text = completion.choices[0].message.content.strip().upper()
            return response_text
        except Exception as e:
            logger.error(f"Error detecting document structure: {str(e)}")
            return "UNSTRUCTURED"  # Default to unstructured if analysis fails
 
    def generate_topics_from_content(self, text: str, client) -> List[str]:
        """Generate key topics from document content when no clear headings exist"""
        try:
            system_message = "You are an expert at identifying main topics and themes from documents. Generate concise, specific topic phrases."
           
            user_prompt = f"""
            Analyze this document and identify 5-7 main topics or themes discussed:
           
            Text: {text[:6000]}
           
            Generate concise topic phrases (3-5 words each) that represent the main concepts, ideas, or subjects covered in this document.
           
            IMPORTANT:
            - Focus on the most important themes and concepts
            - Use clear, descriptive phrases
            - Make topics specific to the document content
            - Avoid generic terms
            - Each topic should be 3-5 words
           
            Format as a simple numbered list:
            1. Topic phrase
            2. Topic phrase
            ...
            """
           
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.4,
                max_tokens=500
            )
           
            response_text = completion.choices[0].message.content.strip()
            topics = []
           
            # Parse the numbered list
            lines = response_text.split('\n')
            for line in lines:
                line = line.strip()
                if line and (line[0].isdigit() or line.startswith('•') or line.startswith('-')):
                    # Extract topic after number/bullet
                    topic = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
                    if topic:
                        topics.append(topic)
           
            return topics[:7]  # Limit to 7 topics
        except Exception as e:
            logger.error(f"Error generating topics: {str(e)}")
            return []
 
    def extract_headings_from_structured_doc(self, text: str, client) -> List[str]:
        """Extract actual headings from structured documents"""
        try:
            system_message = "You are an expert at extracting headings and section titles from structured documents. Extract only exact headings that appear in the document."
           
            user_prompt = f"""
            Extract the main headings and section titles from this structured document:
           
            Text: {text[:8000]}
           
            IMPORTANT:
            - Extract ONLY actual headings, titles, or section names that appear in the document
            - Use exact words from the document (3-6 words per heading)
            - Look for numbered sections, bold text, capitalized headings
            - Do not create new phrases - use exact text from document
           
            Format as a simple numbered list:
            1. Exact heading from document
            2. Exact heading from document
            ...
            """
           
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.2,
                max_tokens=800
            )
           
            response_text = completion.choices[0].message.content.strip()
            headings = []
           
            # Parse the numbered list
            lines = response_text.split('\n')
            for line in lines:
                line = line.strip()
                if line and (line[0].isdigit() or line.startswith('•') or line.startswith('-')):
                    # Extract heading after number/bullet
                    heading = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
                    if heading:
                        headings.append(heading)
           
            return headings[:7]  # Limit to 7 headings
        except Exception as e:
            logger.error(f"Error extracting headings: {str(e)}")
            return []
 
    def generate_summary(self, text, file_name, user=None):
        """
        REPLACED: Generate summary and key points using GPT-4 with smart topic detection
        Returns: (summary_text, key_points_list)
        """
        if not text.strip():
            return "No extractable text found in the document.", []
       
        try:
           
            # First, detect if document has structure
            doc_structure = self.detect_document_structure(text, client)
           
            # Generate summary
            summary_system_message = """You are an expert document analyzer. Provide comprehensive document summaries with proper HTML-like formatting."""
           
            summary_user_prompt = f"""
            Please analyze the following document '{file_name}' and provide a comprehensive quick summary (8-10 lines):
           
            Text: {text[:8000]}
           
            IMPORTANT FOR SUMMARY:
            - Must be 8-10 lines providing comprehensive coverage
            - FOCUS ON EXECUTION AND VALIDATION: Emphasize practical implementation steps, validation methods, testing approaches, and measurable outcomes
            - DESCRIBE MORE CLEARLY: Use specific language, avoid vague terms, provide clear explanations of processes and concepts
            - MENTION SPECIFIC SECTORS: Identify and name the exact industries, market segments, business sectors, or domains discussed
            - BE MORE CONCRETE: Include specific numbers, percentages, timeframes, locations, company names, product names, or measurable data points
            - COMPREHENSIVE AND APPLIED PICTURE: Show how the analysis translates to real-world applications, practical benefits, and actionable insights
            - Include key findings, main themes, important data/statistics if present
            - Cover methodology, results, conclusions, and recommendations if applicable
            - Explain the context, significance, and implications with specific examples
            - Make it thorough but still readable and well-structured
           
            ### Expected Response Format:
            <b>Summary Overview</b>
            <p>High-level introduction to the document's main theme</p>
 
            <b>Key Highlights</b>
            <ul>
                <li>First major insight</li>
                <li>Second major insight</li>
                <li>Third major insight</li>
            </ul>
 
            <b>Detailed Insights</b>
            <p>Expanded explanation of the document's core content and significance</p>
            """
           
            summary_completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": summary_system_message},
                    {"role": "user", "content": summary_user_prompt}
                ],
                temperature=0.4,
                max_tokens=2000
            )
           
            summary = summary_completion.choices[0].message.content.strip()
           
            # Format the summary response to ensure proper HTML-like formatting
            summary = self.format_summary_response(summary)
           
            # Generate key points based on document structure
            if doc_structure == "STRUCTURED":
                key_points = self.extract_headings_from_structured_doc(text, client)
            else:
                key_points = self.generate_topics_from_content(text, client)
           
            # Fallback if no key points were generated
            if not key_points:
                fallback_system_message = "You are an expert at identifying main topics and themes from documents."
               
                fallback_user_prompt = f"""
                From this document, identify 5 main topics or themes:
               
                Text: {text[:4000]}
               
                Provide 5 short phrases (3-4 words each) that capture the main ideas.
                Format as:
                1. Topic phrase
                2. Topic phrase
                ...
                """
               
                fallback_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": fallback_system_message},
                        {"role": "user", "content": fallback_user_prompt}
                    ],
                    temperature=0.4,
                    max_tokens=500
                )
               
                fallback_response = fallback_completion.choices[0].message.content.strip()
                lines = fallback_response.split('\n')
                key_points = []
                for line in lines:
                    line = line.strip()
                    if line and (line[0].isdigit() or line.startswith('•')):
                        topic = line.split('.', 1)[-1].strip() if '.' in line else line[1:].strip()
                        if topic:
                            key_points.append(topic)
           
            logger.info(f"Successfully generated summary and {len(key_points)} key points for {file_name}")
            print("!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!", key_points)
            return summary, key_points
           
        except Exception as e:
            logger.error(f"Error generating content with GPT-4: {str(e)}")
            return self.format_error_message(file_name), []
    
    def format_summary_response(self, response_text):
        """Ensures proper HTML-like formatting in the summary response."""
        import re
        
        # Wrap paragraphs in <p> tags if not already wrapped
        response_text = re.sub(r'^(?!<[p|b|u|ul|li])(.*?)$', r'<p>\1</p>', response_text, flags=re.MULTILINE)
        
        # Ensure bold tags for section headers
        section_headers = ['Summary Overview', 'Key Highlights', 'Detailed Insights']
        for header in section_headers:
            response_text = response_text.replace(header, f'<b>{header}</b>')
        
        return response_text

    def format_error_message(self, file_name):
        """Returns a formatted error message."""
        return f"""
        <b>Summary Generation Error</b>
        <p>Unable to generate a comprehensive summary for {file_name}</p>
        
        <b>Possible Reasons</b>
        <ul>
            <li>Document may be too complex</li>
            <li>Parsing issues encountered</li>
            <li>Insufficient context extracted</li>
        </ul>
        """


    def convert_video_to_audio(self, video_path):
        """
        Convert video files to audio format using MoviePy.
        Returns the path to the created audio file.
        
        Args:
            video_path: Path to the video file
        """
        
        from moviepy import VideoFileClip
        import logging
        
        logger = logging.getLogger(__name__)
        
        # Set up the audio directory
        AUDIO_DIR = "audio_files"
        os.makedirs(AUDIO_DIR, exist_ok=True)
        
        try:
            # Generate output path for the audio file
            base_name = os.path.splitext(os.path.basename(video_path))[0]
            output_path = os.path.join(AUDIO_DIR, f"{base_name}.mp3")
            
            logger.info(f"Loading video: {video_path}")
            video = VideoFileClip(video_path)
            
            logger.info("Extracting audio...")
            audio = video.audio
            
            if audio is None:
                logger.warning(f"No audio stream found in video: {video_path}")
                return None
            
            logger.info(f"Converting to MP3: {output_path}")
            audio.write_audiofile(output_path, codec='mp3')
            
            # Close files
            audio.close()
            video.close()
            
            logger.info(f"Video-to-audio conversion complete: {output_path}")
            return output_path
        
        except Exception as e:
            logger.error(f"Error during video-to-audio conversion: {str(e)}")
            return None
    def extract_text_from_video(self, file_path, user=None):
        """
        Extract text from video files by first converting to audio, then transcribing.
        Returns the transcribed text.
        
        Args:
            file_path: Path to the video file
            user: Django user object to get API token
        """
        
        logger = logging.getLogger(__name__)
        print(f"Starting video-to-text extraction for: {file_path}")
        
        try:
            # First convert video to audio
            print("Attempting to convert video to audio...")
            audio_path = self.convert_video_to_audio(file_path)
            
            if not audio_path:
                print(f"Failed to convert video to audio: {file_path}")
                logger.error(f"Failed to convert video to audio: {file_path}")
                return f"Error: Failed to extract audio from video file."
            
            print(f"Successfully converted video to audio: {audio_path}")
            
            # Then extract text from the resulting audio file
            print(f"Starting audio transcription...")
            extracted_text = self.extract_text_from_audio(audio_path, user=user)
            
            print(f"Audio transcription completed, text length: {len(extracted_text) if extracted_text else 0}")
            
            #Clean up the temporary audio file if needed
            #Uncomment if you want to delete the audio file after processing
            if os.path.exists(audio_path):
                os.remove(audio_path)
            
            return extracted_text
        
        except Exception as e:
            print(f"Exception in extract_text_from_video: {str(e)}")
            logger.error(f"Error extracting text from video: {str(e)}")
            return f"Error transcribing video: {str(e)}"

    def process_document_from_text(self, text, filename):
        """
        Process document directly from provided text.
        For use with transcripts that are already extracted.
        """
        
        try:
            # Clean the text
            cleaned_text = self.clean_text(text)
            
            # Create document record for processing
            doc = {
                'text': cleaned_text,
                'name': filename
            }
            
            all_chunks = []
            # Split text into chunks, with smaller size for transcripts
            chunks = self.split_text_into_chunks(doc['text'], chunk_size=800, chunk_overlap=200)
            
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    'text': chunk,
                    'source': doc['name'],
                    'source_file': doc['name'],
                    'chunk_id': i,
                    'is_transcript': True  # Mark this as a transcript
                })
            
            # Get embeddings for all chunks
            text_chunks = [chunk['text'] for chunk in all_chunks]
            print(f"Getting embeddings for {len(text_chunks)} text chunks")
            embeddings = self.get_embeddings(text_chunks)
            
            if not embeddings:
                raise ValueError("Failed to generate embeddings for document chunks")
            
            # Create FAISS index
            index = self.create_faiss_index(embeddings)
            
            # Generate a unique session ID for this document
            session_id = uuid.uuid4().hex
            
            # Save the index and chunks
            index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
            
            print(f"Transcript processed successfully: {len(all_chunks)} chunks created")
            
            return {
                'index_path': index_file,
                'metadata_path': pickle_file,
                'full_text': doc['text']
            }
            
        except Exception as e:
            print(f"Error in process_document_from_text: {str(e)}")
            raise


    def extract_text_from_audio(self, file_path, user=None):
        """
        Extract text from audio files using Google's Gemini API.
        For large files, chunks the audio before transcription.
        Returns the transcribed text.
        
        Args:
            file_path: Path to the audio file
            user: Django user object to get API token
        """

        logger = logging.getLogger(__name__)
        
        # Set up the transcripts directory
        TRANSCRIPT_DIR = "transcripts"
        os.makedirs(TRANSCRIPT_DIR, exist_ok=True)
        
        # Audio chunking settings
        CHUNK_LENGTH_MINUTES = 5
        CHUNK_THRESHOLD_MB = 20
        DURATION_THRESHOLD_MINUTES = 10
        MAX_RETRIES = 3
        
        try:
            # Initialize Gemini client
            client = self.init_gemini(user)
            if not client:
                raise ValueError("Failed to initialize Gemini client")
            
            # Generate a unique base filename for this transcription
            base_filename = f"audio_transcript_{uuid.uuid4().hex}"
            full_transcript = ""
            temp_files = []
            
            # Check if audio file needs chunking
            file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
            
            # Load audio file to check duration using MoviePy
            audio_clip = AudioFileClip(file_path)
            duration_minutes = audio_clip.duration / 60  # Convert seconds to minutes
            
            logger.info(f"Audio file: {file_path}, Size: {file_size_mb:.2f}MB, Duration: {duration_minutes:.2f} minutes")
            
            # Determine if chunking is needed
            need_chunking = False
            if file_size_mb > CHUNK_THRESHOLD_MB or duration_minutes > DURATION_THRESHOLD_MINUTES:
                need_chunking = True
                logger.info(f"File exceeds size/duration thresholds, will chunk for processing")
            
            if need_chunking:
                # Calculate parameters for chunking
                chunk_length_seconds = CHUNK_LENGTH_MINUTES * 60
                overlap_seconds = 5  # 5 second overlap
                total_chunks = int(audio_clip.duration / chunk_length_seconds) + (1 if audio_clip.duration % chunk_length_seconds > 0 else 0)
                
                logger.info(f"Splitting audio into {total_chunks} chunks of {CHUNK_LENGTH_MINUTES} minutes each")
                
                for i in range(total_chunks):
                    # Calculate start and end times with overlap
                    start_time = max(0, i * chunk_length_seconds - overlap_seconds) if i > 0 else 0
                    end_time = min(audio_clip.duration, (i + 1) * chunk_length_seconds + overlap_seconds)
                    
                    # Create a subclip
                    chunk = audio_clip.subclipped(start_time, end_time)
                    
                    # Export the chunk to a temporary file
                    chunk_filename = f"{os.path.splitext(file_path)[0]}_chunk_{i}.mp3"
                    chunk.write_audiofile(chunk_filename, codec='mp3')
                    temp_files.append(chunk_filename)
                    chunk.close()  # Close the subclip to free resources
                    
                    # Process the chunk
                    logger.info(f"Processing chunk {i + 1}/{total_chunks}")
                    
                    # Transcribe each chunk with retry logic
                    retry_count = 0
                    success = False
                    
                    while retry_count < MAX_RETRIES and not success:
                        try:
                            myfile = client.files.upload(file=chunk_filename)
                            
                            
                            

                            prompt = """You are a transcription and translation assistant. Your job is to:
                            1. Transcribe spoken content from audio.
                            2. Translate everything fully into English.
                            3. Output only the translated English version.
                            4. If other languages are spoken, translate them fully into English.
                            5. Identify intelligently the Moderator and Respondent(s) in the conversation.
                            6. If a speaker is continously speaking, do not repeat the speaker label.

                            **Formatting Instructions:**
                            - Format each line like this:
                            M: [Moderator's translated English speech]
                            R: [Respondent's translated English speech]
                            - Use "M:" for the moderator and "R:" for any respondent. Do not attempt to label or distinguish individual respondents.
                            - Do not include names or speaker numbers.
                            - Do not include the original language (e.g., Hindi).
                            - Do not paraphrase or summarize. Translate every spoken sentence accurately.
                            - Output should be clean, easy-to-read translated English dialogue.
                            - Do not use JSON or structured formats.

                            **Example Output:**
                            M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
                            R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
                            R: I live with my parents and younger brother.

                            Only output the above format. Nothing else."""
                                
                            # Generate the transcription using new API
                            response = client.models.generate_content(
                                model="gemini-2.5-flash",
                                contents=[prompt, myfile]
                            )
                            chunk_transcription = response.text
                            
                            # Add a separator between chunks
                            if i > 0:
                                full_transcript += "\n\n--- NEXT SEGMENT ---\n\n"
                            
                            full_transcript += chunk_transcription
                            success = True
                            
                        except Exception as e:
                            retry_count += 1
                            logger.error(f"Error in chunk transcription (attempt {retry_count}): {str(e)}")
                            if retry_count >= MAX_RETRIES:
                                full_transcript += f"\n\n[Error transcribing segment {i + 1}: {str(e)}]\n\n"
                            else:
                                time.sleep(2)  # Wait before retry
                
                # Close the main audio clip
                audio_clip.close()
                    
            else:
                # Process the entire file at once if it's small enough
                logger.info("Processing audio file in a single pass")
                audio_clip.close()  # Close the clip since we're not using it for chunking
                
                myfile = client.files.upload(file=file_path)
                
                # Define the prompt for transcription
                prompt = """You are a transcription and translation assistant. Your job is to:
                            1. Transcribe spoken content from audio.
                            2. Translate everything fully into English.
                            3. Output only the translated English version.
                            4. If other languages are spoken, translate them fully into English.
                            5. Identify intelligently the Moderator and Respondent(s) in the conversation.
                            6. If a speaker is continously speaking, do not repeat the speaker label.
 
                            **Formatting Instructions:**
                            - Format each line like this:
                            M: [Moderator's translated English speech]
                            R: [Respondent's translated English speech]
                            - Use "M:" for the moderator and "R:" for any respondent. Do not attempt to label or distinguish individual respondents.
                            - Do not include names or speaker numbers.
                            - Do not include the original language (e.g., Hindi).
                            - Do not paraphrase or summarize. Translate every spoken sentence accurately.
                            - Output should be clean, easy-to-read translated English dialogue.
                            - Do not use JSON or structured formats.
 
                            **Example Output:**
                            M: Thank you all for joining today's session. Let's begin by getting to know each other a little.
                            R: Sure. My name is Khushboo. I live in Lucknow and work as a school teacher.
                            R: I live with my parents and younger brother.
 
                            Only output the above format. Nothing else.
                            """
                
                # Generate the transcription
                response = client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[prompt, myfile]
                )
                full_transcript = response.text
            
            # Save the full transcription as a text file
            txt_path = os.path.join(TRANSCRIPT_DIR, f"{base_filename}.txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(full_transcript)
            
            # Clean up temporary chunk files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                        logger.info(f"Removed temporary file: {temp_file}")
                except Exception as e:
                    logger.error(f"Error removing file {temp_file}: {str(e)}")
            
            # Extract the text from the saved text file using the existing method
            extracted_text = self.extract_text_from_txt(txt_path)
            
            return extracted_text
        
        except Exception as e:
            logger.error(f"Error transcribing audio file: {str(e)}")
            return f"Error transcribing audio: {str(e)}"

    def extract_text_from_file(self, file_path, user=None):
        """Extract text from different file types."""
        from pathlib import Path
        ext = Path(file_path).suffix.lower()
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif ext == '.txt':
            return self.extract_text_from_txt(file_path)
        
        elif ext in ['.mp3', '.wav', '.mpeg']:
            return self.extract_text_from_audio(file_path, user)
        else:
            return ""

    def extract_text_from_pdf(self, file_path):
        import fitz
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    def extract_text_from_docx(self, file_path):
        import docx
        doc = docx.Document(file_path)
        text = [para.text for para in doc.paragraphs]
        return "\n".join(text)

    def extract_text_from_pptx(self, file_path):
        import pptx
        pres = pptx.Presentation(file_path)
        text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_text_from_xlsx(self, file_path):
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                text.append(" | ".join(row_text))
        return "\n".join(text)
    
    def extract_text_from_txt(self, file_path):
        """Extract text from a plain text file."""
        try:
            # Try multiple encodings in case UTF-8 fails
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    return text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try binary mode and decode with errors='replace'
            with open(file_path, 'rb') as file:
                binary_data = file.read()
                return binary_data.decode('utf-8', errors='replace')
                
        except Exception as e:
            print(f"Error extracting text from txt file: {str(e)}")
            return f"Error extracting text: {str(e)}"

    def clean_text(self, text):
        """Clean and normalize extracted text."""
        import re
        import unicodedata
        
        # Remove excessive whitespace while preserving paragraph breaks
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = text.strip()
        
        # Remove any control characters
        text = ''.join(char for char in text if unicodedata.category(char)[0] != 'C')
        
        return text




 
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def generate_question_from_topic(request):
    topic = request.data.get('topic', '')
    user = request.user
    if not topic:
        return Response({'error': 'Topic is required'}, status=400)
    try:
        prompt = f"""
        Convert the following topic into a clear, insightful question that would help someone explore this subject in depth.
        Topic: {topic}
        Guidelines:
        - Make the question specific and engaging
        - 8-15 words
        - End with a question mark
        - Avoid generic questions
        Format: Just the question text.
        """
        completion = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert at converting topics into insightful questions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=50
        )
        question = completion.choices[0].message.content.strip()
        if not question.endswith('?'):
            question += '?'
        return Response({'question': question}, status=200)
    except Exception as e:
        logger.error(f"Error generating question from topic: {str(e)}")
        return Response({'error': str(e)}, status=500)
class ConsolidatedSummaryView(DocumentProcessingMixin, APIView):
    parser_classes = (JSONParser,)



    def load_faiss_index(self, session_id):
    
        try:
            index_file = f"faiss_indexes/{session_id}_index.faiss"
            pickle_file = f"faiss_indexes/{session_id}_chunks.pkl"
            index = faiss.read_index(index_file)
            with open(pickle_file, "rb") as f:
                chunks = pickle.load(f)
            return index, chunks
        except Exception as e:
            print(f"Error loading FAISS index: {str(e)}")
            return None, []

    def _parse_json_response(self, response_content):
        """Parse JSON response with fallback handling"""
        try:
            import json
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {str(e)}")
            # Return a fallback structure
            return {
                "summary": response_content,  # Use raw content as summary
                "error": "JSON parsing failed"
            }

    def format_summary_response(self, response_text):
        """Ensures proper HTML-like formatting in the summary response."""
        # Wrap paragraphs in <p> tags if not already wrapped
        response_text = re.sub(r'^(?!<[p|b|u|ul|li])(.*?)$', r'<p>\1</p>', response_text, flags=re.MULTILINE)
        
        # Ensure bold tags for section headers
        section_headers = ['Summary Overview', 'Key Highlights', 'Detailed Insights']
        for header in section_headers:
            response_text = response_text.replace(header, f'<b>{header}</b>')
        
        return response_text

    def hierarchical_summary(self, chunks, group_size=10):
        """
        Generates a consolidated summary by grouping the chunk texts, summarizing each group,
        and then summarizing the aggregated group summaries.
        Enhanced to track document sources for better consolidated summaries.
        """
        import tiktoken
        encoding = tiktoken.get_encoding("cl100k_base")
            
        # Organize chunks by source document first
        chunks_by_source = {}
        for chunk in chunks:
            source = chunk.get('source', 'Unknown Source')
            if source not in chunks_by_source:
                chunks_by_source[source] = []
            chunks_by_source[source].append(chunk)
            
        # Process each document separately at first
        all_group_summaries = []
        for source, source_chunks in chunks_by_source.items():
            group_summaries = []
            # Group chunks to avoid exceeding token limits
            for i in range(0, len(source_chunks), group_size):
                group_text = "\n\n".join(chunk['text'] for chunk in source_chunks[i:i+group_size])
                    
                # Ensure we don't exceed token limits for the API call
                tokens = encoding.encode(group_text)
                if len(tokens) > 8000:
                    group_text = encoding.decode(tokens[:8000])
                    
                # Generate summary for this group
                try:
                    # Updated system message for JSON response
                    system_message = """You are a document summarization expert.
                    You must respond in JSON format with the following structure:
                    {
                        "summary": "Your summary with HTML formatting",
                        "source_document": "document_name",
                        "group_number": "group_identifier"
                    }"""

                    group_prompt = f"""
                    Summarize the following text from document '{source}'. Cover all key details:

                    {group_text}
                    
                    CRITICAL: Your response MUST be a valid JSON object with this structure:
                    {{
                        "summary": "Your summary covering all key details",
                        "source_document": "{source}",
                        "group_number": "group_{i//group_size}"
                    }}
                    """

                    completion = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": system_message},
                            {"role": "user", "content": group_prompt}
                        ],
                        response_format={"type": "json_object"},  # Force JSON response
                        temperature=0.3,
                        max_tokens=500
                    )

                    # Parse JSON response
                    json_response = self._parse_json_response(completion.choices[0].message.content)
                    group_summary = json_response.get("summary", f"Summary generation error for group {i//group_size}.")
                    group_summaries.append(f"[From {source}] {group_summary}")
                except Exception as e:
                    print(f"Error generating group summary for {source}: {str(e)}")
                    # Add a placeholder if summarization fails
                    group_summaries.append(f"[From {source}] Summary generation error for chunk group {i//group_size}.")
            
            # Get document-level summary first
            document_summary = "\n\n".join(group_summaries)
            try:
                # Updated system message for document summary
                system_message = """You are a document summarization expert.
                You must respond in JSON format with the following structure:
                {
                    "summary": "Your cohesive document summary",
                    "document_name": "source_document",
                    "summary_type": "document_level"
                }"""

                doc_prompt = f"""
                Combine these partial summaries into a cohesive summary for document '{source}':

                {document_summary}
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "summary": "Your cohesive summary for the document",
                    "document_name": "{source}",
                    "summary_type": "document_level"
                }}
                """

                doc_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": doc_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.3,
                    max_tokens=800
                )

                # Parse JSON response
                json_response = self._parse_json_response(doc_completion.choices[0].message.content)
                doc_final_summary = f"<b>Document: {source}</b>\n{json_response.get('summary', document_summary)}"
                all_group_summaries.append(doc_final_summary)
            except Exception as e:
                print(f"Error generating document summary for {source}: {str(e)}")
                all_group_summaries.append(f"<b>Document: {source}</b>\n{document_summary}")
        
        # Combine all document summaries for a final consolidated summary
        aggregated_summary = "\n\n".join(all_group_summaries)
        
        # Generate final consolidated summary
        try:
            # Updated system message for final consolidated summary
            system_message = """You are a document summarization expert.
            You must respond in JSON format with the following structure:
            {
                "summary": "Your comprehensive consolidated summary with HTML formatting",
                "documents_analyzed": ["list", "of", "documents"],
                "summary_type": "consolidated"
            }"""

            final_prompt = f"""
            Combine these document summaries into a cohesive, comprehensive consolidated summary.
            First provide a high-level overview of all documents, then highlight key points across documents,
            and finally discuss any important similarities, differences, or connections between the documents.
            
            {aggregated_summary}
            
            Format your response with HTML-like tags:
            <b>Consolidated Summary</b>
            <p>Overall summary of all documents together</p>
            
            <b>Key Points Across Documents</b>
            <ul>
                <li>First key point across documents</li>
                <li>Second key point across documents</li>
                <li>Third key point across documents</li>
            </ul>
            
            <b>Document Relationships</b>
            <p>Analysis of how the documents relate to each other, including any contradictions, complementary information, or overlapping themes.</p>
            
            CRITICAL: Your response MUST be a valid JSON object with this structure:
            {{
                "summary": "Your consolidated summary with HTML formatting",
                "documents_analyzed": [list of document names],
                "summary_type": "consolidated"
            }}
            """

            final_completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": final_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
                temperature=0.3,
                max_tokens=1200
            )

            # Parse JSON response
            json_response = self._parse_json_response(final_completion.choices[0].message.content)
            final_summary = json_response.get("summary", "")
        except Exception as e:
            print(f"Error generating final consolidated summary: {str(e)}")
            # Fallback to a simpler message with the aggregated summaries
            final_summary = f"""
            <b>Consolidated Summary</b>
            <p>The following summarizes content from multiple documents:</p>
            
            {aggregated_summary}
            """
        
        return final_summary

    def post(self, request):
        document_ids = request.data.get('document_ids', [])
        main_project_id = request.data.get('main_project_id')
        user = request.user

        if not document_ids or len(document_ids) < 2:
            return Response({
                'error': 'Please select at least two documents for a consolidated summary'
            }, status=status.HTTP_400_BAD_REQUEST)

        try:
            # Verify project access
            main_project = Project.objects.get(id=main_project_id, user=user)
            
            # Get all selected documents
            documents = Document.objects.filter(
                id__in=document_ids,
                user=user,
                main_project=main_project
            )
            
            if not documents or documents.count() < 2:
                return Response({
                    'error': 'Could not find requested documents'
                }, status=status.HTTP_404_NOT_FOUND)
            
            # Collect document chunks for hierarchical summarization
            all_chunks = []
            
            for document in documents:
                try:
                    processed_index = ProcessedIndex.objects.get(document=document)
                    
                    # Load FAISS index chunks for this document
                    _, chunks = self.load_faiss_index(processed_index.metadata.split('/')[-1].split('_')[0])
                    
                    # Add these chunks to our collection
                    all_chunks.extend(chunks)
                    
                except ProcessedIndex.DoesNotExist:
                    # Skip documents that haven't been processed
                    continue
            
            if not all_chunks:
                return Response({
                    'error': 'No processed content found for the selected documents'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Generate hierarchical summary across all documents
            consolidated_summary = self.hierarchical_summary(all_chunks)
            
            # Format the summary with HTML-like tags
            formatted_summary = self.format_summary_response(consolidated_summary)

            # Update project timestamp to track activity
            update_project_timestamp(main_project_id, user)
            
            return Response({
                'consolidated_summary': formatted_summary
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"Error generating consolidated summary: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while generating the consolidated summary'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

#Updated docUploadView with pgvector + blob storage

class DocumentUploadView(DocumentProcessingMixin, APIView):
    parser_classes = (MultiPartParser, FormParser)
   
    def post(self, request):
            files = request.FILES.getlist('files')
            user = request.user
            main_project_id = request.data.get('main_project_id')
            target_user_id = request.data.get('target_user_id')

            if target_user_id and request.user.username == 'admin':
                try:
                    user = User.objects.get(id=target_user_id)
                except User.DoesNotExist:
                    return Response({'error': 'Target user not found'}, status=status.HTTP_404_NOT_FOUND)
            else:
                user = request.user

            try:
                permissions = UserModulePermissions.objects.get(user=user)
                if permissions.disabled_modules.get('document-upload', False):
                    return Response({'error': 'Document uploads are disabled for this user'}, status=status.HTTP_403_FORBIDDEN)
            except UserModulePermissions.DoesNotExist:
                pass

            try:
                upload_permissions = UserUploadPermissions.objects.get(user=user)
                if not upload_permissions.can_upload:
                    return Response({'error': 'Document uploads are disabled for this user'}, status=status.HTTP_403_FORBIDDEN)
            except UserUploadPermissions.DoesNotExist:
                pass

            try:
                main_project = Project.objects.get(id=main_project_id, user=user)
                uploaded_docs = []
                last_processed_doc_id = None
                processed_data = None

                for file in files:
                    file_ext = os.path.splitext(file.name)[1].lower()
                    is_audio_file = file_ext in ['.mp3', '.wav', '.mpeg', '.m4a', '.aac', '.flac']
                    is_video_file = file_ext in ['.mp4', '.avi', '.mov', '.wmv', '.mkv', '.flv', '.webm', '.mts']
                    is_media_file = is_audio_file or is_video_file

                    # --- Status: Uploading ---
                    update_doc_status(user, file.name, "uploading", 10, "Uploading file...")

                    existing_doc = Document.objects.filter(
                        user=user,
                        filename=file.name,
                        main_project=main_project
                    ).first()

                    # --- Status: Uploaded ---
                    update_doc_status(user, file.name, "uploaded", 20, "File uploaded, awaiting processing...")

                    try:
                        if is_media_file:
                            # --- Status: Processing Media ---
                            update_doc_status(user, file.name, "processing", 30, "Processing media file...")

                            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                                for chunk in file.chunks():
                                    tmp_file.write(chunk)
                                file_path = tmp_file.name

                            try:
                                if is_video_file:
                                    update_doc_status(user, file.name, "processing", 35, "Extracting audio from video...")
                                    extracted_text = self.extract_text_from_video(file_path, user=user)
                                else:
                                    update_doc_status(user, file.name, "processing", 35, "Transcribing audio file...")
                                    extracted_text = self.extract_text_from_audio(file_path, user=user)

                                if not extracted_text or (isinstance(extracted_text, str) and extracted_text.startswith("Error")):
                                    update_doc_status(user, file.name, "error", 100, "Failed to extract text from media file.")
                                    return Response({'error': f'Failed to extract text from media file: {file.name}'}, status=status.HTTP_400_BAD_REQUEST)

                                base_name = os.path.splitext(file.name)[0]
                                file_type = "video" if is_video_file else "audio"
                                transcript_filename = f"{base_name}_{file_type}_transcript.txt"

                                update_doc_status(user, file.name, "processing", 45, "Processing transcript...")

                                # Create temporary transcript file for local storage
                                transcript_file_path = os.path.join(tempfile.gettempdir(), transcript_filename)
                                with open(transcript_file_path, 'w', encoding='utf-8') as f:
                                    f.write(extracted_text)

                                try:
                                    with open(transcript_file_path, 'rb') as f:
                                        from django.core.files import File
                                        django_file = File(f, name=transcript_filename)

                                        if existing_doc:
                                            existing_doc.filename = transcript_filename
                                            existing_doc.file = django_file  # Store the local file
                                            existing_doc.save()
                                            document = existing_doc
                                            # Clear old embeddings
                                            DocumentEmbedding.objects.filter(document=document).delete()
                                        else:
                                            document = Document.objects.create(
                                                user=user,
                                                file=django_file,  # Store the local file
                                                filename=transcript_filename,
                                                main_project=main_project
                                            )

                                        # Process document with pgvector
                                        processed_data = self.process_document_from_text_pgvector(extracted_text, transcript_filename, document)

                                        # Create transaction record with page count
                                        file_size = file.size if hasattr(file, 'size') else None
                                        upload_method = 'video_transcript' if is_video_file else 'audio_transcript'
                                        page_count = 1  # Transcripts are always single page
                                        self.create_document_transaction(user, document, main_project, upload_method=upload_method, file_size=file_size, page_count=page_count)

                                finally:
                                    if os.path.exists(transcript_file_path):
                                        os.unlink(transcript_file_path)

                                # --- Status: Indexing ---
                                update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                                # Create or update ProcessedIndex with pgvector
                                ProcessedIndex.objects.update_or_create(
                                    document=document,
                                    defaults={
                                        'storage_type': 'pgvector',
                                        'chunks_count': processed_data['chunks_count'],
                                        'summary': "",
                                        'markdown_path': processed_data.get('markdown_path', '')
                                    }
                                )

                                uploaded_docs.append({
                                    'id': document.id,
                                    'filename': transcript_filename,
                                    'original_media_type': file_type
                                })
                                last_processed_doc_id = document.id

                                # --- Status: Complete ---
                                update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=document.id)

                            finally:
                                if os.path.exists(file_path):
                                    os.unlink(file_path)

                        else:
                            # --- Status: Processing Document ---
                            update_doc_status(user, file.name, "processing", 30, "Processing document...")

                            if existing_doc:
                                try:
                                    processed_index = ProcessedIndex.objects.get(document=existing_doc)
                                    
                                    # If it's already processed with pgvector, just return it
                                    if processed_index.storage_type == 'pgvector' and \
                                    DocumentEmbedding.objects.filter(document=existing_doc).exists():
                                        uploaded_docs.append({
                                            'id': existing_doc.id,
                                            'filename': existing_doc.filename
                                        })
                                        last_processed_doc_id = existing_doc.id
                                        update_doc_status(user, file.name, "complete", 100, "Already processed!", doc_id=existing_doc.id)
                                        continue
                                        
                                except ProcessedIndex.DoesNotExist:
                                    pass
                                    
                                # Update existing document with local file storage
                                existing_doc.file = file
                                existing_doc.save()
                                
                                # Clear old embeddings and reprocess with optimized processing
                                DocumentEmbedding.objects.filter(document=existing_doc).delete()
                                processed_data = self.process_document_pgvector_optimized(file, user, existing_doc)

                                # Get page count from processed data
                                page_count = processed_data.get('page_count', 1)
                                file_size = file.size if hasattr(file, 'size') else None
                                
                                # Update transaction if it exists or create new one
                                try:
                                    # Try to find existing transaction
                                    existing_transaction = DocumentTransaction.objects.filter(
                                        user_transaction__document_id=existing_doc.id
                                    ).first()
                                    
                                    if existing_transaction:
                                        existing_transaction.no_pages = page_count
                                        existing_transaction.file_size = file_size
                                        existing_transaction.save()
                                    else:
                                        # Create new transaction if none exists
                                        self.create_document_transaction(user, existing_doc, main_project, upload_method='regular', file_size=file_size, page_count=page_count)
                                except Exception as e:
                                    print(f"Error updating transaction: {str(e)}")

                                # --- Status: Indexing ---
                                update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                                # Update ProcessedIndex with pgvector
                                ProcessedIndex.objects.update_or_create(
                                    document=existing_doc,
                                    defaults={
                                        'storage_type': 'pgvector',
                                        'chunks_count': processed_data['chunks_count'],
                                        'summary': "",
                                        'markdown_path': processed_data.get('markdown_path', '')
                                    }
                                )
                                
                                uploaded_docs.append({
                                    'id': existing_doc.id,
                                    'filename': existing_doc.filename
                                })
                                last_processed_doc_id = existing_doc.id

                                # --- Status: Complete ---
                                update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=existing_doc.id)
                            else:
                                # Create new document with local file storage
                                document = Document.objects.create(
                                    user=user,
                                    file=file,  # Store the local file directly
                                    filename=file.name,
                                    main_project=main_project
                                )

                             

                                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                                    for chunk in file.chunks():
                                        tmp_file.write(chunk)
                                    file_path = tmp_file.name

                                file_ext = os.path.splitext(file_path)[1].lower()
                
                                # Count pages based on file type
                                page_count = self.count_document_pages(file_path, file_ext)
                                print(f"Detected {page_count} pages in {file.name}")

                                file_size = file.size if hasattr(file, 'size') else None

                                # CHECK PAGE LIMIT BEFORE PROCESSING - MOVE THIS UP
                                transaction_result = self.create_document_transaction(
                                    user, document, main_project, 
                                    upload_method='regular', 
                                    file_size=file_size, 
                                    page_count=page_count
                                )
                                
                                # Check if page limit was exceeded
                                if isinstance(transaction_result, dict) and transaction_result.get('error'):
                                    # Delete the document we just created since we can't proceed
                                    document.delete()
                                    
                                    # Clean up temp file
                                    if os.path.exists(file_path):
                                        os.unlink(file_path)
                                    
                                    # Update status to failed
                                    update_doc_status(user, file.name, "failed", 0, f"Upload failed: {transaction_result.get('message', 'Page limit exceeded')}")
                                    
                                    print(f"Page limit exceeded for ########################### {file.name}")
                                    # Return error response - THIS WILL STOP THE PROCESSING
                                    return Response(transaction_result, status=status.HTTP_429_TOO_MANY_REQUESTS)
                                                
                                # ONLY PROCESS IF PAGE LIMIT CHECK PASSED
                                # Process the document with optimized pgvector processing
                                processed_data = self.process_document_pgvector_optimized(file, user, document)
                                
                                # Clean up temp file
                                if os.path.exists(file_path):
                                    os.unlink(file_path)
                                # --- Status: Indexing ---
                                update_doc_status(user, file.name, "indexing", 60, "Indexing document...")

                                # Create ProcessedIndex with pgvector
                                ProcessedIndex.objects.create(
                                    document=document,
                                    storage_type='pgvector',
                                    chunks_count=processed_data['chunks_count'],
                                    summary="",
                                    markdown_path=processed_data.get('markdown_path', '')
                                )
                                
                                uploaded_docs.append({
                                    'id': document.id,
                                    'filename': document.filename
                                })
                                last_processed_doc_id = document.id

                                # --- Status: Complete ---
                                update_doc_status(user, file.name, "complete", 100, "Processing complete!", doc_id=document.id)

                    except Exception as e:
                        update_doc_status(user, file.name, "error", 100, f"Error: {str(e)}")
                        print(f"Error processing document: {str(e)}")
                        return Response({
                            'error': str(e),
                            'detail': 'An error occurred while processing the document'
                        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

                request.session['active_document_id'] = last_processed_doc_id

                if processed_data is None:
                    return Response({'error': 'No documents were successfully processed'}, status=status.HTTP_400_BAD_REQUEST)

                update_project_timestamp(main_project_id, user)

                return Response({
                    'message': 'Documents uploaded successfully',
                    'documents': uploaded_docs,
                    'active_document_id': last_processed_doc_id
                }, status=status.HTTP_201_CREATED)

            except Exception as e:
                print(f"Error processing document: {str(e)}")
                return Response({
                    'error': str(e),
                    'detail': 'An error occurred while processing the document'
                }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def create_document_transaction(self, user, document, main_project, upload_method, file_size=None, page_count=None):
        """Create transaction record for document upload with page limit validation"""
        try:
            print("&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&&",page_count)
            
            # Get or create user API tokens record
            user_api_tokens, created = UserAPITokens.objects.get_or_create(
                user=user,
                defaults={
                    'page_limit': 20,
                    'token_limit': 1_000_000
                }
            )
            
            # Calculate current page usage for this user
            current_total_pages = self.get_user_total_page_usage(user)
            
            # Get pages for this request
            request_pages = page_count or 1  # Default to 1 if no page count provided
            
            # Check if adding these pages would exceed the limit
            new_total_pages = current_total_pages + request_pages
            
            print(f"Page usage check - Current: {current_total_pages}, Request: {request_pages}, New Total: {new_total_pages}, Limit: {user_api_tokens.page_limit}")
            
            if new_total_pages > user_api_tokens.page_limit:
                # Return error data instead of creating transaction
                remaining_pages = max(0, user_api_tokens.page_limit - current_total_pages)
                error_data = {
                    'error': 'page_limit_exceeded',
                    'message': f"You have reached your document page limit of {user_api_tokens.page_limit} pages. Current usage: {current_total_pages} pages. Remaining: {remaining_pages} pages. This document has {request_pages} pages.",
                    'current_usage': current_total_pages,
                    'page_limit': user_api_tokens.page_limit,
                    'remaining_pages': remaining_pages,
                    'requested_pages': request_pages,
                    'document_name': document.filename
                }
                print(f"Page limit exceeded for user {user.id}: {new_total_pages} > {user_api_tokens.page_limit}")
                return error_data
            
            # Create main transaction record
            user_transaction = UserTransaction.objects.create(
                user=user,
                transaction_type=TransactionType.DOCUMENT_UPLOAD,
                document_name=document.filename,
                document_id=document.id,
                main_project=main_project,
                metadata={
                    'upload_timestamp': timezone.now().isoformat(),
                    'upload_method': upload_method,
                    'file_size': file_size or getattr(document.file, 'size', None),
                    'page_count': page_count,
                    'page_limit_check': {
                        'previous_total': current_total_pages,
                        'request_pages': request_pages,
                        'new_total': new_total_pages,
                        'limit': user_api_tokens.page_limit
                    }
                }
            )
        
            # Create detailed document transaction with page count
            DocumentTransaction.objects.create(
                user_transaction=user_transaction,
                original_filename=document.filename,
                file_size=file_size or getattr(document.file, 'size', None),
                file_type=os.path.splitext(document.filename)[1].lower(),
                upload_method=upload_method,
                no_pages=page_count  # Store page count in the new field
            )
        
            print(f"Transaction recorded for document: {document.filename} with {page_count} pages")
            print(f"User total page usage: {new_total_pages}/{user_api_tokens.page_limit}")
            
            # Return success data
            return {
                'success': True,
                'current_usage': new_total_pages,
                'page_limit': user_api_tokens.page_limit,
                'remaining_pages': user_api_tokens.page_limit - new_total_pages,
                'document_pages': request_pages
            }
        
        except Exception as e:
            print(f"Error creating document transaction: {str(e)}")
            return {
                'error': 'transaction_error',
                'message': f'Error creating document transaction: {str(e)}'
            }

    def get_user_total_page_usage(self, user):
        """Calculate total page usage for a user across all documents"""
        from django.db.models import Sum,Count
        try:
            # Get all document transactions for this user
            total_pages = DocumentTransaction.objects.filter(
                user_transaction__user=user,
                user_transaction__is_active=True
            ).aggregate(
                total_pages=Sum('no_pages')
            )
            
            pages_total = total_pages.get('total_pages') or 0
            
            return pages_total
            
        except Exception as e:
            print(f"Error calculating user page usage: {str(e)}")
        return 0
    # Keep the complexity detection from original code
    # -------------------------------
    def detect_document_complexity(self, file_path):
        """Detect if PDF contains images."""
        import fitz
        try:
            doc = fitz.open(file_path)
            for page in doc:
                images = page.get_images()
                if images:
                    return True
            return False
        except Exception as e:
            print(f"Error detecting images in PDF: {e}")
            return False
   
    # Add LlamaParse integration from Streamlit code
    def split_text_into_chunks(self, text, chunk_size=1500, chunk_overlap=300):
        """
        Splits text into chunks while attempting to preserve sentence and paragraph structure.
        If the text is very short, a single chunk is returned.
        """
        if len(text) <= chunk_size:
            return [text]
        chunks = []
        start = 0
        text_length = len(text)
        while start < text_length:
            end = min(start + chunk_size, text_length)
            # Try to find a paragraph break to end the chunk
            if end < text_length:
                break_point = text.rfind("\n\n", start, end)
                if break_point != -1 and break_point > start + chunk_size // 2:
                    end = break_point + 2
                else:
                    # Fallback to sentence break
                    sentence_break = text.rfind(". ", start, end)
                    if sentence_break != -1 and sentence_break > start + chunk_size // 2:
                        end = sentence_break + 2
                    else:
                        # Fallback to space
                        space_break = text.rfind(" ", start, end)
                        if space_break != -1 and space_break > start + chunk_size // 2:
                            end = space_break + 1
            chunks.append(text[start:end])
            start = end - chunk_overlap if end < text_length else text_length
        return chunks
    
    
    # Clean text function (from Streamlit)
    def clean_text(self, text):
        """Clean extracted text by removing extra whitespace and special characters."""
        # Remove extra whitespace
        cleaned = re.sub(r'\s+', ' ', text)
        # Remove special characters that might cause issues
        cleaned = re.sub(r'[^\w\s.,;:!?"\'\-()]', ' ', cleaned)
        return cleaned.strip()
   
    def get_embeddings(self, texts):
        """
        Gets embeddings using OpenAI's text-embedding-3-small model.
        """
        try:
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small"
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            print(f"Error getting embeddings: {str(e)}")
            return []

    # -------------------------------
    # Process documents function that integrates with Django database
    import tempfile
    import uuid
    import numpy as np
    import fitz  # PyMuPDF
    import os
    from pathlib import Path
    from concurrent.futures import ThreadPoolExecutor, as_completed
    import threading
    # -------------------------------
    # Replace process_pdf_document_optimized() with:
    def process_pdf_document_optimized_pgvector(self, file, user, document):
            """
            Optimized PDF processing with pgvector storage instead of FAISS
            """
            import fitz 
            try:
                # Save uploaded file to temporary location
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                    for chunk in file.chunks():
                        tmp_file.write(chunk)
                    pdf_path = tmp_file.name
                
                try:
                    # ... Keep all your existing optimization logic ...
                    pdf_doc = fitz.open(pdf_path)
                    total_pages = len(pdf_doc)
                    print(f"Processing PDF with {total_pages} pages: {file.name}")
                    
                    # Keep all your parallel processing logic
                    pages_with_images, pages_without_images = self._bulk_analyze_pages_for_images(pdf_doc)
                    
                    combined_markdown = ""
                    page_results = {}
                    
                    # Process pages without images in parallel (fast)
                    if pages_without_images:
                        print("Processing text-only pages in parallel...")
                        text_results = self._process_text_pages_parallel(pdf_doc, pages_without_images)
                        page_results.update(text_results)
                    
                    # OPTIMIZATION 3: Batch process image pages if there are many
                    if pages_with_images:
                        if len(pages_with_images) > 5:  # If many image pages, use batch processing
                            print(f"Batch processing {len(pages_with_images)} image pages...")
                            image_results = self._batch_process_image_pages_pgvector(pdf_doc, pages_with_images, file.name, user, document)
                        else:
                            print(f"Processing {len(pages_with_images)} image pages individually...")
                            image_results = self._process_image_pages_parallel_pgvector(pdf_doc, pages_with_images, file.name, user, document)
                        page_results.update(image_results)
                    
                    pdf_doc.close()
                    
                    # Combine results in page order
                    for page_num in range(total_pages):
                        page_text = page_results.get(page_num, "[No extractable content]")
                        combined_markdown += f"\n\n## Page {page_num + 1}\n\n{page_text.strip()}\n"
                    
                    if not combined_markdown.strip():
                        raise ValueError("No content could be extracted from any page of the PDF")
                    
                    # Save markdown file locally instead of blob storage
                    import uuid
                    session_id = uuid.uuid4().hex
                    markdown_filename = f"document_{session_id}.md"
                    markdown_path = os.path.join(tempfile.gettempdir(), markdown_filename)
                    
                    with open(markdown_path, 'w', encoding='utf-8') as f:
                        f.write(combined_markdown)
                    
                    print(f"Complete markdown saved to: {markdown_path}")
                    
                    # Process with pgvector instead of FAISS
                    cleaned_text = self.clean_text(combined_markdown)
                    chunks = self._smart_chunk_text(cleaned_text, target_chunk_size=2000)
                    
                    all_chunks = []
                    for i, chunk in enumerate(chunks):
                        all_chunks.append({
                            'text': chunk,
                            'source': file.name,
                            'source_file': file.name,
                            'chunk_id': i
                        })
                    
                    # Batch embeddings
                    text_chunks = [chunk['text'] for chunk in all_chunks]
                    print(f"Getting embeddings for {len(text_chunks)} optimized chunks")
                    embeddings = self._get_embeddings_batch(text_chunks, batch_size=50)
                    
                    if not embeddings:
                        raise ValueError("Failed to generate embeddings for document chunks")
                    
                    # Save to pgvector instead of FAISS
                    self.save_embeddings_to_pgvector(all_chunks, embeddings, document)
                    
                    print(f"PDF processed successfully: {len(all_chunks)} chunks saved to pgvector from {total_pages} pages")
                    
                    return {
                        'storage_type': 'pgvector',
                        'chunks_count': len(all_chunks),
                        'full_text': cleaned_text,
                        'markdown_path': markdown_path,
                        'page_count': total_pages
                    }
                    
                finally:
                    if os.path.exists(pdf_path):
                        os.unlink(pdf_path)
                        
            except Exception as e:
                print(f"Error in process_pdf_document_optimized_pgvector: {str(e)}")
                raise

    def _bulk_analyze_pages_for_images(self, pdf_doc):
        """
        OPTIMIZATION: Analyze all pages for images in one pass
        """
        pages_with_images = []
        pages_without_images = []
        
        for page_num in range(len(pdf_doc)):
            try:
                page = pdf_doc.load_page(page_num)
                image_list = page.get_images()
                if len(image_list) > 0:
                    pages_with_images.append(page_num)
                else:
                    pages_without_images.append(page_num)
            except Exception as e:
                print(f"Error analyzing page {page_num}: {str(e)}")
                pages_without_images.append(page_num)  # Default to text extraction
        
        return pages_with_images, pages_without_images



    def _process_text_pages_parallel(self, pdf_doc, page_numbers):
        """
        OPTIMIZATION: Process text-only pages in parallel
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import threading
        def extract_text_from_page(page_num):
            try:
                page = pdf_doc.load_page(page_num)
                text = page.get_text()
                print(f"Extracted text from page {page_num}: {text}")
                return page_num, text if text and text.strip() else ""
            
            except Exception as e:
                print(f"Error extracting text from page {page_num}: {str(e)}")
                return page_num, ""
        
        results = {}
        # Use ThreadPoolExecutor for I/O bound operations
        with ThreadPoolExecutor(max_workers=min(8, len(page_numbers))) as executor:
            future_to_page = {executor.submit(extract_text_from_page, page_num): page_num 
                            for page_num in page_numbers}
            
            for future in as_completed(future_to_page):
                page_num, text = future.result()
                results[page_num] = text
        
        return results

    def _batch_process_image_pages_pgvector(self, pdf_doc, page_numbers, original_filename, user, document):
        """
        OPTIMIZATION: Process multiple image pages together using LlamaParse
        """
        import tempfile
        import fitz  # PyMuPDF
        
        try:
            # Create a single PDF with all image pages
            batch_pdf = fitz.open()
            page_mapping = {}  # Maps batch page index to original page number
            
            for i, original_page_num in enumerate(page_numbers):
                batch_pdf.insert_pdf(pdf_doc, from_page=original_page_num, to_page=original_page_num)
                page_mapping[i] = original_page_num
            
            # Save batch PDF
            batch_filename = f"batch_{uuid.uuid4().hex[:8]}.pdf"
            batch_path = os.path.join(tempfile.gettempdir(), batch_filename)
            batch_pdf.save(batch_path)
            batch_pdf.close()
            
            # Process with LlamaParse
            try:
                batch_result = self.process_complex_document_with_llamaparse_pgvector_blob(
                    batch_path, f"{original_filename}_batch", user, document
                )
                
                # Extract text and try to split by pages
                if isinstance(batch_result, dict) and 'full_text' in batch_result:
                    full_text = batch_result['full_text']
                else:
                    full_text = str(batch_result)
                
                # Try to split the result back to individual pages
                # This is approximate since LlamaParse may not preserve exact page boundaries
                results = self._split_batch_result_to_pages(full_text, page_mapping)
                print(f"Batch LlamaParse processed successfully: {len(results)} chunks saved to pgvector from {len(page_numbers)} pages")
                
            except Exception as e:
                print(f"Batch LlamaParse failed, falling back to individual processing: {str(e)}")
                # Fallback to individual processing
                results = self._process_image_pages_parallel_pgvector(pdf_doc, page_numbers, original_filename, user, document)
            
            finally:
                # Cleanup
                if os.path.exists(batch_path):
                    os.unlink(batch_path)
            
            return results
            
        except Exception as e:
            print(f"Error in batch processing: {str(e)}")
            # Fallback to individual processing
            return self._process_image_pages_parallel_pgvector(pdf_doc, page_numbers, original_filename, user, document)

    def _process_image_pages_parallel_pgvector(self, pdf_doc, page_numbers, original_filename, user, document):
        """
        Process image pages individually but in parallel
        """
        import tempfile
        import fitz  # PyMuPDF
        from concurrent.futures import ThreadPoolExecutor, as_completed
        def process_single_image_page(page_num):

            
            try:
                # Extract single page
                single_pdf = fitz.open()
                single_pdf.insert_pdf(pdf_doc, from_page=page_num, to_page=page_num)
                
                temp_filename = f"page_{page_num}_{uuid.uuid4().hex[:8]}.pdf"
                temp_path = os.path.join(tempfile.gettempdir(), temp_filename)
                single_pdf.save(temp_path)
                single_pdf.close()
                
                try:
                    # Process with LlamaParse
                    result = self.process_complex_document_with_llamaparse_pgvector_blob(
                        temp_path, f"{original_filename}_page_{page_num}", user, document
                    )
                    
                    if isinstance(result, dict) and 'full_text' in result:
                        text = result['full_text']
                    else:
                        text = str(result)
                    
                    return page_num, text
                    
                finally:
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
                        
            except Exception as e:
                print(f"Error processing image page {page_num}: {str(e)}")
                # Fallback to text extraction
                try:
                    page = pdf_doc.load_page(page_num)
                    text = page.get_text()
                    return page_num, text if text and text.strip() else ""
                except:
                    return page_num, ""
        
        results = {}
        # Limit concurrent LlamaParse calls to avoid rate limiting
        with ThreadPoolExecutor(max_workers=min(3, len(page_numbers))) as executor:
            future_to_page = {executor.submit(process_single_image_page, page_num): page_num 
                            for page_num in page_numbers}
            
            for future in as_completed(future_to_page):
                page_num, text = future.result()
                results[page_num] = text
        
        return results

    def _split_batch_result_to_pages(self, full_text, page_mapping):
        """
        Try to split batch LlamaParse result back to individual pages
        """
        results = {}
        
        # Simple heuristic: split text roughly equally among pages
        lines = full_text.split('\n')
        lines_per_page = max(1, len(lines) // len(page_mapping))
        
        for i, original_page_num in page_mapping.items():
            start_line = i * lines_per_page
            end_line = start_line + lines_per_page if i < len(page_mapping) - 1 else len(lines)
            page_text = '\n'.join(lines[start_line:end_line]).strip()
            results[original_page_num] = page_text if page_text else "[Content from batch processing]"
        
        return results

    def _smart_chunk_text(self, text, target_chunk_size=2000):
        """
        OPTIMIZATION: Create larger, smarter chunks to reduce embedding calls
        """
        # Use existing chunking but with larger target size
        return self.split_text_into_chunks(text)

    def _get_embeddings_batch(self, text_chunks, batch_size=50):
        """
        OPTIMIZATION: Process embeddings in batches for better performance
        """
        all_embeddings = []
        
        for i in range(0, len(text_chunks), batch_size):
            batch = text_chunks[i:i + batch_size]
            print(f"Processing embedding batch {i//batch_size + 1}/{(len(text_chunks) + batch_size - 1)//batch_size}")
            
            try:
                batch_embeddings = self.get_embeddings(batch)
                if batch_embeddings:
                    all_embeddings.extend(batch_embeddings)
                else:
                    print(f"Warning: Failed to get embeddings for batch starting at index {i}")
            except Exception as e:
                print(f"Error processing embedding batch {i}: {str(e)}")
                continue
        
        return all_embeddings if all_embeddings else None

    # Replace this function name:
    # process_document() -> process_document_pgvector_optimized()

    def process_document_pgvector_optimized(self, file, user, document):
        """
        Enhanced process_document that routes to optimized PDF processing with pgvector
        """
        import tempfile
        
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                for chunk in file.chunks():
                    tmp_file.write(chunk)
                file_path = tmp_file.name

            try:
                file_ext = os.path.splitext(file_path)[1].lower()
                
                # Count pages based on file type
                page_count = self.count_document_pages(file_path, file_ext)
                print(f"Detected {page_count} pages in {file.name}")
                
                # Route PDF files to optimized page-by-page processing
                if file_ext == '.pdf':
                    print(f"PDF detected: {file.name}. Using optimized page-by-page processing...")
                    result = self.process_pdf_document_optimized_pgvector(file, user, document)
                    result['page_count'] = page_count
                    return result

                # Keep existing logic for other file types but update to pgvector
                if file_ext == '.pptx':
                    result = self.process_complex_document_with_llamaparse_pgvector_blob(file_path, file.name, user, document)
                    result['page_count'] = page_count
                    return result

                if file_ext in ['.jpg', '.jpeg', '.png', '.bmp']:
                    result = self.process_complex_document_with_llamaparse_pgvector_blob(file_path, file.name, user, document)
                    result['page_count'] = page_count
                    return result
                
                # Check document complexity for other file types
                is_complex = self.detect_document_complexity(file_path)

                if is_complex:
                    print(f"Complex document detected: {file.name}. Using LlamaParse...")
                    result = self.process_complex_document_with_llamaparse_pgvector_blob(file_path, file.name, user, document)
                    result['page_count'] = page_count
                    return result
                else:
                    # Simple document processing
                    extracted_text = self.extract_text_from_file(file_path, user=user)
                
                if not extracted_text:
                    raise ValueError("No content could be extracted from the document")
                
                result = self.process_document_from_text_pgvector(extracted_text, file.name, document)
                result['page_count'] = page_count
                return result
                
            finally:
                if os.path.exists(file_path):
                    os.unlink(file_path)
                    
        except Exception as e:
            print(f"Error in process_document_pgvector_optimized: {str(e)}")
            raise

    def save_embeddings_to_pgvector(self, chunks_with_metadata, embeddings, document):
        """
        Save chunks and their embeddings to PostgreSQL using pgvector
        """
        try:
            # Clear existing embeddings for this document
            DocumentEmbedding.objects.filter(document=document).delete()
            
            # Create new embeddings
            embedding_objects = []
            for i, (chunk_data, embedding) in enumerate(zip(chunks_with_metadata, embeddings)):
                embedding_obj = DocumentEmbedding(
                    document=document,
                    content=chunk_data['text'],
                    embedding=embedding,
                    chunk_id=chunk_data['chunk_id'],
                    source=chunk_data.get('source', ''),
                    source_file=chunk_data.get('source_file', '')
                )
                embedding_objects.append(embedding_obj)
            
            # Bulk create for efficiency
            DocumentEmbedding.objects.bulk_create(embedding_objects, batch_size=100)
            
            print(f"Saved {len(embedding_objects)} embeddings to pgvector database")
            
        except Exception as e:
            print(f"Error saving embeddings to pgvector: {str(e)}")
            raise

    def process_document_from_text_pgvector(self, extracted_text, filename, document):
        """
        Process extracted text and save to pgvector
        """
        try:
            # Clean the text
            cleaned_text = self.clean_text(extracted_text)
            
            # Split text into chunks
            chunks = self.split_text_into_chunks(cleaned_text)
            print(f"Created {len(chunks)} chunks from {filename}")
            
            # Prepare chunks with metadata
            all_chunks = []
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    'text': chunk,
                    'source': filename,
                    'source_file': filename,
                    'chunk_id': i
                })
            
            # Get embeddings for all chunks
            text_chunks = [chunk['text'] for chunk in all_chunks]
            print(f"Getting embeddings for {len(text_chunks)} chunks")
            embeddings = self.get_embeddings(text_chunks)
            
            if not embeddings:
                raise ValueError("Failed to generate embeddings for document chunks")
            
            # Save to pgvector database
            self.save_embeddings_to_pgvector(all_chunks, embeddings, document)
            
            print(f"Document processed successfully: {len(all_chunks)} chunks saved to pgvector")
            
            return {
                'storage_type': 'pgvector',
                'chunks_count': len(all_chunks),
                'full_text': cleaned_text,
                'markdown_path': ''
            }
            
        except Exception as e:
            print(f"Error in process_document_from_text_pgvector: {str(e)}")
            raise


    def process_complex_document_with_llamaparse_pgvector_blob(self, file_path, file_name, user, document):
        """
        Process complex document using LlamaParse with user's API token and save to pgvector.
        """
        import tempfile
        from llama_parse import LlamaParse
       
        try:
            # Get the user's Llama API token
            try:
                user_api_tokens = UserAPITokens.objects.get(user=user)
                llama_api_key = user_api_tokens.llama_token
               
                # If no token is saved for the user, use a fallback mechanism or raise an error
                if not llama_api_key:
                    logger.warning(f"No Llama API token found for user {user.username}")
                    # Optional: You could implement a fallback or raise an error
                    raise ValueError("Llama API token is required for processing complex documents")
                   
            except UserAPITokens.DoesNotExist:
                logger.error(f"No API tokens record found for user {user.username}")
                raise ValueError("User API tokens not configured")
           
            parser = LlamaParse(
                api_key=llama_api_key,  # Use the user's Llama API token
                result_type="markdown",
                verbose=True,
                images=True,
                premium_mode=True
            )
           
            parsed_documents = parser.load_data(file_path)
            full_text = '\n'.join([doc.text for doc in parsed_documents])
           
            # Check if extracted text is empty or only whitespace
            if not full_text or not full_text.strip():
                logger.error(f"No text could be extracted from document: {file_name}")
                raise ValueError(f"Failed to extract any text content from the document '{file_name}'. The document may be corrupted, password-protected, or in an unsupported format.")
           
            # Additional check for very short content (optional - adjust threshold as needed)
            if len(full_text.strip()) < 5:  # Less than 5 characters
                logger.warning(f"Very little text extracted from document: {file_name} (only {len(full_text.strip())} characters)")
                raise ValueError(f"Insufficient text content extracted from document '{file_name}'. Only {len(full_text.strip())} characters were extracted.")
           
            logger.info(f"Successfully extracted {len(full_text)} characters from document: {file_name}")
           
            # Save markdown (adapted from Streamlit implementation)
            base_name = os.path.splitext(file_name)[0]
            safe_name = re.sub(r'[^\w\-_.]', '_', base_name)
            md_path = os.path.join("markdown_files", f"{safe_name}.md")
            os.makedirs("markdown_files", exist_ok=True)
            with open(md_path, "w", encoding='utf-8') as f:
                f.write(full_text)
           
            # Create chunks with overlap for better retrieval
            chunked_texts = []
            all_chunks = []
            chunk_id = 0
            
            for doc in parsed_documents:
                # Skip documents with no meaningful content
                if not doc.text or not doc.text.strip():
                    continue
                   
                # Original document as a chunk
                chunked_texts.append(doc.text)
                all_chunks.append({
                    'text': doc.text,
                    'source': file_name,
                    'source_file': file_name,
                    'chunk_id': chunk_id
                })
                chunk_id += 1
               
                # Create additional smaller chunks for better retrieval
                words = doc.text.split()
                chunk_size = 200  # Smaller chunk size
                stride = 100      # With overlap
               
                if len(words) > chunk_size:
                    for i in range(0, len(words) - chunk_size, stride):
                        chunk = " ".join(words[i:i+chunk_size])
                        if len(chunk.split()) > 50:  # Ensure chunk has substantial content
                            chunked_texts.append(chunk)
                            all_chunks.append({
                                'text': chunk,
                                'source': file_name,
                                'source_file': file_name,
                                'chunk_id': chunk_id
                            })
                            chunk_id += 1
           
            # Final check to ensure we have chunks to process
            if not chunked_texts:
                logger.error(f"No valid text chunks created from document: {file_name}")
                raise ValueError(f"Unable to create text chunks from document '{file_name}'. The document content may be insufficient or invalid.")
           
            # Get embeddings for all chunks
            print(f"Getting embeddings for {len(chunked_texts)} chunks")
            embeddings = self.get_embeddings(chunked_texts)
            
            if not embeddings:
                raise ValueError("Failed to generate embeddings for document chunks")
            
            # Save to pgvector database
            self.save_embeddings_to_pgvector(all_chunks, embeddings, document)
           
            return {
                'storage_type': 'pgvector',
                'chunks_count': len(all_chunks),
                'full_text': full_text,
                'markdown_path': md_path  # Include the markdown path
            }
           
        except Exception as e:
            logger.error(f"Error in complex document processing for file '{file_name}': {str(e)}")
            print(f"Error in complex document processing: {str(e)}")
            raise

    def count_document_pages(self, file_path, file_ext):
        """
        Count the number of pages in a document based on file type
        For PDF and PPTX files, only count pages that contain images
        """
        try:
            if file_ext == '.pdf':
                # PDF files - count only pages with images
                import fitz
                pdf_doc = fitz.open(file_path)
                pages_with_images = 0
                
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    # Get list of images on the page
                    image_list = page.get_images()
                    if image_list:  # If page has images
                        pages_with_images += 1
                
                pdf_doc.close()
                return pages_with_images
                
            elif file_ext == '.pptx':
                # PowerPoint presentations - count only slides with images
                import fitz
                try:
                    # Convert PPTX to PDF in memory for image detection
                    pptx_doc = fitz.open(file_path)
                    slides_with_images = 0
                    
                    for page_num in range(len(pptx_doc)):
                        page = pptx_doc[page_num]
                        # Get list of images on the page/slide
                        image_list = page.get_images()
                        if image_list:  # If slide has images
                            slides_with_images += 1
                    
                    pptx_doc.close()
                    return slides_with_images
                    
                except Exception as pptx_error:
                    # Fallback to python-pptx if fitz fails
                    try:
                        from pptx import Presentation
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        
                        prs = Presentation(file_path)
                        slides_with_images = 0
                        
                        for slide in prs.slides:
                            has_image = False
                            for shape in slide.shapes:
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    has_image = True
                                    break
                            if has_image:
                                slides_with_images += 1
                        
                        return slides_with_images
                    except:
                        return 0  # Return 0 if we can't determine
                    
            elif file_ext in ['.docx', '.doc']:
                # Word documents
                try:
                    import docx
                    doc = docx.Document(file_path)
                    # For Word docs, we'll count sections or estimate based on content
                    # This is an approximation as Word doesn't have explicit page numbers
                    total_paragraphs = len(doc.paragraphs)
                    # Rough estimation: ~25-30 paragraphs per page (adjustable)
                    estimated_pages = max(1, total_paragraphs // 25)
                    return estimated_pages
                except:
                    return 1  # Default to 1 if we can't determine
                    
            elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff']:
                # Image files - always 1 page
                return 1
                
            elif file_ext == '.txt':
                # Text files - estimate based on content
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        # Rough estimation: ~3000 characters per page
                        estimated_pages = max(1, len(content) // 3000)
                        return estimated_pages
                except:
                    return 1
                    
            elif file_ext in ['.xlsx', '.xls']:
                # Excel files - count worksheets
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(file_path)
                    return len(wb.worksheets)
                except:
                    return 1
                    
            else:
                # For other file types, default to 1
                return 1
                
        except Exception as e:
            print(f"Error counting pages for {file_path}: {str(e)}")
            return 1  # Default to 1 page if there's an error
    
def update_doc_status(user, doc_name, status, progress=0, message="", doc_id=None):
    obj, _ = DocumentProcessingStatus.objects.get_or_create(
        user=user, document_name=doc_name,
        defaults={"status": status, "progress": progress, "message": message, "document_id": doc_id}
    )
    obj.status = status
    obj.progress = progress
    obj.message = message
    if doc_id:
        obj.document_id = doc_id
    obj.save()

class DocumentProcessingStatusView(APIView):
    permission_classes = [IsAuthenticated]
 
    def get(self, request):
        statuses = DocumentProcessingStatus.objects.filter(user=request.user)
        return Response([
            {
                "filename": s.document_name,
                "status": s.status,
                "progress": s.progress,
                "message": s.message,
                "document_id": s.document_id,
            }
            for s in statuses.order_by("created_at")
        ])


class CheckUploadPermissionsView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request):
        try:
            user = request.user
            
            # Skip permission check for admin users
            if user.username == 'admin':
                return Response({
                    'can_upload': True
                }, status=status.HTTP_200_OK)
            
            # Check if the user has upload permissions
            try:
                permissions = UserUploadPermissions.objects.get(user=user)
                can_upload = permissions.can_upload
            except UserUploadPermissions.DoesNotExist:
                # Default to allowed if no permissions are explicitly set
                can_upload = True
            
            return Response({
                'can_upload': can_upload
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            logger.error(f"Error in CheckUploadPermissionsView: {str(e)}")
            return Response({
                'error': f'Failed to check upload permissions: {str(e)}',
                'can_upload': False  # Default to disallowing upload on error
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class GenerateDocumentSummaryView(DocumentProcessingMixin, APIView):
    def post(self, request):
        document_ids = request.data.get('document_ids', [])
        user = request.user
        main_project_id = request.data.get('main_project_id')

        try:
            # Import your models here (adjust based on your project structure)
            from .models import Document, ProcessedIndex
                      # Adjust import as needed
            
            documents = Document.objects.filter(
                id__in=document_ids,
                user=user,
                main_project=main_project_id
            )

            if not documents:
                return Response({
                    'error': 'No valid documents found'
                }, status=status.HTTP_404_NOT_FOUND)

            all_summaries = []
            for document in documents:
                try:
                    # Extract text from the document (your existing flow)
                    text = self.extract_text_from_file(document.file.path, user)
                    cleaned_text = self.clean_text(text)
                    
                    # Generate summary and key points using NEW Gemini-based function
                    summary, key_points = self.generate_summary(cleaned_text, document.filename, user)
                    
                    # Save to database - update existing ProcessedIndex
                    processed_index, created = ProcessedIndex.objects.get_or_create(
                        document=document,
                        defaults={
                            'summary': summary,
                            'key_points': key_points,  # Store key points in new field
                            'faiss_index': '',  # Set appropriate defaults for required fields
                            'metadata': '',
                            'markdown_path': ''
                        }
                    )
                    if not created:
                        processed_index.summary = summary
                        processed_index.key_points = key_points  # Update key points
                        processed_index.save()

                    all_summaries.append({
                        'document_id': document.id,
                        'filename': document.filename,
                        'summary': summary,
                        'key_points': key_points,  # Include key points in response
                        'success': True
                    })
                        
                except Exception as e:
                    logger.error(f"Error processing document {document.id}: {str(e)}")
                    all_summaries.append({
                        'document_id': document.id,
                        'filename': document.filename,
                        'summary': f'Error processing document: {str(e)}',
                        'key_points': [],
                        'success': False,
                        'error': str(e)
                    })

            if main_project_id:
                update_project_timestamp(main_project_id, user)        

            return Response({
                'message': 'Summaries generated successfully',
                'summaries': all_summaries
            }, status=status.HTTP_200_OK)

        except Exception as e:
            logger.error(f"Error in GenerateDocumentSummaryView: {str(e)}")
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)     


class DeleteDocumentView(APIView):
    permission_classes = [IsAuthenticated]
    
    def delete(self, request, document_id):
        try:
            # Find the document and ensure it belongs to the current user
            document = Document.objects.get(
                id=document_id, 
                user=request.user
            )

            # Get main_project_id from the document
            main_project_id = document.main_project_id
            
            # Optional: Delete associated ProcessedIndex
            try:
                processed_index = ProcessedIndex.objects.get(document=document)
                processed_index.delete()
            except ProcessedIndex.DoesNotExist:
                pass
            
            # Delete the document
            document.delete()

            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
            
            return Response(
                {'message': 'Document deleted successfully'}, 
                status=status.HTTP_200_OK
            )
        
        except Document.DoesNotExist:
            return Response(
                {'error': 'Document not found'}, 
                status=status.HTTP_404_NOT_FOUND
            )
        except Exception as e:
            return Response(
                {'error': f'Failed to delete document: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


class GetChatHistoryView(APIView):
    def get(self, request):
        try:
            user = request.user
            main_project_id = request.query_params.get('main_project_id')
            
            print(f"Getting chat history for user {user.username} and project {main_project_id}")
            
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)

            show_archived = request.query_params.get('archived', 'false').lower() == 'true'
 
            # Filter conversations by both user and main_project_id
            conversations = ChatHistory.objects.filter(
                user=user,
                main_project_id=main_project_id,

                is_archived=show_archived

            ).order_by('-created_at')
            
            history = []
            for conversation in conversations:
                messages = conversation.messages_NB.all().order_by('created_at')
                message_list = []
                for message in messages:
                    msg_data = {
                        'id': message.id,
                        'role': message.role,
                        'content': message.content,
                        'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                        'citations': message.citations or []
                    }
                    if message.role == 'assistant':
                        msg_data['use_web_knowledge'] = getattr(message, 'use_web_knowledge', False)
                        msg_data['response_length'] = getattr(message, 'response_length', 'comprehensive')
                        msg_data['response_format'] = getattr(message, 'response_format', 'natural')
                        msg_data['sources_info'] = getattr(message, 'sources', "")
                        # ADD THIS LOGIC:
                        if getattr(message, 'sources', '') == "Image Analysis":
                            msg_data['context_mode'] = "image"
                    message_list.append(msg_data)
            
                history.append({
                    'conversation_id': str(conversation.conversation_id),
                    'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                    'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                    'messages': message_list,
                    'preview': message_list[0]['content'] if message_list else "",
                    'follow_up_questions': conversation.follow_up_questions or [],

                    'selected_documents': [str(doc.id) for doc in conversation.documents.all()],

                    'is_archived': conversation.is_archived

                })
            #print("##############################", history)
            return Response(history, status=status.HTTP_200_OK)
            
        except Exception as e:
            print(f"Error in GetChatHistoryView: {str(e)}")
            return Response(
                {'error': f'Failed to fetch chat history: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class ArchiveConversationView(APIView):
    permission_classes = [IsAuthenticated]
 
    def patch(self, request, conversation_id):
        try:
            conversation = ChatHistory.objects.get(conversation_id=conversation_id, user=request.user)
            archive = request.data.get('archive', True)
            conversation.is_archived = archive
            conversation.save()
            return Response({'success': True, 'is_archived': conversation.is_archived}, status=status.HTTP_200_OK)
        except ChatHistory.DoesNotExist:
            return Response({'error': 'Conversation not found.'}, status=status.HTTP_404_NOT_FOUND)


class GenerateTitleView(APIView):
    permission_classes = [IsAuthenticated]
   
    def post(self, request):
        message = request.data.get('message')
        max_length = request.data.get('max_length', 40)
       
        try:
            # Prompt engineering for better titles
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{
                    "role": "system",
                    "content": "Generate a concise, descriptive title (2-5 words) that captures the main topic or question. Make it clear and specific, like a document heading."
                }, {
                    "role": "user",
                    "content": f"Generate a short title for this: {message}"
                }],
                max_tokens=max_length,
                temperature=0.3
            )
           
            title = response.choices[0].message.content.strip()
           
            return Response({
                "success": True,
                "title": title
            })
           
        except Exception as e:
            logger.error(f"Error generating title: {str(e)}")
            return Response({
                "success": False,
                "error": "Failed to generate title"
            }, status=500)
 


class SetActiveDocumentView(APIView):
    def post(self, request):
        try:
            document_id = request.data.get('document_id')
            
            if not document_id:
                return Response(
                    {'error': 'Document ID is required'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Verify the document exists and belongs to the user
            try:
                document = Document.objects.get(
                    id=document_id, 
                    user=request.user
                )
                
                # Check if document is processed
                ProcessedIndex.objects.get(document=document)
                # Get the main_project_id from the document
                main_project_id = document.main_project_id
                
                # Update project timestamp to track activity
                if main_project_id:
                    update_project_timestamp(main_project_id, request.user)
            except Document.DoesNotExist:
                return Response(
                    {'error': 'Document not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
            except ProcessedIndex.DoesNotExist:
                return Response(
                    {'error': 'Document not processed'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Set the active document in the session
            request.session['active_document_id'] = document_id

            return Response({
                'message': 'Active document set successfully',
                'active_document_id': document_id,
                'filename': document.filename
            }, status=status.HTTP_200_OK)
        
        except Document.DoesNotExist:
            return Response(
                {'error': 'Document not found'},
                status=status.HTTP_404_NOT_FOUND
            )
        
        except Exception as e:
            import traceback
            print(f"Detailed Error: {str(e)}")
            print(traceback.format_exc())  # This will print the full stack trace
            return Response(
                {'error': f'An unexpected error occurred: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
def post_process_response(response_text):
    """
    Post-process the generated response to improve formatting and readability
    but preserve the original meaning and content quality.
    
    Args:
        response_text (str): Raw generated response text
    
    Returns:
        str: Cleaned and formatted response
    """
    try:
        import re
        # Create a copy of the original text before processing
        original_text = response_text
        
        # Only remove explicit section headers - not content
        response_text = re.sub(r'^[IVX]+\.\s*[\w\s]+:', '', response_text, flags=re.MULTILINE)
        response_text = re.sub(r'^\d+\.\s*[\w\s]+:', '', response_text, flags=re.MULTILINE)

        # Convert markdown-style formatting to HTML
        response_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', response_text)
        response_text = re.sub(r'\*(.*?)\*', r'<em>\1</em>', response_text)
        
        # Remove specific section headers but keep their content
        section_headers = [
            'Contextual Insight', 
            'Structured Response', 
            'Analytical Depth', 
            'Interactive Engagement'
        ]
        for header in section_headers:
            response_text = response_text.replace(header + ':', '')
        
        # Clean up extra whitespace while preserving paragraph structure
        response_text = re.sub(r'\n{3,}', '\n\n', response_text)
        
        # If the processed text is significantly shorter or might have lost content,
        # revert to the original
        if len(response_text) < len(original_text) * 0.9:
            logger.warning("Post-processing reduced content significantly, reverting to original")
            return original_text
            
        return response_text.strip()
    
    except Exception as e:
        logger.error(f"Error in post-processing response: {str(e)}", exc_info=True)
        return response_text  # Return original if any error occurs


import uuid
import os
import pickle
import numpy as np
import faiss
import re
from datetime import datetime
import logging
from sklearn.feature_extraction.text import TfidfVectorizer
from agno.agent import Agent
from agno.models.openai import OpenAIChat
from agno.tools.duckduckgo import DuckDuckGoTools

logger = logging.getLogger(__name__)

from google import genai
from google.genai import types
import logging
import os

class ChatView(APIView):
    permission_classes = [IsAuthenticated]
 
    def __init__(self, post_process_func=post_process_response):
        self.post_process_func = post_process_func
        self.agent = Agent(
                model=OpenAIChat(id="gpt-4o"),
                description="""You are a helpful, friendly AI assistant providing informative and thoughtful responses.
                Use semantic HTML tags for structure (<b>, <p>, <ul>, <li>) to enhance readability.
                Maintain a natural, conversational tone.
                Provide comprehensive and clear explanations, breaking down complex topics into easily understandable sections.
                When using web search, prioritize recent, reliable, and relevant information to support your responses.""",
                tools=[DuckDuckGoTools()],
                show_tool_calls=True,
                markdown=True
            )
    def normalize_citation_markers(self, text):
        """
        Convert all citation formats to consistent [n] format:
        - [Source: 1, 2] -> [1][2]
        - [Sources: 1, 2] -> [1][2]
        - Source: 1, 2 -> [1][2]
        - (Source 1, 2) -> [1][2]
        """
        if not text:
            return text
        # Replace [Source: n] or [Sources: n] with [n]
        text = re.sub(r'\[(?:Source|Sources?)\s*:?\s*(\d+)\]', r'[\1]', text, flags=re.IGNORECASE)
        # Replace [Source: n, m] or [Sources: n, m] with [n][m]
        def bracketize(match):
            nums = re.findall(r'\d+', match.group(0))
            return ''.join([f'[{n}]' for n in nums])
        text = re.sub(r'\[(?:Source|Sources?)\s*:?\s*[\d,\s]+\]', bracketize, text, flags=re.IGNORECASE)
        # Replace Source: n, m or Sources: n, m (not in brackets) with [n][m]
        text = re.sub(r'(?:Source|Sources?)\s*:?\s*([\d,\s]+)',
                    lambda m: ''.join([f'[{n}]' for n in re.findall(r'\d+', m.group(1))]),
                    text, flags=re.IGNORECASE)
        # Replace (Source n, m) with [n][m]
        text = re.sub(r'\((?:Source|Sources?)\s*:?\s*([\d,\s]+)\)',
                    lambda m: ''.join([f'[{n}]' for n in re.findall(r'\d+', m.group(1))]),
                    text, flags=re.IGNORECASE)
        # Handle grouped citations [1, 2, 3] -> [1][2][3]
        text = re.sub(r'\[(\d+(?:\s*,\s*\d+)+)\]',
                    lambda m: ''.join([f'[{n.strip()}]' for n in m.group(1).split(',')]),
                    text)
        return text
 
    def post(self, request):
        user = request.user
        try:
            # Extract data with more robust handling
            message = request.data.get('message')
            conversation_id = request.data.get('conversation_id')
            main_project_id = request.data.get('main_project_id')
            selected_documents = request.data.get('selected_documents', [])
           
            # Extract web_mode flag from request
            use_web_knowledge = request.data.get('use_web_knowledge', False)
           
            # Extract general_chat_mode flag
            general_chat_mode = request.data.get('general_chat_mode', False)
 
            # Extract response_length preference
            response_length = request.data.get('response_length', 'comprehensive')
           
            # Extract response_format parameter or default to 'natural'
            response_format = request.data.get('response_format', 'natural')
 
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Validate message
            if not message:
                return Response(
                    {'error': 'Message is required'},
                    status=status.HTTP_400_BAD_REQUEST
                )
 
            # Process URLs in the query - new feature
            modified_message, url_context, extracted_urls = self.process_urls_in_query(message)
            has_url_content = bool(url_context)
           
            # Log URL processing results
            if has_url_content:
                print(f"URLs detected in query: {extracted_urls}")
                print(f"Modified query: {modified_message}")
                print(f"URL context size: {len(url_context)} characters")
 
            conversation_context = ""
            if conversation_id:
                try:
                    # Try to get existing conversation
                    conversation = ChatHistory.objects.get(
                        conversation_id=conversation_id,
                        user=user
                    )
                    
                    # Get last 10 conversation messages (5 user + 5 assistant = 5 complete turns)
                    recent_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('-created_at')[:10]
                    
                    # Format for context - get full messages, not truncated
                    if recent_messages:
                        context_messages = []
                        for msg in reversed(recent_messages):  # Reverse to get chronological order
                            prefix = "User: " if msg.role == 'user' else "Assistant: "
                            # Don't truncate the messages - get full context
                            context_messages.append(f"{prefix}{msg.content}")
                        
                        conversation_context = "\n\n".join(context_messages)
                        print(f"Retrieved conversation context with {len(recent_messages)} messages")
                    
                    # Try to get memory buffer for additional context summary
                    try:
                        memory_buffer = ConversationMemoryBuffer.objects.get(conversation=conversation)
                        if memory_buffer.context_summary:
                            conversation_context += f"\n\n--- Conversation Summary ---\n{memory_buffer.context_summary}"
                    except ConversationMemoryBuffer.DoesNotExist:
                        # No memory buffer yet, that's fine
                        pass
                        
                except ChatHistory.DoesNotExist:
                    # No conversation yet, that's fine
                    pass
           
            # Skip document validation if in general chat mode or if URLs are present
            if not general_chat_mode and not has_url_content:
                # First, check for active document in session
                active_document_id = request.session.get('active_document_id')
               
                if not selected_documents:
                    active_document_id = request.session.get('active_document_id')
                    if active_document_id:
                        selected_documents = [active_document_id]
               
                # Validate document selection only when not in general chat mode and no URLs are provided
                if not selected_documents:
                    return Response(
                        {'error': 'Please select at least one document or set an active document'},
                        status=status.HTTP_400_BAD_REQUEST
                    )
 
            # Log the inputs for debugging
            print(f"Chat request received")
            print(f"Message: {message}")
            print(f"General chat mode: {general_chat_mode}")
            print(f"Use web knowledge: {use_web_knowledge}")
            print(f"Selected documents: {selected_documents}")
            print(f"Response format: {response_format}")
            print(f"Response length: {response_length}")
            print(f"Has conversation context: {bool(conversation_context)}")
            print(f"Has URL content: {has_url_content}")
 
            # Process document search early to get context for format detection
            all_chunks = []
            content_sources = []
            similar_contents = []  # Store document content separately
            citation_sources = {}  # Initialize citation sources dict
           
            if not general_chat_mode:
                try:
                    processed_docs = ProcessedIndex.objects.filter(
                        document_id__in=selected_documents,
                        document__user=user
                    )
                   
                    if not processed_docs.exists() and not has_url_content:
                        return Response(
                            {'error': 'No valid documents found'},
                            status=status.HTTP_404_NOT_FOUND
                        )
                       
                    # Log the processed docs for debugging
                    print(f"Found {processed_docs.count()} processed documents")
                    for doc in processed_docs:
                        print(f"Document: {doc.document.filename}")
                        print(f"Storage type: {doc.storage_type}")
                        print(f"Chunks count: {doc.chunks_count}")
                        if doc.storage_type == 'faiss':
                            print(f"FAISS index path: {doc.faiss_index}")
                            print(f"Metadata path: {doc.metadata}")
                        print(f"Markdown path: {doc.markdown_path}")
                       
                except Exception as e:
                    print(f"Error fetching documents: {str(e)}")
                    if not has_url_content:
                        return Response(
                            {'error': f'Document retrieval error: {str(e)}'},
                            status=status.HTTP_400_BAD_REQUEST
                        )
 
                # Updated document processing for pgvector compatibility
                all_metadata_store = []
               
                # Check storage types and process accordingly
                pgvector_docs = []
                faiss_docs = []
               
                for proc_doc in processed_docs:
                    if proc_doc.storage_type == 'pgvector':
                        # Check if embeddings exist in database
                        embeddings_exist = DocumentEmbedding.objects.filter(document=proc_doc.document).exists()
                        if embeddings_exist:
                            pgvector_docs.append(proc_doc)
                            print(f"Added pgvector document: {proc_doc.document.filename}")
                        else:
                            print(f"Warning: pgvector document {proc_doc.document.filename} has no embeddings")
                    else:
                        # FAISS document - check for file existence
                        if (proc_doc.faiss_index and proc_doc.metadata and
                            os.path.exists(proc_doc.faiss_index) and os.path.exists(proc_doc.metadata)):
                            faiss_docs.append(proc_doc)
                            print(f"Added FAISS document: {proc_doc.document.filename}")
                        else:
                            print(f"Warning: FAISS files missing for {proc_doc.document.filename}")
               
                # Process search based on available storage types
                if pgvector_docs or faiss_docs:
                    print(f"Searching in {len(pgvector_docs)} pgvector docs and {len(faiss_docs)} FAISS docs")
                   
                    # Use the updated search function that handles both storage types
                    similar_contents, content_sources, chunk_citation_sources = self.search_similar_content(
                        modified_message if has_url_content else message,
                        processed_docs,  # Pass all processed docs, function will handle different storage types
                        None  # metadata_store not needed for pgvector
                    )
                   
                    # Store citation sources
                    citation_sources = chunk_citation_sources
                   
                    if similar_contents:
                        print(f"Found {len(similar_contents)} relevant content chunks")
                        all_chunks = [{'text': content, 'source': source} for content, source in zip(similar_contents, content_sources)]
                    else:
                        print("No relevant content found in documents")
                        all_chunks = []
                       
                        # If pgvector search failed, try fallback text search
                        if pgvector_docs:
                            print("Attempting fallback text search in pgvector documents...")
                            fallback_results = self.fallback_text_search_pgvector(
                                modified_message if has_url_content else message,
                                pgvector_docs
                            )
                            if fallback_results:
                                similar_contents, content_sources, citation_sources = fallback_results
                                all_chunks = [{'text': content, 'source': source} for content, source in zip(similar_contents, content_sources)]
                                print(f"Fallback search found {len(similar_contents)} results")
                else:
                    print("No valid documents found for search")
                    all_chunks = []
           
            # Check if we should auto-detect the response format
            if response_format == 'auto_detect':
                context_snippets = [chunk.get('text', '') for chunk in all_chunks[:5]] if all_chunks else []
                detected_format, secondary_format, confidence = self.detect_question_format(modified_message if has_url_content else message, context_snippets)
               
                # Use the detected format if confidence is reasonable, otherwise default to 'natural'
                if confidence >= 5.0:
                    response_format = detected_format
                    print(f"Auto-detected format: {response_format} (confidence: {confidence})")
                else:
                    response_format = 'natural'
                    print(f"Low confidence in format detection ({confidence}), defaulting to 'natural'")
           
            # Get answer based on mode
            web_knowledge_response = None
            web_sources = []
            doc_citation_sources = {}
           
            # Get answer based on mode
            if general_chat_mode:
                # For general chat mode, consider URL content if available
                if has_url_content:
                    # Create a hybrid prompt that includes the URL content
                    augmented_query = f"""
                    Please answer this question: {modified_message}
                   
                    The question references content from URLs, which is provided here:
                   
                    {url_context}
                   
                    Based on this URL content, answer the user's question.
                    """
                   
                    # Use the existing general chat function but with our augmented query
                    answer, token_usage = self.get_general_chat_answer(augmented_query, use_web_knowledge, response_length, response_format,conversation_context, user=user)
                    print("Generated response using general chat mode with URL content")
                else:
                    # Standard general chat response without URL content
                    answer, token_usage = self.get_general_chat_answer(message, use_web_knowledge, response_length, response_format, conversation_context, user=user)
                    print("Generated response using general chat mode")
            else:
                # Document chat mode - now handles URL content too
                document_content_exists = bool(all_chunks)

                # Initialize token usage tracking
                token_usage = {'input_tokens': 0, 'output_tokens': 0, 'total_tokens': 0}
               
                if document_content_exists or has_url_content:
                    # If we have URL content, add it to the context
                    if has_url_content:
                        # Add URL context to similar_contents and content_sources
                        combined_contents = similar_contents.copy() if similar_contents else []
                        combined_contents.append(url_context)
                       
                        combined_sources = content_sources.copy() if content_sources else []
                        combined_sources.append("URL Content")
                       
                        # Add URL to citation sources
                        if citation_sources:
                            url_citation_key = max(citation_sources.keys()) + 1 if citation_sources else 1
                            citation_sources[url_citation_key] = {
                                'source_file': 'URL Content',
                                'text': url_context[:500] + "..." if len(url_context) > 500 else url_context,
                                'document_id': 'url_content',
                                'score': 0.9,  # High score for URL content
                                'display_num': url_citation_key
                            }
                       
                        # Generate response using combined context with citation sources
                        if response_length == 'short':
                            document_answer, doc_citation_sources = self.generate_short_response(
                                modified_message,
                                combined_contents,
                                combined_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                        else:  # Default to comprehensive
                            document_answer, doc_citation_sources = self.generate_response(
                                modified_message,
                                combined_contents,
                                combined_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                       
                        print(f"Generated document-based answer with URL content using {response_length} response length")
                    else:
                        # Standard document-based response without URL content
                        if response_length == 'short':
                            document_answer, doc_citation_sources, token_usage = self.generate_short_response(
                                message,
                                similar_contents,
                                content_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                        else:  # Default to comprehensive
                            document_answer, doc_citation_sources, token_usage = self.generate_response(
                                message,
                                similar_contents,
                                content_sources,
                                False,
                                response_format,
                                conversation_context
                            )
                       
                        print(f"Generated document-based answer using {response_length} response length")
                        print(f"Token usage - Input: {token_usage.get('input_tokens', 0)}, Output: {token_usage.get('output_tokens', 0)}")
                                    
                    # If web knowledge is requested, get web response using document context to enhance the query
                    if use_web_knowledge:
                        web_response_data = self.get_web_knowledge_response(
                            modified_message if has_url_content else message,
                            user=user,
                            document_context=similar_contents  # Pass document context to enhance web search  
                        )
 
                        # Extract the components from the JSON response
                        if isinstance(web_response_data, dict):
                            web_knowledge_response = web_response_data.get("content", "No web response received")
                            web_sources = web_response_data.get("web_sources", [])
                        else:
                            # Fallback for backward compatibility
                            web_knowledge_response = str(web_response_data)
                            web_sources = []
                        print(f"Web knowledge response received with document context, source count: {len(web_sources)}")
                       
                        # Combine document and web responses
                   
                        combined_response_data, combined_citation_sources, token_usage = self.combine_document_and_web_responses(
                            modified_message if has_url_content else message,
                            document_answer,
                            web_knowledge_response,
                            combined_sources if has_url_content else content_sources,
                            web_sources,
                            response_format,
                            conversation_context,
                            original_doc_context=similar_contents,
                            doc_token_usage=token_usage
                        )
 
                        # Extract the answer from the JSON response
                        if isinstance(combined_response_data, dict):
                            answer = combined_response_data.get("content", "Error combining responses")
                        else:
                            answer = str(combined_response_data)
                        print("Combined document and web responses")
                       
                        # Update citation sources with combined sources
                        doc_citation_sources = combined_citation_sources
                    else:
                        # Just use document answer
                        answer = document_answer
                else:
                    # No document content found
                    if use_web_knowledge:
                        # Only web knowledge available
                        web_response_data = self.get_web_knowledge_response(
                            modified_message if has_url_content else message,
                            user=user
                        )
 
                        if isinstance(web_response_data, dict):
                            answer = web_response_data.get("content", "No web response received")
                            web_sources = web_response_data.get("web_sources", [])
                        else:
                            answer = str(web_response_data)
                            web_sources = []
                        print("Using web knowledge response as document search returned no results")
                    else:
                        print("No document content found and web knowledge not enabled")
                        answer = "I couldn't find any relevant information in the documents."
           
            # Extract main response and citation info (if any)
            if "\n\n*Sources:" in answer:
                parts = answer.split("\n\n*Sources:")
                clean_response = parts[0]
                source_info = parts[1].split('*\n')[0] if len(parts) > 1 else ""
            else:
                clean_response = answer
                source_info = "Web search results" if use_web_knowledge else ""
               
                # Add URL sources if applicable
                if has_url_content:
                    urls_domains = [self.extract_domain(url) for url in extracted_urls]
                    url_sources = ", ".join(urls_domains)
                    if source_info:
                        source_info += f", URLs: {url_sources}"
                    else:
                        source_info = f"URLs: {url_sources}"
 
            clean_response = self.normalize_citation_markers(clean_response)
           
            # Generate follow-up questions (either from documents or general chat)
            if general_chat_mode:
                follow_up_questions = self.generate_general_follow_up_questions(
                    modified_message if has_url_content else message,
                    clean_response,
                    user=user
                )
            else:
                if all_chunks:
                    context_texts = [chunk.get('text', '') for chunk in all_chunks[:3]]
                    print("Generating follow-up questions...")
                    follow_up_questions = self.generate_follow_up_questions(context_texts, user=user)
                    print("Follow-up questions generated:", follow_up_questions)
                else:
                    follow_up_questions = [
                        "What else would you like to know about this content?",
                        "Would you like me to elaborate on any specific point?",
                        "Do you have any other questions I can help with?"
                    ]
           
            # Create citations from chunks and citation sources
            citations = []
            if not general_chat_mode:
                # First add citations from doc_citation_sources (our enhanced citations)
                for citation_num, citation_info in doc_citation_sources.items():
                    # Format the snippet to properly handle tables
                    formatted_snippet = self.format_citation_content(citation_info.get('snippet', ''))
                    citations.append({
                        'source_file': citation_info.get('source_file', 'Unknown'),
                        'page_number': citation_info.get('page_number', 'Unknown'),
                        'section_title': citation_info.get('section_title', 'Unknown'),
                        'snippet': formatted_snippet,
                        'document_id': citation_info.get('document_id', 'Unknown')
                    })
               
                # If no enhanced citations, fall back to chunk-based citations (original method)
                if not citations:
                    for chunk in all_chunks:
                        chunk_text = chunk.get('text', '')
                        formatted_snippet = self.format_citation_content(chunk_text)
                        citations.append({
                            'source_file': chunk.get('source', 'Unknown'),
                            'page_number': chunk.get('chunk_id', 'Unknown'),
                            'section_title': 'Unknown',
                            'snippet': formatted_snippet[:200] + "..." if len(formatted_snippet) > 200 else formatted_snippet,
                            'document_id': next((str(doc.document.id) for doc in processed_docs if doc.document.filename == chunk.get('source')), 'Unknown')
                        })
               
                # Add URL citations if applicable
                if has_url_content:
                    for idx, url in enumerate(extracted_urls):
                        domain = self.extract_domain(url)
                        citations.append({
                            'source_file': domain,
                            'page_number': 'URL',
                            'section_title': 'URL Content',
                            'snippet': f"Content from {domain}",
                            'document_id': f"url_{idx}",
                            'url': url
                        })
               
                # Add web citations if applicable
                if use_web_knowledge and web_sources:
                    for idx, source in enumerate(web_sources):
                        citations.append({
                            'source_file': source.get('title', 'Web Source'),
                            'page_number': 'Web',
                            'section_title': 'Web Search Result',
                            'snippet': source.get('snippet', '')[:200] + "..." if source.get('snippet', '') else "Web content",
                            'document_id': f"web_{idx}",
                            'url': source.get('url', '')
                        })
 
            # Prepare conversation details
            conversation_id = conversation_id or str(uuid.uuid4())
            title = f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"
 
            # Create or retrieve conversation
            if conversation_id:
                conversation, created = ChatHistory.objects.get_or_create(
                    user=user,
                    conversation_id=conversation_id,
                    main_project_id=main_project_id,
                    defaults={
                        'title': f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                        'is_active': True,
                        'follow_up_questions': follow_up_questions,
                    }
                )
 
            if created:
                # Create conversation transaction with token limit check
                transaction_result = self.create_conversation_transaction(
                    user, conversation, main_project_id,
                    use_web_knowledge, response_format, response_length,
                    token_usage  # Pass token usage
                )
                
                # Check if token limit was exceeded
                if isinstance(transaction_result, dict) and transaction_result.get('error'):
                    # Instead of deleting conversation and returning HTTP error,
                    # create user message and error response message
                    
                    # Create user message
                    user_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='user',
                        content=message
                    )
                    
                    # Create AI error response message 
                    error_message = transaction_result.get('message', 'Token limit exceeded')
                    ai_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='assistant',
                        content=error_message,
                        sources='',
                        citations=[],
                        use_web_knowledge=use_web_knowledge,
                        response_length=response_length,
                        response_format=response_format,
                    )
                    
                    # Update memory buffer
                    memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                        conversation=conversation
                    )
                    
                    # Get all messages for this conversation
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                    
                    # Update memory with all messages
                    memory_buffer.update_memory(all_messages)
                    
                    # Update conversation details
                    conversation.title = title
                    conversation.follow_up_questions = []  # No follow-up questions for error
                    conversation.save()
                    
                    # Prepare response data
                    response_data = {
                        'response': error_message,
                        'follow_up_questions': [],
                        'conversation_id': str(conversation.conversation_id),
                        'citations': [],
                        'active_document_id': request.session.get('active_document_id') if not general_chat_mode else None,
                        'sources_info': '',
                        'use_web_knowledge': use_web_knowledge,
                        'general_chat_mode': general_chat_mode,
                        'response_length': response_length,
                        'response_format': response_format,
                        'url_content_used': False,
                        'extracted_urls': [],
                        'storage_types_used': []
                    }
                    
                    # Get all messages for this conversation (with IDs)
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                    
                    message_list = []
                    for msg in all_messages:
                        message_data = {
                            'id': msg.id,
                            'role': msg.role,
                            'content': msg.content,
                            'created_at': msg.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': msg.citations or [],
                            'sources_info': getattr(msg, 'sources', ''),
                            'extracted_urls': getattr(msg, 'extracted_urls', []),
                        }
                        # Add badge properties for assistant messages
                        if msg.role == 'assistant':
                            message_data['use_web_knowledge'] = getattr(msg, 'use_web_knowledge', False)
                            message_data['response_length'] = getattr(msg, 'response_length', 'comprehensive')
                            message_data['response_format'] = getattr(msg, 'response_format', 'natural')
                        message_list.append(message_data)
                    
                    response_data['messages'] = message_list
                    
                    # Return successful response with error message as content
                    return Response(response_data, status=status.HTTP_200_OK)
                    
            else:
                # Update existing conversation transaction with token limit check
                transaction_result = self.update_conversation_transaction(
                    conversation, use_web_knowledge, token_usage
                )
                
                # Check if token limit was exceeded
                if isinstance(transaction_result, dict) and transaction_result.get('error'):
                    # Create user message
                    user_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='user',
                        content=message
                    )
                    
                    # Create AI error response message
                    error_message = transaction_result.get('message', 'Token limit exceeded')
                    ai_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='assistant',
                        content=error_message,
                        sources='',
                        citations=[],
                        use_web_knowledge=use_web_knowledge,
                        response_length=response_length,
                        response_format=response_format,
                    )
                    
                    # Update memory buffer
                    memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                        conversation=conversation
                    )
                    
                    # Get all messages for this conversation
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                    
                    # Update memory with all messages
                    memory_buffer.update_memory(all_messages)
                    
                    # Update conversation details
                    conversation.title = title
                    conversation.follow_up_questions = []  # No follow-up questions for error
                    conversation.save()
                    
                    # Prepare response data
                    response_data = {
                        'response': error_message,
                        'follow_up_questions': [],
                        'conversation_id': str(conversation.conversation_id),
                        'citations': [],
                        'active_document_id': request.session.get('active_document_id') if not general_chat_mode else None,
                        'sources_info': '',
                        'use_web_knowledge': use_web_knowledge,
                        'general_chat_mode': general_chat_mode,
                        'response_length': response_length,
                        'response_format': response_format,
                        'url_content_used': False,
                        'extracted_urls': [],
                        'storage_types_used': []
                    }
                    
                    # Get all messages for this conversation (with IDs)
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                    
                    message_list = []
                    for msg in all_messages:
                        message_data = {
                            'id': msg.id,
                            'role': msg.role,
                            'content': msg.content,
                            'created_at': msg.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': msg.citations or [],
                            'sources_info': getattr(msg, 'sources', ''),
                            'extracted_urls': getattr(msg, 'extracted_urls', []),
                        }
                        # Add badge properties for assistant messages
                        if msg.role == 'assistant':
                            message_data['use_web_knowledge'] = getattr(msg, 'use_web_knowledge', False)
                            message_data['response_length'] = getattr(msg, 'response_length', 'comprehensive')
                            message_data['response_format'] = getattr(msg, 'response_format', 'natural')
                        message_list.append(message_data)
                    
                    response_data['messages'] = message_list
                    
                    # Return successful response with error message as content
                    return Response(response_data, status=status.HTTP_200_OK)
 
            # Create user message
            user_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='user',
                content=message
            )
 
            # Create AI response message
            ai_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='assistant',
                content=clean_response,
                sources=source_info,
                citations=citations,
                use_web_knowledge=use_web_knowledge,
                response_length=response_length,
                response_format=response_format,
            )
 
            # Update or create memory buffer
            memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                conversation=conversation
            )
           
            # Get all messages for this conversation
            all_messages = ChatMessage.objects.filter(
                chat_history=conversation
            ).order_by('created_at')
           
            # Update memory with all messages
            memory_buffer.update_memory(all_messages)
 
            # Add selected documents to the conversation (skip if general chat mode)
            if selected_documents and not general_chat_mode:
                documents = Document.objects.filter(
                    id__in=selected_documents,
                    user=user
                )
                conversation.documents.set(documents)
 
            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
 
            # Update conversation details
            conversation.title = title
            conversation.follow_up_questions = follow_up_questions
            conversation.save()
 
            # Use citation_sources in the response data if available
            response_data = {
                'response': clean_response,
                'follow_up_questions': follow_up_questions,
                'conversation_id': str(conversation.conversation_id),
                'citations': citations,
                'active_document_id': request.session.get('active_document_id') if not general_chat_mode else None,
                'sources_info': source_info,
                'use_web_knowledge': use_web_knowledge,
                'general_chat_mode': general_chat_mode,
                'response_length': response_length,  # Include in response for tracking
                'response_format': response_format,   # Include detected/used format in response
                'url_content_used': has_url_content,  # Flag if URL content was used
                'extracted_urls': extracted_urls if has_url_content else [],  # List of extracted URLs
                'storage_types_used': self.get_storage_types_used(processed_docs) if not general_chat_mode else []  # Debug info
            }
 
            # Print detailed chat response information
            print("\n--- Chat Interaction Logged ---")
            print(f"User Question: {message}")
            print(f"Mode: {'Web Knowledge' if use_web_knowledge else 'General Chat' if general_chat_mode else 'Document Chat'}")
            print(f"Format: {response_format}")
            print(f"Length: {response_length}")
            print(f"URL content used: {has_url_content}")
            if has_url_content:
                print(f"Extracted URLs: {extracted_urls}")
            if not general_chat_mode:
                storage_types = self.get_storage_types_used(processed_docs)
                print(f"Storage types used: {storage_types}")
            print(f"Assistant Response: {clean_response[:500]}...")  # First 500 chars
            print(f"Citation count: {len(citations)}")
            print(f"Follow-up Questions: {len(follow_up_questions)}")
            print("-----------------------------\n")
 
            # Get all messages for this conversation (with IDs)
            all_messages = ChatMessage.objects.filter(
                chat_history=conversation
            ).order_by('created_at')
 
            message_list = []
            for message in all_messages:
                message_data = {
                    'id': message.id,
                    'role': message.role,
                    'content': message.content,
                    'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                    'citations': message.citations or [],
                    'sources_info': getattr(message, 'sources', ''),  # Add this line
                    'extracted_urls': getattr(message, 'extracted_urls', []),  # Add this line
                }
                # Add badge properties for assistant messages
                if message.role == 'assistant':
                    message_data['use_web_knowledge'] = getattr(message, 'use_web_knowledge', False)
                    message_data['response_length'] = getattr(message, 'response_length', 'comprehensive')
                    message_data['response_format'] = getattr(message, 'response_format', 'natural')
                message_list.append(message_data)
 
            response_data['messages'] = message_list
 
            return Response(response_data, status=status.HTTP_200_OK)
 
        except Exception as e:
            logger.error(f"Unexpected error in ChatView: {str(e)}", exc_info=True)
            return Response(
                {'error': f'An unexpected error occurred: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
 
    # ===================================================================
    # Helper methods for pgvector integration
    # ===================================================================
   
    def get_storage_types_used(self, processed_docs):
        """Get list of storage types used in the current query for debugging"""
        storage_types = []
        for doc in processed_docs:
            if doc.storage_type not in storage_types:
                storage_types.append(doc.storage_type)
        return storage_types
   
    def fallback_text_search_pgvector(self, query, pgvector_docs, limit=10):
        """
        Fallback text search when pgvector similarity search fails
        """
        try:
            from django.db.models import Q
           
            query_lower = query.lower()
            all_results = []
            all_sources = []
            citation_sources = {}
           
            for proc_doc in pgvector_docs:
                # Search for text matches in document embeddings
                text_matches = DocumentEmbedding.objects.filter(
                    document=proc_doc.document,
                    content__icontains=query_lower
                )[:limit]
               
                for idx, embedding_obj in enumerate(text_matches):
                    citation_key = len(all_results) + 1
                   
                    citation_sources[citation_key] = {
                        'source': proc_doc.document.filename,
                        'text': embedding_obj.content,
                        'relevance_score': 0.5,  # Default score for text match
                        'document_id': proc_doc.document.id,
                        'chunk_idx': embedding_obj.chunk_id,
                        'display_num': citation_key
                    }
                   
                    all_results.append(embedding_obj.content)
                    all_sources.append(proc_doc.document.filename)
           
            if all_results:
                print(f"Fallback text search found {len(all_results)} results")
                return all_results, all_sources, citation_sources
            else:
                return None
               
        except Exception as e:
            print(f"Error in fallback text search: {str(e)}")
            return None
    def create_conversation_transaction(self, user, conversation, main_project_id, use_web_knowledge, response_format, response_length, token_usage=None):
        """Create transaction record for conversation with token limit validation"""
        try:
            # Get or create user API tokens record
            user_api_tokens, created = UserAPITokens.objects.get_or_create(
                user=user,
                defaults={
                    'page_limit': 20,
                    'token_limit': 1_000_000
                }
            )
            
            # Calculate current token usage for this user
            current_total_tokens = self.get_user_total_token_usage(user)
            
            # Get tokens for this request
            request_tokens = 0
            if token_usage:
                request_tokens = token_usage.get('total_tokens', 0)
                if request_tokens == 0:  # Fallback calculation
                    request_tokens = token_usage.get('input_tokens', 0) + token_usage.get('output_tokens', 0)
            
            # Check if adding these tokens would exceed the limit
            new_total = current_total_tokens + request_tokens
            
            print(f"Token usage check - Current: {current_total_tokens}, Request: {request_tokens}, New Total: {new_total}, Limit: {user_api_tokens.token_limit}")
            
            if new_total > user_api_tokens.token_limit:
                # Return error data instead of creating transaction
                remaining_tokens = max(0, user_api_tokens.token_limit - current_total_tokens)
                error_data = {
                    'error': 'token_limit_exceeded',
                    'message': f"You have reached your conversation token limit of {user_api_tokens.token_limit:,} tokens. Current usage: {current_total_tokens:,} tokens. Remaining: {remaining_tokens:,} tokens. The current message exceeds your remaining token quota and cannot be processed. Please limit your usage or request for reset at - contact@klarifai.ai",
                    'current_usage': current_total_tokens,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': remaining_tokens,
                    'requested_tokens': request_tokens
                }
                print(f"Token limit exceeded for user {user.id}: {new_total} > {user_api_tokens.token_limit}")
                return error_data
            
            # Get the project
            main_project = Project.objects.get(id=main_project_id, user=user)
            
            # Create main transaction record
            user_transaction = UserTransaction.objects.create(
                user=user,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                conversation_title=conversation.title,
                conversation_id=conversation.conversation_id,
                main_project=main_project,
                metadata={
                    'creation_timestamp': timezone.now().isoformat(),
                    'web_knowledge_used': use_web_knowledge,
                    'response_format': response_format,
                    'response_length': response_length,
                    'token_usage': token_usage or {},
                    'token_limit_check': {
                        'previous_total': current_total_tokens,
                        'request_tokens': request_tokens,
                        'new_total': new_total,
                        'limit': user_api_tokens.token_limit
                    }
                }
            )
            
            # Extract token information if available
            input_tokens = token_usage.get('input_tokens', 0) if token_usage else 0
            output_tokens = token_usage.get('output_tokens', 0) if token_usage else 0
            
            # Create detailed conversation transaction
            ConversationTransaction.objects.create(
                user_transaction=user_transaction,
                message_count=1,  # Initial message
                question_count=1,  # First question
                input_api_tokens=input_tokens,
                output_api_tokens=output_tokens,
                web_knowledge_used=use_web_knowledge,
                response_format=response_format,
                response_length=response_length
            )
            
            print(f"Transaction recorded for conversation: {conversation.title}")
            print(f"Token usage - Input: {input_tokens}, Output: {output_tokens}, Total: {input_tokens + output_tokens}")
            print(f"User total token usage: {new_total}/{user_api_tokens.token_limit}")
            
            # Return success data
            return {
                'success': True,
                'current_usage': new_total,
                'token_limit': user_api_tokens.token_limit,
                'remaining_tokens': user_api_tokens.token_limit - new_total
            }
            
        except Exception as e:
            print(f"Error creating conversation transaction: {str(e)}")
            return {
                'error': 'transaction_error',
                'message': f'Error creating conversation transaction: {str(e)}'
            }

    def get_user_total_token_usage(self, user):
        """Calculate total token usage for a user across all conversations"""
        from django.db.models import Sum,Count
        try:
            from django.db.models import Sum,Count
            # Get all conversation transactions for this user
            total_tokens = ConversationTransaction.objects.filter(
                user_transaction__user=user,
                user_transaction__is_active=True
            ).aggregate(
                total_input=Sum('input_api_tokens'),
                total_output=Sum('output_api_tokens')
            )
            
            input_total = total_tokens.get('total_input') or 0
            output_total = total_tokens.get('total_output') or 0
            
            return input_total + output_total
            
        except Exception as e:
            print(f"Error calculating user token usage: {str(e)}")
            return 0
 
    def update_conversation_transaction(self, conversation, use_web_knowledge, token_usage=None):
        """Update existing conversation transaction when new messages are added with token limit validation"""
        try:
            # Get user from conversation
            user = conversation.user
            
            # Get or create user API tokens record
            user_api_tokens, created = UserAPITokens.objects.get_or_create(
                user=user,
                defaults={
                    'page_limit': 20,
                    'token_limit': 1_000_000
                }
            )
            
            # Calculate current token usage for this user
            current_total_tokens = self.get_user_total_token_usage(user)
            
            # Get tokens for this request
            request_tokens = 0
            if token_usage:
                request_tokens = token_usage.get('total_tokens', 0)
                if request_tokens == 0:  # Fallback calculation
                    request_tokens = token_usage.get('input_tokens', 0) + token_usage.get('output_tokens', 0)
            
            # Check if adding these tokens would exceed the limit
            new_total = current_total_tokens + request_tokens
            
            print(f"Token usage update check - Current: {current_total_tokens}, Request: {request_tokens}, New Total: {new_total}, Limit: {user_api_tokens.token_limit}")
            
            if new_total > user_api_tokens.token_limit:
                # Return error data instead of updating transaction
                remaining_tokens = max(0, user_api_tokens.token_limit - current_total_tokens)
                error_data = {
                    'error': 'token_limit_exceeded',
                    'message': f"You have reached your conversation token limit of {user_api_tokens.token_limit:,} tokens. Current usage: {current_total_tokens:,} tokens. Remaining: {remaining_tokens:,} tokens. The current message exceeds your remaining token quota and cannot be processed. Please limit your usage or request for reset at - contact@klarifai.ai",
                    'current_usage': current_total_tokens,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': remaining_tokens,
                    'requested_tokens': request_tokens
                }
                print(f"Token limit exceeded during update for user {user.id}: {new_total} > {user_api_tokens.token_limit}")
                return error_data
            
            # Find existing transaction for this conversation
            user_transaction = UserTransaction.objects.filter(
                conversation_id=conversation.conversation_id,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                is_active=True
            ).first()
            
            if user_transaction and hasattr(user_transaction, 'notebook_conversation_details'):
                # Get the conversation details
                conv_details = user_transaction.notebook_conversation_details
                
                # Update message count and question count
                conv_details.message_count += 1
                conv_details.question_count += 1  # Each update represents a new question
                
                # Update token usage if provided
                if token_usage:
                    conv_details.input_api_tokens += token_usage.get('input_tokens', 0)
                    conv_details.output_api_tokens += token_usage.get('output_tokens', 0)
                
                if use_web_knowledge:
                    conv_details.web_knowledge_used = True
                    
                conv_details.save()
                
                # Update metadata with token usage
                if token_usage:
                    if 'token_history' not in user_transaction.metadata:
                        user_transaction.metadata['token_history'] = []
                    user_transaction.metadata['token_history'].append({
                        'timestamp': timezone.now().isoformat(),
                        'input_tokens': token_usage.get('input_tokens', 0),
                        'output_tokens': token_usage.get('output_tokens', 0),
                        'total_tokens': token_usage.get('total_tokens', 0)
                    })
                
                user_transaction.metadata['last_message_timestamp'] = timezone.now().isoformat()
                user_transaction.metadata['token_limit_check'] = {
                    'previous_total': current_total_tokens,
                    'request_tokens': request_tokens,
                    'new_total': new_total,
                    'limit': user_api_tokens.token_limit
                }
                user_transaction.save()
                
                print(f"Updated conversation transaction - Questions: {conv_details.question_count}, Total tokens: {conv_details.total_tokens_used}")
                print(f"User total token usage: {new_total}/{user_api_tokens.token_limit}")
                
                # Return success data
                return {
                    'success': True,
                    'current_usage': new_total,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': user_api_tokens.token_limit - new_total
                }
                
        except Exception as e:
            print(f"Error updating conversation transaction: {str(e)}")
            return {
                'error': 'transaction_error',
                'message': f'Error updating conversation transaction: {str(e)}'
            }

    def _parse_json_response(self, response_content, fallback_message="Error processing response"):
        """Helper method to safely parse JSON responses from LLM"""
        try:
            return json.loads(response_content)
        except json.JSONDecodeError as e:
            logger.error(f"JSON parsing failed: {str(e)}")
            return {
                "content": fallback_message,
                "error": True,
                "format_type": "error"
            }
   

    def format_citation_content(self, snippet):
        """
        Format citation snippets to optimize display of complex tables with long content
        Designed to work with any table structure from any file type
        """
        if not snippet:
            return snippet
            
        # Generic table pattern detection - using structural indicators, not content-specific terms
        has_table_pattern = (
            # Standard markdown table patterns
            ('|' in snippet and any(separator in snippet for separator in ['|---', '| ---', '|----', '------'])) or
            # Table rows with pipes at beginning and end
            ('|' in snippet and any(row.strip().startswith('|') and row.strip().endswith('|') for row in snippet.split('\n'))) or
            # Any structured data that resembles a table (multiple lines with consistent pipe counts)
            self._detect_tabular_structure(snippet)
        )
        
        if has_table_pattern:
            try:
                # Start with basic cleanup
                content = snippet
                
                # Replace emoji arrows with text versions for better compatibility
                content = content.replace('🔻', '↓').replace('🔼', '↑')
                
                # Step 1: Process section headers to ensure they're properly formatted
                # Match any line starting with # and make sure there's a newline after it
                content = re.sub(r'(#+\s*[^\n]+)(?![\n])', r'\1\n', content)
                
                # Step 2: Make sure all headers have proper spacing
                content = re.sub(r'(#+)(\S)', r'\1 \2', content)
                
                # Step 3: Detect and format tabular structures
                lines = content.split('\n')
                formatted_lines = []
                
                in_table = False
                table_lines = []
                table_header_detected = False
                pending_table_header = None
                
                for i, line in enumerate(lines):
                    stripped_line = line.strip()
                    
                    # Check if this is a header followed by a potential table
                    is_header = stripped_line.startswith('#')
                    next_line_is_table = (i < len(lines) - 1 and 
                                        '|' in lines[i+1] and 
                                        lines[i+1].strip().startswith('|') and 
                                        lines[i+1].strip().endswith('|'))
                    
                    # Handle headers that might be associated with tables
                    if is_header and next_line_is_table:
                        formatted_lines.append(line)
                        continue
                    
                    # Check if this is a table line using structural patterns
                    is_table_line = '|' in stripped_line and stripped_line.startswith('|') and stripped_line.endswith('|')
                    
                    # Check if this is a separator line
                    is_separator = is_table_line and any(c in stripped_line for c in ['-', '='])
                    
                    # Check for potential table headers based on structure (first row of pipe-delimited content)
                    is_potential_header = (is_table_line and 
                                        not is_separator and 
                                        not in_table and
                                        i < len(lines) - 1 and
                                        '|' in lines[i+1])
                    
                    # Fix incomplete tables missing separators
                    if is_potential_header and i < len(lines) - 1:
                        next_line = lines[i+1].strip()
                        is_next_separator = '|' in next_line and any(c in next_line for c in ['-', '='])
                        
                        if '|' in next_line and not is_next_separator:
                            # This looks like a header row without a separator, let's add one
                            pending_table_header = line
                            table_header_detected = True
                            continue
                    
                    if table_header_detected and is_table_line and not is_separator:
                        # We found a data row after a header without separator, add the header and create a separator
                        if pending_table_header:
                            formatted_lines.append(pending_table_header)
                            # Create a separator row based on the header
                            parts = [p.strip() for p in pending_table_header.split('|')]
                            separator = '|'
                            for part in parts[1:-1]:  # Skip first and last empty parts
                                separator += ' ' + '-' * len(part) + ' |'
                            formatted_lines.append(separator)
                            pending_table_header = None
                            table_header_detected = False
                            in_table = True
                        formatted_lines.append(line)
                        continue
                    
                    if is_table_line:
                        if not in_table:
                            # Start of a new table
                            in_table = True
                            # If we have pending lines, add them
                            if formatted_lines and not formatted_lines[-1].strip() == '':
                                formatted_lines.append('')  # Empty line before table
                        
                        # Add to current table
                        table_lines.append(line)
                    else:
                        if in_table:
                            # End of table, process it
                            in_table = False
                            
                            # Process and optimize the table
                            optimized_table = self.optimize_table_structure(table_lines)
                            
                            # Add the optimized table
                            formatted_lines.extend(optimized_table)
                            
                            if not stripped_line:  # If this is an empty line
                                formatted_lines.append('')  # Empty line after table
                            else:
                                # If next line has content, add a separator
                                formatted_lines.append('')  # Empty line after table
                                formatted_lines.append(line)  # Add the current non-table line
                            
                            table_lines = []
                            continue
                        
                        # Add non-table line
                        formatted_lines.append(line)
                
                # Don't forget to process any table that's at the end
                if in_table and table_lines:
                    optimized_table = self.optimize_table_structure(table_lines)
                    formatted_lines.extend(optimized_table)
                
                # Join everything back together
                content = '\n'.join(formatted_lines)
                
                # Final cleanup, remove excessive newlines
                content = re.sub(r'\n{3,}', '\n\n', content)
                
                return content
                
            except Exception as e:
                print(f"Error formatting complex table content: {str(e)}")
                # Return original content on error
                return snippet
        
        # If not a table or processing failed, return the original content
        return snippet

    def _detect_tabular_structure(self, text):
        """
        Generic detection of tabular structures in text.
        Looks for consistent patterns of delimiter characters that indicate a table.
        """
        if not text or '|' not in text:
            return False
            
        lines = text.split('\n')
        pipe_counts = []
        
        # Count pipes in consecutive non-empty lines
        consecutive_lines = 0
        for line in lines:
            if not line.strip():
                consecutive_lines = 0
                continue
                
            if '|' in line:
                pipe_count = line.count('|')
                pipe_counts.append(pipe_count)
                consecutive_lines += 1
                
                # If we have at least 3 consecutive lines with pipes, check pattern
                if consecutive_lines >= 3:
                    # Check if pipe counts are consistent (allowing small variations)
                    unique_counts = set(pipe_counts[-3:])
                    if len(unique_counts) <= 2 and max(unique_counts) - min(unique_counts) <= 2:
                        return True
        
        # Alternative detection: Check if we have a header-like row followed by a separator-like row
        for i in range(len(lines) - 1):
            current_line = lines[i].strip()
            next_line = lines[i+1].strip()
            
            if (current_line.startswith('|') and current_line.endswith('|') and
                next_line.startswith('|') and next_line.endswith('|') and
                '-' in next_line):
                return True
        
        return False

    def optimize_table_structure(self, table_lines):
        """
        Process and optimize a table's structure for better display
        Works with any table format from any file type
        """
        if not table_lines:
            return []
        
        # Find header and separator rows
        header_idx = -1
        separator_idx = -1
        
        for i, line in enumerate(table_lines):
            if '-' in line or '=' in line:
                separator_idx = i
                header_idx = max(0, i - 1)
                break
        
        # If no separator found, try to identify the header row and create a separator
        if separator_idx == -1 and len(table_lines) >= 2:
            # Assume first row is header if we have at least two rows with same number of columns
            header_row = table_lines[0]
            header_parts = header_row.split('|')
            
            first_data_row = table_lines[1]
            data_parts = first_data_row.split('|')
            
            # If first and second row have similar structure, assume first is header
            if abs(len(header_parts) - len(data_parts)) <= 1:
                header_idx = 0
                # Create a separator row
                parts = header_row.split('|')
                separator = '|'
                for part in parts[1:-1]:  # Skip first and last empty parts
                    separator += ' ' + '-' * len(part.strip()) + ' |'
                
                # Insert the separator row
                table_lines.insert(1, separator)
                separator_idx = 1
        
        # If we still don't have a separator, check if we have consistent columns
        if separator_idx == -1:
            # Try to convert any table-like structure to proper format
            has_consistent_columns = True
            if len(table_lines) >= 2:
                col_count = len(table_lines[0].split('|'))
                for line in table_lines[1:]:
                    if len(line.split('|')) != col_count:
                        has_consistent_columns = False
                        break
                
                if has_consistent_columns:
                    # Create a proper table with first row as header
                    header_idx = 0
                    separator = '|'
                    header_parts = table_lines[0].split('|')
                    for part in header_parts[1:-1]:
                        separator += ' ' + '-' * len(part.strip()) + ' |'
                    
                    table_lines.insert(1, separator)
                    separator_idx = 1
        
        # If still no structure found, create a default structure
        if header_idx == -1 or separator_idx == -1:
            # The simplest approach: consider the first row as header
            if len(table_lines) >= 1:
                header_idx = 0
                
                # Create and insert a separator row
                header_parts = table_lines[0].split('|')
                separator = '|'
                for part in header_parts[1:-1]:  # Skip first and last empty parts
                    separator += ' ' + '-' * len(part.strip()) + ' |'
                
                if len(table_lines) > 1:
                    table_lines.insert(1, separator)
                    separator_idx = 1
                else:
                    table_lines.append(separator)
                    separator_idx = 1
        
        # Enhancement: Format cells with indicators, preserving the original indicator characters
        for i in range(len(table_lines)):
            # Skip header and separator
            if i != header_idx and i != separator_idx:
                line = table_lines[i]
                parts = line.split('|')
                formatted_parts = []
                
                for part in parts:
                    trimmed = part.strip()
                    
                    # Detect any trend indicators without hardcoding specific symbols
                    # This looks for common directional symbols in a generic way
                    has_up_indicator = any(up_symbol in trimmed for up_symbol in ['↑', '▲', '🔼', '⬆'])
                    has_down_indicator = any(down_symbol in trimmed for down_symbol in ['↓', '▼', '🔻', '⬇'])
                    
                    if has_up_indicator or has_down_indicator:
                        # Keep original format but ensure consistent spacing
                        formatted_parts.append(' ' + trimmed + ' ')
                    else:
                        formatted_parts.append(part)
                
                table_lines[i] = '|'.join(formatted_parts)
        
        # Analyze table for column consistency
        column_counts = [len(line.split('|')) - 1 for line in table_lines]  # -1 because n+1 pipes make n columns
        
        if len(set(column_counts)) > 1:
            # Inconsistent column counts - normalize
            max_columns = max(column_counts)
            
            for i in range(len(table_lines)):
                parts = table_lines[i].split('|')
                current_columns = len(parts) - 1
                
                if current_columns < max_columns:
                    # Add missing columns
                    if table_lines[i].endswith('|'):
                        # Add columns before the trailing |
                        table_lines[i] = table_lines[i][:-1] + ('| ' * (max_columns - current_columns)) + '|'
                    else:
                        # Add columns including trailing |
                        table_lines[i] = table_lines[i] + ('| ' * (max_columns - current_columns)) + '|'
        
        # Check for long content cells and analyze column widths
        has_long_cells = False
        max_col_lengths = []
        
        # Initialize with header column lengths if header exists
        if header_idx >= 0 and header_idx < len(table_lines):
            header_parts = table_lines[header_idx].split('|')[1:-1]  # Skip first and last
            max_col_lengths = [len(part.strip()) for part in header_parts]
        else:
            # Initialize with zeros if no header
            max_columns = max(len(line.split('|')) - 1 for line in table_lines)
            max_col_lengths = [0] * max_columns
        
        # Update with actual content lengths
        for line in table_lines:
            if separator_idx >= 0 and line == table_lines[separator_idx]:
                continue
                
            parts = line.split('|')[1:-1]  # Skip first and last
            
            for i, part in enumerate(parts):
                if i < len(max_col_lengths):
                    max_col_lengths[i] = max(max_col_lengths[i], len(part.strip()))
                    if len(part.strip()) > 20:
                        has_long_cells = True
                elif i >= len(max_col_lengths) and len(parts) > len(max_col_lengths):
                    # Extend max_col_lengths if needed
                    max_col_lengths.extend([0] * (len(parts) - len(max_col_lengths)))
                    max_col_lengths[i] = len(part.strip())
                    if len(part.strip()) > 20:
                        has_long_cells = True
        
        # Optimize the formatting if we have long cells or inconsistent structure
        if has_long_cells or len(set(column_counts)) > 1:
            # Format the table with consistent column widths
            formatted_table = []
            
            for i, line in enumerate(table_lines):
                if i == separator_idx:
                    # Special handling for separator
                    parts = ['']  # Start with empty part for leading pipe
                    for j, max_length in enumerate(max_col_lengths):
                        parts.append(' ' + '-' * max_length + ' ')
                    parts.append('')  # End with empty part for trailing pipe
                    formatted_table.append('|'.join(parts))
                else:
                    parts = line.split('|')
                    formatted_parts = ['']  # Start with empty part for leading pipe
                    
                    for j in range(len(max_col_lengths)):
                        if j + 1 < len(parts):
                            # Part exists in current row
                            part = parts[j+1]
                            # Clean up whitespace but maintain padding
                            cleaned = part.strip()
                            # Ensure consistent width formatting
                            formatted_parts.append(' ' + cleaned.ljust(max_col_lengths[j]) + ' ')
                        else:
                            # Missing part in current row
                            formatted_parts.append(' ' + ' ' * max_col_lengths[j] + ' ')
                    
                    formatted_parts.append('')  # End with empty part for trailing pipe
                    formatted_table.append('|'.join(formatted_parts))
            
            return formatted_table
        
        # If no special formatting needed, just return the original table lines
        return table_lines

    
    def scrape_webpage(self, url):
        """Scrape content from a webpage"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36 Edg/91.0.864.59'
            }

            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # Remove script and style elements
            for script_or_style in soup(['script', 'style', 'header', 'footer', 'nav']):
                script_or_style.decompose()
                
            # Extract title
            title = soup.title.string if soup.title else "No title found"
            
            # Extract main content - focus on article, main, or div with content
            content_elements = soup.select('article, main, .content, #content, .post, .entry')
            
            if not content_elements:
                # If no specific content elements, get all paragraphs
                paragraphs = soup.find_all('p')
                content = ' '.join([p.get_text() for p in paragraphs])
            else:
                # Use the first content element found
                content = content_elements[0].get_text()
            
            # Clean the content
            content = re.sub(r'\s+', ' ', content).strip()
            
            return {
                'title': title,
                'content': content[:15000],  # Limit content length
                'url': url
            }
        except Exception as e:
            logger.error(f"Error scraping webpage {url}: {str(e)}")
            return {
                'title': 'Error scraping page',
                'content': f"Could not retrieve content: {str(e)}",
                'url': url
            }

    def extract_urls_from_text(self, text):
        """
        Extract URLs from text using regex pattern matching.
        
        Args:
            text (str): The text to extract URLs from
            
        Returns:
            list: List of extracted URLs
        """
        import re
        
        # Pattern to match URLs - handles http, https, www formats
        url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!./?=&#+]*)*'
        www_pattern = r'www\.(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!./?=&#+]*)*'
        
        # Find all matches
        http_urls = re.findall(url_pattern, text)
        www_urls = re.findall(www_pattern, text)
        
        # Add http:// prefix to www URLs if they aren't already prefixed
        processed_www_urls = []
        for url in www_urls:
            if not any(url in http_url for http_url in http_urls):
                processed_www_urls.append(f"http://{url}")
        
        # Combine and return all URLs
        return http_urls + processed_www_urls

    def process_urls_in_query(self, query):
        """
        Process query to extract URLs, scrape their content, and return enhanced context.
        
        Args:
            query (str): The user's query that may contain URLs
            
        Returns:
            tuple: (modified_query, url_context, extracted_urls)
        """
        # Extract URLs from the query
        extracted_urls = self.extract_urls_from_text(query)
        
        if not extracted_urls:
            return query, "", []
        
        # Log found URLs
        logger.info(f"Found {len(extracted_urls)} URLs in query: {extracted_urls}")
        
        # Scrape content from each URL
        url_contents = []
        successful_urls = []
        
        for url in extracted_urls:
            try:
                scraped_content = self.scrape_webpage(url)
                if scraped_content and scraped_content.get('content'):
                    url_contents.append({
                        'url': url,
                        'title': scraped_content.get('title', 'No title'),
                        'content': scraped_content.get('content')
                    })
                    successful_urls.append(url)
                    logger.info(f"Successfully scraped content from {url}")
                else:
                    logger.warning(f"No content retrieved from {url}")
            except Exception as e:
                logger.error(f"Error scraping URL {url}: {str(e)}")
        
        if not url_contents:
            return query, "", extracted_urls
        
        # Build context from scraped content
        url_context = "CONTENT FROM URLS IN QUERY:\n\n"
        
        for content in url_contents:
            domain = self.extract_domain(content['url'])
            url_context += f"Source: {domain}\n"
            url_context += f"Title: {content['title']}\n"
            # Limit content length to avoid exceeding token limits
            url_context += f"Content: {content['content'][:5000]}...\n\n"
        
        # Create a modified query that references the URLs without including them
        modified_query = query
        for url in extracted_urls:
            # Replace URLs with a more readable reference
            domain = self.extract_domain(url)
            modified_query = modified_query.replace(url, f"[content from {domain}]")
        
        return modified_query, url_context, extracted_urls


    def get_web_knowledge_response(self, query, user=None, document_context=None):
        """
        Search the web for information and return direct search results using enhanced Gemini search.
        Maintains original functionality while using new Gemini approach with proper citations.
   
        Args:
            query (str): The user's query
            document_context (list, optional): List of context chunks from document search
            user: User object for API token access
       
        Returns:
            tuple: (response, sources)
        """
        try:
            # Step 1: Extract keywords and context from document context if available
            enhanced_query = query
            doc_context_text = ""
       
            if document_context and len(document_context) > 0:
                # Combine all document context into a single text for analysis
                combined_context = " ".join(document_context[:3])  # Use top 3 context chunks
           
                # Extract keywords using a smarter prompt that identifies entity types
                keyword_prompt = f"""
                Analyze the following document context and the user query to extract key information for a web search.
 
                USER QUERY: "{query}"
           
                DOCUMENT CONTEXT:
                {combined_context[:2500]}
           
                Follow these steps:
                1. Identify the main entities (brands, products, organizations, etc.) in the document
                2. For each entity, determine its type (e.g., "cosmetic brand", "financial service", "software company")
                3. Extract distinctive attributes or features of these entities
                4. Identify industry-specific terminology that might help with search precision
                5. Include geographical context if relevant
           
                Return ONLY a comma-separated list of 5-8 search terms that would help find the most relevant information online.
           
                Note: Focus on terms that will disambiguate search results and prevent confusion with similarly-named but unrelated entities.
                """
           
                try:
                    # Extract context-aware keywords using OpenAI
                    keyword_response = client.chat.completions.create(
                        model="gpt-4o",  # Using a smaller model for efficiency
                        messages=[
                            {"role": "system", "content": "You are a search optimization specialist who identifies the most effective search terms based on document context."},
                            {"role": "user", "content": keyword_prompt}
                        ],
                        temperature=0.3,
                        max_tokens=150
                    )
               
                    search_terms = keyword_response.choices[0].message.content.strip()
                    logger.info(f"Extracted search terms: {search_terms}")
               
                    # Create a more precise enhanced query that includes both the original query and the extracted search terms
                    enhanced_query = f"{query} {search_terms}"
               
                    # Also extract key passages for the context section of our prompt
                    context_prompt = f"""
                    Analyze the following document context and select 2-3 short passages (2-3 sentences each) that contain the most important information related to this query: "{query}"
               
                    DOCUMENT CONTEXT:
                    {combined_context[:3000]}
               
                    Return ONLY the extracted passages with no additional commentary.
                    """
               
                    context_response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "system", "content": "You extract the most relevant passages from documents."},
                            {"role": "user", "content": context_prompt}
                        ],
                        temperature=0.3,
                        max_tokens=300
                    )
               
                    # Get key passages for later use in the main prompt
                    doc_context_text = context_response.choices[0].message.content.strip()
                    logger.info(f"Extracted key passages for context")
               
                except Exception as ke:
                    logger.warning(f"Error extracting keywords: {str(ke)}, proceeding with original query")
                    enhanced_query = query
       
            # Log the queries for comparison
            logger.info(f"Original query: {query}")
            logger.info(f"Enhanced query: {enhanced_query}")
       
            # Step 2: PRIMARY METHOD - Use new Gemini search_web function to get structured search results
            logger.info("PRIMARY: Using enhanced Gemini search_web function to get search results")
            response_text, metadata = self.search_web(enhanced_query, max_results=8, user=user)
       
            if not response_text and not metadata.get('sources'):
                logger.warning("No search results found")
                return "I couldn't find relevant information on the web for your query.", []
       
            # Step 3: Process all sources from metadata
            web_sources = []
            domains_seen = set()
            source_domains = []
           
            # Get sources from metadata (these are from Gemini's grounding_chunks)
            gemini_sources = metadata.get('sources', [])
            search_queries = metadata.get('search_queries', [])
           
            # Process Gemini sources with complete information
            # Process Gemini sources with complete information
            for source in gemini_sources:
                title = source.get('title', 'Unknown Source')
                gemini_redirect_url = source.get('uri', '')  # This is the actual working URL
                
                if not gemini_redirect_url:
                    continue
                
                # Extract clean domain from title for display
                if title and '.' in title and not title.startswith('http') and len(title.split('.')) >= 2:
                    display_domain = title.lower().strip()
                else:
                    # Fallback: extract from the redirect URL (though this will show vertexaisearch.cloud.google.com)
                    display_domain = self.extract_domain(gemini_redirect_url)
                
                # Add to sources if domain is valid and not seen
                if display_domain not in domains_seen and display_domain != "unknown.com":
                    domains_seen.add(display_domain)
                    source_domains.append(display_domain)
                    
                    # Add complete source information
                    web_sources.append({
                        'title': title,
                        'url': gemini_redirect_url,  # Keep the actual working Gemini redirect URL
                        'snippet': response_text[:300] if response_text else '',
                        'domain': display_domain,  # Clean domain for display (e.g., "nasa.gov")
                        'source_type': 'gemini_grounding',
                        'search_queries': search_queries,
                        'actual_url': gemini_redirect_url  # Add this for frontend use
                    })
           
            # If we have additional results from the old web_results format, process those too
            # (This handles the case where search_web might return both response_text and web_results)
            if hasattr(self, '_last_web_results') and self._last_web_results:
                for i, result in enumerate(self._last_web_results):
                    title = result.get('title', 'Unknown Source')
                    url = result.get('href', '')
                    snippet = result.get('body', '')
                   
                    if not url:
                        continue
                       
                    domain = self.extract_domain(url)
                    if domain not in domains_seen and domain != "unknown.com":
                        domains_seen.add(domain)
                        source_domains.append(domain)
                       
                        web_sources.append({
                            'title': title,
                            'url': url,
                            'snippet': snippet[:300],
                            'domain': domain,
                            'source_type': 'web_result',
                            'search_queries': search_queries
                        })
       
            # Format source information for display
            source_info = ", ".join(source_domains) if source_domains else ""
           
            logger.info("Successfully processed Gemini search results")
            logger.info(f"Search queries used: {search_queries}")
            logger.info(f"Total sources returned: {len(web_sources)}")
           
            # Step 4: Return comprehensive response with all source information
            if response_text:
                # If we have document context, enhance the response
                if doc_context_text:
                    enhanced_response = f"{response_text}\n\n*Document Context: {doc_context_text[:200]}...*"
                else:
                    enhanced_response = response_text
                   
                # Add source attribution
                if web_sources:
                    # Create source info with actual URLs
                    source_urls = [f"{source['domain']}|{source['url']}" for source in web_sources if source.get('url')]
                    source_info_with_urls = ", ".join(source_urls)
                    final_response = f"{enhanced_response}\n\n*Sources: {source_info_with_urls}*"
                else:
                    final_response = enhanced_response
                   
                return final_response, web_sources
            else:
                # Fallback: create response from source titles and snippets if no direct response
                if web_sources:
                    formatted_results = []
                    for source in web_sources[:5]:  # Limit display to top 5
                        title = source['title']
                        snippet = source['snippet']
                        formatted_results.append(f"**{title}**\n{snippet}")
                   
                    combined_response = "\n\n".join(formatted_results)
                    final_response = f"{combined_response}\n\n*Sources: {source_info}*"
                    return final_response, web_sources
                else:
                    return "I couldn't find relevant information on the web for your query.", []
       
        except Exception as e:
            logger.error(f"Error in web knowledge search: {str(e)}", exc_info=True)
            return f"An error occurred while searching the web: {str(e)}", []

    def search_web(self, query, max_results=5, user=None):
        """Search the web using Google Gemini and return response text + metadata with sources"""
        try:
            # Get API key from user tokens
            user_api_tokens = UserAPITokens.objects.get(user=user)
            GOOGLE_API_KEY = user_api_tokens.gemini_token
            if not GOOGLE_API_KEY:
                logger.error("GOOGLE_API_KEY not found in user tokens.")
                return "", {'error': 'GOOGLE_API_KEY not found'}
 
            # Initialize Google Gemini client
            client = genai.Client(api_key=GOOGLE_API_KEY)
            model_id = "gemini-2.5-flash"
 
            # Call Gemini with search enabled
            response = client.models.generate_content(
                model=model_id,
                contents=query,
                config={"tools": [{"google_search": {}}]},
            )
            print("\nRaw Gemini Response:")
            print(response.text if hasattr(response, "text") else "No text in response")
            print("\nResponse Metadata:")
            print(f"Response type: {type(response)}")
            if hasattr(response, 'candidates'):
                print(f"Number of candidates: {len(response.candidates)}")
            if hasattr(response, 'sources'):
                print("\nSource URLs:")
                for source in response.sources:
                    print(f"- {source.uri if hasattr(source, 'uri') else source}")

            print("=== GEMINI SEARCH DEBUG END ===\n")



 
            # Extract response text
            response_text = response.text if response.text else ""
            logger.info(f"Generated response text length: {len(response_text)}")
 
            metadata = {
                'response_text': response_text,
                'search_queries': [],
                'sources': [],
                'error': None
            }
 
            try:
                if hasattr(response, 'candidates') and response.candidates:
                    candidate = response.candidates[0]
                    if hasattr(candidate, 'grounding_metadata') and candidate.grounding_metadata:
                        grounding = candidate.grounding_metadata
 
                        # Search queries
                        if hasattr(grounding, 'web_search_queries') and grounding.web_search_queries:
                            metadata['search_queries'] = list(grounding.web_search_queries)
                            logger.info(f"Search queries extracted: {metadata['search_queries']}")
 
                        # Source pages
                        if hasattr(grounding, 'grounding_chunks') and grounding.grounding_chunks:
                            sources = []
                            for i, chunk in enumerate(grounding.grounding_chunks):
                                if hasattr(chunk, 'web') and chunk.web:
                                    title = getattr(chunk.web, 'title', f"Source {i+1}")
                                    uri = getattr(chunk.web, 'uri', '')
                                    if not uri:
                                        continue
 
                                    sources.append({
                                        'title': title,
                                        'uri': uri
                                    })
                                    if len(sources) >= max_results:
                                        break
                            metadata['sources'] = sources
            except Exception as e:
                logger.error(f"Error extracting grounding metadata: {e}", exc_info=True)
                metadata['error'] = str(e)
            print("metadata", metadata)
            return response_text, metadata
 
        except UserAPITokens.DoesNotExist:
            logger.error(f"User API tokens not found for user: {user}")
            return "", {'error': 'User API tokens not found'}
        except Exception as e:
            logger.error(f"Error searching the web via Google Gemini: {str(e)}", exc_info=True)
            return "", {'error': str(e)}
    def extract_domain(self, url):
        """Extract domain from URL - helper method"""
        try:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            return parsed.netloc.replace('www.', '')
        except Exception as e:
            logger.warning(f"Error extracting domain from {url}: {e}")
            return 'unknown.com'

    def resolve_redirect_url(self, redirect_url):
        """Resolve a redirect URL to get the actual destination URL"""
        try:
            import requests
           
            # Set a reasonable timeout and don't follow too many redirects
            response = requests.head(redirect_url, allow_redirects=True, timeout=5,
                                headers={'User-Agent': 'Mozilla/5.0 (compatible; SearchBot)'})
           
            # Return the final URL after following redirects
            final_url = response.url
            logger.info(f"Resolved redirect: {redirect_url} -> {final_url}")
            return final_url
           
        except Exception as e:
            logger.warning(f"Failed to resolve redirect URL {redirect_url}: {e}")
            return None
 
   
 
    
    def combine_document_and_web_responses(self, query, document_response, web_response, doc_sources, web_sources, response_format, conversation_context, original_doc_context=None, doc_citation_sources=None, doc_token_usage=None):
        """
        Combine document-based response and web-based response into a single coherent answer.
        Returns JSON format with enhanced citation handling.
        
        Args:
            query (str): The original user query
            document_response (str): Response generated from document context
            web_response (str or dict): Response generated from web search (can be JSON)
            doc_sources (list): Document sources
            web_sources (list): Web sources
            response_format (str): Desired format for the response
            conversation_context (str): Previous conversation context
            original_doc_context (list, optional): The original document context chunks
            doc_citation_sources (dict, optional): Document citation sources
            
        Returns:
            tuple: (combined_response_dict, combined_citation_sources)
        """

        combined_token_usage = doc_token_usage.copy() if doc_token_usage else {'input_tokens': 0, 'output_tokens': 0, 'total_tokens': 0}
        # Check if the query is a simple greeting
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya",
            
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
        
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
        
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
        
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
        
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
        
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            
            return {
                "content": greeting_responses[response_index],
                "sources": "",
                "citation_sources": {},
                "response_type": "greeting",
                "json_response": True,
                "token_usage": combined_token_usage  # Add this line
            }, {}, combined_token_usage
            

        # Initialize combined citation sources with document citations
        combined_citation_sources = doc_citation_sources.copy() if doc_citation_sources else {}
        
        # Handle web response format (could be string or dict)
        if isinstance(web_response, dict):
            web_content = web_response.get("content", str(web_response))
            web_sources_from_response = web_response.get("sources", [])
        else:
            web_content = str(web_response)
            web_sources_from_response = []
        
        # If not a greeting, proceed with the regular flow
        # Clean up responses by removing source sections
        if "\n\n*Sources:" in document_response:
            document_response = document_response.split("\n\n*Sources:")[0]
            
        if "\n\n*Sources:" in web_content:
            web_response_clean = web_content.split("\n\n*Sources:")[0]
        else:
            web_response_clean = web_content
        
        # Extract web source domains for final source list
        web_source_domains = []
        for source in web_sources:
            domain = self.extract_domain(source.get('url', ''))
            if domain and domain != "unknown.com":
                web_source_domains.append(domain)
        
        # Detect if this is a comparison query that needs special handling
        comparison_words = ["compare", "comparison", "versus", "vs", "difference", "similarities",
                            "better than", "worse than", "alternative", "competitors", "competition",
                            "similar to", "different from", "compared to"]
        
        is_comparison_query = any(word in query_lower for word in comparison_words)
        
        # Analyze the document response and original context to extract key entities for comparison
        # Prepare document context if available
        doc_context_text = ""
        if original_doc_context and len(original_doc_context) > 0:
            # Use up to 5 most relevant chunks from the original document context
            doc_context_text = "\n\n".join(original_doc_context[:5])
        
        # Create analysis prompt for JSON response
        analyze_prompt = f"""
        Analyze the following information to identify key details about the topic being discussed.
        
        USER QUERY: {query}
        
        DOCUMENT RESPONSE:
        {document_response}
        
        {f"ORIGINAL DOCUMENT CONTEXT:{doc_context_text}" if doc_context_text else ""}
        
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "main_entities": ["list", "of", "exact", "entity", "names"],
            "domain": "specific domain or category",
            "key_attributes": ["important", "attributes", "mentioned"],
            "mentioned_competitors": ["any", "competitors", "mentioned"],
            "is_comparison": true/false,
            "comparison_focus": "what aspects are being compared if applicable",
            "industry_context": "specific industry or market context",
            "missing_information": ["what", "important", "info", "seems", "missing"]
        }}
        """
        
        try:
            # Extract key information from document response to enhance web results integration
            analysis_response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert analyst who extracts structured information from text. Respond only in JSON format."},
                    {"role": "user", "content": analyze_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON
                temperature=0.2,
                max_tokens=500
            )

            if hasattr(analysis_response, 'usage') and analysis_response.usage:
                combined_token_usage['input_tokens'] += analysis_response.usage.prompt_tokens
                combined_token_usage['output_tokens'] += analysis_response.usage.completion_tokens
                combined_token_usage['total_tokens'] += analysis_response.usage.total_tokens
            
            # Parse JSON analysis
            analysis_json = self._parse_json_response(analysis_response.choices[0].message.content)
            logger.info(f"Document analysis: {analysis_json}")
            
            # Check if web response seems insufficient based on analysis
            web_info_lacking = False
            if "Unfortunately, there was no additional relevant information found on the web" in web_response_clean or "no relevant information" in web_response_clean.lower():
                web_info_lacking = True
                logger.info("Web information appears to be lacking")
            
            # Determine if we need to use model knowledge for a better response
            use_model_knowledge = False
            
            # If this is a comparison query and web info is lacking, use model knowledge
            if is_comparison_query and (web_info_lacking or not web_sources):
                use_model_knowledge = True
                logger.info("Using model knowledge to supplement missing web information for comparison")
            
            # If web is lacking information about mentioned entities, use model knowledge
            if web_info_lacking and analysis_json.get('main_entities'):
                use_model_knowledge = True
                logger.info("Using model knowledge to supplement information about entities")
            
            # Include conversation context if available (same as above)
            conversation_prompt = ""
            if conversation_context and conversation_context.strip():
                conversation_prompt = f"""
                RECENT CONVERSATION HISTORY:
                {conversation_context}
                
                The above messages are the previous parts of the same ongoing conversation with the user.
                When generating your response:
                - Consider the conversation flow and any references to previous topics
                - If the user is asking follow-up questions or clarifications, refer back to the relevant previous context
                - Maintain consistency with earlier responses and user preferences shown in the conversation
                - If the current question builds upon or relates to previous discussions, acknowledge that connection
                - Avoid repeating information already covered unless specifically asked for clarification
                - Keep the conversation natural and flowing, showing awareness of the ongoing dialogue
                
                Use this conversation history to provide more contextually aware and relevant responses.
                """
                print(f"Including conversation context in prompt (length: {len(conversation_context)} chars)")
            else:
                print("No conversation context available for this response")
            # Create system message for JSON response combination
            system_message = """You are an expert at combining information from multiple sources and using your knowledge to provide comprehensive answers.
            You must respond in JSON format with the following structure:
            {
                "content": "Your comprehensive response with proper HTML formatting and citations",
                "sources_combined": ["list", "of", "all", "sources", "used"],
                "response_sections": {
                    "document_info": "key points from documents",
                    "web_info": "key points from web",
                    "additional_knowledge": "any additional knowledge provided"
                },
                "combination_type": "standard/knowledge_enhanced",
                "citations_included": true/false,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }
            
            CRITICAL INSTRUCTION FOR WEB SOURCES - MUST BE FOLLOWED:
            In your JSON response, include a "web_sources_extracted" field that contains an array of actual website domains mentioned in the web response content. Extract the actual website domains mentioned in the web response content. Do not include generic domains. Look for domain names that appear in the web source content."""
            
            # Create enhanced prompt based on the analysis and whether to use model knowledge
            if use_model_knowledge:
                # Extract key information from the analysis
                main_entities = analysis_json.get('main_entities', [])
                domain = analysis_json.get('domain', '')
                mentioned_competitors = analysis_json.get('mentioned_competitors', [])
                industry_context = analysis_json.get('industry_context', '')
                
                # Enhanced prompt without any explicit brand examples
                enhanced_prompt = f"""
                You have a document response and a web response to the user's query. Your task is to combine these into a coherent response.

                DETAILED CONTEXT:
                - Domain: {domain if domain else "Not clearly identified"}
                - Industry: {industry_context if industry_context else "Not specified"}
                - Main entities: {', '.join(main_entities) if main_entities else "Not clearly identified"}
                - Mentioned competitors: {', '.join(mentioned_competitors) if mentioned_competitors else "None explicitly mentioned"}

                USER QUERY: {query}

                DOCUMENT RESPONSE:
                {document_response}

                WEB RESPONSE:
                {web_response_clean}

                SPECIAL INSTRUCTIONS:
                1. The web search did not return sufficiently helpful information, especially about additional competitors or alternatives not mentioned in the document.

                2. Use your built-in knowledge to SIGNIFICANTLY EXPAND the information, focusing on:
                - ADDITIONAL competitors or alternatives NOT mentioned in the document
                - The broader competitive landscape beyond just the entities mentioned in the document
                - Different market segments and their leading brands relevant to this industry
                - Comparing the main entities with other market players not covered in the document or web results
                
                3. When using your built-in knowledge:
                - Clearly label it as "Additional Market Information"
                - Focus on providing NEW information not found in either the document or web results
                - Include major competing products/services relevant to this specific industry
                - Include international alternatives if relevant
                - Mention different positioning of these alternatives (by demographic, price point, features, etc.)

                4. Format your response with these sections:
                - <h3>Information from Documents</h3> (prioritize this information)
                - <h3>Information from the Web</h3> (include relevant web information)
                - <h3>Additional Market Information</h3> (for supplementary information from your knowledge, focusing on NEW competitors and alternatives)
                
                5. In cases where you are adding information about competitors or alternatives:
                - Focus on major players NOT mentioned in either the document or web response
                - Include at least 3-5 additional competing entities if relevant to the domain
                - Highlight how they differ from the main entities in the document
                - Include details about:
                    * Target audience or use case
                    * Distinctive features or characteristics
                    * Market positioning (premium, mid-range, budget)
                    * Distribution approach
                    * Unique selling propositions
                - Organize the information in a clear, structured way (tables or bullet points if appropriate)

                6. The final response should:
                - Begin with document information (it's specific to the user's context)
                - Include relevant web information when available
                - Add substantial new knowledge about additional competitors and market alternatives
                - Be clearly structured and easy to follow
                - Use proper HTML formatting (<h3>, <b>, <p>, <ul>, <li>, <table> if appropriate)

                The response format should be: {response_format.replace('_', ' ').title()}
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "content": "Your comprehensive response with HTML formatting",
                    "sources_combined": ["document sources", "web sources", "knowledge base"],
                    "response_sections": {{
                        "document_info": "summary of document insights",
                        "web_info": "summary of web insights", 
                        "additional_knowledge": "summary of additional knowledge provided"
                    }},
                    "combination_type": "knowledge_enhanced",
                    "citations_included": false,
                    "web_sources_extracted": ["domain1.com", "domain2.com"]
                }}
                """
                
            else:
                # Update the prompt to include citation information and source JSON extraction
                citation_instructions = """
                CITATION GUIDELINES:
                When referencing information from the document sources, cite the source using [n] format.
                When referencing information from web sources, cite the source using [Wn] format (e.g., [W1], [W2]).
                Every statement that comes directly from a source should include a citation.
                Place citations immediately after the information they support.
                """
                
                # Standard combination prompt if we don't need model knowledge
                enhanced_prompt = f"""
                You have two responses to the user's query: one based on their documents and one based on web search.
                Your task is to combine these into a single coherent response that prioritizes the document information
                (since that is more specific to the user) but supplements with web information when valuable.

                {citation_instructions}

                The response format should be: {response_format.replace('_', ' ').title()}

                Previous conversation context:
                {conversation_context}

                User query: {query}

                conversation history:
                {conversation_prompt}

                RESPONSE FROM DOCUMENTS:
                {document_response}

                RESPONSE FROM WEB SEARCH:
                {web_response_clean}

                GUIDELINES FOR COMBINATION:
                1. When information appears in both sources, prioritize the document information.

                2. Clearly distinguish when you're providing information from the web vs. from documents.

                3. The web information should EXPAND BEYOND what's in the documents - focus on adding NEW information:
                - Additional competitors or alternatives NOT mentioned in the documents
                - Information about different market segments not covered in the documents
                - Broader competitive landscape and market positioning
                - If the query is about competitors, aim to mention additional major players not in the document

                4. Maintain a natural tone and logical flow between document and web information.

                5. If there are contradictions, note them and explain the different perspectives.

                6. Use HTML tags for structure: <p>, <ul>, <li>, <b>, etc.

                7. Structure your response with proper headings and sections for readability.

                8. Use <h3> tags for different sources:
                - <h3>Information from Documents</h3> for document-based information.
                - <h3>Information from the Web</h3> for web-based information.

                9. For document information, maintain the existing citation numbers [n].
                10. For web information, use [W1], [W2], etc. format.

                11. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.

                Create a single, coherent, well-structured response that combines both sources of information.
                Make sure your response is contextually relevant to the ongoing conversation.

                IMPORTANT: If the web search doesn't provide much additional information beyond what's in the documents,
                use your background knowledge to expand the response with additional relevant information about competitors,
                alternatives, or market players that aren't mentioned in either source.
                
                CRITICAL: Your response MUST be a valid JSON object with this structure:
                {{
                    "content": "Your comprehensive response with HTML formatting",
                    "sources_combined": ["list", "of", "all", "sources"],
                    "response_sections": {{
                        "document_info": "key insights from documents",
                        "web_info": "key insights from web"
                    }},
                    "combination_type": "standard",
                    "citations_included": true,
                    "web_sources_extracted": ["domain1.com", "domain2.com"]
                }}
                """
            
            # Generate combined response with JSON format
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": enhanced_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON
                temperature=0.3,
                max_tokens=2500
            )

            # Add token tracking:
            if hasattr(completion, 'usage') and completion.usage:
                combined_token_usage['input_tokens'] += completion.usage.prompt_tokens
                combined_token_usage['output_tokens'] += completion.usage.completion_tokens
                combined_token_usage['total_tokens'] += completion.usage.total_tokens
            
            # Parse the JSON response
            combined_json = self._parse_json_response(completion.choices[0].message.content)
            combined_response = combined_json.get("content", "Error generating combined response")
            
            # Debug: Print the full response to see what GPT returned
            print("=== FULL COMBINED RESPONSE ===")
            print(combined_response)
            print("=== END FULL RESPONSE ===")
            
            # Extract web sources from JSON response
            web_sources_extracted = combined_json.get("web_sources_extracted", [])
            
            print(f"=== EXTRACTED WEB SOURCES FROM JSON ===")
            print(f"Web Sources: {web_sources_extracted}")
            print("=== END EXTRACTED WEB SOURCES ===")
            
            # If no web sources extracted from JSON, fall back to manual extraction
            if not web_sources_extracted:
                import re
                # Look for domain patterns in the response using regex
                domain_pattern = r'([a-zA-Z0-9-]+\.(?:com|in|org|net|co\.in|co\.uk))'
                domains_found = re.findall(domain_pattern, combined_response.lower())
                
                # Filter out common generic domains
                generic_domains = ['google.com', 'search.com', 'example.com', 'website.com']
                for domain in domains_found:
                    if domain not in generic_domains and domain not in web_sources_extracted:
                        web_sources_extracted.append(domain)
                
                print(f"=== FALLBACK EXTRACTION RESULTS ===")
                print(f"Fallback extracted web sources: {web_sources_extracted}")
                print("=== END FALLBACK EXTRACTION ===")
            
            # Process citations in the combined response to include web sources
            web_citation_pattern = r'\[W(\d+)\]'
            web_citations = re.findall(web_citation_pattern, combined_response)

            # Add web sources to combined_citation_sources
            if web_citations and web_sources:
                web_citation_base = max(combined_citation_sources.keys()) + 1 if combined_citation_sources else 1
                
                for i, web_citation in enumerate(set(web_citations)):
                    web_idx = int(web_citation) - 1
                    if web_idx < len(web_sources):
                        web_source = web_sources[web_idx]
                        citation_num = web_citation_base + i
                        
                        # Add web source to citation sources
                        combined_citation_sources[citation_num] = {
                            'source_file': web_source.get('title', 'Web Source'),
                            'page_number': 'Web',
                            'section_title': 'Web Search Result',
                            'snippet': web_source.get('snippet', '')[:200] + "..." if web_source.get('snippet', '') else "Web content",
                            'document_id': f"web_{web_idx}",
                            'url': web_source.get('url', '')
                        }
                        
                        # Replace [Wn] with [citation_num]
                        combined_response = combined_response.replace(f"[W{web_citation}]", f"[{citation_num}]")

            # Create combined sources list for final display
            all_doc_sources = list(set(doc_sources))
            
            # Use extracted web sources if available, otherwise fall back to web_source_domains
            if web_sources_extracted:
                final_web_sources = web_sources_extracted
            elif web_source_domains:
                final_web_sources = web_source_domains
            else:
                final_web_sources = []
            
            # Include web sources from JSON response if available
            if isinstance(web_response, dict) and web_response.get("sources"):
                final_web_sources.extend(web_response["sources"])
            
            all_sources = all_doc_sources 
            
            # If we used model knowledge, add a note
            if use_model_knowledge:
                if all_sources:
                    all_sources_str = ", ".join(all_sources)
                else:
                    all_sources_str = "Document context and general knowledge"
            else:
                all_sources_str = ", ".join(all_sources) if all_sources else "Document context only"

            # Clean up and process citations in the final response
            combined_response = self.normalize_citation_markers(combined_response)
            processed_response, processed_citation_sources = self._format_citations_for_response(combined_response, combined_citation_sources)
            citations_list = [
                processed_citation_sources[k]
                for k in sorted(processed_citation_sources.keys())
            ]
            
            # Return JSON structure instead of plain text
            print("$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$$ FINAL COMBINED RESPONSE", f"{processed_response}\n\n*Sources: {all_sources_str}*")
            return {
                "content": f"{processed_response}\n\n*Sources: {all_sources_str}*",
                "sources": all_sources_str,
                "citation_sources": processed_citation_sources,
                "citations": citations_list,
                "document_sources": all_doc_sources,
                "web_sources": "",
                "combination_type": combined_json.get("combination_type", "standard"),
                "response_sections": combined_json.get("response_sections", {}),
                "used_model_knowledge": use_model_knowledge,
                "is_comparison_query": is_comparison_query,
                "json_response": True,
                "success": True,
                "token_usage": combined_token_usage  # Add this line
            }, processed_citation_sources, combined_token_usage
            
        except Exception as e:
            logger.error(f"Error combining responses: {str(e)}", exc_info=True)
            
            # Fallback: Simply concatenate the responses with a divider in JSON format
            fallback_system_message = """You are an assistant that combines information from multiple sources.
            You must respond in JSON format with this structure:
            {
                "content": "Your combined response with HTML formatting",
                "sources": "list of sources used",
                "combination_type": "fallback",
                "error_handled": true,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }"""
            
            fallback_prompt = f"""
            Combine the following information sources for the user's query: "{query}"
            
            Document Information:
            {document_response}
            
            Web Information:
            {web_response_clean}
            
            Create a coherent response that includes both sources of information.
            Use <h3>Information from Your Documents:</h3> and <h3>Additional Information from Web Search:</h3> as section headers.
            
            CRITICAL: Your response MUST be a valid JSON object with this structure:
            {{
                "content": "Your combined response with HTML formatting",
                "sources": "list of sources",
                "combination_type": "fallback",
                "error_handled": true,
                "web_sources_extracted": ["domain1.com", "domain2.com"]
            }}
            """
            
            try:
                fallback_response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": fallback_system_message},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON
                    temperature=0.3,
                    max_tokens=1500
                )

                # Add token tracking:
                if hasattr(fallback_response, 'usage') and fallback_response.usage:
                    combined_token_usage['input_tokens'] += fallback_response.usage.prompt_tokens
                    combined_token_usage['output_tokens'] += fallback_response.usage.completion_tokens
                    combined_token_usage['total_tokens'] += fallback_response.usage.total_tokens
                
                fallback_json = self._parse_json_response(fallback_response.choices[0].message.content)
                fallback_content = fallback_json.get("content", "Error generating fallback response")
                fallback_web_sources = fallback_json.get("web_sources_extracted", [])
                
            except Exception as fallback_error:
                logger.error(f"Fallback response generation failed: {str(fallback_error)}")
                # Ultimate fallback - create JSON manually
                fallback_content = f"""
                <h3>Information from Your Documents:</h3>
                <div>
                {document_response}
                </div>
                
                <h3>Additional Information from Web Search:</h3>
                <div>
                {web_response_clean}
                </div>
                """
                fallback_web_sources = []
            
            # Add sources for fallback
            all_doc_sources = list(set(doc_sources))
            
            # Use extracted web sources if available from fallback, otherwise fall back to web_source_domains
            if fallback_web_sources:
                final_web_sources = fallback_web_sources
            elif web_source_domains:
                final_web_sources = web_source_domains
            else:
                final_web_sources = []
                
            # Include web sources from JSON response if available
            if isinstance(web_response, dict) and web_response.get("sources"):
                final_web_sources.extend(web_response["sources"])
                
            all_sources = all_doc_sources
            all_sources_str = ", ".join(all_sources) if all_sources else "Document and web sources"
            
            fallback_content = self.normalize_citation_markers(fallback_content)
            return {
                "content": f"{fallback_content}\n\n*Sources: {all_sources_str}*",
                "sources": all_sources_str,
                "citation_sources": doc_citation_sources if doc_citation_sources else {},
                "document_sources": all_doc_sources,
                "web_sources": "",
                "combination_type": "fallback",
                "error": str(e),
                "error_handled": True,
                "json_response": True,
                "success": False,
                "token_usage": combined_token_usage  # Add this line
            }, doc_citation_sources if doc_citation_sources else {}, combined_token_usage 

    

    def detect_question_format(self, query, context_snippets=None):
        """
        Analyzes the question and detects the most appropriate response format
        based on the query content and any available document context.
        
        Args:
            query (str): The user's question
            context_snippets (list): Optional list of document context snippets
            
        Returns:
            tuple: (primary_format, secondary_format, confidence)
        """
        # Implementation remains the same as in your original code
        print("###########################################################")
        
        # Normalize query for keyword matching
        query_lower = query.lower()
        
        # Format detection keywords dictionary
        format_keywords = {
            "executive_summary": [
                "summarize", "summary", "overview", "high level", "highlights", "key points", 
                "brief overview", "executive summary", "main points", "takeaways", "tldr",
                "in a nutshell", "brief summary", "key highlights", "summarise", "recap",
                "give me the gist", "summarization", "outline the main", "key takeaways",
                "bullet points of", "boil down", "distill", "condense"
            ],
            "comparative_analysis": [
                "compare", "comparison", "versus", "vs", "difference", "differences", 
                "similarities", "similarity", "better", "worse", "advantage", "disadvantage",
                "pros and cons", "strengths and weaknesses", "contrast", "choose between",
                "which is better", "how does", "stack up", "relative to", "compared to",
                "weighing", "juxtapose", "differentiate between", "distinction", "compare and contrast",
                "side by side", "match up against", "measure against"
            ],
            "detailed_analysis": [
                "analyze", "analysis", "deep dive", "in-depth", "detailed", "breakdown",
                "examine", "investigate", "elaborate on", "explain in detail", "thorough",
                "comprehensive", "drill down", "what factors", "analyze in detail",
                "dissect", "unpack", "scrutinize", "evaluate", "assessment", "diagnosis",
                "root cause", "deconstructing", "implications of", "examine carefully", 
                "look closely at", "study"
            ],
            "strategic_recommendation": [
                "recommend", "recommendation", "suggest", "advice", "should", "best course",
                "action plan", "strategic", "next steps", "how should", "what would you advise",
                "strategy", "approach", "what should", "best way to", "how to best",
                "action items", "propose", "guidance", "direction", "advise on",
                "what's the right approach", "strategic guidance", "roadmap", "plan for",
                "steps to take", "best practice", "optimal strategy"
            ],
            "market_insights": [
                "market trends", "industry", "segment", "consumer behavior", "market data", 
                "competitive landscape", "market share", "industry trend", "consumer preference",
                "market analysis", "market research", "market performance", "market outlook",
                "sector analysis", "competitors", "customer demographics", "forecast",
                "target audience", "customer segment", "market growth", "consumer insights",
                "purchase behavior", "audience analysis", "industry position"
            ],
            "factual_brief": [
                "what is", "when was", "where is", "who is", "facts about", "define", 
                "information on", "tell me about", "data on", "evidence of", "statistics on",
                "key facts", "figures on", "numbers for", "metrics on", "how many",
                "percentage of", "rate of", "frequency of", "occurrence of", "incidents of",
                "quantity of", "value of", "amount of", "total number", "statistics for",
                "list the", "enumerate", "identify", "reference"
            ],
            "technical_deep_dive": [
                "how does it work", "technical", "technology", "mechanism", "process", "methodology",
                "framework", "system", "architecture", "technical details", "underlying", "step by step",
                "explain the process", "technical explanation", "in technical terms",
                "inner workings", "components", "operation", "workflow", "algorithm",
                "explain technically", "implementation", "procedure", "specification",
                "design", "function", "operation of", "engineering behind"
            ]
        }
        
        # Check for format-specific keywords
        matched_formats = []
        for format_key, keywords in format_keywords.items():
            # Count matches and weight longer phrases higher
            match_score = 0
            for keyword in keywords:
                if keyword in query_lower:
                    # Weight multi-word phrases higher
                    words_in_keyword = len(keyword.split())
                    match_score += words_in_keyword  
            
            if match_score > 0:
                matched_formats.append((format_key, match_score))
        
        # Sort by match score (higher score = stronger signal)
        matched_formats.sort(key=lambda x: x[1], reverse=True)
        
        # If we have keyword matches, use the best match
        if matched_formats and matched_formats[0][1] >= 1:
            primary_format = matched_formats[0][0]
            # Set confidence based on match score
            base_confidence = 5.0
            # Boost confidence based on match score
            match_bonus = min(matched_formats[0][1], 4.0)  # Cap the bonus
            confidence = base_confidence + match_bonus
            
            # Get secondary format if available
            secondary_format = matched_formats[1][0] if len(matched_formats) > 1 else None
            
            return primary_format, secondary_format, confidence
        
        # If no keyword matches, we can try using LLM for classification
        # This would integrate with your existing classify_question_intent function
        # but we'll keep it simple for now
        
        # Default to natural if no clear classification
        return "natural", None, 4.0


    # Add this helper method to get format-specific guidance
    def _get_format_guidance(self, response_format, length_preference):
        """
        Get format-specific guidance based on response format and length preference.
        
        Args:
            response_format (str): The requested response format
            length_preference (str): 'short' or 'comprehensive'
            
        Returns:
            str: Format-specific guidance text
        """
        # Format guidance dictionary
        format_guidance = {
            "executive_summary": {

                 """
    This is an EXECUTIVE SUMMARY request. Provide a concise, high-level overview of the key points and implications. Specifically:

    1. Begin with a brief 1-2 sentence overview of the main finding or takeaway
    2. Structure your summary with clear sections for different key points
    3. Prioritize only the most important findings and insights (quality over quantity)
    4. Include critical data points and evidence that support the key messages
    5. Ensure every statement adds value and nothing is redundant
    6. Conclude with the most important implications or next steps
    7. Keep the entire summary focused and concise - aim for brevity while maintaining completeness

    Your summary should give the reader a clear understanding of the essential points without needing to read a longer document.
    """
            },
            "detailed_analysis": {

                 """
    This is a DETAILED ANALYSIS request. Provide a comprehensive analysis with in-depth examination of the information and its implications. Specifically:

    1. Begin with an executive summary of your key findings
    2. Include a rigorous analysis of the data, breaking it down into logical sections
    3. Examine patterns, trends, anomalies, and relationships in the data
    4. Support your analysis with specific evidence and examples from the documents
    5. Consider multiple perspectives and interpretations of the information
    6. Discuss implications of your findings for decision-making
    7. Include relevant charts or tables to visualize key data points
    8. Conclude with synthesis of the most important insights

    Your analysis should be thorough and demonstrate critical thinking throughout.
    """
            },
            "strategic_recommendation": {

               """
    This is a STRATEGIC RECOMMENDATION request. Provide actionable, prioritized recommendations with clear strategic direction. Specifically:

    1. Begin with a brief overview of the current situation or challenge
    2. Present 3-5 specific, actionable recommendations in priority order
    3. For each recommendation:
    - Clearly state what should be done
    - Explain why this approach is recommended (with evidence)
    - Outline expected outcomes or benefits
    - Address potential implementation challenges
    4. Include a high-level implementation timeline or roadmap if appropriate
    5. Conclude with the critical success factors for implementation

    Your recommendations should be practical, evidence-based, and strategically sound.
    """
            },
            "comparative_analysis": {
                """
    This is a COMPARATIVE question. Organize your response around comparing different options, approaches, or entities mentioned in the query. Specifically:

    1. Begin with a brief introduction of the items being compared
    2. Use a structure that clearly contrasts the different options (side-by-side format where possible)
    3. Create comparison tables that directly align comparable attributes
    4. Identify key differentiating factors between the options
    5. Discuss relative advantages and disadvantages of each
    6. Include direct quantitative comparisons where data is available
    7. Conclude with insights about the relative positioning of the compared items

    When presenting comparative data, use tables with this format when appropriate:

    | Attribute | Option A | Option B | Key Difference |
    |-----------|----------|----------|---------------|
    | Feature 1 | Value    | Value    | Insight       |
    """
            },
            "market_insights": {

               """
    This is a MARKET INSIGHTS request. Focus on market trends, consumer behavior, and competitive dynamics. Specifically:

    1. Begin with an overview of the key market dynamics or trends
    2. Analyze market segments, sizes, and growth rates with supporting data
    3. Examine competitive landscape and relative positioning of key players
    4. Discuss consumer behavior patterns, preferences, and changing demands
    5. Identify market opportunities and potential threats
    6. Present relevant market statistics in tabular format where appropriate
    7. Conclude with strategic implications of these market insights

    Your response should combine data-driven market analysis with strategic interpretation.
    """
            },
            "factual_brief": {

             """
    This is a FACTUAL INQUIRY. Provide direct, fact-based answers with supporting evidence. Specifically:

    1. Begin with a clear, direct answer to the question asked
    2. Present the key facts in a structured, logical order
    3. Support each fact with specific evidence from the documents
    4. Include relevant numbers, statistics, or data points
    5. Use bullet points or numbered lists where appropriate to improve readability
    6. Avoid speculation or unsupported assertions
    7. Clearly indicate if certain requested information is not available in the documents

    Your response should be factual, concise, and directly address the query.
    """
            },
            "technical_deep_dive": {

             """
    This is a TECHNICAL question. Include detailed explanations of processes, methodologies, or procedures. Specifically:

    1. Begin with a high-level overview of the technical concept or process
    2. Break down the technical components or steps in a logical sequence
    3. Explain how each component works or contributes to the whole
    4. Include relevant technical specifications, parameters, or metrics
    5. Use clear examples to illustrate technical concepts
    6. Consider including diagrams or technical illustrations where helpful
    7. Address limitations, constraints, or technical considerations
    8. Conclude with practical applications or implications

    Your explanation should be technically accurate while remaining accessible to a business audience.
    """
            },
            "natural": {
                "short": "Structure your response in a natural way that best addresses the query, while keeping it concise and to-the-point.",
                "comprehensive": "Structure your response in the most natural way that best addresses the query, with appropriate depth and detail."
            }
        }
        
        # Return the appropriate guidance based on format and length preference
        if response_format in format_guidance and length_preference in format_guidance[response_format]:
            return format_guidance[response_format][length_preference]
        
        # Default fallback
        return "Provide a helpful response to the query based on the document context."


    # Add these citation-specific prompt templates to the ChatView class

    CITATION_QA_TEMPLATE = """
    Please provide an answer based solely on the provided sources. 
    When referencing information from a source, cite the appropriate source(s) using their corresponding numbers within square brackets.
    Every statement in your answer should include at least one source citation.
    Only cite a source when you are explicitly referencing it.
    If none of the sources are helpful, you should indicate that.

    Examples of good citations:
    1. 'The sky is red in the evening [1]'
    2. 'Water is wet when the sky is red [2], which occurs in the evening [1]'
    3. 'According to the research [3], both findings [1][2] support this conclusion'

    Your answer should be clear, accurate, and directly tied to the source material.
    Below are several numbered sources of information:

    {context_str}

    Query: {query_str}
    Answer: 
    """

    CITATION_REFINE_TEMPLATE = """
    You are refining an existing answer to make it more accurate and well-cited.
    When referencing information from a source, cite the appropriate source(s) using their corresponding numbers within square brackets.
    Every statement in your answer should include at least one source citation.
    Only cite a source when you are explicitly referencing it.

    We have provided an existing answer: {existing_answer}

    Below are numbered sources of information.
    Use them to refine the existing answer, ensuring all statements have proper citations.

    {context_msg}

    Query: {query_str}
    Refined Answer:
    """


    def _clean_citations(self, text, citation_sources):

        """Clean up citation format to ensure consistent [n][m] format and remove duplicates."""

        # Fix citations that might not have brackets

        text = re.sub(r'(\s)(\d+)(\s)', lambda m: m.group(1) + '[' + m.group(2) + ']' + m.group(3) 

                    if self._is_likely_citation(int(m.group(2)), citation_sources) else m.group(0), text)

        # Fix spaces in citations and normalize format

        text = re.sub(r'\[\s*(\d+)\s*\]', r'[\1]', text)

        # Handle multiple citations in various formats and convert to separate brackets

        # Match patterns like [1, 2], [1,2], [1 2], etc.

        def normalize_multiple_citations(match):

            citation_content = match.group(1)

            # Extract all numbers

            numbers = re.findall(r'\d+', citation_content)

            # Remove duplicates while preserving order

            seen = set()

            unique_numbers = []

            for num in numbers:

                if num not in seen and self._is_likely_citation(int(num), citation_sources):

                    seen.add(num)

                    unique_numbers.append(num)

            # Return as separate brackets

            return ''.join(f'[{num}]' for num in unique_numbers)

        # Apply normalization for multiple citation patterns

        text = re.sub(r'\[([0-9,\s]+)\]', normalize_multiple_citations, text)

        # Remove duplicate citations and clean up

        text = self._remove_consecutive_duplicate_citations(text)

        return text
    


    def _ensure_separate_citation_brackets(self, text):
        """Ensure all citations are in separate brackets [1][2] not [1,2] or [1, 2]."""
        # Handle any remaining grouped citations
        def separate_brackets(match):
            numbers = re.findall(r'\d+', match.group(1))
            # Remove duplicates while preserving order
            seen = set()
            unique_numbers = []
            for num in numbers:
                if num not in seen:
                    seen.add(num)
                    unique_numbers.append(num)
            return ''.join(f'[{num}]' for num in unique_numbers)
        # Match any remaining grouped citations
        text = re.sub(r'\[([0-9,\s]+)\]', separate_brackets, text)
        # Final cleanup - remove any remaining duplicates in the same vicinity
        # Look for patterns like [1][1] within a short span
        words = text.split()
        cleaned_words = []
        for word in words:
            # If word contains multiple identical citations, deduplicate
            if '[' in word and ']' in word:
                citations = re.findall(r'\[(\d+)\]', word)
                if len(citations) != len(set(citations)):
                    # Remove duplicates
                    unique_citations = []
                    seen = set()
                    for cite in citations:
                        if cite not in seen:
                            unique_citations.append(cite)
                            seen.add(cite)
                    # Rebuild the word
                    citation_pattern = r'\[\d+\]'
                    parts = re.split(citation_pattern, word)
                    # Reconstruct with unique citations
                    result = parts[0]
                    for i, cite in enumerate(unique_citations):
                        result += f'[{cite}]'
                        if i + 1 < len(parts):
                            result += parts[i + 1]
                    cleaned_words.append(result)
                else:
                    cleaned_words.append(word)
            else:
                cleaned_words.append(word)
        return ' '.join(cleaned_words)
    
    def _is_likely_citation(self, num, citation_sources):
        """Check if a number is likely to be a citation based on available sources."""
        return num in citation_sources and 1 <= num <= len(citation_sources)
    
    def _format_citations_for_response(self, response_text, citation_sources):
        """Format the citations for response and ensure sequential numbering"""
        # First extract ALL citations from the text
        citation_pattern = r'\[(\d+)\]'
        all_citations = re.findall(citation_pattern, response_text)
        
        # Get unique citations while preserving order of first appearance
        unique_citations = []
        seen = set()
        for cite in all_citations:
            if cite.isdigit() and int(cite) in citation_sources and cite not in seen:
                seen.add(cite)
                unique_citations.append(cite)
        
        # Create mapping from original to sequential numbers (1,2,3...)
        citation_mapping = {int(old): new for new, old in enumerate(unique_citations, 1)}
        
        # Replace ALL citations in the text with sequential numbers
        def replace_citation(match):
            old_num = match.group(1)
            if old_num.isdigit() and int(old_num) in citation_mapping:
                return f'[{citation_mapping[int(old_num)]}]'
            return match.group(0)
        
        processed_text = re.sub(citation_pattern, replace_citation, response_text)
        
        # Create new citation sources with sequential numbers
        display_citation_sources = {}
        for display_num, original_num in enumerate(unique_citations, 1):
            original_num = int(original_num)
            if original_num in citation_sources:
                source_info = citation_sources[original_num]
                snippet = source_info.get('text', '')
                snippet_length = 1500 
                display_citation_sources[display_num] = {
                    'source_file': source_info.get('source_file', 'Unknown'),
                    'page_number': 'Unknown',
                    'section_title': 'Unknown',
                    'snippet': snippet[:snippet_length] + "..." if len(snippet) > snippet_length else snippet,
                    'document_id': source_info.get('document_id', 'Unknown')
                }
        
        return processed_text, display_citation_sources

    # def _format_citations_for_response(self, response_text, citation_sources):
    #     """Format the citations for response and remove duplicates."""
    #     # Extract all citations from the text
    #     response_text = self._ensure_separate_citation_brackets(response_text)
    #     citation_pattern = r'\[(\d+)\]'
    #     citations = re.findall(citation_pattern, response_text)
        
    #     # Use set to get unique citations, then convert back to sorted list
    #     unique_citations = sorted(set([int(c) for c in citations if c.isdigit()]))
        
    #     # Create mapping from original citation numbers to sequential display numbers
    #     citation_mapping = {original_num: display_num for display_num, original_num in enumerate(unique_citations, 1)}
        
    #     # Replace original citation numbers with sequential display numbers
    #     # Process in reverse order to avoid issues with overlapping replacements
    #     processed_text = response_text
    #     for original_num in sorted(citation_mapping.keys(), reverse=True):
    #         # Use word boundaries to avoid partial matches
    #         pattern = r'\[' + str(original_num) + r'\]'
    #         replacement = f'[{citation_mapping[original_num]}]'
    #         processed_text = re.sub(pattern, replacement, processed_text)
        
    #     # Remove duplicate citations that appear consecutively
    #     processed_text = self._remove_consecutive_duplicate_citations(processed_text)
        
    #     # Create new citation sources dictionary with display numbers
    #     display_citation_sources = {}
    #     for display_num, original_num in enumerate(unique_citations, 1):
    #         if original_num in citation_sources:
    #             source_info = citation_sources[original_num]
    #             snippet = source_info.get('text', '')
    #             snippet_length = 1500 
    #             display_citation_sources[display_num] = {
    #                 'source_file': source_info.get('source_file', 'Unknown'),
    #                 'page_number': 'Unknown',
    #                 'section_title': 'Unknown',
    #                 'snippet': snippet[:snippet_length] + "..." if len(snippet) > snippet_length else snippet,
    #                 'document_id': source_info.get('document_id', 'Unknown')
    #             }
        
    #     return processed_text, display_citation_sources

    def _remove_consecutive_duplicate_citations(self, text):
        """Remove consecutive duplicate citations and ensure proper bracket separation."""
        # First, handle grouped citations like [1, 11] and convert to [1][11]
        def separate_grouped_citations(match):
            citation_content = match.group(1)
            # Split by comma and clean up
            citation_nums = [num.strip() for num in citation_content.split(',') if num.strip().isdigit()]
            # Remove duplicates while preserving order
            seen = set()
            unique_citations = []
            for num in citation_nums:
                if num not in seen:
                    seen.add(num)
                    unique_citations.append(num)
            # Return as separate brackets
            return ''.join(f'[{num}]' for num in unique_citations)
        # Pattern to match grouped citations like [1, 11] or [1,11] or [1 11]
        grouped_pattern = r'\[([0-9,\s]+)\]'
        text = re.sub(grouped_pattern, separate_grouped_citations, text)
        # Remove consecutive identical citations like [1][1] -> [1]
        consecutive_pattern = r'\[(\d+)\](?:\s*\[\1\])+' 
        def replace_consecutive_duplicates(match):
            citation_num = match.group(1)
            return f'[{citation_num}]'
        text = re.sub(consecutive_pattern, replace_consecutive_duplicates, text)
        # Remove citations separated by minimal punctuation like [1], [1] -> [1]
        punctuation_pattern = r'\[(\d+)\]\s*[,.;:!?]*\s*\[\1\]'
        text = re.sub(punctuation_pattern, r'[\1]', text)
        # Remove duplicate citations within the same sentence (more aggressive)
        # Find all citations in each sentence and deduplicate
        sentences = text.split('.')
        processed_sentences = []
        for sentence in sentences:
            # Find all citations in this sentence
            citations_in_sentence = re.findall(r'\[(\d+)\]', sentence)
            if len(citations_in_sentence) > len(set(citations_in_sentence)):
                # There are duplicates, remove them
                seen_in_sentence = set()
                def replace_if_seen(match):
                    num = match.group(1)
                    if num in seen_in_sentence:
                        return ''  # Remove duplicate
                    else:
                        seen_in_sentence.add(num)
                        return match.group(0)  # Keep first occurrence
                sentence = re.sub(r'\[(\d+)\]', replace_if_seen, sentence)
            processed_sentences.append(sentence)
        return '.'.join(processed_sentences)

    def generate_response(self, query, context, sources, use_web_knowledge=False, response_format='natural', conversation_context=""):
        """
        Generate a comprehensive response using the provided context and sources with web search capability.
        Enhanced to handle URL content in context and generate detailed, long-form responses.
       
        Args:
            query (str): User's original query
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            use_web_knowledge (bool): Whether to use web search in addition to documents
            response_format (str): Format style for the response
            conversation_context (str): Previous conversation history context
           
        Returns:
            str: Generated comprehensive response based on the context and/or web search
        """
        # Check if the query is a simple greeting
        greeting_keywords = [
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            "hlw", "g'day", "yo", "heya",
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
       
        query_lower = query.lower().strip()
        is_greeting = False
       
        if query_lower in greeting_keywords:
            is_greeting = True
       
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:
                is_greeting = True
                break
       
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index], {}
       
        if not context and not use_web_knowledge:
            return "I cannot answer this question based on the provided documents."
       
        # Enhanced URL and context processing
        url_context = ""
        document_context = []
        document_sources = []
        url_sources = []
       
        for i, content_item in enumerate(context):
            if i < len(sources) and sources[i] == "URL Content" and "CONTENT FROM URLS IN QUERY:" in content_item:
                url_context = content_item
                url_sources.append(sources[i])
            else:
                document_context.append(content_item)
                if i < len(sources):
                    document_sources.append(sources[i])
       
        citation_sources = {}
        # Use MORE context for comprehensive responses - increased limits
        selected_context, selected_sources = self._prepare_context_comprehensive(document_context, document_sources)
       
        contextualized_content = []
        for i, (content, source) in enumerate(zip(selected_context, selected_sources), 1):
            citation_sources[i] = {
                "source_file": source,
                "text": content,
                "document_id": "Unknown",
                "score": 1.0 - (i * 0.02),  # Reduced penalty for more sources
                "display_num": i
            }
            contextualized_content.append(f"Source {i}:\n{content}")
 
        if url_context:
            citation_count = len(contextualized_content) + 1
            contextualized_content.append(f"Source {citation_count}:\n{url_context}")
            selected_sources.append("URL Content")
            citation_sources[citation_count] = {
                "source_file": "URL Content",
                "text": url_context,
                "document_id": "url_content",
                "score": 0.8,
                "display_num": citation_count
            }
 
        full_context = "\n\n".join(contextualized_content)
 
        # Get format-specific guidance for COMPREHENSIVE responses
        format_guidance = self._get_format_guidance(response_format, 'comprehensive')
       
        # Include conversation context if available
        conversation_prompt = ""
        if conversation_context and conversation_context.strip():
            conversation_prompt = f"""
            RECENT CONVERSATION HISTORY:
            {conversation_context}
            
            The above messages are the previous parts of the same ongoing conversation with the user.
            When generating your response:
            - Consider the conversation flow and any references to previous topics
            - If the user is asking follow-up questions or clarifications, refer back to the relevant previous context
            - Maintain consistency with earlier responses and user preferences shown in the conversation
            - If the current question builds upon or relates to previous discussions, acknowledge that connection
            - Avoid repeating information already covered unless specifically asked for clarification
            - Keep the conversation natural and flowing, showing awareness of the ongoing dialogue
            
            Use this conversation history to provide more contextually aware and relevant responses.
            """
            print(f"Including conversation context in prompt (length: {len(conversation_context)} chars)")
        else:
            print("No conversation context available for this response")
 
        # Get project description if available
        project_description = self._get_project_description(query)
       
        project_guidance = ""
        if project_description and len(project_description.strip()) > 5:
            project_guidance = f"""
            PROJECT CONTEXT:
            {project_description}
           
            IMPORTANT: This project context is only to help you understand the domain and purpose of this project. Your responses MUST be derived exclusively from the uploaded documents provided in the document context.
           
            When formulating your response:
            1. Use this project context solely to understand the general domain and focus of the project
            2. Draw ALL factual information, data, and specific content ONLY from the uploaded documents
            3. Shape your response style and tone to align with this project's domain, but never invent content
            4. DO NOT add information from your general knowledge even if relevant to the project description
            5. If the documents don't contain information relevant to the query, clearly state this limitation
           
            Remember: The project description helps you understand the context, but all response content must come from the uploaded documents.
            """
 
        # Updated system message for JSON response
        system_message = """You are a document analysis expert with conversation memory.
        You must respond in JSON format with the following structure:
        {
            "content": "Your detailed response content here with proper HTML formatting and citations",
            "citations_used": [list of citation numbers used in the response],
            "format_type": "response_format_used",
            "sources_referenced": ["list", "of", "source", "names"]
        }
       
        In the content field, use HTML tags for formatting (<b>, <p>, <ul>, <li>, <table>, etc.).
        Use source citations [n] for all information from the documents.
        Provide a comprehensive and detailed answer while maintaining conversation continuity."""
 
        # Enhanced comprehensive response prompt
        user_prompt = f"""
        You are a highly skilled expert in deep document analysis and summarization. Your task is to create a **massive, exhaustive, and detailed answer** of the provided context using rich HTML formatting, prioritizing **tables** whenever asked and **clear breakdowns** for optimal readability.


        <b>OBJECTIVE:</b> Provide a **comprehensive, all-encompassing summary** of the information extracted from multiple documents and URLs, with a focus on **clarity, structured analysis, tabular representation**, and in-depth explanation of every important aspect.
 
        <b>RESPONSE STYLE & FOCUS:</b>
        - Use an **encyclopedic**, **scholarly**, and **analytical tone**.
        - Focus heavily on clarity, comparison, and synthesis.
        - Provide detailed breakdowns, categories, and layered analysis.
        - Aim for maximum **coverage, detail, structure, and readability**.
 
        <b>STRUCTURE & FORMATTING:</b>
        - Begin with a short introduction paragraph summarizing what the content covers.
        - Add supplementary <p> explanations under each table to elaborate with context, implications, and examples.
        - Use HTML formatting only (<b>, <p>, <ul>, <li>, <table>, <tr>, <td>, etc.). Avoid Markdown.
        - Clearly indicate relationships, comparisons, or patterns when analyzing multiple sources.
        - Use sub-headings for logical sections.
 
        <b>COMPREHENSIVE RESPONSE GENERATION GUIDELINES:</b>
        - Ensure complete coverage of all document points and topics.
        - Include detailed analysis, implications, contrasts, and supporting examples.
        - If possible, show timelines, process flows, or hierarchies in structured formats.
        - Integrate and reference all relevant document and URL material without skipping or summarizing vaguely.
        - Provide deep, layered explanations under each table for full contextual clarity.
        - Use citations [n] for all derived content and maintain a scholarly tone.
        - Always refer back to original source content when explaining or summarizing.
 
        DOCUMENT AND URL CONTEXT:
        {full_context}
 
        CONVERSATION HISTORY:
        {conversation_prompt}
        
        USER QUERY:
        {query}
 
        <b>CRITICAL OUTPUT FORMAT:</b> Your output MUST be a valid JSON object with this structure:
        {{
            "content": "Massive, detailed response in rich HTML format with multiple tables, paragraph explanations, structured sections, and full use of provided context.",
            "citations_used": [list of citation numbers used],
            "format_type": "{response_format}",
            "sources_referenced": ["list of source names referenced"]
        }}
 
        <b>STRICT CONSTRAINTS:</b>
        - DO NOT use any knowledge or speculation outside the provided document and URL context.
        - ONLY use details that are verifiably present in the provided context.
        - Provide deep, analytical insights—not high-level summaries.
        - Every important idea, point, or piece of data from the context should appear in the output.
        - Ensure output is substantial, structured, and formatted with maximum clarity.
 
        <b>RESPONSE FORMAT REQUIREMENTS:</b>
        1. Use <b> for section and table headers.
        2. Use <table>, <tr>, <td> as the primary format to organize information.
        3. Provide <p> tag explanations below each table for depth and clarity.
        4. Avoid Markdown; only HTML formatting.
        5. Include detailed, well-supported elaboration on every relevant element.
        6. Structure the summary by themes, categories, comparisons, or timelines where applicable.
        """
 
       
        try:
            # Call the OpenAI chat completion API with JSON format
            completion = client.chat.completions.create(
                model="o3-mini",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
             
               
            )

            token_usage = {
            'input_tokens': completion.usage.prompt_tokens if completion.usage else 0,
            'output_tokens': completion.usage.completion_tokens if completion.usage else 0,
            'total_tokens': completion.usage.total_tokens if completion.usage else 0
            }
           
             # ===== LOG RAW LLM RESPONSE =====
            raw_llm_response = completion.choices[0].message.content
            print("\n" + "="*100)
            print("🤖 RAW LLM RESPONSE FROM OPENAI (generate_short_response)")
            print("="*100)
            print("Model used:", completion.model)
            print("Response ID:", completion.id)
            print("Raw response content:")
            print(raw_llm_response)
            print("Response length:", len(raw_llm_response))
            print("Response type:", type(raw_llm_response))
            print("="*100)
            # ===== END LOG RAW LLM RESPONSE =====
           
            json_response = self._parse_json_response(completion.choices[0].message.content)
           
            if json_response.get("error"):
                # If JSON parsing failed, return the fallback message
                answer = json_response.get("content", "An error occurred while generating the response.")
            else:
                # Extract the content from JSON response
                answer = json_response.get("content", "No response content found.")
 
            # Clean and process citations in the answer
            answer = self._clean_citations(answer, citation_sources)
            answer = self._ensure_separate_citation_brackets(answer)
            processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
           
            # If answer seems too short, try to generate a more detailed one
            if len(answer.split()) < 100:  # If less than ~100 words
                enhanced_system_message = """You are a document analysis expert with conversation memory.
                You must respond in JSON format with the following structure:
                {
                    "content": "Your EXTREMELY DETAILED and COMPREHENSIVE response with proper HTML formatting and citations",
                    "citations_used": [list of citation numbers used],
                    "format_type": "response_format_used",
                    "sources_referenced": ["list", "of", "source", "names"]
                }
               
                Provide an EXTREMELY DETAILED and COMPREHENSIVE response including ALL information from the context while maintaining conversational continuity. Use source citations [n] for all information from the documents."""
               
                completion = client.chat.completions.create(
                    model="o3-mini",
                    messages=[
                        {"role": "system", "content": enhanced_system_message},
                        {"role": "user", "content": user_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                   
                   
                )
                token_usage = {
                'input_tokens': completion.usage.prompt_tokens if completion.usage else 0,
                'output_tokens': completion.usage.completion_tokens if completion.usage else 0,
                'total_tokens': completion.usage.total_tokens if completion.usage else 0
                }
               
                # Parse the enhanced JSON response
                json_response = self._parse_json_response(completion.choices[0].message.content)
                answer = json_response.get("content", answer)  # Use original if parsing fails
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
               
        except Exception as e:
            # If there's an error (e.g., token limit), try with fewer context chunks
            reduced_context = "\n\n".join(contextualized_content[:8])
            # Create a shorter version of conversation context if needed
            short_conv_context = ""
            if conversation_context:
                short_conv_context = f"Previous conversation context: {conversation_context[:300]}..."
               
            fallback_prompt = f"""
            Based ONLY on the following context, provide a comprehensive answer to the question: {query}
           
            {short_conv_context}
           
            CONTEXT:
            {reduced_context}
           
            CRITICAL: Respond in JSON format:
            {{
                "content": "Your detailed response with HTML formatting",
                "citations_used": [],
                "format_type": "fallback",
                "sources_referenced": []
            }}
            """
           
            try:
                fallback_completion = client.chat.completions.create(
                    model="o3-mini",
                    messages=[
                        {"role": "system", "content": "You are a document analysis expert. Respond in JSON format."},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                   
                 
                )
                token_usage = {
                'input_tokens': fallback_completion.usage.prompt_tokens if fallback_completion.usage else 0,
                'output_tokens': fallback_completion.usage.completion_tokens if fallback_completion.usage else 0,
                'total_tokens': fallback_completion.usage.total_tokens if fallback_completion.usage else 0
                }
               
                json_response = self._parse_json_response(fallback_completion.choices[0].message.content)
                answer = json_response.get("content", "An error occurred while generating the response.")
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
            except Exception as nested_e:
                # Last resort fallback
                answer = f"An error occurred while generating the response: {str(e)}. Please try a more specific question or with fewer documents."
                processed_answer = answer
                processed_citations = {}
 
        # Add source information
        source_list = list(set(selected_sources))
        source_info = ", ".join(source_list)
        return f"{processed_answer}\n\n*Sources: {source_info}*", processed_citations, token_usage
    
    #  generate_short_response function with this fixed version:

    def generate_short_response(self, query, context, sources, use_web_knowledge=False, response_format='natural', conversation_context=""):
        """
        Generate a shorter, concise response using the provided context with web search capability.
        Enhanced to handle URL content in context using Gemini Flash 2.0.
       
        Args:
            query (str): User's original query
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            use_web_knowledge (bool): Whether to use web search in addition to documents
            response_format (str): Format style for the response
            conversation_context (str): Previous conversation history context
           
        Returns:
            str: Generated response based on the context and/or web search
        """
        # Check if the query is a simple greeting (same logic as above)
        greeting_keywords = [
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
            "hlw", "g'day", "yo", "heya",
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
            "morning", "afternoon", "evening", "gm", "ga", "ge",
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
       
        query_lower = query.lower().strip()
        is_greeting = False
       
        if query_lower in greeting_keywords:
            is_greeting = True
       
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:
                is_greeting = True
                break
       
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index], {}
       
        if not context and not use_web_knowledge:
            return "I cannot answer this question based on the provided documents."
       
        # [Same URL and context processing logic as in generate_response]
        url_context = ""
        document_context = []
        document_sources = []
        url_sources = []
       
        for i, content_item in enumerate(context):
            if i < len(sources) and sources[i] == "URL Content" and "CONTENT FROM URLS IN QUERY:" in content_item:
                url_context = content_item
                url_sources.append(sources[i])
            else:
                document_context.append(content_item)
                if i < len(sources):
                    document_sources.append(sources[i])
       
        citation_sources = {}
        selected_context, selected_sources = self._prepare_context_comprehensive(document_context, document_sources)
       
        contextualized_content = []
        for i, (content, source) in enumerate(zip(selected_context, selected_sources), 1):
            citation_sources[i] = {
                "source_file": source,
                "text": content,
                "document_id": "Unknown",
                "score": 1.0 - (i * 0.05),
                "display_num": i
            }
            contextualized_content.append(f"Source {i}:\n{content}")
 
        if url_context:
            citation_count = len(contextualized_content) + 1
            contextualized_content.append(f"Source {citation_count}:\n{url_context}")
            selected_sources.append("URL Content")
            citation_sources[citation_count] = {
                "source_file": "URL Content",
                "text": url_context,
                "document_id": "url_content",
                "score": 0.8,
                "display_num": citation_count
            }
 
        full_context = "\n\n".join(contextualized_content)
 
        # Get format-specific guidance for short responses
        format_guidance = self._get_format_guidance(response_format, 'short')
       
        # Include conversation context if available (same as above)
        conversation_prompt = ""
        if conversation_context and conversation_context.strip():
            conversation_prompt = f"""
            RECENT CONVERSATION HISTORY:
            {conversation_context}
            
            The above messages are the previous parts of the same ongoing conversation with the user.
            When generating your response:
            - Consider the conversation flow and any references to previous topics
            - If the user is asking follow-up questions or clarifications, refer back to the relevant previous context
            - Maintain consistency with earlier responses and user preferences shown in the conversation
            - If the current question builds upon or relates to previous discussions, acknowledge that connection
            - Avoid repeating information already covered unless specifically asked for clarification
            - Keep the conversation natural and flowing, showing awareness of the ongoing dialogue
            
            Use this conversation history to provide more contextually aware and relevant responses.
            """
            print(f"Including conversation context in prompt (length: {len(conversation_context)} chars)")
        else:
            print("No conversation context available for this response")
 
        # Get project description if available
        project_description = self._get_project_description(query)
       
        project_guidance = ""
        if project_description and len(project_description.strip()) > 5:
            project_guidance = f"""
            PROJECT CONTEXT:
            {project_description}
           
            IMPORTANT: This project context is only to help you understand the domain and purpose of this project. Your responses MUST be derived exclusively from the uploaded documents provided in the document context.
           
            When formulating your response:
            1. Use this project context solely to understand the general domain and focus of the project
            2. Draw ALL factual information, data, and specific content ONLY from the uploaded documents
            3. Shape your response style and tone to align with this project's domain, but never invent content
            4. DO NOT add information from your general knowledge even if relevant to the project description
            5. If the documents don't contain information relevant to the query, clearly state this limitation
           
            Remember: The project description helps you understand the context, but all response content must come from the uploaded documents.
            """
 
        # Updated system message for JSON response
        system_message = """You are a document analysis expert with conversation memory.
        You must respond in JSON format with the following structure:
        {
            "content": "Your concise yet informative response with proper HTML formatting and citations",
            "citations_used": [list of citation numbers used],
            "format_type": "response_format_used",
            "sources_referenced": ["list", "of", "source", "names"]
        }
       
        Provide a concise yet informative response that maintains conversation continuity. Use source citations [n] for information from the documents."""
        # Define the user prompt for short response with conversation context
        user_prompt = f"""
        You are an expert document analyst tasked with providing focused, accurate, and complete answers based on the provided context. Your goal is to deliver precise responses that directly address the user's question while ensuring no important details are missed.
 
        Based ONLY on the following context from multiple documents and URLs, provide a focused and comprehensive answer to the question. Your response should be directly relevant to the query, well-structured, and include all pertinent details without unnecessary elaboration.
 
        RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
 
        {format_guidance}
 
        {conversation_prompt}
 
        {project_guidance}
 
        FOCUSED RESPONSE GENERATION GUIDELINES:
        - DIRECTLY address the user's specific question as the primary focus
        - Include ALL relevant details and data points that relate to the query
        - Structure information logically with clear organization and proper HTML formatting
        - Provide complete answers - don't leave out important context or supporting details
        - Include specific examples, data, quotes, and evidence from the documents when relevant
        - Use a clear, professional tone that balances conciseness with completeness
        - Ensure every piece of relevant information from the context is utilized appropriately
        - Organize content with proper headings, lists, and formatting for clarity
        - Include quantitative data, statistics, and specific metrics when available
        - Provide sufficient context for understanding while avoiding unnecessary detail
        - Previous conversation context should inform your response only when directly relevant
        - Aim for concise, informative responses that can be quickly read and understood
 
        CONTENT COMPLETENESS REQUIREMENTS:
        - Cover all aspects of the question that are addressed in the provided context
        - Include supporting details that enhance understanding of the main answer
        - Provide specific examples and evidence from the documents
        - Include relevant background information that directly supports the answer
        - Mention related information that adds value to understanding the topic
        - Ensure all important data points and findings are included
        - Use proper citations to clearly attribute information to sources
        - Prioritize clarity and brevity - eliminate redundant or overly detailed explanations
 
        DOCUMENT AND URL CONTEXT:
        {full_context}
 
        USER QUERY: {query}
 
        CRITICAL: Your response MUST be a valid JSON object with this exact structure:
        {{
            "content": "Your focused, complete response with proper HTML formatting, clear structure, and [n] citations. Include all relevant details while maintaining direct relevance to the query.",
            "citations_used": [list of citation numbers used],
            "format_type": "{response_format}",
            "sources_referenced": ["list of source names referenced"]
        }}
 
        CRITICAL CONSTRAINTS:
        - Use ONLY information from the provided document and URL context
        - DO NOT generate speculative or external information beyond the provided context
        - Your response must be based EXCLUSIVELY on the document and URL context provided
        - Provide a complete answer that includes all relevant details from the context
        - Structure your response with proper HTML formatting (h3, h4, p, ul, ol, strong, em tags)
        - Include comprehensive coverage of the topic while maintaining focus on the specific query
        - Use source citations [n] for all information derived from the documents
        - Ensure accuracy and completeness while maintaining conciseness - don't omit important details but avoid verbose explanations
        - Balance thoroughness with relevance - include complete information without unnecessary tangents
        - If the context contains URL content, integrate and properly attribute information from URLs
        - If specific information is missing from the documents, clearly state what is not available
        - Maintain professional clarity while ensuring no relevant detail is overlooked
 
        RESPONSE FORMAT REQUIREMENTS:
            1. Follow with detailed elaboration on each relevant point.
            2. Include specific details, examples, and support from the source content.
            3. Use clear, logical sections when needed.
            4. Conclude with any other contextually relevant insights.
            5. Use <b> for section headers.
            6. Use <p> for body text.
            7. Use <ul> / <li> for lists.
            8. Use <table>, <tr>, <td> for tables, if applicable.
            9. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.
        """
       
        try:
            # Call the OpenAI chat completion API for short response with JSON format
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},  # Force JSON response
                temperature=0.4,
                max_tokens=2000
            )
            token_usage = {
            'input_tokens': completion.usage.prompt_tokens if completion.usage else 0,
            'output_tokens': completion.usage.completion_tokens if completion.usage else 0,
            'total_tokens': completion.usage.total_tokens if completion.usage else 0
            }
                # ===== LOG RAW LLM RESPONSE =====
            raw_llm_response = completion.choices[0].message.content
            print("\n" + "="*100)
            print("🤖 RAW LLM RESPONSE FROM OPENAI (generate_short_response)")
            print("="*100)
            print("Model used:", completion.model)
            print("Response ID:", completion.id)
            print("Raw response content:")
            print(raw_llm_response)
            print("Response length:", len(raw_llm_response))
            print("Response type:", type(raw_llm_response))
            print("="*100)
            # ===== END LOG RAW LLM RESPONSE =====
            # Parse the JSON response
            json_response = self._parse_json_response(completion.choices[0].message.content)
            answer = json_response.get("content", "No response content found.")
            answer = self._clean_citations(answer, citation_sources)
            answer = self._ensure_separate_citation_brackets(answer)
            processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
               
        except Exception as e:
            # If there's an error, try with even fewer context chunks
            reduced_context = "\n\n".join(contextualized_content[:5])
            fallback_prompt = f"""
            Based ONLY on the following context, provide a brief answer to the question: {query}
           
            CONTEXT:
            {reduced_context}
           
            CRITICAL: Respond in JSON format:
            {{
                "content": "Your brief response with HTML formatting",
                "citations_used": [],
                "format_type": "fallback",
                "sources_referenced": []
            }}
            """
           
            try:
                fallback_completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": "You are a document analysis expert. Be brief. Respond in JSON format."},
                        {"role": "user", "content": fallback_prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.4,
                    max_tokens=1024
                )
                token_usage = {
                'input_tokens': fallback_completion.usage.prompt_tokens if fallback_completion.usage else 0,
                'output_tokens': fallback_completion.usage.completion_tokens if fallback_completion.usage else 0,
                'total_tokens': fallback_completion.usage.total_tokens if fallback_completion.usage else 0
                }

                json_response = self._parse_json_response(fallback_completion.choices[0].message.content)
                answer = json_response.get("content", "An error occurred while generating the response.")
                answer = self._clean_citations(answer, citation_sources)
                answer = self._ensure_separate_citation_brackets(answer)
                processed_answer, processed_citations = self._format_citations_for_response(answer, citation_sources)
            except Exception as nested_e:
                # Last resort fallback
                answer = f"An error occurred while generating the response. Please try a more specific question."
                processed_answer = answer
                processed_citations = {}
 
        # Add source information
        source_list = list(set(selected_sources))
        source_info = ", ".join(source_list)
        return f"{processed_answer}\n\n*Sources: {source_info}*", processed_citations, token_usage


    def _get_project_description(self, query):
        """
        Get the project description for the current main project
        
        Returns:
            str: Project description or empty string if not found
        """
        try:
            # Get the main_project_id from the current request
            main_project_id = self.request.data.get('main_project_id')
            
            if not main_project_id:
                return ""
                
            # Get the project from the database
            project = Project.objects.filter(id=main_project_id).first()
            
            if project and project.description:
                return project.description
                
            return ""
        except Exception as e:
            logger.error(f"Error getting project description: {str(e)}")
            return ""

    def get_general_chat_answer(self, query, use_web_knowledge=False, response_length='comprehensive', response_format='natural', conversation_context="", user=None):
        """
        Generate an answer for general chat mode, optionally using web knowledge.
        Enhanced with new Gemini search approach while preserving all original functionality.
 
        Args:
            query (str): The user's query
            use_web_knowledge (bool): Whether to use web search
            response_length (str): 'short' or 'comprehensive'
            response_format (str): The format to use for the response
            conversation_context (str): Previous conversation context
            user: User object for API access
       
        Returns:
            tuple: (response_text, token_usage_dict)
        """
        token_usage = {'input_tokens': 0, 'output_tokens': 0, 'total_tokens': 0}
 
        # Check if the query is a simple greeting - Complete original logic
        greeting_keywords = [
            # Basic greetings
            "hi", "hello", "hey", "greetings", "howdy", "hola", "good morning",
            "good afternoon", "good evening", "what's up", "sup", "hiya",
       
            # Variations with repeated letters
            "hiii", "hiiii", "hiiiii", "helloooo", "hellooo", "heyyy", "heyyyy",
       
            # Short forms/abbreviations
            "hlw", "g'day", "yo", "heya",
       
            # Informal greetings
            "wassup", "whats up", "whatcha up to", "how's it going", "how are you",
            "how r u", "how r you", "how are ya", "how ya doin", "how you doing",
       
            # Time-specific with variations
            "morning", "afternoon", "evening", "gm", "ga", "ge",
       
            # Other languages and cultural greetings
            "namaste", "bonjour", "ciao", "konnichiwa", "aloha", "salut", "hallo",
       
            # Casual text-style greetings
            "sup", "yo yo", "hiya", "hai", "hullo"
        ]
 
        # Convert to lowercase and strip to properly detect greetings
        query_lower = query.lower().strip()
 
        # Check if the query is just a greeting or a greeting followed by a name/simple phrase
        is_greeting = False
 
        # Exact matches
        if query_lower in greeting_keywords:
            is_greeting = True
 
        # Starting with a greeting
        for greeting in greeting_keywords:
            if query_lower.startswith(greeting) and len(query_lower) < len(greeting) + 15:  # Limit to short phrases
                is_greeting = True
                break
 
        # If it's a greeting, respond directly without using more complex logic
        if is_greeting:
            greeting_responses = [
                f"Hello! How can I assist you today?",
                f"Hi there! What can I help you with?",
                f"Greetings! How may I be of service?",
                f"Hello! I'm here to help. What do you need?",
                f"Hey! What questions do you have for me today?"
            ]
            # Select a response using a simple hash of the query to ensure consistent responses
            response_index = hash(query_lower) % len(greeting_responses)
            return greeting_responses[response_index], token_usage
 
        # If not a greeting, proceed with the regular flow
        # If web knowledge is enabled, implement the full web search flow
        if use_web_knowledge:
            try:
                # Step 1: Search the web using enhanced Gemini approach
                print(f"Searching web for: {query}")
                web_response, web_sources = self.get_web_knowledge_response(query, user=user, document_context=None)
 
                if not web_response or not web_sources:
                    return "I couldn't find relevant information on the web for your query.", token_usage
               
                # Extract web source domains for display - FIXED TO USE PROPER DOMAINS
                web_source_domains = []
                for source in web_sources:
                    domain = source.get('domain', '')  # Use the processed domain from get_web_knowledge_response
                    if domain and domain != "unknown.com" and domain != "vertexaisearch.cloud.google.com":
                        web_source_domains.append(domain)
           
                # Remove existing sources from web_response to avoid duplication in GPT processing
                if "*Sources:" in web_response:
                    web_response_clean = web_response.split("*Sources:")[0].strip()
                else:
                    web_response_clean = web_response
 
                # Get format-specific guidance
                format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
 
                # Include conversation context if available
                conversation_prompt = ""
                if conversation_context and conversation_context.strip():
                    conversation_prompt = f"""
                    RECENT CONVERSATION HISTORY:
                    {conversation_context}
                   
                    The above messages are the previous parts of the same ongoing conversation with the user.
                    When generating your response:
                    - Consider the conversation flow and any references to previous topics
                    - If the user is asking follow-up questions or clarifications, refer back to the relevant previous context
                    - Maintain consistency with earlier responses and user preferences shown in the conversation
                    - If the current question builds upon or relates to previous discussions, acknowledge that connection
                    - Avoid repeating information already covered unless specifically asked for clarification
                    - Keep the conversation natural and flowing, showing awareness of the ongoing dialogue
                   
                    Use this conversation history to provide more contextually aware and relevant responses.
                    """
                    print(f"Including conversation context in prompt (length: {len(conversation_context)} chars)")
                else:
                    print("No conversation context available for this response")
           
                # Updated system message for JSON response
                system_message = """You are a web research assistant.
                You must respond in JSON format with the following structure:
                {
                    "content": "Your response with proper HTML formatting",
                    "web_sources_used": true,
                    "response_type": "web_enhanced",
                    "format_type": "response_format_used"
                }
               
                Provide detailed, comprehensive responses based on web sources."""
 
                # Create the prompt for OpenAI - REMOVED THE SOURCES_JSON EXTRACTION REQUIREMENT
                prompt = f"""
                You are a web research assistant. Based ONLY on the following information from multiple web sources, answer the user question.
 
                If relevant details are missing or contradictory, clearly acknowledge this and summarize available related content. Be helpful by highlighting potentially useful information even if it doesn't fully resolve the question. Provide quantitative details where available.
 
                RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
 
                {format_guidance}
 
                RESPONSE GENERATION GUIDELINES:
                - Provide a DETAILED, COMPREHENSIVE, and THOROUGH answer.
                - Include ALL relevant information from the web sources below.
                - Prioritize completeness over brevity — include important details.
                - If multiple perspectives or data points exist, include all of them.
                - When appropriate, structure the information clearly with headings.
                - Maintain a natural, conversational tone.
                - If the web sources do NOT contain relevant information, clearly state that and summarize any related or tangential insights.
 
                QUESTION: {query}
 
                conversation history:
                {conversation_prompt}
 
                INFORMATION FROM WEB SOURCES:
                {web_response_clean}
 
                CRITICAL: Your response MUST be a valid JSON object with this exact structure:
                {{
                    "content": "Your response with HTML tags for structure and formatting",
                    "web_sources_used": true,
                    "response_type": "web_enhanced",
                    "format_type": "{response_format}"
                }}
 
                RESPONSE FORMAT REQUIREMENTS:
                1. Follow with detailed elaboration on each relevant point.
                2. Include specific details, examples, and support from the source content.
                3. Use clear, logical sections when needed.
                4. Conclude with any other contextually relevant insights.
                5. Use <b> for section headers.
                6. Use <p> for body text.
                7. Use <ul> / <li> for lists.
                8. Use <table>, <tr>, <td> for tables, if applicable.
                9. Avoid using Markdown formatting (**bold** and etc.). Use HTML only.
 
                CRITICAL CONSTRAINTS:
                - Ensure clarity, accuracy, and full coverage of the relevant content.
 
                RESPONSE LENGTH: {'Provide a focused, concise response prioritizing the most important information.' if response_length == 'short' else 'Provide a comprehensive, detailed response that thoroughly covers the topic.'}
                """
 
                temperature = 0.3 if response_format in ['factual_brief', 'technical_deep_dive'] else 0.5
                max_tokens = 2000 if response_length == 'short' else 2000
 
                print(f"Calling OpenAI API with temperature={temperature}, max_tokens={max_tokens}")
                response = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=temperature,
                    max_tokens=max_tokens
                )
 
                # Add token tracking after the API call:
                if hasattr(response, 'usage') and response.usage:
                    token_usage['input_tokens'] += response.usage.prompt_tokens
                    token_usage['output_tokens'] += response.usage.completion_tokens
                    token_usage['total_tokens'] += response.usage.total_tokens
 
                # Parse the JSON response
                json_response = self._parse_json_response(response.choices[0].message.content)
                formatted_web_response = json_response.get("content", "No response content found.")
 
                # Debug: Print the full response to see what GPT returned
                print("=== FULL GPT RESPONSE ===")
                print(formatted_web_response)
                print("=== END FULL RESPONSE ===")
 
                # FIXED: Use the properly processed domains from web_sources instead of trying to extract from content
                if web_source_domains:
                    source_info = ", ".join(web_source_domains)
                    final_response = f"{formatted_web_response}\n\n*Sources: {source_info}*"
                    print("=== FINAL RESPONSE WITH SOURCES ===")
                    print(final_response)
                    print("=== END FINAL RESPONSE ===")
                    return final_response, token_usage
                else:
                    return formatted_web_response, token_usage
           
            except Exception as e:
                logger.error(f"Error in web knowledge general chat: {str(e)}", exc_info=True)
                return f"I encountered an error while searching the web: {str(e)}. Please try a different question or try again later.", token_usage
 
        # If no web knowledge requested, use standard chat completion
        else:
            # Format-specific guidance
            format_guidance = self._get_format_guidance(response_format, 'short' if response_length == 'short' else 'comprehensive')
       
            # Configure conciseness based on response length
            conciseness = "Be concise and to-the-point." if response_length == 'short' else "Provide a comprehensive and detailed response."
       
            # Updated system message for JSON response
            system_message = """You are a helpful, friendly AI assistant.
            You must respond in JSON format with the following structure:
            {
                "content": "Your response with proper HTML formatting",
                "response_type": "general_chat",
                "format_type": "response_format_used",
                "web_sources_used": false
            }
           
            Provide informative and thoughtful responses using HTML tags for structure."""
 
            # Create a prompt for the general chat with JSON requirement
            prompt = f"""
            Please answer the following question in a helpful, informative way.
           
            RESPONSE FORMAT: {response_format.replace('_', ' ').title()}
           
            {format_guidance}
           
            {conciseness}
           
            Use HTML tags for structure (<b>, <p>, <ul>, <li>) to enhance readability.
           
            USER QUERY: {query}
           
            CRITICAL: Your response MUST be a valid JSON object with this exact structure:
            {{
                "content": "Your response with HTML formatting",
                "response_type": "general_chat",
                "format_type": "{response_format}",
                "web_sources_used": false
            }}
            """
       
            try:
                # Use the OpenAI API with JSON format
                completion = client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_message},
                        {"role": "user", "content": prompt}
                    ],
                    response_format={"type": "json_object"},  # Force JSON response
                    temperature=0.5,
                    max_tokens=2000 if response_length == 'comprehensive' else 800
                )
                # Add token tracking after the API call:
                if hasattr(completion, 'usage') and completion.usage:
                    token_usage['input_tokens'] += completion.usage.prompt_tokens
                    token_usage['output_tokens'] += completion.usage.completion_tokens
                    token_usage['total_tokens'] += completion.usage.total_tokens
               
                # Parse the JSON response
                json_response = self._parse_json_response(completion.choices[0].message.content)
                return json_response.get("content", "No response content found."), token_usage
           
            except Exception as e:
                logger.error(f"Error in general chat: {str(e)}", exc_info=True)
                return f"I'm sorry, I encountered an error while processing your request. Please try again or rephrase your question.", token_usage


    def get_embeddings(self, texts):
        """
        Gets embeddings using OpenAI's text-embedding-3-small model.
        """
        try:
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small"
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            logger.error(f"Error getting embeddings: {str(e)}")
            return []
    
    
    def search_similar_content(self, query, processed_docs, metadata_store=None, k=80):
            """
            Enhanced search function using pgvector for similarity search
            Maintains the same return structure as the original FAISS version
            """
            print(f"\n🔍 SEARCH DEBUG: Starting pgvector search for query: '{query}'")
            print(f"🔍 SEARCH DEBUG: Number of processed docs: {len(processed_docs)}")
       
            # Get embeddings for the query
            query_embedding = self.get_embeddings([query])
            if not query_embedding:
                print("❌ SEARCH DEBUG: Failed to get query embedding")
                return [], [], {}
       
            print("✅ SEARCH DEBUG: Got query embedding")
       
            # Search results storage
            all_results = []
            all_distances = []
            all_sources = []
            citation_mapping = {}
            citation_count = 0
            seen_content_hashes = set()
            valid_docs_found = False
       
            for i, proc_doc in enumerate(processed_docs):
                print(f"\n🔍 SEARCH DEBUG: Processing document {i+1}: {proc_doc.document.filename}")
           
                # Skip if not using pgvector or no embeddings exist
                if proc_doc.storage_type != 'pgvector':
                    print(f"❌ SEARCH DEBUG: Document {proc_doc.document.filename} not using pgvector storage")
                    continue
           
                # Check if embeddings exist for this document
                embeddings_exist = DocumentEmbedding.objects.filter(document=proc_doc.document).exists()
                if not embeddings_exist:
                    print(f"❌ SEARCH DEBUG: No embeddings found for {proc_doc.document.filename}")
                    continue
           
                valid_docs_found = True
           
                try:
                    # Perform pgvector similarity search
                    from pgvector.django import CosineDistance
               
                    # Calculate search limit per document
                    search_k = min(k // len(processed_docs) + 10, proc_doc.chunks_count or 50)
               
                    # Query embeddings using pgvector
                    similar_embeddings = DocumentEmbedding.objects.filter(
                        document=proc_doc.document
                    ).annotate(
                        distance=CosineDistance('embedding', query_embedding[0])
                    ).order_by('distance')[:search_k]
               
                    print(f"✅ SEARCH DEBUG: pgvector search returned {len(similar_embeddings)} results")
               
                    # Process results
                    doc_results_count = 0
                    for embedding_obj in similar_embeddings:
                        distance = float(embedding_obj.distance)
                   
                        # Filter by distance threshold
                        if distance < 0.8:  # Cosine distance threshold (lower is better)
                            content = embedding_obj.content
                       
                            if content and content.strip():
                                # Duplicate detection
                                content_hash = hash(content[:150])
                           
                                if content_hash not in seen_content_hashes:
                                    seen_content_hashes.add(content_hash)
                                    doc_results_count += 1
                                    citation_count += 1
                               
                                    # Store citation mapping
                                    citation_mapping[citation_count] = {
                                        'source': proc_doc.document.filename,
                                        'text': content,
                                        'relevance_score': distance,
                                        'document_id': proc_doc.document.id,
                                        'chunk_idx': embedding_obj.chunk_id,
                                        'display_num': citation_count
                                    }
                               
                                    all_results.append(content)
                                    all_distances.append(distance)
                                    all_sources.append(proc_doc.document.filename)
                               
                                    print(f"   ✅ Added result {citation_count}: distance={distance:.4f}, content='{content[:100]}...'")
               
                    print(f"✅ SEARCH DEBUG: Added {doc_results_count} results from {proc_doc.document.filename}")
               
                except Exception as e:
                    print(f"❌ SEARCH DEBUG: Error searching embeddings for {proc_doc.document.filename}: {str(e)}")
               
                    # Fallback to text search
                    try:
                        print(f"🔄 SEARCH DEBUG: Falling back to text search for {proc_doc.document.filename}")
                        query_lower = query.lower()
                   
                        text_matches = DocumentEmbedding.objects.filter(
                            document=proc_doc.document,
                            content__icontains=query_lower
                        )[:10]
                   
                        print(f"   Found {len(text_matches)} text matches")
                   
                        for embedding_obj in text_matches:
                            content = embedding_obj.content
                            content_hash = hash(content[:150])
                       
                            if content_hash not in seen_content_hashes:
                                seen_content_hashes.add(content_hash)
                                citation_count += 1
                           
                                citation_mapping[citation_count] = {
                                    'source': proc_doc.document.filename,
                                    'text': content,
                                    'relevance_score': 0.5,
                                    'document_id': proc_doc.document.id,
                                    'chunk_idx': embedding_obj.chunk_id,
                                    'display_num': citation_count
                                }
                           
                                all_results.append(content)
                                all_distances.append(0.5)
                                all_sources.append(proc_doc.document.filename)
                           
                                print(f"   ✅ Added text match {citation_count}: '{content[:100]}...'")
                           
                    except Exception as fallback_error:
                        print(f"❌ SEARCH DEBUG: Fallback search failed: {str(fallback_error)}")
                        continue
       
            print(f"\n🔍 SEARCH DEBUG: Search completed. Found {len(all_results)} total results")
       
            if not all_results:
                if valid_docs_found:
                    print(f"⚠️ SEARCH DEBUG: No matching content found in documents for query: {query}")
                else:
                    print(f"❌ SEARCH DEBUG: No valid documents found to search")
                return [], [], {}
       
            # Sort results by distance (lower is better for cosine distance)
            combined_results = list(zip(all_results, all_sources, all_distances, range(1, len(all_results) + 1)))
            sorted_results = sorted(combined_results, key=lambda x: x[2])
       
            print(f"🔍 SEARCH DEBUG: Sorted results by distance")
            for i, (content, source, distance, _) in enumerate(sorted_results[:5]):
                print(f"   Result {i+1}: distance={distance:.4f}, source='{source}', content='{content[:100]}...'")
       
            # Extract sorted results
            results = [res[0] for res in sorted_results]
            sources = [res[1] for res in sorted_results]
       
            # Update citation mapping with new ranking
            reordered_citation_mapping = {}
            for new_idx, (_, _, _, old_key) in enumerate(sorted_results, 1):
                if old_key in citation_mapping:
                    reordered_citation_mapping[new_idx] = citation_mapping[old_key].copy()
                    reordered_citation_mapping[new_idx]['display_num'] = new_idx
       
            citation_mapping = reordered_citation_mapping
       
            # Final deduplication and filtering
            seen_content = set()
            filtered_results = []
            filtered_sources = []
            filtered_citation_mapping = {}
            new_citation_count = 0
       
            for i, (content, source) in enumerate(zip(results, sources)):
                content_hash = content[:80] if content else ""
                if content_hash and content_hash not in seen_content:
                    seen_content.add(content_hash)
                    filtered_results.append(content)
                    filtered_sources.append(source)
   
                    new_citation_count += 1
                    old_idx = i + 1
                    if old_idx in citation_mapping:
                        filtered_citation_mapping[new_citation_count] = citation_mapping[old_idx].copy()
                        filtered_citation_mapping[new_citation_count]['display_num'] = new_citation_count
       
            # Return final results (limit to 30 for comprehensive responses)
            max_results = min(len(filtered_results), 30)
            final_results = filtered_results[:max_results]
            final_sources = filtered_sources[:max_results]
            final_citations = {k: v for k, v in filtered_citation_mapping.items() if k <= max_results}
       
            print(f"\n✅ SEARCH DEBUG: Final results: {len(final_results)} items")
            for i, (content, source) in enumerate(zip(final_results[:5], final_sources[:5])):
                print(f"   Final {i+1}: source='{source}', content='{content[:100]}...'")
       
            return final_results, final_sources, final_citations


    def _prepare_context_comprehensive(self, context, sources):
        """
        Enhanced context preparation for comprehensive responses
        
        Args:
            context (list): List of context chunks
            sources (list): List of source documents for each context chunk
            
        Returns:
            tuple: (selected_context, selected_sources)
        """
        selected_context = []
        selected_sources = []
        seen_content_hashes = set()
        
        # Significantly increased context limits for comprehensive responses
        initial_count = min(20, len(context))  # Increased from 8 to 20
        for i in range(initial_count):
            selected_context.append(context[i])
            selected_sources.append(sources[i])
            seen_content_hashes.add(context[i][:80])  # Reduced hash length for less strict deduplication
        
        # Add more diverse chunks from the remaining context
        for i in range(initial_count, len(context)):
            content_hash = context[i][:80]  # Reduced from 100 to 80
            if content_hash not in seen_content_hashes:
                selected_context.append(context[i])
                selected_sources.append(sources[i])
                seen_content_hashes.add(content_hash)
                # Significantly increased context limit from 15 to 40
                if len(selected_context) >= 40:
                    break
        
        print(f"🔍 CONTEXT DEBUG: Selected {len(selected_context)} context chunks for comprehensive response")
        return selected_context, selected_sources
            
    # Keep general follow-up question generation as is
    def generate_general_follow_up_questions(self, question, answer, user=None):
        system_message = "You are an expert at generating relevant follow-up questions. Always respond with valid JSON format."
       
        user_prompt = f"""
        Based on this user question and your answer, suggest 3 relevant follow-up questions that the user might want to ask next.
        The questions should be short, interesting, and directly related to the topic.
       
        User Question: {question}
        Your Answer (abbreviated): {answer[:500]}...
       
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "questions": ["question1", "question2", "question3"],
            "topic": "main_topic_discussed"
        }}
        """
       
        try:            
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.3,
                max_tokens=500
            )
           
            # Parse JSON response
            json_response = json.loads(completion.choices[0].message.content)
            questions = json_response.get("questions", [])
            return questions[:3]
           
        except Exception as e:
            logger.error(f"Error generating follow-up questions: {str(e)}")
            return [
                "What else would you like to know about this topic?",
                "Would you like me to elaborate on any specific point?",
                "Do you have any other questions I can help with?"
            ]
 
    def generate_follow_up_questions(self, context, user=None):
        context_sample = "\n".join(context[:3]) if context else ""
       
        system_message = "You are an expert at generating relevant follow-up questions based on context. Always respond with valid JSON format."
       
        user_prompt = f"""
        Based on this context, suggest 3 relevant follow-up questions. The questions should be short and directly related to the content.
       
        Context: {context_sample}
       
        CRITICAL: Your response MUST be a valid JSON object with this structure:
        {{
            "questions": ["question1", "question2", "question3"],
            "context_topic": "main_topic_from_context"
        }}
        """
       
        try:
            completion = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.3,
                max_tokens=500
            )
           
            # Parse JSON response
            json_response = json.loads(completion.choices[0].message.content)
            questions = json_response.get("questions", [])
            return questions[:3]
           
        except Exception as e:
            logger.error(f"Error generating follow-up questions: {str(e)}")
            return [
                "What additional information would you like to know?",
                "Would you like me to elaborate on any specific point?",
                "How can I help clarify this information further?"
            ]

class DeleteMessagePairView(APIView):
    permission_classes = [IsAuthenticated]
 
    def post(self, request):
        print("DeleteMessagePairView POST data:", request.data)
        user = request.user
        user_message_id = request.data.get('user_message_id')
        assistant_message_id = request.data.get('assistant_message_id')
        conversation_id = request.data.get('conversation_id')
 
 
        if not (user_message_id and assistant_message_id and conversation_id):
            return Response({"error": "Missing required parameters."}, status=status.HTTP_400_BAD_REQUEST)
 
        try:
            with transaction.atomic():
                # Ensure both messages belong to the user and conversation
                chat_history = ChatHistory.objects.get(conversation_id=conversation_id, user=user)
                user_msg = ChatMessage.objects.get(id=user_message_id, chat_history=chat_history, role="user")
                assistant_msg = ChatMessage.objects.get(id=assistant_message_id, chat_history=chat_history, role="assistant")
 
                user_msg.delete()
                assistant_msg.delete()
 
            return Response({"success": True}, status=status.HTTP_200_OK)
        except ChatHistory.DoesNotExist:
            return Response({"error": "Conversation not found."}, status=status.HTTP_404_NOT_FOUND)
        except ChatMessage.DoesNotExist:
            return Response({"error": "Message not found."}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class GetConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request, conversation_id=None):
        try:
            user = request.user
            
            if conversation_id:
                # Fetch specific conversation
                try:
                    conversation = ChatHistory.objects.get(
                        conversation_id=conversation_id, 
                        user=user
                    )
                    
                    # Retrieve all messages for this conversation
                    messages = conversation.messages_NB.all().order_by('created_at')
                    
                    # Prepare message list - UPDATED to include all badge properties
                    message_list = []
                    for message in messages:
                        message_data = {
                            'id': message.id,
                            'role': message.role,
                            'content': message.content,
                            'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': message.citations or []
                        }
                        if message.role == 'assistant':
                            message_data['use_web_knowledge'] = getattr(message, 'use_web_knowledge', False)
                            message_data['response_length'] = getattr(message, 'response_length', 'comprehensive')
                            message_data['response_format'] = getattr(message, 'response_format', 'natural')
                            message_data['sources_info'] = getattr(message, 'sources', "")
                            # ADD THIS LOGIC:
                            if getattr(message, 'sources', '') == "Image Analysis":
                                message_data['context_mode'] = "image"
                        message_list.append(message_data)
                    
                    return Response({
                        'conversation_id': str(conversation.conversation_id),
                        'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                        'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                        'messages': message_list,
                        'follow_up_questions': conversation.follow_up_questions or [],
                        'selected_documents': [
                            str(doc.id) for doc in conversation.documents.all()
                        ]
                    }, status=status.HTTP_200_OK)
                
                except ChatHistory.DoesNotExist:
                    return Response(
                        {'error': 'Conversation not found'}, 
                        status=status.HTTP_404_NOT_FOUND
                    )
            
            else:
                # Fetch all conversations for the user
                conversations = ChatHistory.objects.filter(user=user) \
                    .filter(is_active=True) \
                    .order_by('-created_at')
                
                conversation_list = []
                for conversation in conversations:
                    # Get the first and last messages
                    messages = conversation.messages_NB.all().order_by('created_at')
                    
                    if messages:
                        first_message = messages.first()
                        last_message = messages.last()
                        
                        conversation_list.append({
                            'conversation_id': str(conversation.conversation_id),
                            'title': conversation.title or f"Chat from {conversation.created_at.strftime('%Y-%m-%d %H:%M')}",
                            'created_at': conversation.created_at.strftime('%Y-%m-%d %H:%M'),
                            'first_message': first_message.content,
                            'last_message': last_message.content,
                            'message_count': messages.count(),
                            'follow_up_questions': conversation.follow_up_questions or []
                        })
                
                return Response(conversation_list, status=status.HTTP_200_OK)
        
        except Exception as e:
            return Response(
                {'error': f'Failed to fetch conversations: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    # Add PATCH method to handle title updates
    def patch(self, request, conversation_id):
        try:
            # Log incoming request data for debugging
            print(f"PATCH request for conversation {conversation_id}")
            print(f"Request data: {request.data}")

            # Validate input
            new_title = request.data.get('title')
            is_active = request.data.get('is_active', True)

            if not new_title or not new_title.strip():
                return Response(
                    {'error': 'Title cannot be empty'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            # Find the conversation
            try:
                conversation = ChatHistory.objects.get(
                    conversation_id=conversation_id, 
                    user=request.user
                )
            except ChatHistory.DoesNotExist:
                return Response(
                    {'error': 'Conversation not found'},
                    status=status.HTTP_404_NOT_FOUND
                )

            # Update the title and active status
            conversation.title = new_title.strip()
            conversation.is_active = is_active
            conversation.save()

            # Log successful update
            print(f"Conversation {conversation_id} updated successfully")
            print(f"New title: {conversation.title}")

            return Response({
                'message': 'Conversation title updated successfully',
                'conversation_id': str(conversation.conversation_id),
                'new_title': conversation.title,
                'is_active': conversation.is_active
            }, status=status.HTTP_200_OK)

        except Exception as e:
            # Comprehensive error logging
            print(f"Error updating conversation title: {str(e)}")
            logger.error(f"Conversation title update error: {str(e)}")

            return Response(
                {'error': 'Failed to update conversation title', 'details': str(e)},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    # Also implement PUT for completeness
    def put(self, request, conversation_id):
        return self.patch(request, conversation_id)

# Optional: Add a view to delete a conversation
class DeleteConversationView(APIView):
    permission_classes = [IsAuthenticated]
    
    def delete(self, request, conversation_id):
        try:
            conversation = ChatHistory.objects.get(
                conversation_id=conversation_id, 
                user=request.user
            )
            
            conversation.delete()
            return Response(
                {'message': 'Conversation deleted successfully'}, 
                status=status.HTTP_200_OK
            )
        
        except ChatHistory.DoesNotExist:
            return Response(
                {'error': 'Conversation not found'}, 
                status=status.HTTP_404_NOT_FOUND
            )
        except Exception as e:
            return Response(
                {'error': f'Failed to delete conversation: {str(e)}'}, 
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


class OriginalDocumentView(APIView):
    permission_classes = [IsAuthenticated]
    
    def get(self, request, document_id):
        try:
            # Get the document making sure it belongs to the user
            document = get_object_or_404(Document, id=document_id, user=request.user)
            
            # Get the file path
            file_path = document.file.path
            
            # Check if file exists
            if not os.path.exists(file_path):
                return Response(
                    {'error': 'Document file not found'},
                    status=status.HTTP_404_NOT_FOUND
                )
                
            # Determine content type based on file extension
            content_type = self.get_content_type(file_path)
            
            # Create a file response
            response = FileResponse(
                open(file_path, 'rb'),
                content_type=content_type
            )
            
            # Set content disposition to attachment with the original filename
            response['Content-Disposition'] = f'inline; filename="{document.filename}"'
            
            # Log the file access
            self.log_document_view(request.user, document)
            
            return response
            
        except Exception as e:
            print(f"Error serving original document: {str(e)}")
            return Response(
                {'error': f'Failed to retrieve document: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
    
    def get_content_type(self, file_path):
        """Determine content type based on file extension"""
        extension = os.path.splitext(file_path)[1].lower()
        
        content_types = {
            '.pdf': 'application/pdf',
            '.doc': 'application/msword',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xls': 'application/vnd.ms-excel',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.txt': 'text/plain',
            '.csv': 'text/csv',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.md': 'text/markdown'
        }
        
        return content_types.get(extension, 'application/octet-stream')
    
    def log_document_view(self, user, document):
        """Log document view for analytics (optional)"""
        try:
            # You can create a DocumentViewLog model to track this
            # For now, we'll just log it
            logger.info(f"User {user.username} viewed document {document.id}: {document.filename}")
            
        except Exception as e:
            logger.error(f"Error logging document view: {str(e)}")
            # Don't raise the exception - this is a non-critical feature

class DocumentViewLogView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request, document_id):
        """Log when a user views a document"""
        try:
            # Get the document ensuring it belongs to the user
            document = get_object_or_404(Document, id=document_id, user=request.user)
    
            
            # For simplicity, just update the document's view count
            document.view_count = F('view_count') + 1  # This requires adding view_count field to Document model
            document.last_viewed_at = timezone.now()  # This requires adding last_viewed_at field to Document model
            document.save(update_fields=['view_count', 'last_viewed_at'])

            # Get main_project_id from the document
            main_project_id = document.main_project_id
            
            # Update project timestamp to track activity
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
            
            return Response({
                'success': True,
                'message': 'Document view logged successfully'
            }, status=status.HTTP_200_OK)
            
        except Document.DoesNotExist:
            return Response({
                'error': 'Document not found'
            }, status=status.HTTP_404_NOT_FOUND)
            
        except Exception as e:
            print(f"Error logging document view: {str(e)}")
            # Return success anyway - this is a non-critical feature
            return Response({
                'success': True,
                'message': 'Continued without logging view'
            }, status=status.HTTP_200_OK)




class DocumentContentSearchView(APIView):
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        try:
            user = request.user
            query = request.data.get('query')
            main_project_id = request.data.get('main_project_id')
            
            if not query or not main_project_id:
                return Response({
                    'error': 'Query and project ID are required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Get all documents for this user and project
            documents = Document.objects.filter(
                user=user, 
                main_project_id=main_project_id
            )
            
            search_results = []
            
            # For each document, search for query in the content
            for doc in documents:
                try:
                    # Get the processed index for this document
                    processed_index = ProcessedIndex.objects.get(document=doc)
                    
                    # First, check if this is a LlamaParse document
                    if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                        # This is a LlamaParse document, read the markdown content directly
                        try:
                            with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                                markdown_content = f.read()
                            
                            # Convert query to lowercase for case-insensitive search
                            query_lower = query.lower()
                            
                            # Check if query exists in the markdown content
                            if query_lower in markdown_content.lower():
                                # Found a match - get the surrounding context
                                index = markdown_content.lower().find(query_lower)
                                start = max(0, index - 100)
                                end = min(len(markdown_content), index + len(query) + 100)
                                match_context = markdown_content[start:end]
                                
                                # Count occurrences
                                match_count = markdown_content.lower().count(query_lower)
                                
                                search_results.append({
                                    'id': doc.id,
                                    'filename': doc.filename,
                                    'match_count': match_count,
                                    'match_context': match_context,
                                    'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M')
                                })
                                
                                # Continue to the next document
                                continue
                        except Exception as e:
                            logger.error(f"Error searching in markdown file for document {doc.id}: {str(e)}")
                            # Continue with the FAISS approach as fallback
                    
                    # Standard FAISS approach for non-LlamaParse documents
                    if processed_index.faiss_index and processed_index.metadata:
                        # Check if the files exist
                        if not os.path.exists(processed_index.faiss_index) or not os.path.exists(processed_index.metadata):
                            continue
                            
                        # Load the metadata (chunks)
                        with open(processed_index.metadata, 'rb') as f:
                            chunks = pickle.load(f)
                        
                        # Convert query to lowercase for case-insensitive search
                        query_lower = query.lower()
                        matches = []
                        
                        # Search in each chunk
                        for chunk in chunks:
                            chunk_text = chunk.get('text', '') if isinstance(chunk, dict) else str(chunk)
                            if chunk_text and query_lower in chunk_text.lower():
                                # Found a match, add the context
                                matches.append(chunk_text)
                        
                        # If we found matches, add to results
                        if matches:
                            # Use the first match for preview, but count all matches
                            search_results.append({
                                'id': doc.id,
                                'filename': doc.filename,
                                'match_count': len(matches),
                                'match_context': matches[0] if matches else "",
                                'uploaded_at': doc.uploaded_at.strftime('%Y-%m-%d %H:%M'),
                                'source_type': 'chunks'
                            })
                            
                except ProcessedIndex.DoesNotExist:
                    # Skip documents that haven't been processed
                    logger.info(f"No processed index found for document {doc.id}")
                    continue
                except Exception as e:
                    logger.error(f"Error searching document {doc.id}: {str(e)}")
                    continue
            
            # Sort results by number of matches (most relevant first)
            search_results.sort(key=lambda x: x['match_count'], reverse=True)
            
            # Update project timestamp
            if main_project_id:
                update_project_timestamp(main_project_id, user)
            
            return Response({
                'results': search_results,
                'total_count': len(search_results),
                'query': query
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            logger.error(f"Error in document content search: {str(e)}", exc_info=True)
            return Response({
                'error': f'An error occurred: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 

class ProcessCitationsView(APIView):
    """
    API endpoint to process a response and map citations to relevant parts of the text.
    This can be called by the frontend when displaying a message.
    """
    permission_classes = [IsAuthenticated]
    
    def post(self, request):
        try:
            response_text = request.data.get('response_text')
            citations = request.data.get('citations')
            
            if not response_text or not citations:
                return Response({
                    'error': 'Both response_text and citations are required'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Process the response text with citations
            processed_text, enhanced_citations = self.process_response_with_citations(response_text, citations)
            
            return Response({
                'processed_text': processed_text,
                'enhanced_citations': enhanced_citations
            }, status=status.HTTP_200_OK)
            
        except Exception as e:
            import traceback
            print(f"Citation processing error: {str(e)}")
            print(traceback.format_exc())
            return Response({
                'error': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
    def extract_citation_metadata(self, citation_text, source_file):
        """
        Extract metadata like page numbers and section titles from citation text.
        
        Args:
            citation_text (str): The citation text from the document
            source_file (str): The source file name
            
        Returns:
            dict: Enhanced citation metadata
        """
        # Default citation structure
        citation = {
            'source_file': source_file,
            'page_number': 'Unknown',
            'section_title': 'Unknown',
            'snippet': citation_text,
        }
        
        # Try to extract page numbers - common formats
        page_matches = re.findall(r'(page|p\.?)\s*(\d+)', citation_text, re.IGNORECASE)
        if page_matches:
            citation['page_number'] = page_matches[0][1]
        
        # Try to extract section titles - look for section headings
        section_matches = re.findall(r'(?:section|heading):\s*["\'"]([^\'"\']+)[\'"\']', citation_text, re.IGNORECASE)
        if section_matches:
            citation['section_title'] = section_matches[0]
        
        # Alternative: Look for text in quotes that might be section titles
        if citation['section_title'] == 'Unknown':
            quote_matches = re.findall(r'[\'"\']([^\'"\']{5,50})[\'"\']', citation_text)
            if quote_matches:
                citation['section_title'] = quote_matches[0]
        
        # Look for date information
        date_matches = re.findall(r'\b(\d{4}[-/]\d{1,2}[-/]\d{1,2})\b', citation_text)
        if date_matches:
            citation['date'] = date_matches[0]
        
        return citation
    
    def process_response_with_citations(self, response_text, citations):
        """
        Maps citations to specific parts of the response text based on content similarity.
        
        Args:
            response_text (str): The response text from the LLM
            citations (list): List of citation objects with source information
            
        Returns:
            tuple: (processed_text, enhanced_citations)
        """
        if not citations or not response_text:
            return response_text, citations
        
        # First, try to improve citation metadata
        enhanced_citations = []

        for idx, citation in enumerate(citations):
            # Get basic info
            source_file = citation.get('source_file', 'Unknown Document')
            snippet = citation.get('snippet', '')
            
            # If metadata is missing, try to extract it
            if citation.get('page_number') in [None, 'Unknown'] or citation.get('section_title') in [None, 'Unknown']:
                enhanced_citation = self.extract_citation_metadata(snippet, source_file)
                # Preserve original data where available
                for key, value in citation.items():
                    if value and value != 'Unknown':
                        enhanced_citation[key] = value
                enhanced_citations.append(enhanced_citation)
            else:
                enhanced_citations.append(citation)
        
        # Remove HTML tags for better text processing
        clean_text = strip_tags(response_text)
        
        # Step 1: Split the response into sentences
        sentences = sent_tokenize(clean_text)
        
        # Step 2: For each citation, find the most relevant sentence
        citation_mapping = {}  # Maps citation index to sentence index
        
        for citation_idx, citation in enumerate(enhanced_citations):
            snippet = citation.get('snippet', '').lower()
            if not snippet:
                continue
                
            # Generate keywords from the snippet
            keywords = set(re.findall(r'\b\w+\b', snippet.lower()))
            keywords = {k for k in keywords if len(k) > 3}  # Filter out short words
            
            # Score each sentence based on keyword overlap
            best_score = 0
            best_sentence_idx = -1
            
            for sent_idx, sentence in enumerate(sentences):
                sentence_lower = sentence.lower()
                sentence_words = set(re.findall(r'\b\w+\b', sentence_lower))
                
                # Calculate overlap score
                overlap = keywords.intersection(sentence_words)
                if overlap:
                    score = len(overlap) / max(len(keywords), 1)  # Avoid division by zero
                    
                    # Give higher score for exact phrase matches
                    for i in range(0, len(snippet), 10):
                        if i + 20 <= len(snippet):
                            phrase = snippet[i:i+20]
                            if phrase in sentence_lower:
                                score += 0.5
                    
                    if score > best_score:
                        best_score = score
                        best_sentence_idx = sent_idx
            
            # Map citation to sentence if we found a good match
            if best_score > 0.2 and best_sentence_idx >= 0:  # Threshold to ensure relevance
                if best_sentence_idx not in citation_mapping:
                    citation_mapping[best_sentence_idx] = []
                citation_mapping[best_sentence_idx].append(citation_idx)
        
        # Step 3: Reinsert the HTML and add citation markers
        # Start by marking positions in the original response_text where sentences end
        sentence_positions = []
        current_pos = 0
        
        for sentence in sentences:
            # Find this sentence in the original text
            sentence_pos = response_text.find(sentence, current_pos)
            if sentence_pos >= 0:
                sentence_end = sentence_pos + len(sentence)
                sentence_positions.append(sentence_end)
                current_pos = sentence_end
        
        # Now insert citation markers at the end of relevant sentences
        result = response_text
        offset = 0  # Track position offset as we insert markers
        
        for sent_idx, end_pos in enumerate(sentence_positions):
            if sent_idx in citation_mapping:
                # Create citation markers
                citation_markers = ""
                for citation_idx in citation_mapping[sent_idx]:
                    citation_markers += f'<citation id="{citation_idx}"></citation>'
                
                # Insert at the adjusted position
                insert_pos = end_pos + offset
                if insert_pos <= len(result):
                    result = result[:insert_pos] + citation_markers + result[insert_pos:]
                    offset += len(citation_markers)
        
        return result, enhanced_citations


# codes for chat download


class ChatExportView(APIView):
    permission_classes = [IsAuthenticated]
 
    def remove_citations_from_html(self, html):
        """
        Remove all [n], [n][m], [n, m], [Source: n], [Sources: n, m], and 'Sources: n, m' style citations from HTML/text content.
        """
        if not html:
            return html
 
        # Remove [Source: n], [Sources: n, m], [Source: n, m], etc.
        html = re.sub(r'\[Sources?:\s*[\d,\s]+\]', '', html, flags=re.IGNORECASE)
        # Remove 'Sources: n, m' or 'Source: n' at the end of sentences or lines
        html = re.sub(r'Sources?:\s*[\d,\s]+\.?', '', html, flags=re.IGNORECASE)
        # Remove [number] and [number][number] patterns
        html = re.sub(r'(\[\d+(?:,\s*\d+)*\])+', '', html)
        # Optionally, remove citations inside <sup> or <span> tags (if any)
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup.find_all(['sup', 'span']):
            if tag.string and re.match(r'^\[\d+(?:,\s*\d+)*\]$', tag.string.strip()):
                tag.decompose()
        return str(soup)
   
    def post(self, request):
        try:
            # Extract parameters from request
            chats = request.data.get('chats', [])
            date_range = request.data.get('date_range')
            main_project_id = request.data.get('main_project_id')
            options = request.data.get('options', {})
 
            print()
           
            # Validate inputs
            if not chats and not (date_range and main_project_id):
                return Response({
                    'error': 'Either chat data or date_range with main_project_id is required'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Default options
            include_timestamps = options.get('includeTimestamps', True)
            include_metadata = options.get('includeChatMetadata', True)
            include_follow_ups = options.get('includeFollowUpQuestions', True)
            format_code = options.get('formatCode', True),
            include_sources = options.get('includeSources', True)
           
            # If we have date_range but no chats, fetch the chats from database
            if date_range and main_project_id and not chats:
                try:
                    start_date = datetime.fromisoformat(date_range.get('startDate').replace('Z', '+00:00'))
                    end_date = datetime.fromisoformat(date_range.get('endDate').replace('Z', '+00:00'))
                    end_date = end_date + timedelta(days=1)  # Include the entire end day
                   
                    # Fetch conversations from database
                    conversations = ChatHistory.objects.filter(
                        user=request.user,
                        main_project_id=main_project_id,
                        created_at__gte=start_date,
                        created_at__lt=end_date
                    ).order_by('-created_at')
                   
                    if not conversations:
                        return Response({
                            'error': 'No conversations found in the specified date range'
                        }, status=status.HTTP_404_NOT_FOUND)
         
                    # Convert to the same format as frontend-processed chats
                    chats = [{
                        'conversation_id': conv.conversation_id,
                        'title': conv.title,
                        'created_at': conv.created_at.strftime('%Y-%m-%d %H:%M'),
                        'messages': [{
                            'role': msg.role,
                            'content': msg.content,
                            'created_at': msg.created_at.strftime('%Y-%m-%d %H:%M'),
                            'sources_info': msg.sources or ''  # ✅ FIXED: Use 'sources' field from model
                        } for msg in conv.messages_NB.all().order_by('created_at')],
                        'follow_up_questions': conv.follow_up_questions or []
                    } for conv in conversations]
                    print('###########inside',chats)
                   
                except Exception as e:
                    return Response({
                        'error': f'Invalid date format: {str(e)}'
                    }, status=status.HTTP_400_BAD_REQUEST)
 
            print('###########outside', chats)
           
            # Create a DOCX file in memory
            doc = DocxDocument()
           
            # Set document styles
            styles = doc.styles
           
            # Heading styles
            heading1_style = styles['Heading 1']
            heading1_style.font.size = Pt(16)
            heading1_style.font.bold = True
            heading1_style.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue
           
            heading2_style = styles['Heading 2']
            heading2_style.font.size = Pt(14)
            heading2_style.font.bold = True
            heading2_style.font.color.rgb = RGBColor(0, 0, 100)  # Slightly darker blue
           
            # Normal style
            normal_style = styles['Normal']
            normal_style.font.size = Pt(11)
            normal_style.font.name = 'Calibri'
           
            # List styles
            if 'List Bullet' in styles:
                bullet_style = styles['List Bullet']
                bullet_style.paragraph_format.left_indent = Inches(0.25)
                bullet_style.paragraph_format.first_line_indent = Inches(-0.25)
               
            if 'List Number' in styles:
                number_style = styles['List Number']
                number_style.paragraph_format.left_indent = Inches(0.25)
                number_style.paragraph_format.first_line_indent = Inches(-0.25)
           
            # Code style
            if 'Code' not in styles:
                code_style = styles.add_style('Code', WD_STYLE_TYPE.PARAGRAPH)
                code_style.font.name = 'Consolas'
                code_style.font.size = Pt(10)
                code_style.paragraph_format.space_before = Pt(6)
                code_style.paragraph_format.space_after = Pt(6)
                code_style.paragraph_format.left_indent = Inches(0.25)
                code_style.font.color.rgb = RGBColor(50, 50, 50)
           
            # Set document margins
            sections = doc.sections
            for section in sections:
                section.left_margin = Inches(1.0)
                section.right_margin = Inches(1.0)
                section.top_margin = Inches(1.0)
                section.bottom_margin = Inches(1.0)
           
            # Document title
            if len(chats) == 1:
                doc.add_heading(f"Chat Conversation: {chats[0]['title']}", level=0)
            else:
                # Get the start date and add one day
                start_date_str = date_range.get('startDate').split('T')[0] if date_range else chats[-1]['created_at'].split(' ')[0]
                start_date = datetime.strptime(start_date_str, '%Y-%m-%d') + timedelta(days=1)
                start_date_str = start_date.strftime('%Y-%m-%d')
               
                # Get the end date and add one day
                end_date_str = date_range.get('endDate').split('T')[0] if date_range else chats[0]['created_at'].split(' ')[0]
                end_date = datetime.strptime(end_date_str, '%Y-%m-%d') + timedelta(days=1)
                end_date_str = end_date.strftime('%Y-%m-%d')
               
                doc.add_heading(f"Chat Conversations: {start_date_str} to {end_date_str}", level=0)
           
            # Process each conversation
            for i, chat in enumerate(chats):
                # Add conversation title
                section_heading = doc.add_heading(chat['title'], level=1)
               
                # Add metadata if requested
                if include_metadata:
                    metadata_para = doc.add_paragraph(style='Intense Quote')
                    metadata_para.add_run(f"Created: {chat['created_at']}")
               
                # Process each message
                for message in chat['messages']:
                    # Message role as heading
                    role_heading = "User Question:" if message['role'] == 'user' else "Assistant Answer:"
                    doc.add_heading(role_heading, level=2)
                   
                    # Add timestamp if requested
                    if include_timestamps:
                        timestamp = doc.add_paragraph(style='Subtitle')
                        timestamp.add_run(f"{message['created_at']}")
                   
                    content = message['content']
                    if message['role'] == 'assistant':
                        content = self.remove_citations_from_html(content)
 
                    # Add the content (already processed by frontend)
                    self.add_html_to_docx(doc, content, format_code)
 
 
                    if message['role'] == 'assistant' and include_sources:
                        # Try sources_info first
                        sources_content = message.get('sources_info', '').strip()
                       
                        # Fallback to citation source_files if sources_info is missing
                        citation_sources = []
                        if not sources_content:
                            for citation in message.get('citations', []):
                                if 'source_file' in citation:
                                    citation_sources.append(citation['source_file'])
 
                            # Remove duplicates
                            citation_sources = list(set(citation_sources))
 
                        # Decide which to use
                        final_sources = []
                        if sources_content:
                            final_sources = [s.strip() for s in sources_content.split('\n') if s.strip()]
                        elif citation_sources:
                            final_sources = citation_sources
 
                        if final_sources:
                            doc.add_heading("Sources:", level=3)
                            for src in final_sources:
                                p = doc.add_paragraph(style='List Bullet')
                                p.add_run(src).italic = True
 
 
               
               
                # Add follow-up questions if requested
                if include_follow_ups and chat.get('follow_up_questions'):
                    doc.add_heading("Suggested Follow-up Questions:", level=2)
                    for question in chat['follow_up_questions']:
                        p = doc.add_paragraph(style='List Bullet')
                        p.add_run(question)
               
                # Add page break between conversations if multiple
                if i < len(chats) - 1:
                    doc.add_page_break()
           
            # Save document to a BytesIO object
            docx_file = BytesIO()
            doc.save(docx_file)
            docx_file.seek(0)
 
            # Update project timestamp if needed
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
           
            # Create response with file
            if len(chats) == 1:
                safe_title = ''.join(c if c.isalnum() or c in ' _-' else '_' for c in chat['title'])
                safe_title = '_'.join(safe_title.split()[:3])
                filename = f"Chat_{safe_title}.docx"
            else:
                start_date = date_range.get('startDate').split('T')[0] if date_range else chats[-1]['created_at'].split(' ')[0]
                end_date = date_range.get('endDate').split('T')[0] if date_range else chats[0]['created_at'].split(' ')[0]
                filename = f"Chats_{start_date.replace('-', '')}_to_{end_date.replace('-', '')}.docx"
           
            response = HttpResponse(
                docx_file.getvalue(),
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
           
            return response
           
        except Exception as e:
            logger.error(f"Error exporting chat: {str(e)}", exc_info=True)
            return Response({
                'error': f'Failed to export chat: {str(e)}'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    def add_html_to_docx(self, doc, html_content, format_code=True):
        """Parse HTML content and add it to the document with proper formatting."""
        soup = BeautifulSoup(html_content, "html.parser")
       
        # Process each top-level element
        for elem in soup.children:
            if isinstance(elem, NavigableString) and elem.strip():
                # Plain text
                doc.add_paragraph(elem.strip())
            elif elem.name:
                self.process_element(doc, elem, format_code, list_level=0)
   
    def process_element(self, doc, element, format_code=True, list_level=0):
        """Process a single HTML element and add it to the document."""
        # Handle different element types
        if element.name == 'p':
            p = doc.add_paragraph()
            self.process_inline_elements(p, element)
       
        elif element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])
            doc.add_heading(element.get_text(), level=min(level, 9))
       
        elif element.name == 'ul':
            for li in element.find_all('li', recursive=False):
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.left_indent = Inches(0.25 + (0.25 * list_level))
                p.paragraph_format.first_line_indent = Inches(-0.25)
               
                has_nested_list = False
                nested_list = None
               
                for child in li.children:
                    if child.name in ['ul', 'ol']:
                        has_nested_list = True
                        nested_list = child
                        break
               
                if has_nested_list and nested_list:
                    nested_list.extract()
               
                self.process_inline_elements(p, li)
               
                if has_nested_list and nested_list:
                    self.process_element(doc, nested_list, format_code, list_level + 1)
       
        elif element.name == 'ol':
            for i, li in enumerate(element.find_all('li', recursive=False), 1):
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.left_indent = Inches(0.25 + (0.25 * list_level))
                p.paragraph_format.first_line_indent = Inches(-0.25)
               
                has_nested_list = False
                nested_list = None
               
                for child in li.children:
                    if child.name in ['ul', 'ol']:
                        has_nested_list = True
                        nested_list = child
                        break
               
                if has_nested_list and nested_list:
                    nested_list.extract()
               
                self.process_inline_elements(p, li)
               
                if has_nested_list and nested_list:
                    self.process_element(doc, nested_list, format_code, list_level + 1)
       
        elif element.name == 'pre' and format_code:
            code_element = element.find('code')
            code_text = code_element.get_text() if code_element else element.get_text()
            doc.add_paragraph(code_text, style='Code')
       
        elif element.name == 'code' and element.parent.name != 'pre' and format_code:
            p = doc.add_paragraph()
            code_run = p.add_run(element.get_text())
            code_run.font.name = 'Consolas'
            code_run.font.size = Pt(10)
           
        elif element.name == 'table':
            rows = element.find_all('tr')
            if rows:
                max_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
                if max_cols > 0:
                    table = doc.add_table(rows=len(rows), cols=max_cols)
                    table.style = 'Table Grid'
                   
                    for i, row in enumerate(rows):
                        cells = row.find_all(['td', 'th'])
                        for j, cell in enumerate(cells):
                            if j < max_cols:
                                cell_para = table.cell(i, j).paragraphs[0]
                                if cell.name == 'th':
                                    run = cell_para.add_run(cell.get_text())
                                    run.bold = True
                                else:
                                    # --- FIX: Handle block elements inside table cells ---
                                    has_block = False
                                    for child in cell.children:
                                        if getattr(child, "name", None) in ['ul', 'ol', 'p', 'div']:
                                            has_block = True
                                            # For block elements, process them inside the cell
                                            self.process_element(table.cell(i, j), child, format_code)
                                    if not has_block:
                                        self.process_inline_elements(cell_para, cell)
 
        elif element.name == 'blockquote':
            p = doc.add_paragraph(style='Quote')
            self.process_inline_elements(p, element)
       
        elif element.name == 'div':
            for child in element.children:
                if isinstance(child, NavigableString) and child.strip():
                    doc.add_paragraph(child.strip())
                elif child.name:
                    self.process_element(doc, child, format_code)
       
        else:
            text = element.get_text(strip=True)
            if text:
                doc.add_paragraph(text)
   
    def process_inline_elements(self, paragraph, element):
        """Process inline elements like bold, italic, etc. within a paragraph."""
        if element is None:
            return
           
        for child in element.children:
            if isinstance(child, NavigableString):
                if str(child).strip():
                    paragraph.add_run(str(child))
            elif child.name == 'b' or child.name == 'strong':
                run = paragraph.add_run(child.get_text())
                run.bold = True
            elif child.name == 'i' or child.name == 'em':
                run = paragraph.add_run(child.get_text())
                run.italic = True
            elif child.name == 'u':
                run = paragraph.add_run(child.get_text())
                run.underline = True
            elif child.name == 'a':
                run = paragraph.add_run(child.get_text())
                run.underline = True
                run.font.color.rgb = RGBColor(0, 0, 255)
            elif child.name == 'code':
                run = paragraph.add_run(child.get_text())
                run.font.name = 'Consolas'
                run.font.size = Pt(10)
            elif child.name == 'br':
                paragraph.add_run("\n")
            elif child.name not in ['ul', 'ol']:
                self.process_inline_elements(paragraph, child)
 
class YouTubeUploadView(DocumentProcessingMixin, APIView):
    parser_classes = (JSONParser,)
    def sanitize_filename(self, filename):
        """Sanitize filename to make it safe for file system"""
        import re
       
        # Remove or replace characters that are not allowed in filenames
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Remove leading/trailing dots and spaces
        filename = filename.strip('. ')
        # Ensure the filename is not empty
        if not filename:
            filename = 'untitled'
        # Limit length to avoid file system issues
        if len(filename) > 200:
            filename = filename[:200]
       
        return filename
   
    def is_valid_youtube_url(self, url):
        """Check if the provided URL is a valid YouTube URL."""
        youtube_patterns = [
            r'(?:https?://)?(?:www\.)?(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/v/)([a-zA-Z0-9_-]{11})',
            r'(?:https?://)?(?:www\.)?youtube\.com/shorts/([a-zA-Z0-9_-]{11})',
        ]
       
        for pattern in youtube_patterns:
            if re.search(pattern, url):
                return True
        return False
   
    def download_youtube_audio(self, url):
        """
        Download YouTube video as audio using yt_dlp.
        Returns the path to the downloaded audio file and video title.
        """
        try:
            if not self.is_valid_youtube_url(url):
                raise ValueError("Invalid YouTube URL provided")
           
            # Set up download directory
            DOWNLOAD_DIR = "youtube_downloads"
            os.makedirs(DOWNLOAD_DIR, exist_ok=True)
           
            # Configure yt_dlp options for audio only
            ydl_opts = {
                'format': 'bestaudio/best',
                'outtmpl': os.path.join(DOWNLOAD_DIR, '%(title)s.%(ext)s'),
                'extractaudio': True,
                'audioformat': 'mp3',
                'audioquality': '192K',
                'noplaylist': True,
                'quiet': False,
                'no_warnings': False,
            }
           
            logger.info(f"Starting download from: {url}")
           
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                # Get video info first
                info = ydl.extract_info(url, download=False)
                title = info.get('title', 'Unknown')
                duration = info.get('duration', 0)
               
                # Check duration (optional limit)
                if duration > 7200:  # 2 hour limit
                    raise ValueError(f"Video is too long ({duration//60} minutes). Maximum allowed is 120 minutes.")
               
                logger.info(f"Title: {title}, Duration: {duration//60}:{duration%60:02d}")
               
                # Download the audio
                ydl.download([url])
               
                # Find the downloaded file
                expected_filename = ydl.prepare_filename(info)
               
                # Handle different possible extensions
                possible_files = []
                base_name = os.path.splitext(expected_filename)[0]
               
                # Common audio extensions when downloading audio
                audio_extensions = ['.mp3', '.m4a', '.webm', '.opus']
                for ext in audio_extensions:
                    possible_files.append(base_name + ext)
               
                # Find which file actually exists
                downloaded_file = None
                for file_path in possible_files:
                    if os.path.exists(file_path):
                        downloaded_file = file_path
                        break
               
                # If not found in expected locations, search the download directory
                if not downloaded_file:
                    for file in os.listdir(DOWNLOAD_DIR):
                        if title.replace('/', '_').replace('\\', '_') in file:
                            downloaded_file = os.path.join(DOWNLOAD_DIR, file)
                            break
               
                if not downloaded_file:
                    raise Exception("Downloaded file not found")
               
                logger.info(f"Successfully downloaded: {downloaded_file}")
                return downloaded_file, title
       
        except Exception as e:
            logger.error(f"Error downloading YouTube video: {str(e)}")
            raise Exception(f"Failed to download YouTube video: {str(e)}")
   
    def process_document_from_text(self, text, filename, user=None):
        """
        Process extracted text directly without file operations.
        Creates FAISS index and metadata from the provided text.
        """
        import numpy as np
        import faiss
        import pickle
       
        try:
            # Clean the text
            cleaned_text = self.clean_text(text)
           
            all_chunks = []
            # Split text into chunks
            chunks = self.split_text_into_chunks(cleaned_text)
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    'text': chunk,  # Use 'text' key for consistency
                    'source_file': filename,
                    'chunk_id': i
                })
           
            # Get embeddings for all chunks
            text_chunks = [chunk['text'] for chunk in all_chunks]
            logger.info(f"Getting embeddings for {len(text_chunks)} text chunks")
            embeddings = self.get_embeddings(text_chunks)
           
            if not embeddings:
                raise ValueError("Failed to generate embeddings for document chunks")
           
            # Create FAISS index
            index = self.create_faiss_index(embeddings)
           
            # Generate a unique session ID for this document
            session_id = uuid.uuid4().hex
           
            # Save the index and chunks
            index_file, pickle_file = self.save_faiss_index(index, all_chunks, session_id)
           
            logger.info(f"Text processed successfully: {len(all_chunks)} chunks created")
           
            return {
                'index_path': index_file,
                'metadata_path': pickle_file,
                'full_text': cleaned_text
            }
           
        except Exception as e:
            logger.error(f"Error in process_document_from_text: {str(e)}")
            raise
   
    def post(self, request):
        youtube_url = request.data.get('youtube_url')
        user = request.user
        main_project_id = request.data.get('main_project_id')
        target_user_id = request.data.get('target_user_id')
       
        if not youtube_url:
            return Response({
                'error': 'YouTube URL is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        if not main_project_id:
            return Response({
                'error': 'Main project ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        if target_user_id and request.user.username == 'admin':
            # Admin uploading for another user
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
        else:
            # Regular upload
            user = request.user
 
        try:
            permissions = UserModulePermissions.objects.get(user=user)
            if permissions.disabled_modules.get('document-upload', False):
                return Response({
                    'error': 'Document uploads are disabled for this user'
                }, status=status.HTTP_403_FORBIDDEN)
        except UserModulePermissions.DoesNotExist:
            pass
 
        try:
            main_project = Project.objects.get(id=main_project_id, user=user)
           
            # Validate YouTube URL
            if not self.is_valid_youtube_url(youtube_url):
                return Response({
                    'error': 'Invalid YouTube URL provided'
                }, status=status.HTTP_400_BAD_REQUEST)
           
            # Download audio from YouTube
            logger.info(f"Processing YouTube URL: {youtube_url}")
            downloaded_file, video_title = self.download_youtube_audio(youtube_url)
           
            try:
                # Extract text from downloaded audio
                logger.info(f"Extracting text from downloaded audio: {video_title}")
                extracted_text = self.extract_text_from_youtube_audio(downloaded_file, user=user)
               
                if not extracted_text or (isinstance(extracted_text, str) and extracted_text.startswith("Error")):
                    return Response({
                        'error': f'Failed to extract text from YouTube video: {video_title}'
                    }, status=status.HTTP_400_BAD_REQUEST)
               
                # Create transcript filename based on video title
                # Clean the title for filename use
                safe_title = re.sub(r'[^\w\s-]', '', video_title)
                safe_title = re.sub(r'[-\s]+', '_', safe_title)
                transcript_filename = f"{safe_title}_youtube_transcript.txt"
               
                # Check for existing document with similar name
                existing_doc = Document.objects.filter(
                    user=user,
                    filename__icontains=safe_title,
                    main_project=main_project
                ).first()
               
                # Create a temporary file with the transcript
                transcript_file_path = os.path.join(tempfile.gettempdir(), transcript_filename)
                with open(transcript_file_path, 'w', encoding='utf-8') as f:
                    f.write(extracted_text)
               
                # Create a Django File object from the transcript
                with open(transcript_file_path, 'rb') as f:
                    from django.core.files import File
                    django_file = File(f, name=transcript_filename)
                   
                    # Create or update document with transcript
                    if existing_doc:
                        # Update existing document
                        existing_doc.file = django_file
                        existing_doc.filename = transcript_filename
                        existing_doc.save()
                        document = existing_doc
                        logger.info(f"Updated existing document: {document.id}")
                    else:
                        # Create new document with transcript
                        document = Document.objects.create(
                            user=user,
                            file=django_file,
                            filename=transcript_filename,
                            main_project=main_project
                        )
                        logger.info(f"Created new document: {document.id}")
               
                # Delete the temporary transcript file
                os.unlink(transcript_file_path)
               
                # Process the document for RAG
                processed_data = self.process_document_from_text(extracted_text, transcript_filename, user)
               
                # Create or update ProcessedIndex
                processed_index, created = ProcessedIndex.objects.get_or_create(
                    document=document,
                    defaults={
                        'faiss_index': processed_data['index_path'],
                        'metadata': processed_data['metadata_path'],
                        'summary': "",
                        'markdown_path': processed_data.get('markdown_path', '')
                    }
                )
               
                if not created:
                    # Update existing ProcessedIndex
                    processed_index.faiss_index = processed_data['index_path']
                    processed_index.metadata = processed_data['metadata_path']
                    processed_index.markdown_path = processed_data.get('markdown_path', '')
                    processed_index.save()
               
                # Store the document ID in the session
                request.session['active_document_id'] = document.id
               
                # Update project timestamp to track activity
                update_project_timestamp(main_project_id, user)
               
                return Response({
                    'message': 'YouTube video processed successfully',
                    'document': {
                        'id': document.id,
                        'filename': transcript_filename,
                        'original_video_title': video_title,
                        'youtube_url': youtube_url
                    },
                    'active_document_id': document.id
                }, status=status.HTTP_201_CREATED)
           
            finally:
                # Clean up downloaded audio file
                try:
                    if os.path.exists(downloaded_file):
                        os.unlink(downloaded_file)
                        logger.info(f"Cleaned up downloaded file: {downloaded_file}")
                except Exception as cleanup_error:
                    logger.warning(f"Cleanup error: {str(cleanup_error)}")
 
        except Exception as e:
            logger.error(f"Error processing YouTube video: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while processing the YouTube video'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            
    # Add this method to your YouTubeUploadView class

    # def extract_text_from_youtube_audio(self, audio_file_path, user=None):
    #     """
    #     Extract text from audio file specifically for YouTube uploads using Google Gemini API.
    #     Includes proper file state checking and fallback mechanisms.
    #     """
    #     import google.generativeai as genai
    #     import time
    #     import os
    #     from pathlib import Path
        
    #     try:
    #         # Get the user's Gemini API token
    #         gemini_api_key = None
    #         if user:
    #             try:
    #                 user_api_tokens = UserAPITokens.objects.get(user=user)
    #                 gemini_api_key = user_api_tokens.gemini_token
                    
    #                 if not gemini_api_key:
    #                     logger.warning(f"No Gemini API token found for user {user.username}")
    #                     # Fallback to environment variable
    #                     gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
    #                     if not gemini_api_key:
    #                         raise ValueError("Gemini API token is required for processing")
                        
    #             except UserAPITokens.DoesNotExist:
    #                 logger.error(f"No API tokens record found for user {user.username}")
    #                 # Fallback to environment variable
    #                 gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
    #                 if not gemini_api_key:
    #                     raise ValueError("Gemini API token is required for processing")
    #         else:
    #             # Fallback to environment variable if no user provided
    #             gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
    #             if not gemini_api_key:
    #                 raise ValueError("Gemini API token is required for processing")
            
    #         # Configure the API
    #         genai.configure(api_key=gemini_api_key)
            
    #         # Get file info
    #         file_size = os.path.getsize(audio_file_path)
    #         duration_minutes = file_size / (1024 * 1024 * 0.1)  # Rough estimate
            
    #         logger.info(f"Audio file: {audio_file_path}, Size: {file_size/(1024*1024):.2f}MB, Duration: {duration_minutes:.2f} minutes")
            
    #         # Check file format and size limitations
    #         file_extension = Path(audio_file_path).suffix.lower()
    #         if file_size > 500 * 1024 * 1024:  # 20MB limit
    #             raise Exception(f"File too large: {file_size/(1024*1024):.2f}MB. Maximum size is 500MB.")
            
    #         # Upload the file with metadata
    #         logger.info("Uploading audio file to Google API...")
    #         try:
    #             audio_file = genai.upload_file(
    #                 path=audio_file_path,
    #                 display_name=os.path.basename(audio_file_path),
    #                 mime_type=f"audio/{file_extension.lstrip('.')}" if file_extension in ['.mp3', '.wav', '.m4a'] else 'audio/webm'
    #             )
    #         except Exception as upload_error:
    #             logger.error(f"Upload failed: {upload_error}")
    #             # Try different approach or return error
    #             raise Exception(f"Failed to upload file to Google API: {upload_error}")
            
    #         logger.info(f"File uploaded successfully: {audio_file.name}")
            
    #         # Wait for the file to be processed with better error handling
    #         logger.info(f"Waiting for file {audio_file.name} to be processed...")
    #         max_wait_time = 180  # 3 minutes max wait (reduced from 5)
    #         wait_interval = 3    # Check every 3 seconds
    #         waited_time = 0
            
    #         while waited_time < max_wait_time:
    #             try:
    #                 current_file = genai.get_file(audio_file.name)
    #                 state_name = getattr(current_file.state, 'name', str(current_file.state))
    #                 logger.info(f"File state: {state_name}, waited: {waited_time}s")
                    
    #                 if state_name == "ACTIVE":
    #                     logger.info(f"File {audio_file.name} is now ACTIVE and ready for processing")
    #                     break
    #                 elif state_name == "FAILED" or str(current_file.state) == "10":
    #                     # Clean up failed file
    #                     try:
    #                         genai.delete_file(audio_file.name)
    #                     except:
    #                         pass
    #                     raise Exception(f"Google API failed to process the audio file. This might be due to: unsupported format, corrupted audio, or content policy restrictions.")
    #                 elif state_name == "PROCESSING":
    #                     time.sleep(wait_interval)
    #                     waited_time += wait_interval
    #                 else:
    #                     logger.warning(f"Unknown file state: {state_name}")
    #                     time.sleep(wait_interval)
    #                     waited_time += wait_interval
                        
    #             except Exception as status_error:
    #                 logger.error(f"Error checking file status: {status_error}")
    #                 time.sleep(wait_interval)
    #                 waited_time += wait_interval
            
    #         # Final check
    #         try:
    #             final_file = genai.get_file(audio_file.name)
    #             final_state = getattr(final_file.state, 'name', str(final_file.state))
    #             if final_state != "ACTIVE":
    #                 raise Exception(f"File processing timeout or failed. Final state: {final_state}")
    #         except Exception as final_check_error:
    #             raise Exception(f"File not ready for transcription: {final_check_error}")
            
    #         # Configure the model
    #         try:
    #             model = genai.GenerativeModel("gemini-2.0-flash")
    #         except:
    #             model = genai.GenerativeModel("gemini-2.5-pro")  # Fallback
            
    #         # Create the prompt for transcription
    #         prompt = """
    #         Please transcribe this audio file accurately. 
    #         Return only the spoken text content without any additional commentary.
    #         Provide proper punctuation and formatting.
    #         """
            
    #         # Generate the transcription with timeout
    #         logger.info("Starting transcription...")
    #         try:
    #             response = model.generate_content([prompt, audio_file])
    #             logger.info("Transcription completed")
    #         except Exception as transcription_error:
    #             logger.error(f"Transcription failed: {transcription_error}")
    #             raise Exception(f"Failed to transcribe audio: {transcription_error}")
            
    #         # Clean up the uploaded file from Google's servers
    #         try:
    #             genai.delete_file(audio_file.name)
    #             logger.info(f"Cleaned up file {audio_file.name} from Google's servers")
    #         except Exception as cleanup_error:
    #             logger.warning(f"Could not clean up file from Google's servers: {cleanup_error}")
            
    #         if response and response.text:
    #             transcribed_text = response.text.strip()
    #             if len(transcribed_text) < 10:  # Very short transcription might indicate an issue
    #                 logger.warning(f"Transcription seems unusually short: '{transcribed_text}'")
    #             return transcribed_text
    #         else:
    #             raise Exception("No transcription text was generated")
                
    #     except Exception as e:
    #         logger.error(f"Error in extract_text_from_youtube_audio: {str(e)}")
    #         return f"Error transcribing audio: {str(e)}"
    
    # New Code

    def extract_text_from_youtube_audio(self, audio_file_path, user=None):
        """
        Extract text from audio file specifically for YouTube uploads using Google Gemini API.
        Updated to use gemini-2.5-flash with new syntax while maintaining file upload workflow.
        """
        from google import genai
        from google.genai import types
        import time
        import os
        from pathlib import Path
       
        try:
            # Get the user's Gemini API token
            gemini_api_key = None
            if user:
                try:
                    user_api_tokens = UserAPITokens.objects.get(user=user)
                    gemini_api_key = user_api_tokens.gemini_token
                   
                    if not gemini_api_key:
                        logger.warning(f"No Gemini API token found for user {user.username}")
                        # Fallback to environment variable
                        gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                        if not gemini_api_key:
                            raise ValueError("Gemini API token is required for processing")
                       
                except UserAPITokens.DoesNotExist:
                    logger.error(f"No API tokens record found for user {user.username}")
                    # Fallback to environment variable
                    gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                    if not gemini_api_key:
                        raise ValueError("Gemini API token is required for processing")
            else:
                # Fallback to environment variable if no user provided
                gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                if not gemini_api_key:
                    raise ValueError("Gemini API token is required for processing")
           
            # Configure the Gemini client
            gemini_client = genai.Client(api_key=gemini_api_key)
           
            # Get file info
            file_size = os.path.getsize(audio_file_path)
            duration_minutes = file_size / (1024 * 1024 * 0.1)  # Rough estimate
           
            logger.info(f"Audio file: {audio_file_path}, Size: {file_size/(1024*1024):.2f}MB, Duration: {duration_minutes:.2f} minutes")
           
            # Check file format and size limitations
            file_extension = Path(audio_file_path).suffix.lower()
            if file_size > 500 * 1024 * 1024:  # 500MB limit
                raise Exception(f"File too large: {file_size/(1024*1024):.2f}MB. Maximum size is 500MB.")
           
            # Upload the file with metadata using new client
            logger.info("Uploading audio file to Google API...")
            try:
                audio_file = gemini_client.files.upload(
                    path=audio_file_path,
                    config=types.FileUploadConfig(
                        display_name=os.path.basename(audio_file_path),
                        mime_type=f"audio/{file_extension.lstrip('.')}" if file_extension in ['.mp3', '.wav', '.m4a'] else 'audio/webm'
                    )
                )
            except Exception as upload_error:
                logger.error(f"Upload failed: {upload_error}")
                # Try different approach or return error
                raise Exception(f"Failed to upload file to Google API: {upload_error}")
           
            logger.info(f"File uploaded successfully: {audio_file.name}")
           
            # Wait for the file to be processed with better error handling
            logger.info(f"Waiting for file {audio_file.name} to be processed...")
            max_wait_time = 180  # 3 minutes max wait (reduced from 5)
            wait_interval = 3    # Check every 3 seconds
            waited_time = 0
           
            while waited_time < max_wait_time:
                try:
                    current_file = gemini_client.files.get(audio_file.name)
                    state_name = getattr(current_file.state, 'name', str(current_file.state))
                    logger.info(f"File state: {state_name}, waited: {waited_time}s")
                   
                    if state_name == "ACTIVE":
                        logger.info(f"File {audio_file.name} is now ACTIVE and ready for processing")
                        break
                    elif state_name == "FAILED" or str(current_file.state) == "10":
                        # Clean up failed file
                        try:
                            gemini_client.files.delete(audio_file.name)
                        except:
                            pass
                        raise Exception(f"Google API failed to process the audio file. This might be due to: unsupported format, corrupted audio, or content policy restrictions.")
                    elif state_name == "PROCESSING":
                        time.sleep(wait_interval)
                        waited_time += wait_interval
                    else:
                        logger.warning(f"Unknown file state: {state_name}")
                        time.sleep(wait_interval)
                        waited_time += wait_interval
                       
                except Exception as status_error:
                    logger.error(f"Error checking file status: {status_error}")
                    time.sleep(wait_interval)
                    waited_time += wait_interval
           
            # Final check
            try:
                final_file = gemini_client.files.get(audio_file.name)
                final_state = getattr(final_file.state, 'name', str(final_file.state))
                if final_state != "ACTIVE":
                    raise Exception(f"File processing timeout or failed. Final state: {final_state}")
            except Exception as final_check_error:
                raise Exception(f"File not ready for transcription: {final_check_error}")
           
            # Create the prompt for transcription
            prompt = """
            Please transcribe this audio file accurately.
            Return only the spoken text content without any additional commentary.
            Provide proper punctuation and formatting.
            """
           
            # Generate the transcription using new API syntax with uploaded file
            logger.info("Starting transcription...")
            try:
                response = gemini_client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[prompt, audio_file],
                    config=types.GenerateContentConfig(
                        temperature=0.1
                    )
                )
                logger.info("Transcription completed")
            except Exception as transcription_error:
                logger.error(f"Transcription failed: {transcription_error}")
                raise Exception(f"Failed to transcribe audio: {transcription_error}")
           
            # Clean up the uploaded file from Google's servers
            try:
                gemini_client.files.delete(audio_file.name)
                logger.info(f"Cleaned up file {audio_file.name} from Google's servers")
            except Exception as cleanup_error:
                logger.warning(f"Could not clean up file from Google's servers: {cleanup_error}")
           
            if response and response.text:
                transcribed_text = response.text.strip()
                if len(transcribed_text) < 10:  # Very short transcription might indicate an issue
                    logger.warning(f"Transcription seems unusually short: '{transcribed_text}'")
                return transcribed_text
            else:
                raise Exception("No transcription text was generated")
               
        except Exception as e:
            logger.error(f"Error in extract_text_from_youtube_audio: {str(e)}")
            return f"Error transcribing audio: {str(e)}"
    
    
    def split_text_into_chunks(self, text, chunk_size=1000, overlap=200):
        """
        Split text into overlapping chunks for processing.
        """
        if not text or len(text.strip()) == 0:
            return []
        
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk_words = words[i:i + chunk_size]
            chunk_text = ' '.join(chunk_words)
            if chunk_text.strip():
                chunks.append(chunk_text)
        
        return chunks

    def clean_text(self, text):
        """
        Clean and normalize text content.
        """
        import re
        
        if not text:
            return ""
        
        # Remove extra whitespace and normalize line breaks
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n+', '\n', text)
        
        # Remove special characters that might cause issues
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]\"\'\n]', '', text)
        
        return text.strip()

    def get_embeddings(self, texts):
        """
        Generate embeddings for text chunks using OpenAI or similar service.
        """
        try:
            # You'll need to implement this based on your embedding service
            # This is a placeholder - replace with your actual embedding logic
            
            import openai
            from django.conf import settings
            
            if hasattr(settings, 'OPENAI_API_KEY'):
                openai.api_key = settings.OPENAI_API_KEY
                
                embeddings = []
                for text in texts:
                    try:
                        response = openai.embeddings.create(
                            model="text-embedding-ada-002",
                            input=text
                        )
                        embeddings.append(response.data[0].embedding)
                    except Exception as e:
                        logger.error(f"Error getting embedding: {e}")
                        # Create a dummy embedding if API fails
                        embeddings.append([0.0] * 1536)  # Ada-002 embedding size
                
                return embeddings
            else:
                # Fallback: create dummy embeddings
                logger.warning("No OpenAI API key found, creating dummy embeddings")
                return [[0.0] * 1536 for _ in texts]  # Dummy embeddings
                
        except Exception as e:
            logger.error(f"Error in get_embeddings: {e}")
            # Return dummy embeddings as fallback
            return [[0.0] * 1536 for _ in texts]

    def create_faiss_index(self, embeddings):
        """
        Create FAISS index from embeddings.
        """
        import numpy as np
        import faiss
        
        if not embeddings:
            raise ValueError("No embeddings provided")
        
        # Convert to numpy array
        embeddings_array = np.array(embeddings, dtype=np.float32)
        
        # Create FAISS index
        dimension = embeddings_array.shape[1]
        index = faiss.IndexFlatL2(dimension)
        index.add(embeddings_array)
        
        return index

    def save_faiss_index(self, index, chunks, session_id):
        """
        Save FAISS index and metadata to files.
        """
        import pickle
        import os
        from django.conf import settings
        
        try:
            # Create directory for indexes
            index_dir = os.path.join(settings.MEDIA_ROOT, 'faiss_indexes')
            os.makedirs(index_dir, exist_ok=True)
            
            # File paths
            index_file = os.path.join(index_dir, f'index_{session_id}.faiss')
            pickle_file = os.path.join(index_dir, f'metadata_{session_id}.pkl')
            
            # Save FAISS index
            import faiss
            faiss.write_index(index, index_file)
            
            # Save metadata
            with open(pickle_file, 'wb') as f:
                pickle.dump(chunks, f)
            
            logger.info(f"Saved FAISS index to {index_file}")
            logger.info(f"Saved metadata to {pickle_file}")
            
            return index_file, pickle_file
            
        except Exception as e:
            logger.error(f"Error in save_faiss_index: {str(e)}")
            raise Exception(f"Failed to save FAISS index: {str(e)}")
  

class NoteManagementView(DocumentUploadView, APIView):
    parser_classes = (JSONParser,)
   
    def post(self, request):
        """Handle note saving and conversion to document"""
        action = request.data.get('action')  # 'save' or 'convert'
       
        if action == 'save':
            return self.save_note(request)
        elif action == 'convert':
            return self.convert_note_to_document(request)
        else:
            return Response({
                'error': 'Invalid action. Use "save" or "convert"'
            }, status=status.HTTP_400_BAD_REQUEST)
   
    def save_note(self, request):
        """Save or update a note"""
        user = request.user
        note_id = request.data.get('note_id')  # For updating existing note
        title = request.data.get('title', '').strip()
        content = request.data.get('content', '').strip()
        main_project_id = request.data.get('main_project_id')
        target_user_id = request.data.get('target_user_id')
       
        # Validate required fields
        if not content:
            return Response({
                'error': 'Note content is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        if not main_project_id:
            return Response({
                'error': 'Main project ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        # Handle admin uploading for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
       
        # Check user permissions
        try:
            permissions = UserModulePermissions.objects.get(user=user)
            if permissions.disabled_modules.get('note-management', False):
                return Response({
                    'error': 'Note management is disabled for this user'
                }, status=status.HTTP_403_FORBIDDEN)
        except UserModulePermissions.DoesNotExist:
            pass
       
        try:
            # Verify project exists and belongs to user
            main_project = Project.objects.get(id=main_project_id, user=user)
           
            # Generate title if not provided
            if not title:
                title = self.generate_note_title(content)
           
            # Save or update note
            if note_id:
                # Update existing note
                try:
                    note = Note.objects.get(id=note_id, user=user, main_project=main_project)
                    note.title = title
                    note.content = content
                    note.save()
                    logger.info(f"Updated note: {note.id}")
                except Note.DoesNotExist:
                    return Response({
                        'error': 'Note not found'
                    }, status=status.HTTP_404_NOT_FOUND)
            else:
                # Create new note
                note = Note.objects.create(
                    user=user,
                    title=title,
                    content=content,
                    main_project=main_project
                )
                logger.info(f"Created new note: {note.id}")
           
            # Update project timestamp
            update_project_timestamp(main_project_id, user)
           
            return Response({
                'message': 'Note saved successfully',
                'note': {
                    'id': note.id,
                    'title': note.title,
                    'content': note.content,
                    'created_at': note.created_at,
                    'updated_at': note.updated_at,
                    'is_converted_to_document': note.is_converted_to_document
                }
            }, status=status.HTTP_201_CREATED if not note_id else status.HTTP_200_OK)
           
        except Project.DoesNotExist:
            return Response({
                'error': 'Project not found'
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"Error saving note: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while saving the note'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    def convert_note_to_document(self, request):
        """Convert a note to a document and process it for RAG"""
        user = request.user
        note_id = request.data.get('note_id')
        target_user_id = request.data.get('target_user_id')
 
        if not note_id:
            return Response({
                'error': 'Note ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
 
        # Handle admin converting for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
 
        # Check user permissions
        try:
            permissions = UserModulePermissions.objects.get(user=user)
            if permissions.disabled_modules.get('document-upload', False):
                return Response({
                    'error': 'Document uploads are disabled for this user'
                }, status=status.HTTP_403_FORBIDDEN)
        except UserModulePermissions.DoesNotExist:
            pass
 
        try:
            # Get the note
            note = Note.objects.get(id=note_id, user=user)
       
            # Check if already converted
            if note.is_converted_to_document and note.converted_document:
                return Response({
                    'error': 'This note has already been converted to a document',
                    'document': {
                        'id': note.converted_document.id,
                        'filename': note.converted_document.filename
                    }
                }, status=status.HTTP_400_BAD_REQUEST)
       
            # Generate filename from note title (internal use with .txt extension)
            safe_title = self.sanitize_filename(note.title or f"Note_{note.id}")
            filename = f"{safe_title}.txt"
            # Display filename without extension for user
            display_filename = safe_title
       
            # Clean and convert the content
            cleaned_content = self.clean_html_content(note.content)
           
            # Create temporary file with cleaned note content
            temp_file_path = os.path.join(tempfile.gettempdir(), filename)
       
            try:
                with open(temp_file_path, 'w', encoding='utf-8') as f:
                    # Add title as header if available
                    if note.title:
                        f.write(f"Title: {note.title}\n")
                        f.write("=" * (len(note.title) + 7) + "\n\n")
                    f.write(cleaned_content)
           
                # Create Django File object
                with open(temp_file_path, 'rb') as f:
                    django_file = File(f, name=filename)
               
                    # Use transaction to ensure data consistency
                    with transaction.atomic():
                        # Create document
                        document = Document.objects.create(
                            user=user,
                            file=django_file,
                            filename=filename,
                            main_project=note.main_project
                        )
                   
                        logger.info(f"Created document from note: {document.id}")
                   
                        # Process the document for RAG using the inherited method
                        try:
                            # Prepare the content with title
                            full_content = ""
                            if note.title:
                                full_content += f"Title: {note.title}\n"
                                full_content += "=" * (len(note.title) + 7) + "\n\n"
                            full_content += cleaned_content
                           
                            # Use the process_document_from_text method from YouTubeUploadView
                            processed_data = self.process_document_from_text_pgvector(full_content, filename, document)
                           
                            logger.info(f"Note processed successfully for RAG")
                           
                        except Exception as processing_error:
                            logger.error(f"Error processing note for RAG: {str(processing_error)}")
                            # If processing fails, still create the document but without processing
                            processed_data = {
                                'index_path': '',
                                'metadata_path': '',
                                'full_text': full_content
                            }
                   
                        # Create ProcessedIndex with pgvector metadata
                        processed_index = ProcessedIndex.objects.create(
                            document=document,
                            storage_type='pgvector',
                            chunks_count=processed_data.get('chunks_count', 0),
                            summary="",
                            markdown_path=processed_data.get('markdown_path', '')
                        )
                        logger.info(f"Created ProcessedIndex for document: {document.id} with pgvector storage")
                   
                        # Update note to mark as converted
                        note.is_converted_to_document = True
                        note.converted_document = document
                        note.save()
                   
                        # Store the document ID in the session
                        request.session['active_document_id'] = document.id
                   
                        # Update project timestamp
                        update_project_timestamp(note.main_project.id, user)
           
                return Response({
                    'message': 'Note converted to document successfully',
                    'note': {
                        'id': note.id,
                        'title': note.title,
                        'is_converted_to_document': True
                    },
                    'document': {
                        'id': document.id,
                        'filename': display_filename,  # Show without .txt extension
                        'uploaded_at': document.uploaded_at,
                        'file_type': 'text'  # Added file type for frontend
                    },
                    'active_document_id': document.id
                }, status=status.HTTP_201_CREATED)
           
            finally:
                # Clean up temporary file
                try:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                        logger.info(f"Cleaned up temporary file: {temp_file_path}")
                except Exception as cleanup_error:
                    logger.warning(f"Cleanup error: {str(cleanup_error)}")
       
        except Note.DoesNotExist:
            return Response({
                'error': 'Note not found'
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"Error converting note to document: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while converting the note to document'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
       
   
    def get(self, request):
        """Get notes for a user and project"""
        user = request.user
        main_project_id = request.GET.get('main_project_id')
        target_user_id = request.GET.get('target_user_id')
       
        # Handle admin viewing notes for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
       
        try:
            # Filter notes
            notes_query = Note.objects.filter(user=user)
           
            if main_project_id:
                notes_query = notes_query.filter(main_project_id=main_project_id)
           
            notes = notes_query.order_by('-updated_at')
           
            notes_data = []
            for note in notes:
                note_data = {
                    'id': note.id,
                    'title': note.title,
                    'content': note.content,
                    'created_at': note.created_at,
                    'updated_at': note.updated_at,
                    'is_converted_to_document': note.is_converted_to_document,
                    'main_project_id': note.main_project.id if note.main_project else None
                }
               
                if note.converted_document:
                    note_data['converted_document'] = {
                        'id': note.converted_document.id,
                        'filename': os.path.splitext(note.converted_document.filename)[0]  # Remove .txt extension
                    }
               
                notes_data.append(note_data)
           
            return Response({
                'notes': notes_data,
                'count': len(notes_data)
            }, status=status.HTTP_200_OK)
           
        except Exception as e:
            logger.error(f"Error retrieving notes: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while retrieving notes'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    def delete(self, request):
        """Delete a note"""
        user = request.user
        note_id = request.data.get('note_id')
        target_user_id = request.data.get('target_user_id')
       
        if not note_id:
            return Response({
                'error': 'Note ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        # Handle admin deleting for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
       
        try:
            note = Note.objects.get(id=note_id, user=user)
            note_title = note.title or f"Note {note.id}"
            note.delete()
           
            logger.info(f"Deleted note: {note_id}")
           
            return Response({
                'message': f'Note "{note_title}" deleted successfully'
            }, status=status.HTTP_200_OK)
           
        except Note.DoesNotExist:
            return Response({
                'error': 'Note not found'
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"Error deleting note: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while deleting the note'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    # Helper methods
    def clean_html_content(self, html_content):
        """Clean HTML content and convert to readable text format"""
        if not html_content:
            return ""
       
        try:
            # Parse HTML with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
           
            # Handle different HTML elements
           
            # Convert headings
            for i in range(1, 7):
                for heading in soup.find_all(f'h{i}'):
                    heading_text = heading.get_text().strip()
                    if heading_text:
                        heading.replace_with(f"\n{heading_text.upper()}\n{'=' * len(heading_text)}\n\n")
           
            # Convert paragraphs
            for p in soup.find_all('p'):
                p_text = p.get_text().strip()
                if p_text:
                    p.replace_with(f"{p_text}\n\n")
           
            # Convert line breaks
            for br in soup.find_all('br'):
                br.replace_with('\n')
           
            # Convert bold and strong
            for tag in soup.find_all(['b', 'strong']):
                tag_text = tag.get_text().strip()
                if tag_text:
                    tag.replace_with(tag_text)
           
            # Convert italic and emphasis
            for tag in soup.find_all(['i', 'em']):
                tag_text = tag.get_text().strip()
                if tag_text:
                    tag.replace_with(tag_text)
           
            # Convert code blocks
            for code in soup.find_all('code'):
                code_text = code.get_text().strip()
                if code_text:
                    code.replace_with(f"`{code_text}`")
           
            # Convert unordered lists
            for ul in soup.find_all('ul'):
                list_items = []
                for li in ul.find_all('li'):
                    li_text = li.get_text().strip()
                    if li_text:
                        list_items.append(f"• {li_text}")
                if list_items:
                    ul.replace_with('\n' + '\n'.join(list_items) + '\n\n')
           
            # Convert ordered lists
            for ol in soup.find_all('ol'):
                list_items = []
                for idx, li in enumerate(ol.find_all('li'), 1):
                    li_text = li.get_text().strip()
                    if li_text:
                        list_items.append(f"{idx}. {li_text}")
                if list_items:
                    ol.replace_with('\n' + '\n'.join(list_items) + '\n\n')
           
            # Convert links
            for a in soup.find_all('a'):
                href = a.get('href', '')
                text = a.get_text().strip()
                if href and text:
                    a.replace_with(f"{text} ({href})")
                elif text:
                    a.replace_with(text)
           
            # Convert pre blocks
            for pre in soup.find_all('pre'):
                pre_text = pre.get_text().strip()
                if pre_text:
                    pre.replace_with(f"\n{pre_text}\n\n")
           
            # Convert blockquotes
            for blockquote in soup.find_all('blockquote'):
                quote_text = blockquote.get_text().strip()
                if quote_text:
                    # Add > to each line
                    quoted_lines = [f"    {line}" for line in quote_text.split('\n') if line.strip()]
                    blockquote.replace_with('\n' + '\n'.join(quoted_lines) + '\n\n')
           
            # Convert tables (basic conversion)
            for table in soup.find_all('table'):
                table_text = []
                for row in table.find_all('tr'):
                    cells = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                    if cells:
                        table_text.append(' | '.join(cells))
                if table_text:
                    table.replace_with('\n' + '\n'.join(table_text) + '\n\n')
           
            # Get the final cleaned text
            cleaned_text = soup.get_text()
           
            # Clean up extra whitespace
            cleaned_text = re.sub(r'\n\s*\n\s*\n', '\n\n', cleaned_text)  # Multiple newlines to double
            cleaned_text = re.sub(r'[ \t]+', ' ', cleaned_text)  # Multiple spaces to single
            cleaned_text = cleaned_text.strip()
           
            return cleaned_text
           
        except Exception as e:
            logger.warning(f"Error cleaning HTML content: {str(e)}")
            # Fallback: just strip HTML tags
            return re.sub(r'<[^>]+>', '', html_content).strip()
   
    def generate_note_title(self, content):
        """Generate a title from note content"""
        # First clean the content if it's HTML
        cleaned_content = self.clean_html_content(content)
       
        # Take first 50 characters and clean them up
        title = cleaned_content[:50].strip()
        # Remove newlines and extra spaces
        title = ' '.join(title.split())
        # Add ellipsis if truncated
        if len(cleaned_content) > 50:
            title += "..."
        return title or "Untitled Note"
   
    def sanitize_filename(self, filename):
        """Sanitize filename to make it safe for file system"""
        import re
       
        # Remove or replace characters that are not allowed in filenames
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Remove leading/trailing dots and spaces
        filename = filename.strip('. ')
        # Ensure the filename is not empty
        if not filename:
            filename = 'untitled'
        # Limit length to avoid file system issues
        if len(filename) > 200:
            filename = filename[:200]
       
        return filename

from django.http import HttpResponse, JsonResponse
from django.shortcuts import get_object_or_404
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated
from .models import Document
import mimetypes
import os
import logging

logger = logging.getLogger(__name__)


class DocumentContentView(APIView):
    """
    API endpoints for retrieving document content in different formats
    FIXED: Now properly handles Azure Blob Storage
    """
    permission_classes = [IsAuthenticated]
   
    def get_document_metadata(self, request, document_id):
        """Get document metadata including file type"""
        try:
            document = get_object_or_404(Document, id=document_id, user=request.user)
           
            # Determine file type
            filename = document.filename or ''
            file_extension = os.path.splitext(filename)[1].lower()
           
            if file_extension == '.pdf':
                file_type = 'pdf'
            elif file_extension in ['.txt', '.text']:
                file_type = 'text'
            elif file_extension in ['.doc', '.docx']:
                file_type = 'document'
            elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                file_type = 'image'
            else:
                # Default based on content or filename
                file_type = 'text' if filename.endswith('.txt') else 'unknown'
           
            # FIXED: Get file size from blob storage if available
            file_size = 0
            if document.file:
                try:
                    if hasattr(document.file, 'size'):
                        file_size = document.file.size
                    elif hasattr(document.file, 'name'):
                        # For blob storage, we might need to get size differently
                        from notebook.utils import get_blob_content
                        blob_content = get_blob_content(document.file.name)
                        if blob_content:
                            file_size = len(blob_content)
                except Exception as e:
                    logger.warning(f"Could not get file size for document {document_id}: {str(e)}")
           
            return JsonResponse({
                'id': document.id,
                'filename': document.filename,
                'file_type': file_type,
                'file_size': file_size,
                'uploaded_at': document.uploaded_at.isoformat() if document.uploaded_at else None,
                'main_project_id': document.main_project.id if document.main_project else None,
            })
           
        except Exception as e:
            logger.error(f"Error getting document metadata for {document_id}: {str(e)}")
            return JsonResponse({
                'error': str(e)
            }, status=500)
   
    def get_document_content(self, request, document_id):
        """
        Get document content as text (for text files converted from notes)
        FIXED: Now properly handles Azure Blob Storage
        """
        try:
            document = get_object_or_404(Document, id=document_id, user=request.user)
           
            if not document.file:
                return HttpResponse("Document file not found", status=404)
           
            # FIXED: Handle both local files and blob storage
            content = None
            
            # Method 1: Try to get content from blob storage
            if hasattr(document.file, 'name') and hasattr(document.file, 'storage'):
                # This is a Django FieldFile with Azure Blob Storage
                blob_path = document.file.name
                logger.info(f"Getting text content from blob: {blob_path}")
                
                from notebook.utils import get_blob_content
                blob_content = get_blob_content(blob_path)
                
                if blob_content:
                    # Try different encodings to decode the content
                    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                    for encoding in encodings:
                        try:
                            content = blob_content.decode(encoding)
                            logger.info(f"Successfully decoded content with {encoding}")
                            break
                        except UnicodeDecodeError:
                            continue
                    
                    if not content:
                        logger.error("Unable to decode blob content with any encoding")
                        return HttpResponse("Unable to decode file content", status=400)
            
            # Method 2: Try local file access (fallback)
            elif hasattr(document.file, 'path') and os.path.exists(document.file.path):
                logger.info(f"Getting text content from local file: {document.file.path}")
                
                encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                for encoding in encodings:
                    try:
                        with open(document.file.path, 'r', encoding=encoding) as f:
                            content = f.read()
                        logger.info(f"Successfully read local file with {encoding}")
                        break
                    except (UnicodeDecodeError, IOError):
                        continue
            
            # Method 3: Try Django file interface (another fallback)
            else:
                logger.info("Trying Django file interface")
                try:
                    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                    for encoding in encodings:
                        try:
                            with document.file.open('r', encoding=encoding) as f:
                                content = f.read()
                            break
                        except UnicodeDecodeError:
                            continue
                except Exception as e:
                    logger.error(f"Django file interface failed: {str(e)}")
            
            if content is not None:
                return HttpResponse(content, content_type='text/plain; charset=utf-8')
            else:
                return HttpResponse("Unable to read file content", status=400)
               
        except Exception as e:
            logger.error(f"Error reading document {document_id}: {str(e)}")
            return HttpResponse(f"Error reading document: {str(e)}", status=500)
   
    def get_original_document(self, request, document_id):
        """
        Get original document file (for PDFs and other binary files)
        FIXED: Now properly handles Azure Blob Storage
        """
        try:
            document = get_object_or_404(Document, id=document_id, user=request.user)
           
            if not document.file:
                return HttpResponse("Document file not found", status=404)
           
            # FIXED: Handle both local files and blob storage
            file_content = None
            
            # Method 1: Try to get content from blob storage
            if hasattr(document.file, 'name') and hasattr(document.file, 'storage'):
                # This is a Django FieldFile with Azure Blob Storage
                blob_path = document.file.name
                logger.info(f"Getting original document from blob: {blob_path}")
                
                from notebook.utils import get_blob_content
                file_content = get_blob_content(blob_path)
                
                if not file_content:
                    logger.error(f"Failed to get blob content for: {blob_path}")
                    return HttpResponse("Document content not found in storage", status=404)
            
            # Method 2: Try local file access (fallback)
            elif hasattr(document.file, 'path') and os.path.exists(document.file.path):
                logger.info(f"Getting original document from local file: {document.file.path}")
                with open(document.file.path, 'rb') as f:
                    file_content = f.read()
            
            # Method 3: Try Django file interface (another fallback)
            else:
                logger.info("Trying Django file interface for original document")
                try:
                    file_content = document.file.read()
                except Exception as e:
                    logger.error(f"Django file interface failed: {str(e)}")
                    return HttpResponse("Unable to access file", status=500)
           
            if file_content is None:
                return HttpResponse("Unable to read file content", status=500)
           
            # Determine content type
            content_type, _ = mimetypes.guess_type(document.filename)
            if not content_type:
                # Fallback content types based on file extension
                file_extension = os.path.splitext(document.filename)[1].lower()
                content_type_map = {
                    '.pdf': 'application/pdf',
                    '.doc': 'application/msword',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.xls': 'application/vnd.ms-excel',
                    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    '.ppt': 'application/vnd.ms-powerpoint',
                    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    '.txt': 'text/plain',
                    '.csv': 'text/csv',
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.png': 'image/png',
                    '.gif': 'image/gif',
                }
                content_type = content_type_map.get(file_extension, 'application/octet-stream')
           
            # Create response
            response = HttpResponse(file_content, content_type=content_type)
            response['Content-Disposition'] = f'inline; filename="{document.filename}"'
            response['Content-Length'] = len(file_content)
            response['Cache-Control'] = 'no-cache'
            
            # Set CORS headers if needed for frontend
            response['Access-Control-Allow-Origin'] = '*'
            response['Access-Control-Allow-Methods'] = 'GET'
            response['Access-Control-Allow-Headers'] = 'Content-Type'
           
            logger.info(f"Successfully served document {document_id}: {document.filename}")
            return response
           
        except Exception as e:
            logger.error(f"Error retrieving document {document_id}: {str(e)}")
            return HttpResponse(f"Error retrieving document: {str(e)}", status=500)


import requests
from bs4 import BeautifulSoup
import tempfile
import os
import re
from urllib.parse import urlparse, urljoin
from django.core.files import File
from django.db import transaction
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.parsers import JSONParser
from .models import Document, ProcessedIndex, UserModulePermissions
from core.models import Project
from core.utils import update_project_timestamp
import logging

logger = logging.getLogger(__name__)

class WebsiteLinkUploadView(YouTubeUploadView):
    """
    Inherits from YouTubeUploadView to reuse document processing methods.
    Extracts text content from website URLs using BeautifulSoup.
    """
    parser_classes = (JSONParser,)
    
    def is_valid_url(self, url):
        """Check if the provided URL is valid"""
        try:
            parsed = urlparse(url)
            return bool(parsed.netloc) and bool(parsed.scheme)
        except:
            return False
    
    def extract_text_from_website(self, url, timeout=30):
        """
        Extract text content from a website URL using BeautifulSoup.
        Returns the extracted text and page title.
        """
        try:
            # Validate URL
            if not self.is_valid_url(url):
                raise ValueError("Invalid URL provided")
            
            # Add protocol if missing
            if not url.startswith(('http://', 'https://')):
                url = 'https://' + url
            
            logger.info(f"Fetching content from: {url}")
            
            # Set up headers to mimic a real browser
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'en-US,en;q=0.5',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
            }
            
            # Make request with timeout
            response = requests.get(url, headers=headers, timeout=timeout, verify=True)
            response.raise_for_status()
            
            # Check content type
            content_type = response.headers.get('content-type', '').lower()
            if 'text/html' not in content_type:
                raise ValueError(f"URL does not contain HTML content. Content-Type: {content_type}")
            
            # Parse HTML content
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Extract title
            title_tag = soup.find('title')
            page_title = title_tag.get_text().strip() if title_tag else 'Untitled Page'
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "header", "footer", "aside"]):
                script.decompose()
            
            # Try to find main content areas first
            main_content = None
            content_selectors = [
                'main', 'article', '[role="main"]', '.main-content', 
                '.content', '.post-content', '.entry-content', '#content',
                '.article-body', '.story-body', '.post-body'
            ]
            
            for selector in content_selectors:
                main_content = soup.select_one(selector)
                if main_content:
                    logger.info(f"Found main content using selector: {selector}")
                    break
            
            # If no main content found, use body
            if not main_content:
                main_content = soup.find('body')
                logger.info("Using body tag as main content")
            
            if not main_content:
                raise ValueError("Could not find any content in the webpage")
            
            # Extract text content
            text_content = main_content.get_text(separator='\n', strip=True)
            
            # Clean up the text
            lines = text_content.split('\n')
            cleaned_lines = []
            
            for line in lines:
                line = line.strip()
                # Skip empty lines and very short lines (likely navigation/UI elements)
                if len(line) > 3 and not line.isdigit():
                    cleaned_lines.append(line)
            
            cleaned_text = '\n'.join(cleaned_lines)
            
            # Remove excessive whitespace
            cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
            cleaned_text = re.sub(r' {2,}', ' ', cleaned_text)
            
            if len(cleaned_text.strip()) < 50:
                raise ValueError("Extracted content is too short (less than 50 characters)")
            
            logger.info(f"Successfully extracted {len(cleaned_text)} characters from {url}")
            logger.info(f"Page title: {page_title}")
            
            return cleaned_text, page_title
            
        except requests.exceptions.Timeout:
            raise Exception(f"Request timeout while fetching URL: {url}")
        except requests.exceptions.ConnectionError:
            raise Exception(f"Connection error while fetching URL: {url}")
        except requests.exceptions.HTTPError as e:
            raise Exception(f"HTTP error {e.response.status_code} while fetching URL: {url}")
        except Exception as e:
            logger.error(f"Error extracting text from website: {str(e)}")
            raise Exception(f"Failed to extract text from website: {str(e)}")
    
    def post(self, request):
        website_url = request.data.get('website_url')
        user = request.user
        main_project_id = request.data.get('main_project_id')
        target_user_id = request.data.get('target_user_id')
        custom_title = request.data.get('custom_title', '')  # Optional custom title
        
        if not website_url:
            return Response({
                'error': 'Website URL is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        if not main_project_id:
            return Response({
                'error': 'Main project ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Handle admin uploading for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
        
        # Check user permissions
        try:
            permissions = UserModulePermissions.objects.get(user=user)
            if permissions.disabled_modules.get('document-upload', False):
                return Response({
                    'error': 'Document uploads are disabled for this user'
                }, status=status.HTTP_403_FORBIDDEN)
        except UserModulePermissions.DoesNotExist:
            pass
        
        try:
            main_project = Project.objects.get(id=main_project_id, user=user)
            
            # Extract text from website
            logger.info(f"Processing website URL: {website_url}")
            extracted_text, page_title = self.extract_text_from_website(website_url)
            
            if not extracted_text or len(extracted_text.strip()) < 10:
                return Response({
                    'error': f'Failed to extract meaningful content from website: {website_url}'
                }, status=status.HTTP_400_BAD_REQUEST)
            
            # Use custom title if provided, otherwise use page title
            final_title = custom_title.strip() if custom_title.strip() else page_title
            
            # Create filename based on title
            safe_title = self.sanitize_filename(final_title)
            filename = f"{safe_title}_website_content.txt"
            
            # Check for existing document with similar name
            existing_doc = Document.objects.filter(
                user=user,
                filename__icontains=safe_title,
                main_project=main_project
            ).first()
            
            # Create temporary file with the extracted content
            temp_file_path = os.path.join(tempfile.gettempdir(), filename)
            
            try:
                # Prepare content with metadata
                full_content = f"Website URL: {website_url}\n"
                full_content += f"Title: {final_title}\n"
                full_content += "=" * (len(final_title) + 7) + "\n\n"
                full_content += extracted_text
                
                with open(temp_file_path, 'w', encoding='utf-8') as f:
                    f.write(full_content)
                
                # Create Django File object
                with open(temp_file_path, 'rb') as f:
                    django_file = File(f, name=filename)
                    
                    # Use transaction to ensure data consistency
                    with transaction.atomic():
                        # Create or update document
                        if existing_doc:
                            existing_doc.file = django_file
                            existing_doc.filename = filename
                            existing_doc.save()
                            document = existing_doc
                            logger.info(f"Updated existing document: {document.id}")
                        else:
                            document = Document.objects.create(
                                user=user,
                                file=django_file,
                                filename=filename,
                                main_project=main_project
                            )
                            logger.info(f"Created new document: {document.id}")
                
                # Process the document for RAG
                processed_data = self.process_document_from_text(full_content, filename, user)
                
                # Create or update ProcessedIndex
                processed_index, created = ProcessedIndex.objects.get_or_create(
                    document=document,
                    defaults={
                        'faiss_index': processed_data['index_path'],
                        'metadata': processed_data['metadata_path'],
                        'summary': "",
                        'markdown_path': processed_data.get('markdown_path', '')
                    }
                )
                
                if not created:
                    processed_index.faiss_index = processed_data['index_path']
                    processed_index.metadata = processed_data['metadata_path']
                    processed_index.markdown_path = processed_data.get('markdown_path', '')
                    processed_index.save()
                
                # Store the document ID in the session
                request.session['active_document_id'] = document.id
                
                # Update project timestamp
                update_project_timestamp(main_project_id, user)
                
                return Response({
                    'message': 'Website content processed successfully',
                    'document': {
                        'id': document.id,
                        'filename': filename,
                        'original_page_title': page_title,
                        'website_url': website_url,
                        'content_length': len(extracted_text)
                    },
                    'active_document_id': document.id
                }, status=status.HTTP_201_CREATED)
                
            finally:
                # Clean up temporary file
                try:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                        logger.info(f"Cleaned up temporary file: {temp_file_path}")
                except Exception as cleanup_error:
                    logger.warning(f"Cleanup error: {str(cleanup_error)}")
        
        except Exception as e:
            logger.error(f"Error processing website: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while processing the website content'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class PlainTextUploadView(YouTubeUploadView):
    """
    Inherits from YouTubeUploadView to reuse document processing methods.
    Processes plain text content directly from frontend.
    """
    parser_classes = (JSONParser,)
    
    def post(self, request):
        text_content = request.data.get('text_content')
        user = request.user
        main_project_id = request.data.get('main_project_id')
        target_user_id = request.data.get('target_user_id')
        custom_title = request.data.get('title', '')  # Optional title for the text
        
        if not text_content or not text_content.strip():
            return Response({
                'error': 'Text content is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        if not main_project_id:
            return Response({
                'error': 'Main project ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Validate text length
        if len(text_content.strip()) < 10:
            return Response({
                'error': 'Text content must be at least 10 characters long'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        if len(text_content) > 1000000:  # 1MB text limit
            return Response({
                'error': 'Text content is too large (maximum 1MB)'
            }, status=status.HTTP_400_BAD_REQUEST)
        
        # Handle admin uploading for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
        
        # Check user permissions
        try:
            permissions = UserModulePermissions.objects.get(user=user)
            if permissions.disabled_modules.get('document-upload', False):
                return Response({
                    'error': 'Document uploads are disabled for this user'
                }, status=status.HTTP_403_FORBIDDEN)
        except UserModulePermissions.DoesNotExist:
            pass
        
        try:
            main_project = Project.objects.get(id=main_project_id, user=user)
            
            # Generate title if not provided
            if not custom_title.strip():
                # Generate title from first 50 characters of content
                first_line = text_content.strip().split('\n')[0]
                custom_title = first_line[:50].strip()
                if len(first_line) > 50:
                    custom_title += "..."
                if not custom_title:
                    custom_title = "Plain Text Document"
            
            # Create filename based on title
            safe_title = self.sanitize_filename(custom_title)
            filename = f"{safe_title}_plain_text.txt"
            
            # Check for existing document with similar name
            existing_doc = Document.objects.filter(
                user=user,
                filename__icontains=safe_title,
                main_project=main_project
            ).first()
            
            # Create temporary file with the text content
            temp_file_path = os.path.join(tempfile.gettempdir(), filename)
            
            try:
                # Prepare content with title header
                full_content = f"Title: {custom_title}\n"
                full_content += "=" * (len(custom_title) + 7) + "\n\n"
                full_content += text_content.strip()
                
                with open(temp_file_path, 'w', encoding='utf-8') as f:
                    f.write(full_content)
                
                # Create Django File object
                with open(temp_file_path, 'rb') as f:
                    django_file = File(f, name=filename)
                    
                    # Use transaction to ensure data consistency
                    with transaction.atomic():
                        # Create or update document
                        if existing_doc:
                            existing_doc.file = django_file
                            existing_doc.filename = filename
                            existing_doc.save()
                            document = existing_doc
                            logger.info(f"Updated existing document: {document.id}")
                        else:
                            document = Document.objects.create(
                                user=user,
                                file=django_file,
                                filename=filename,
                                main_project=main_project
                            )
                            logger.info(f"Created new document: {document.id}")
                
                # Process the document for RAG
                processed_data = self.process_document_from_text(full_content, filename, user)
                
                # Create or update ProcessedIndex
                processed_index, created = ProcessedIndex.objects.get_or_create(
                    document=document,
                    defaults={
                        'faiss_index': processed_data['index_path'],
                        'metadata': processed_data['metadata_path'],
                        'summary': "",
                        'markdown_path': processed_data.get('markdown_path', '')
                    }
                )
                
                if not created:
                    processed_index.faiss_index = processed_data['index_path']
                    processed_index.metadata = processed_data['metadata_path']
                    processed_index.markdown_path = processed_data.get('markdown_path', '')
                    processed_index.save()
                
                # Store the document ID in the session
                request.session['active_document_id'] = document.id
                
                # Update project timestamp
                update_project_timestamp(main_project_id, user)
                
                return Response({
                    'message': 'Plain text processed successfully',
                    'document': {
                        'id': document.id,
                        'filename': filename,
                        'title': custom_title,
                        'content_length': len(text_content),
                        'word_count': len(text_content.split())
                    },
                    'active_document_id': document.id
                }, status=status.HTTP_201_CREATED)
                
            finally:
                # Clean up temporary file
                try:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)
                        logger.info(f"Cleaned up temporary file: {temp_file_path}")
                except Exception as cleanup_error:
                    logger.warning(f"Cleanup error: {str(cleanup_error)}")
        
        except Project.DoesNotExist:
            return Response({
                'error': 'Project not found'
            }, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            logger.error(f"Error processing plain text: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while processing the plain text'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# import json

# import os
# import re
# import time
# import pickle
# from rest_framework.decorators import api_view, parser_classes
# from rest_framework.parsers import MultiPartParser, FormParser
# from rest_framework.response import Response
# from rest_framework import status
# from rest_framework.views import APIView
# from rest_framework.permissions import IsAuthenticated
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# import google.generativeai as genai
# from .models import MindMap, Document, ProcessedIndex
# from core.models import Project
# from core.utils import update_project_timestamp
# import logging

# logger = logging.getLogger(__name__)

# # Configuration for Mindmap
# GEMINI_API_KEY = "AIzaSyCv_41De6hVOtuQ3jYkfniGc_61bJ9bvS4"  # Move this to settings
# GEMINI_MODEL = "gemini-2.0-flash"
# CHUNK_SIZE = 2000
# CHUNK_OVERLAP = 200

# # Initialize Gemini for mindmap
# genai.configure(api_key=GEMINI_API_KEY)
# mindmap_model = genai.GenerativeModel(GEMINI_MODEL)

# class MindMapView(APIView):
#     """Generate mindmap for selected documents - matches your 'generate-mindmap/' endpoint"""
#     permission_classes = [IsAuthenticated]

#     def init_gemini(self, user=None):
#         """Initialize Gemini 2.0 Flash client"""
#             # Get the user's Gemini API token
#         gemini_api_key = None
#         if user:
#             from .models import UserAPITokens  # Adjust import based on your models location
#             user_api_tokens = UserAPITokens.objects.get(user=user)
#             gemini_api_key = user_api_tokens.gemini_token
           
#             if not gemini_api_key:
#                 logger.warning(f"No Gemini API token found for user {user.username}")
#                 # Fallback to environment variable
#                 gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
#                 if not gemini_api_key:
#                     raise ValueError("Gemini API token is required for processing")
       
#         self.gemini_api_key = gemini_api_key


#     def post(self, request):
#         try:
#             user = request.user
#             main_project_id = request.data.get('main_project_id')
#             if not main_project_id:
#                 return Response({'error': 'Main project ID is required'}, status=status.HTTP_400_BAD_REQUEST)
 
#             target_user_id = request.data.get('target_user_id')
#             force_regenerate = request.data.get('force_regenerate', False)
#             selected_documents = request.data.get('selected_documents', [])
           
#             if not isinstance(selected_documents, list):
#                 selected_documents = [selected_documents]
           
#             try:
#                 selected_documents = [int(doc_id) for doc_id in selected_documents]
#             except (ValueError, TypeError):
#                 return Response({'error': 'Invalid document ID format'}, status=status.HTTP_400_BAD_REQUEST)
 
#             # Handle admin uploading for another user
#             if target_user_id and request.user.username == 'admin':
#                 try:
#                     user = User.objects.get(id=target_user_id)
#                 except User.DoesNotExist:
#                     return Response({'error': 'Target user not found'}, status=status.HTTP_404_NOT_FOUND)
 
#             if not selected_documents:
#                 # Try to get active document from session
#                 active_document_id = request.session.get('active_document_id')
#                 if active_document_id:
#                     selected_documents = [active_document_id]
#                 else:
#                     return Response({'error': 'Please select at least one document'}, status=status.HTTP_400_BAD_REQUEST)
 
#             # Verify project access
#             try:
#                 main_project = Project.objects.get(id=main_project_id, user=user)
#             except Project.DoesNotExist:
#                 return Response({'error': 'Project not found'}, status=status.HTTP_404_NOT_FOUND)
 
#             # Check if mindmap already exists for these documents (unless force regenerate)
#             if not force_regenerate:
#                 existing_mindmap = self.find_existing_mindmap(user, main_project_id, selected_documents)
#                 if existing_mindmap:
#                     return Response({
#                         'success': True,
#                         'mindmap': existing_mindmap.data,
#                         'mindmap_id': existing_mindmap.id,
#                         'from_cache': True,
#                         'stats': {
#                             'total_characters': 0,
#                             'number_of_chunks': 0,
#                             'documents_processed': len(existing_mindmap.get_document_sources_list()),
#                             'mindmap_nodes': existing_mindmap.total_nodes,
#                             'model_used': GEMINI_MODEL,
#                             'document_sources': existing_mindmap.get_document_sources_list(),
#                             'created_at': existing_mindmap.created_at.isoformat()
#                         }
#                     }, status=status.HTTP_200_OK)
 
#             # Get documents that belong to the user
#             user_documents = Document.objects.filter(
#                 id__in=selected_documents,
#                 user=user
#             )
 
#             if not user_documents.exists():
#                 return Response({'error': 'No valid documents found'}, status=status.HTTP_404_NOT_FOUND)
 
#             logger.info(f"Found {user_documents.count()} documents for mindmap generation")
 
#             # Step 1: Load chunks from pgvector DocumentEmbedding
#             per_doc_chunks = []
#             document_info = []
 
#             for document in user_documents:
#                 try:
#                     # Get all embeddings/chunks for this document from pgvector
#                     document_embeddings = DocumentEmbedding.objects.filter(
#                         document=document
#                     ).order_by('chunk_id')
 
#                     if not document_embeddings.exists():
#                         logger.warning(f"No embeddings found for document {document.filename}")
#                         continue
 
#                     # Extract text chunks from embeddings
#                     chunk_texts = [embedding.content for embedding in document_embeddings]
                   
#                     # Only add if we have meaningful content
#                     chunk_texts = [t for t in chunk_texts if t and len(t.strip()) > 0]
                   
#                     if chunk_texts:
#                         per_doc_chunks.append({
#                             'document_id': document.id,
#                             'filename': document.filename,
#                             'chunks': chunk_texts
#                         })
#                         document_info.append({
#                             'filename': document.filename,
#                             'text_length': sum(len(t) for t in chunk_texts),
#                             'chunks': len(chunk_texts),
#                             'document_id': document.id
#                         })
                       
#                     logger.info(f"Loaded {len(chunk_texts)} chunks from document {document.filename}")
                   
#                 except Exception as e:
#                     logger.error(f"Error extracting content from {document.filename}: {str(e)}")
#                     continue
 
#             if not per_doc_chunks:
#                 return Response({'error': 'No valid processed documents found with embeddings'}, status=status.HTTP_404_NOT_FOUND)
 
#             # Interleave chunks from all documents for balanced mindmap context
#             interleaved_chunks = []
#             max_len = max(len(doc['chunks']) for doc in per_doc_chunks) if per_doc_chunks else 0
           
#             for i in range(max_len):
#                 for doc in per_doc_chunks:
#                     if i < len(doc['chunks']):
#                         interleaved_chunks.append(doc['chunks'][i])
           
#             all_text_chunks = interleaved_chunks
 
#             if not all_text_chunks:
#                 return Response({
#                     'success': True,
#                     'message': 'No content found in selected documents'
#                 }, status=status.HTTP_200_OK)
 
#             # Generate mindmap from combined content
#             mindmap_data = self.generate_comprehensive_mindmap(all_text_chunks)
#             if not mindmap_data or 'name' not in mindmap_data:
#                 return Response({
#                     'error': 'Failed to generate mindmap'
#                 }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 
#             # Generate heading for the mindmap
#             combined_text = "\n\n".join(all_text_chunks)
#             heading = self.generate_mindmap_heading(selected_documents=selected_documents,document_info=document_info)
#             if heading:
#                 mindmap_data['name'] = heading
#             else:
#                 # Fallback to filenames if LLM fails
#                 combined_filenames = " + ".join([doc['filename'] for doc in document_info])
#                 mindmap_data['name'] = f"Mindmap: {combined_filenames}"
 
#             node_count = self.count_mindmap_nodes(mindmap_data)
 
#             # Save MindMap record
#             mindmap_record = MindMap.objects.create(
#                 user=user,
#                 main_project=main_project,
#                 data=mindmap_data,
#                 document_sources=[doc['filename'] for doc in document_info],
#                 total_nodes=node_count
#             )
           
#             # Update project timestamp
#             update_project_timestamp(main_project_id, user)
 
#             response_data = {
#                 'success': True,
#                 'mindmap': mindmap_data,
#                 'mindmap_id': mindmap_record.id,
#                 'from_cache': False,
#                 'stats': {
#                     'total_characters': sum(doc['text_length'] for doc in document_info),
#                     'number_of_chunks': len(all_text_chunks),
#                     'documents_processed': len(document_info),
#                     'mindmap_nodes': node_count,
#                     'model_used': GEMINI_MODEL,
#                     'document_sources': [doc['filename'] for doc in document_info],
#                     'created_at': mindmap_record.created_at.isoformat()
#                 }
#             }
#             return Response(response_data, status=status.HTTP_200_OK)
 
#         except Exception as e:
#             logger.error(f"Error in MindMapView: {str(e)}")
#             return Response({
#                 'error': f'Internal server error: {str(e)}',
#                 'success': False
#             }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 
 
#     def generate_mindmap_heading(self, combined_text, api_key=None):
#         """
#         Use Gemini API to generate a meaningful mindmap heading based on document content.
#         """
#         try:
#             import google.generativeai as genai
#             if api_key:
#                 genai.configure(api_key=api_key)
#             else:
#                 genai.configure(api_key=self.gemini_api_key)
#             model = genai.GenerativeModel(GEMINI_MODEL)
#             prompt = (
#                 "Given the following combined document content, generate a concise, descriptive title "
#                 "that best summarizes the main subject or theme. Avoid generic titles. "
#                 "Return only the title as a single line.\n\n"
#                 f"DOCUMENT CONTENT:\n{combined_text[:3000]}"
#             )
#             response = model.generate_content(prompt)
#             title = response.text.strip().replace('\n', ' ')
#             # Clean up any quotes or markdown
#             if title.startswith('"') and title.endswith('"'):
#                 title = title[1:-1]
#             return title
#         except Exception as e:
#             logger.error(f"Error generating mindmap heading: {e}")
#             return None
 


#     def find_existing_mindmap(self, user, main_project_id, selected_documents):
#         try:
#             document_filenames = set()
#             for doc_id in selected_documents:
#                 try:
#                     doc = Document.objects.get(id=doc_id, user=user)
#                     document_filenames.add(doc.filename)
#                 except Document.DoesNotExist:
#                     continue
#             if not document_filenames:
#                 return None
#             existing_mindmaps = MindMap.objects.filter(
#                 user=user,
#                 main_project_id=main_project_id
#             ).order_by('-created_at')
#             for mindmap in existing_mindmaps:
#                 mindmap_sources = set(mindmap.get_document_sources_list())
#                 if mindmap_sources == document_filenames:
#                     return mindmap
#             return None
#         except Exception:
#             return None

#     def count_mindmap_nodes(self, data):
#         if not isinstance(data, dict):
#             return 0
#         count = 1
#         if 'children' in data and isinstance(data['children'], list):
#             for child in data['children']:
#                 count += self.count_mindmap_nodes(child)
#         return count

#     def generate_comprehensive_mindmap(self, text_chunks):
#         try:
#             # selected_chunks = text_chunks
#             # comprehensive_text = "\n\n".join(selected_chunks)
            
#             selected_chunks = text_chunks
#             combined_text = "\n\n".join(selected_chunks)
#             if len(combined_text) > 12000:
#                 combined_text = combined_text[:12000]  # truncate to avoid overloading prompt
#             selected_chunks = combined_text.split("\n\n")  # re-chunk safely if needed

#             comprehensive_text = "\n\n".join(selected_chunks)

#             prompt = f"""
#             You are an expert knowledge architect and document analyzer. Create a comprehensive, hierarchical mind map by systematically extracting and organizing ALL meaningful content from the document(s).

#             DOCUMENT TEXT:
#             {comprehensive_text}

#             Create a JSON mind map with this structure:
#             {{
#                 "name": "[Document Title/Main Subject]",
#                 "children": [
#                     {{
#                         "name": "[Major Topic 1]",
#                         "summary": "Brief description",
#                         "children": [
#                             {{
#                                 "name": "[Subtopic A]",
#                                 "summary": "Description",
#                                 "children": [
#                                     {{
#                                         "name": "[Detail 1]",
#                                         "summary": "Specific information",
#                                         "children": []
#                                     }}
#                                 ]
#                             }}
#                         ]
#                     }}
#                 ]
#             }}

#             REQUIREMENTS:
#             1. Extract all major topics
#             2. Each major topic must have possible subtopics
#             3. Go 9-12 possible levels deep
#             4. Use exact terminology from the document
#             5. Include meaningful summaries
#             6. Ensure hierarchical organization

#             Return only valid JSON:
#             """
#             try:
#                 response = mindmap_model.generate_content(
#                     prompt,
#                     generation_config=genai.types.GenerationConfig(
#                         temperature=0,
#                         max_output_tokens=12288,
#                         top_p=0.8
#                     )
#                 )
#                 parsed_data = self.extract_json_from_response(response.text)
#                 if parsed_data:
#                     validated_data = self.validate_mindmap_structure(parsed_data)
#                     logger.info(f"Successfully generated mindmap with {len(validated_data.get('children', []))} main branches")
#                     return validated_data
#                 else:
#                     logger.warning("Failed to parse JSON, creating fallback")
#                     return self.create_enhanced_fallback_mindmap(comprehensive_text)
#             except Exception as api_error:
#                 logger.error(f"API Error: {api_error}")
#                 return self.create_enhanced_fallback_mindmap(comprehensive_text)
#         except Exception as e:
#             logger.error(f"Error generating mindmap: {e}")
#             return self.create_enhanced_fallback_mindmap(text_chunks[0] if text_chunks else "No content")

#     def extract_json_from_response(self, response_text):
#         try:
#             cleaned_text = response_text.strip()
#             if cleaned_text.startswith("```json"):
#                 cleaned_text = cleaned_text[7:]
#             elif cleaned_text.startswith("```"):
#                 cleaned_text = cleaned_text[3:]
#             if cleaned_text.endswith("```"):
#                 cleaned_text = cleaned_text[:-3]
#             cleaned_text = cleaned_text.strip()
#             try:
#                 return json.loads(cleaned_text)
#             except json.JSONDecodeError:
#                 pass
#             json_pattern = r'\{(?:[^{}]|{(?:[^{}]|{[^{}]*})*})*\}'
#             matches = re.findall(json_pattern, cleaned_text, re.DOTALL)
#             for match in matches:
#                 try:
#                     parsed = json.loads(match)
#                     if isinstance(parsed, dict) and ('name' in parsed or 'children' in parsed):
#                         return parsed
#                 except json.JSONDecodeError:
#                     continue
#             return None
#         except Exception as e:
#             logger.error(f"Error extracting JSON: {e}")
#             return None

#     def validate_mindmap_structure(self, data):
#         if not isinstance(data, dict):
#             return {"name": "Document", "children": []}
#         if 'name' not in data:
#             data['name'] = "Document Analysis"
#         if 'children' in data:
#             if not isinstance(data['children'], list):
#                 data['children'] = []
#             else:
#                 validated_children = []
#                 for child in data['children']:
#                     validated_child = self.validate_mindmap_structure(child)
#                     if validated_child['name'].strip():
#                         validated_children.append(validated_child)
#                 data['children'] = validated_children
#         else:
#             data['children'] = []
#         allowed_fields = {'name', 'children', 'description', 'type', 'summary'}
#         cleaned_data = {k: v for k, v in data.items() if k in allowed_fields}
#         return cleaned_data

#     def create_enhanced_fallback_mindmap(self, text):
#         try:
#             sentences = [s.strip() for s in text.split('.') if 15 < len(s.strip()) < 200]
#             paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]
#             words = re.findall(r'\b[A-Za-z]{4,}\b', text.lower())
#             word_freq = {}
#             for word in words:
#                 if word not in ['these', 'those', 'which', 'where', 'there', 'should', 'would']:
#                     word_freq[word] = word_freq.get(word, 0) + 1
#             top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:15]
#             children = [
#                 {
#                     "name": "Document Overview",
#                     "summary": "General overview of the document content",
#                     "children": [
#                         {"name": "Main Topics", "children": [], "summary": "Primary subjects covered"},
#                         {"name": "Key Information", "children": [], "summary": "Important facts and data"},
#                         {"name": "Document Structure", "children": [], "summary": "Organization of content"}
#                     ]
#                 },
#                 {
#                     "name": "Content Analysis",
#                     "summary": "Detailed analysis of document content",
#                     "children": [
#                         {"name": "Core Concepts", "children": [], "summary": "Fundamental ideas presented"},
#                         {"name": "Supporting Details", "children": [], "summary": "Evidence and examples"},
#                         {"name": "Conclusions", "children": [], "summary": "Main findings and outcomes"}
#                     ]
#                 }
#             ]
#             if top_keywords:
#                 term_children = []
#                 for keyword, freq in top_keywords[:8]:
#                     term_children.append({
#                         "name": keyword.title(),
#                         "children": [],
#                         "summary": f"Mentioned {freq} times in document"
#                     })
#                 children.append({
#                     "name": "Key Terms and Concepts",
#                     "summary": "Important terminology from the document",
#                     "children": term_children
#                 })
#             root_name = "Document Analysis"
#             if sentences:
#                 first_sentence = sentences[0]
#                 if len(first_sentence) < 100:
#                     root_name = first_sentence
#             return {
#                 "name": root_name,
#                 "children": children
#             }
#         except Exception as e:
#             logger.error(f"Error in fallback mindmap: {e}")
#             return {
#                 "name": "Document Analysis",
#                 "children": [
#                     {
#                         "name": "Content Overview",
#                         "children": [
#                             {"name": "Main Information", "children": []},
#                             {"name": "Key Details", "children": []},
#                             {"name": "Supporting Context", "children": []}
#                         ]
#                     }
#                 ]
#             }


#     def calculate_max_depth(self,node, current_depth=0):
#         """Calculate the maximum depth of a mindmap structure."""
#         if not isinstance(node, dict) or 'children' not in node:
#             return current_depth
        
#         children = node.get('children', [])
#         if not children:
#             return current_depth
        
#         max_child_depth = current_depth
#         for child in children:
#             child_depth = self.calculate_max_depth(child, current_depth + 1)
#             max_child_depth = max(max_child_depth, child_depth)
        
#         return max_child_depth

#     def create_deep_hierarchical_fallback(self,text, batch_number):
#         """Create a deep hierarchical fallback mind map for a specific batch."""
#         try:
#             # Enhanced text analysis for deeper hierarchy
#             sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 15][:25]
#             paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50][:15]
            
#             # Extract keywords with frequency
#             words = re.findall(r'\b[A-Za-z]{3,}\b', text.lower())
#             word_freq = {}
#             stop_words = {'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'was', 'one', 'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'man', 'new', 'now', 'old', 'see', 'two', 'way', 'who', 'boy', 'did', 'its', 'let', 'put', 'say', 'she', 'too', 'use'}
            
#             for word in words:
#                 if word not in stop_words and len(word) > 3:
#                     word_freq[word] = word_freq.get(word, 0) + 1
            
#             top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:12]
            
#             # Create deep hierarchical structure
#             children = []
            
#             # Group keywords into thematic categories for deeper hierarchy
#             if top_keywords:
#                 # Create 3-4 main categories with deep sub-structures
#                 keyword_groups = [
#                     top_keywords[0:3],   # Primary concepts
#                     top_keywords[3:6],   # Secondary concepts  
#                     top_keywords[6:9],   # Supporting concepts
#                     top_keywords[9:12]   # Additional concepts
#                 ]
                
#                 category_names = [
#                     "Primary Concepts",
#                     "Secondary Concepts", 
#                     "Supporting Elements",
#                     "Additional Information"
#                 ]
                
#                 for i, (category_name, keyword_group) in enumerate(zip(category_names, keyword_groups)):
#                     if not keyword_group:
#                         continue
                        
#                     category_children = []
#                     for keyword, freq in keyword_group:
#                         # Create deep structure for each keyword
#                         keyword_node = {
#                             "name": keyword.title().replace('_', ' '),
#                             "children": [
#                                 {
#                                     "name": "Definition & Context",
#                                     "children": [
#                                         {"name": "Core Meaning", "children": [], "summary": f"Primary definition and significance of {keyword}"},
#                                         {"name": "Usage Context", "children": [], "summary": f"How {keyword} is used in this section"}
#                                     ]
#                                 },
#                                 {
#                                     "name": "Details & Examples",
#                                     "children": [
#                                         {"name": "Specific Details", "children": [], "summary": f"Detailed information about {keyword}"},
#                                         {"name": "Related Examples", "children": [], "summary": f"Examples and instances of {keyword} mentioned {freq} times"}
#                                     ]
#                                 }
#                             ]
#                         }
#                         category_children.append(keyword_node)
                    
#                     if category_children:
#                         children.append({
#                             "name": category_name,
#                             "children": category_children
#                         })
            
#             # Add content structure analysis
#             if sentences:
#                 content_structure = {
#                     "name": "Content Structure",
#                     "children": [
#                         {
#                             "name": "Main Information",
#                             "children": [
#                                 {
#                                     "name": "Key Statements",
#                                     "children": [
#                                         {"name": "Primary Points", "children": [], "summary": "Main ideas expressed in this section"},
#                                         {"name": "Supporting Points", "children": [], "summary": "Additional supporting information"}
#                                     ]
#                                 },
#                                 {
#                                     "name": "Detailed Content",
#                                     "children": [
#                                         {"name": "Explanations", "children": [], "summary": "Detailed explanations and descriptions"},
#                                         {"name": "Specifications", "children": [], "summary": "Specific details and technical information"}
#                                     ]
#                                 }
#                             ]
#                         },
#                         {
#                             "name": "Organization",
#                             "children": [
#                                 {"name": "Information Flow", "children": [], "summary": "How information is organized in this section"},
#                                 {"name": "Topic Progression", "children": [], "summary": "Sequence and development of topics"}
#                             ]
#                         }
#                     ]
#                 }
#                 children.append(content_structure)
            
#             # Determine root name based on content
#             root_name = f"Section {batch_number} Analysis"
#             if sentences and len(sentences[0]) < 100:
#                 clean_sentence = re.sub(r'[^\w\s]', '', sentences[0])
#                 if len(clean_sentence) > 10:
#                     root_name = clean_sentence[:70] + "..." if len(clean_sentence) > 70 else clean_sentence
            
#             return {
#                 "name": root_name,
#                 "children": children if children else [
#                     {
#                         "name": "Content Overview",
#                         "children": [
#                             {
#                                 "name": "Section Information",
#                                 "children": [
#                                     {"name": "Main Content", "children": [], "summary": f"Primary information from section {batch_number}"},
#                                     {"name": "Supporting Details", "children": [], "summary": f"Additional details from section {batch_number}"}
#                                 ]
#                             }
#                         ]
#                     }
#                 ]
#             }
            
#         except Exception as e:
#             print(f"Error in deep hierarchical fallback for batch {batch_number}: {e}")
#             return {
#                 "name": f"Section {batch_number} Content",
#                 "children": [
#                     {
#                         "name": "Document Content",
#                         "children": [
#                             {
#                                 "name": "Information",
#                                 "children": [
#                                     {"name": "Content Details", "children": [], "summary": f"Content from batch {batch_number}"}
#                                 ]
#                             }
#                         ]
#                     }
#                 ]
#             }

#     def normalize_topic_name(name):
#         """Normalize topic name for comparison - case insensitive, strip whitespace, handle variations."""
#         if not name or not isinstance(name, str):
#             return ""
        
#         # Convert to lowercase and strip whitespace
#         normalized = name.strip().lower()
        
#         # Remove common variations in formatting
#         normalized = re.sub(r'\s+', ' ', normalized)  # Multiple spaces to single space
#         normalized = re.sub(r'[^\w\s]', '', normalized)  # Remove special characters
        
#         return normalized

#     def find_matching_topic(self,topic_name, existing_children):
#         """Find if a topic already exists in the children list with fuzzy matching."""
#         normalized_new = self.normalize_topic_name(topic_name)
        
#         for i, child in enumerate(existing_children):
#             if isinstance(child, dict) and 'name' in child:
#                 normalized_existing = self.normalize_topic_name(child['name'])
                
#                 # Exact match
#                 if normalized_new == normalized_existing:
#                     return i, child
                
#                 # Fuzzy match - check if one contains the other (for similar topics)
#                 if len(normalized_new) > 3 and len(normalized_existing) > 3:
#                     if normalized_new in normalized_existing or normalized_existing in normalized_new:
#                         # Additional check - ensure they're actually related topics
#                         words_new = set(normalized_new.split())
#                         words_existing = set(normalized_existing.split())
                        
#                         # If they share significant words, consider them a match
#                         if len(words_new.intersection(words_existing)) >= min(len(words_new), len(words_existing)) * 0.6:
#                             return i, child
        
#         return None, None

#     def merge_hierarchical_topics(self,existing_node, new_node, depth=0):
#         """
#         Advanced hierarchical merging that handles deep topic structures:
#         topic -> subtopic -> sub-subtopic -> continue...
        
#         Logic:
#         1. If topic exists, merge subtopics into it
#         2. If topic doesn't exist, add it completely
#         3. Recursively apply this at all levels
#         """
#         try:
#             if not existing_node or not isinstance(existing_node, dict):
#                 return new_node if new_node else {}
            
#             if not new_node or not isinstance(new_node, dict):
#                 return existing_node
            
#             print(f"{'  ' * depth}Merging at depth {depth}: '{existing_node.get('name', 'Unknown')}' + '{new_node.get('name', 'Unknown')}'")
            
#             # Start with existing node as base
#             merged_node = existing_node.copy()
            
#             # Update summary if new node has better content
#             if 'summary' in new_node:
#                 if 'summary' not in merged_node or not merged_node.get('summary') or len(new_node['summary']) > len(merged_node.get('summary', '')):
#                     merged_node['summary'] = new_node['summary']
            
#             # Get children from both nodes
#             existing_children = existing_node.get('children', [])
#             new_children = new_node.get('children', [])
            
#             if not new_children:
#                 return merged_node
            
#             # Start with copy of existing children
#             merged_children = existing_children.copy()
            
#             # Process each new child
#             for new_child in new_children:
#                 if not isinstance(new_child, dict) or 'name' not in new_child:
#                     continue
                
#                 # Try to find matching existing topic
#                 match_index, matching_child = self.find_matching_topic(new_child['name'], merged_children)
                
#                 if matching_child is not None:
#                     # Topic already exists - merge the subtopics recursively
#                     print(f"{'  ' * depth}  ↳ Merging subtopics into existing topic: '{matching_child['name']}'")
#                     merged_children[match_index] = self.merge_hierarchical_topics(matching_child, new_child, depth + 1)
                    
#                     # Count new subtopics added
#                     existing_subtopic_count = len(matching_child.get('children', []))
#                     new_subtopic_count = len(merged_children[match_index].get('children', []))
#                     if new_subtopic_count > existing_subtopic_count:
#                         print(f"{'  ' * depth}    ✓ Added {new_subtopic_count - existing_subtopic_count} new subtopics")
                        
#                 else:
#                     # Topic doesn't exist - add it completely with all its hierarchy
#                     merged_children.append(new_child)
#                     subtopic_count = self.count_mindmap_nodes(new_child) - 1  # Subtract 1 for the topic itself
#                     print(f"{'  ' * depth}  ✓ Added new topic: '{new_child['name']}' (with {subtopic_count} subtopics)")
            
#             merged_node['children'] = merged_children
#             return merged_node
            
#         except Exception as e:
#             print(f"Error in hierarchical merging at depth {depth}: {e}")
#             return existing_node if existing_node else new_node

#     def generate_mindmap_heading(self, selected_documents=None, document_info=None, api_key=None):
 
#         """
 
#         Use GPT-4o API to generate a meaningful mindmap heading based on selected document filenames.
 
#         Uses only the filenames of selected documents to create a descriptive title.
 
#         """
 
#         try:
 
#             import openai
#             load_dotenv()
#             api_key = os.getenv("OPENAI_API_KEY")
#             # Set API key
 
#             if api_key:
 
#                 openai.api_key = api_key
 
#             else:
 
#                 openai.api_key = OPENAI_API_KEY
 
#             # Extract raw filenames - no cleaning, let LLM handle everything
 
#             filenames = []
 
#             if document_info and isinstance(document_info, list):
 
#                 filenames = [doc.get('filename', '') for doc in document_info if doc.get('filename')]
 
#             elif selected_documents and isinstance(selected_documents, list):
 
#                 filenames = [f"Document_{doc_id}" for doc_id in selected_documents]
 
#             if not filenames:
 
#                 filenames = ["untitled_document"]
 
#             # Just pass raw filenames to LLM
 
#             filenames_text = "\n".join([f"- {filename}" for filename in filenames])
 
#             # Let GPT-4o do everything - but MUST stay related to document names
 
#             if len(filenames) == 1:
 
#                 prompt = f"""Generate a professional 7-8 word mindmap heading that is DIRECTLY RELATED to this document filename. The heading must reflect the actual content suggested by the filename.
   
#     Document filename: {filenames[0]}
   
#     CRITICAL REQUIREMENTS:
 
#     - The heading MUST be directly related to what this filename suggests
 
#     - Use keywords and concepts from the actual filename
 
#     - Create EXACTLY 7 or 8 words
 
#     - Make it professional but stay true to the document's apparent topic
 
#     - If filename mentions "sales" - heading must be about sales
 
#     - If filename mentions "marketing" - heading must be about marketing
 
#     - If filename mentions "technical" - heading must be about technical topics
 
#     - DO NOT create generic headings unrelated to the filename
 
#     - Extract the main subject from the filename and build around it
   
#     Generate the professional mindmap heading that reflects this specific document:"""
   
#             else:
 
#                 prompt = f"""Generate a professional 7-8 word mindmap heading that is DIRECTLY RELATED to these specific document filenames. The heading must reflect the actual content suggested by ALL these filenames.
   
#     Document filenames:
 
#     {filenames_text}
   
#     CRITICAL REQUIREMENTS:
 
#     - The heading MUST be directly related to what these filenames suggest
 
#     - Use actual keywords and concepts from the filenames themselves
 
#     - Find the real common themes from the actual filenames (not generic themes)
 
#     - Create EXACTLY 7 or 8 words
 
#     - If filenames are about "budget" and "finance" - heading must be about financial topics
 
#     - If filenames are about "user guide" and "manual" - heading must be about documentation
 
#     - If filenames are about "research" and "data" - heading must be about research/data
 
#     - DO NOT create generic headings unrelated to the actual filename content
 
#     - The heading should make someone think "yes, this is exactly about those documents"
   
#     Generate the professional mindmap heading that reflects these specific documents:"""
   
#             # Call GPT-4o API with emphasis on staying document-relevant
 
#             response = openai.chat.completions.create(
 
#                 model="gpt-4o",  # Latest GPT-4o model
 
#                 messages=[
 
#                     {
 
#                         "role": "system",
 
#                         "content": "You are an expert at creating professional titles that are DIRECTLY related to document content. Never create generic titles. Always extract and use the actual subject matter from the filenames provided. The title must make it obvious what the documents are about."
 
#                     },
 
#                     {
 
#                         "role": "user",
 
#                         "content": prompt
 
#                     }
 
#                 ],
 
#                 max_tokens=50,
 
#                 temperature=0.6,  # Slightly lower to stay more focused on actual content
 
#                 top_p=0.85,
 
#                 frequency_penalty=0.3,
 
#                 presence_penalty=0.2
 
#             )
 
#             # Extract the generated title - minimal processing
 
#             title = response.choices[0].message.content.strip()
 
#             # Basic word count check - if wrong, try once more
 
#             word_count = len(title.split())
 
#             if word_count < 7 or word_count > 8:
 
#                 # Retry with more specific instruction about staying document-relevant
 
#                 retry_response = openai.chat.completions.create(
 
#                     model="gpt-4o",
 
#                     messages=[
 
#                         {
 
#                             "role": "system",
 
#                             "content": "Generate exactly 7 words that MUST relate to the specific document content. No generic titles allowed."
 
#                         },
 
#                         {
 
#                             "role": "user",
 
#                             "content": f"Based on these specific filenames: {', '.join(filenames[:3])}, create exactly 7 words that someone would recognize as being about these exact documents. Use actual keywords from the filenames."
 
#                         }
 
#                     ],
 
#                     max_tokens=30,
 
#                     temperature=0.4  # Lower temperature for more document-focused results
 
#                 )
 
#                 retry_title = retry_response.choices[0].message.content.strip()
 
#                 retry_word_count = len(retry_title.split())
 
#                 if 7 <= retry_word_count <= 8:
 
#                     title = retry_title
 
#             logger.info(f"Generated {len(title.split())}-word heading: '{title}' from {len(filenames)} files")
 
#             return title
 
#         except Exception as e:
 
#             logger.error(f"Error generating mindmap heading with GPT-4o: {e}")
 
#             # Simple fallback without any processing
 
#             if len(filenames) == 1:
 
#                 return f"Professional Analysis of {filenames[0].split('.')[0]} Research Framework"
 
#             else:
 
#                 return f"Integrated Analysis of {len(filenames)} Document Research Collection"


# # In your backend views.py, update the MindMapQuestionView:

# class MindMapQuestionView(APIView):
#     """Generate question for mindmap node - matches your 'mindmap-question/' endpoint"""
#     permission_classes = [IsAuthenticated]
   
#     def post(self, request):
#         try:
#             user = request.user
#             main_project_id = request.data.get('main_project_id')
#             if not main_project_id:
#                 return Response({'error': 'Main project ID is required'}, status=status.HTTP_400_BAD_REQUEST)
#             mindmap_id = request.data.get('mindmap_id')
#             topic_name = request.data.get('topic_name', '')
#             topic_summary = request.data.get('topic_summary', '')
#             node_path = request.data.get('node_path', '')
#             selected_documents = request.data.get('selected_documents', [])
#             target_user_id = request.data.get('target_user_id')
#             force_new_context = request.data.get('force_new_context', False)
#             current_timestamp = request.data.get('current_timestamp')
#             mindmap_document_sources = request.data.get('mindmap_document_sources', [])

#             # Handle admin operation for another user
#             if target_user_id and request.user.username == 'admin':
#                 try:
#                     user = User.objects.get(id=target_user_id)
#                 except User.DoesNotExist:
#                     return Response({
#                         'error': 'Target user not found'
#                     }, status=status.HTTP_404_NOT_FOUND)

#             if not all([main_project_id, topic_name]):
#                 return Response({
#                     'error': 'Missing required parameters'
#                 }, status=status.HTTP_400_BAD_REQUEST)

#             # If multiple documents are selected, combine their context for question generation
#             if selected_documents:
#                 valid_docs = Document.objects.filter(
#                     id__in=selected_documents,
#                     user=user
#                 ).count()
#                 if valid_docs != len(selected_documents):
#                     valid_doc_ids = list(Document.objects.filter(
#                         id__in=selected_documents,
#                         user=user
#                     ).values_list('id', flat=True))
#                     selected_documents = [str(doc_id) for doc_id in valid_doc_ids]

#             # Generate specific question for the topic using combined context
#             question = self.generate_topic_question(
#                 topic_name,
#                 topic_summary,
#                 node_path,
#                 mindmap_id=mindmap_id,
#                 document_context=selected_documents
#             )

#             return Response({
#                 'success': True,
#                 'question': question,
#                 'topic': topic_name,
#                 'node_path': node_path,
#                 'mindmap_id': mindmap_id,
#                 'selected_documents': selected_documents,
#                 'timestamp': current_timestamp or int(time.time() * 1000)
#             }, status=status.HTTP_200_OK)

#         except Exception as e:
#             logger.error(f"Error in MindMapQuestionView: {str(e)}")
#             return Response({
#                 'error': f'Internal server error: {str(e)}',
#                 'success': False
#             }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)



#     def generate_topic_question(self, topic_name, topic_summary, node_path, mindmap_id=None, document_context=None):
#         """Generate a specific question for the mindmap topic with enhanced context."""
#         try:
#             time.sleep(1.0)  # Rate limiting
            
#             # ✅ NEW: Enhanced prompt with document context awareness
#             context_info = ""
#             if document_context and len(document_context) > 0:
#                 context_info = f"The user has selected {len(document_context)} specific document(s) for this question context. "
            
#             prompt = f"""
#                 Generate ONE simple and helpful question about "{topic_name}" that will make it easier for users to understand this topic.

#                 Topic: {topic_name}  
#                 Summary: {topic_summary}  
#                 Context Path: {node_path}
#                 Mindmap ID: {mindmap_id}
#                 {context_info}

#                 Create a question that:
#                 1. Uses simple language  
#                 2. Helps someone understand the topic better  
#                 3. Is 10–25 words long  
#                 4. Focuses on practical use or basic understanding
#                 5. {"Should be answerable from the selected documents" if document_context else "Can be answered from general knowledge"}

#                 Examples of good questions:
#                 - "How is [topic] used in real life?"  
#                 - "What is the main idea behind [topic]?"  
#                 - "Why is [topic] important or useful?"

#                 Generate ONE question:
#             """
           
#             response = mindmap_model.generate_content(
#                 prompt,
#                 generation_config=genai.types.GenerationConfig(
#                     temperature=0.2,
#                     max_output_tokens=512
#                 )
#             )
           
#             question = response.text.strip()
           
#             # Clean the question
#             if question.startswith('"') and question.endswith('"'):
#                 question = question[1:-1]
                
#             # ✅ NEW: Log the generated question with context
#             logger.info(f"Generated question for mindmap {mindmap_id}: {question}")
#             logger.info(f"Question context: {len(document_context or [])} documents")
           
#             return question
           
#         except Exception as e:
#             logger.error(f"Error generating topic question: {e}")
#             return f"What are the key aspects and details about '{topic_name}' in this document?"


#     """Generate question for mindmap node - matches your 'mindmap-question/' endpoint"""
#     permission_classes = [IsAuthenticated]




import json
import os
import re
import time
import pickle
from rest_framework.decorators import api_view, parser_classes
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework.response import Response
from rest_framework import status
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated
from langchain_text_splitters import RecursiveCharacterTextSplitter
from google import genai
from google.genai import types
from .models import MindMap, Document, ProcessedIndex
from core.models import Project
from core.utils import update_project_timestamp
import logging

logger = logging.getLogger(__name__)

# Configuration for Mindmap
GEMINI_API_KEY = "AIzaSyCv_41De6hVOtuQ3jYkfniGc_61bJ9bvS4"  # Move this to settings
GEMINI_MODEL = "gemini-2.5-flash"
CHUNK_SIZE = 2000
CHUNK_OVERLAP = 200

# Initialize Gemini client for mindmap
mindmap_client = genai.Client(api_key=GEMINI_API_KEY)

class MindMapView(APIView):
    """Generate mindmap for selected documents - matches your 'generate-mindmap/' endpoint"""
    permission_classes = [IsAuthenticated]

    def init_gemini(self, user=None):
        """Initialize Gemini 2.5 Flash client"""
            # Get the user's Gemini API token
        gemini_api_key = None
        if user:
            from .models import UserAPITokens  # Adjust import based on your models location
            user_api_tokens = UserAPITokens.objects.get(user=user)
            gemini_api_key = user_api_tokens.gemini_token
           
            if not gemini_api_key:
                logger.warning(f"No Gemini API token found for user {user.username}")
                # Fallback to environment variable
                gemini_api_key = os.environ.get("GOOGLE_API_KEY", "")
                if not gemini_api_key:
                    raise ValueError("Gemini API token is required for processing")
       
        self.gemini_api_key = gemini_api_key
        self.mindmap_client = genai.Client(api_key=gemini_api_key)


    def post(self, request):
        try:
            user = request.user
            main_project_id = request.data.get('main_project_id')
            if not main_project_id:
                return Response({'error': 'Main project ID is required'}, status=status.HTTP_400_BAD_REQUEST)
 
            target_user_id = request.data.get('target_user_id')
            force_regenerate = request.data.get('force_regenerate', False)
            selected_documents = request.data.get('selected_documents', [])
           
            if not isinstance(selected_documents, list):
                selected_documents = [selected_documents]
           
            try:
                selected_documents = [int(doc_id) for doc_id in selected_documents]
            except (ValueError, TypeError):
                return Response({'error': 'Invalid document ID format'}, status=status.HTTP_400_BAD_REQUEST)
 
            # Handle admin uploading for another user
            if target_user_id and request.user.username == 'admin':
                try:
                    user = User.objects.get(id=target_user_id)
                except User.DoesNotExist:
                    return Response({'error': 'Target user not found'}, status=status.HTTP_404_NOT_FOUND)
 
            if not selected_documents:
                # Try to get active document from session
                active_document_id = request.session.get('active_document_id')
                if active_document_id:
                    selected_documents = [active_document_id]
                else:
                    return Response({'error': 'Please select at least one document'}, status=status.HTTP_400_BAD_REQUEST)
 
            # Verify project access
            try:
                main_project = Project.objects.get(id=main_project_id, user=user)
            except Project.DoesNotExist:
                return Response({'error': 'Project not found'}, status=status.HTTP_404_NOT_FOUND)
 
            # Check if mindmap already exists for these documents (unless force regenerate)
            if not force_regenerate:
                existing_mindmap = self.find_existing_mindmap(user, main_project_id, selected_documents)
                if existing_mindmap:
                    return Response({
                        'success': True,
                        'mindmap': existing_mindmap.data,
                        'mindmap_id': existing_mindmap.id,
                        'from_cache': True,
                        'stats': {
                            'total_characters': 0,
                            'number_of_chunks': 0,
                            'documents_processed': len(existing_mindmap.get_document_sources_list()),
                            'mindmap_nodes': existing_mindmap.total_nodes,
                            'model_used': GEMINI_MODEL,
                            'document_sources': existing_mindmap.get_document_sources_list(),
                            'created_at': existing_mindmap.created_at.isoformat()
                        }
                    }, status=status.HTTP_200_OK)
 
            # Get documents that belong to the user
            user_documents = Document.objects.filter(
                id__in=selected_documents,
                user=user
            )
 
            if not user_documents.exists():
                return Response({'error': 'No valid documents found'}, status=status.HTTP_404_NOT_FOUND)
 
            logger.info(f"Found {user_documents.count()} documents for mindmap generation")
 
            # Step 1: Load chunks from pgvector DocumentEmbedding
            per_doc_chunks = []
            document_info = []
 
            for document in user_documents:
                try:
                    # Get all embeddings/chunks for this document from pgvector
                    document_embeddings = DocumentEmbedding.objects.filter(
                        document=document
                    ).order_by('chunk_id')
 
                    if not document_embeddings.exists():
                        logger.warning(f"No embeddings found for document {document.filename}")
                        continue
 
                    # Extract text chunks from embeddings
                    chunk_texts = [embedding.content for embedding in document_embeddings]
                   
                    # Only add if we have meaningful content
                    chunk_texts = [t for t in chunk_texts if t and len(t.strip()) > 0]
                   
                    if chunk_texts:
                        per_doc_chunks.append({
                            'document_id': document.id,
                            'filename': document.filename,
                            'chunks': chunk_texts
                        })
                        document_info.append({
                            'filename': document.filename,
                            'text_length': sum(len(t) for t in chunk_texts),
                            'chunks': len(chunk_texts),
                            'document_id': document.id
                        })
                       
                    logger.info(f"Loaded {len(chunk_texts)} chunks from document {document.filename}")
                   
                except Exception as e:
                    logger.error(f"Error extracting content from {document.filename}: {str(e)}")
                    continue
 
            if not per_doc_chunks:
                return Response({'error': 'No valid processed documents found with embeddings'}, status=status.HTTP_404_NOT_FOUND)
 
            # Interleave chunks from all documents for balanced mindmap context
            interleaved_chunks = []
            max_len = max(len(doc['chunks']) for doc in per_doc_chunks) if per_doc_chunks else 0
           
            for i in range(max_len):
                for doc in per_doc_chunks:
                    if i < len(doc['chunks']):
                        interleaved_chunks.append(doc['chunks'][i])
           
            all_text_chunks = interleaved_chunks
 
            if not all_text_chunks:
                return Response({
                    'success': True,
                    'message': 'No content found in selected documents'
                }, status=status.HTTP_200_OK)
 
            # Generate mindmap from combined content
            mindmap_data = self.generate_comprehensive_mindmap(all_text_chunks)
            if not mindmap_data or 'name' not in mindmap_data:
                return Response({
                    'error': 'Failed to generate mindmap'
                }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 
            # Generate heading for the mindmap
            combined_text = "\n\n".join(all_text_chunks)
            heading = self.generate_mindmap_heading(selected_documents=selected_documents,document_info=document_info)
            if heading:
                mindmap_data['name'] = heading
            else:
                # Fallback to filenames if LLM fails
                combined_filenames = " + ".join([doc['filename'] for doc in document_info])
                mindmap_data['name'] = f"Mindmap: {combined_filenames}"
 
            node_count = self.count_mindmap_nodes(mindmap_data)
 
            # Save MindMap record
            mindmap_record = MindMap.objects.create(
                user=user,
                main_project=main_project,
                data=mindmap_data,
                document_sources=[doc['filename'] for doc in document_info],
                total_nodes=node_count
            )
           
            # Update project timestamp
            update_project_timestamp(main_project_id, user)
 
            response_data = {
                'success': True,
                'mindmap': mindmap_data,
                'mindmap_id': mindmap_record.id,
                'from_cache': False,
                'stats': {
                    'total_characters': sum(doc['text_length'] for doc in document_info),
                    'number_of_chunks': len(all_text_chunks),
                    'documents_processed': len(document_info),
                    'mindmap_nodes': node_count,
                    'model_used': GEMINI_MODEL,
                    'document_sources': [doc['filename'] for doc in document_info],
                    'created_at': mindmap_record.created_at.isoformat()
                }
            }
            return Response(response_data, status=status.HTTP_200_OK)
 
        except Exception as e:
            logger.error(f"Error in MindMapView: {str(e)}")
            return Response({
                'error': f'Internal server error: {str(e)}',
                'success': False
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 
 

    def find_existing_mindmap(self, user, main_project_id, selected_documents):
        try:
            document_filenames = set()
            for doc_id in selected_documents:
                try:
                    doc = Document.objects.get(id=doc_id, user=user)
                    document_filenames.add(doc.filename)
                except Document.DoesNotExist:
                    continue
            if not document_filenames:
                return None
            existing_mindmaps = MindMap.objects.filter(
                user=user,
                main_project_id=main_project_id
            ).order_by('-created_at')
            for mindmap in existing_mindmaps:
                mindmap_sources = set(mindmap.get_document_sources_list())
                if mindmap_sources == document_filenames:
                    return mindmap
            return None
        except Exception:
            return None

    def count_mindmap_nodes(self, data):
        if not isinstance(data, dict):
            return 0
        count = 1
        if 'children' in data and isinstance(data['children'], list):
            for child in data['children']:
                count += self.count_mindmap_nodes(child)
        return count

    def generate_comprehensive_mindmap(self, text_chunks):
        try:
            selected_chunks = text_chunks
            combined_text = "\n\n".join(selected_chunks)
            if len(combined_text) > 12000:
                combined_text = combined_text[:12000]  # truncate to avoid overloading prompt
            selected_chunks = combined_text.split("\n\n")  # re-chunk safely if needed

            comprehensive_text = "\n\n".join(selected_chunks)

            prompt = f"""
            You are an expert knowledge architect and document analyzer. Create a comprehensive, hierarchical mind map by systematically extracting and organizing ALL meaningful content from the document(s).

            DOCUMENT TEXT:
            {comprehensive_text}

            Create a JSON mind map with this structure:
            {{
                "name": "[Document Title/Main Subject]",
                "children": [
                    {{
                        "name": "[Major Topic 1]",
                        "summary": "Brief description",
                        "children": [
                            {{
                                "name": "[Subtopic A]",
                                "summary": "Description",
                                "children": [
                                    {{
                                        "name": "[Detail 1]",
                                        "summary": "Specific information",
                                        "children": []
                                    }}
                                ]
                            }}
                        ]
                    }}
                ]
            }}

            REQUIREMENTS:
            1. Extract all major topics
            2. Each major topic must have possible subtopics
            3. Go 9-12 possible levels deep
            4. Use exact terminology from the document
            5. Include meaningful summaries
            6. Ensure hierarchical organization

            Return only valid JSON:
            """
            try:
                response = mindmap_client.models.generate_content(
                    model=GEMINI_MODEL,
                    contents=prompt,
                    config=types.GenerateContentConfig(
                        temperature=0,
                        max_output_tokens=12288,
                        top_p=0.8,
                        thinking_config=types.ThinkingConfig(thinking_budget=0)
                    )
                )
                
                parsed_data = self.extract_json_from_response(response.text)
                if parsed_data:
                    validated_data = self.validate_mindmap_structure(parsed_data)
                    logger.info(f"Successfully generated mindmap with {len(validated_data.get('children', []))} main branches")
                    return validated_data
                else:
                    logger.warning("Failed to parse JSON, creating fallback")
                    return self.create_enhanced_fallback_mindmap(comprehensive_text)
            except Exception as api_error:
                logger.error(f"API Error: {api_error}")
                return self.create_enhanced_fallback_mindmap(comprehensive_text)
        except Exception as e:
            logger.error(f"Error generating mindmap: {e}")
            return self.create_enhanced_fallback_mindmap(text_chunks[0] if text_chunks else "No content")

    def extract_json_from_response(self, response_text):
        try:
            cleaned_text = response_text.strip()
            if cleaned_text.startswith("```json"):
                cleaned_text = cleaned_text[7:]
            elif cleaned_text.startswith("```"):
                cleaned_text = cleaned_text[3:]
            if cleaned_text.endswith("```"):
                cleaned_text = cleaned_text[:-3]
            cleaned_text = cleaned_text.strip()
            try:
                return json.loads(cleaned_text)
            except json.JSONDecodeError:
                pass
            json_pattern = r'\{(?:[^{}]|{(?:[^{}]|{[^{}]*})*})*\}'
            matches = re.findall(json_pattern, cleaned_text, re.DOTALL)
            for match in matches:
                try:
                    parsed = json.loads(match)
                    if isinstance(parsed, dict) and ('name' in parsed or 'children' in parsed):
                        return parsed
                except json.JSONDecodeError:
                    continue
            return None
        except Exception as e:
            logger.error(f"Error extracting JSON: {e}")
            return None

    def validate_mindmap_structure(self, data):
        if not isinstance(data, dict):
            return {"name": "Document", "children": []}
        if 'name' not in data:
            data['name'] = "Document Analysis"
        if 'children' in data:
            if not isinstance(data['children'], list):
                data['children'] = []
            else:
                validated_children = []
                for child in data['children']:
                    validated_child = self.validate_mindmap_structure(child)
                    if validated_child['name'].strip():
                        validated_children.append(validated_child)
                data['children'] = validated_children
        else:
            data['children'] = []
        allowed_fields = {'name', 'children', 'description', 'type', 'summary'}
        cleaned_data = {k: v for k, v in data.items() if k in allowed_fields}
        return cleaned_data

    def create_enhanced_fallback_mindmap(self, text):
        try:
            sentences = [s.strip() for s in text.split('.') if 15 < len(s.strip()) < 200]
            paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50]
            words = re.findall(r'\b[A-Za-z]{4,}\b', text.lower())
            word_freq = {}
            for word in words:
                if word not in ['these', 'those', 'which', 'where', 'there', 'should', 'would']:
                    word_freq[word] = word_freq.get(word, 0) + 1
            top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:15]
            children = [
                {
                    "name": "Document Overview",
                    "summary": "General overview of the document content",
                    "children": [
                        {"name": "Main Topics", "children": [], "summary": "Primary subjects covered"},
                        {"name": "Key Information", "children": [], "summary": "Important facts and data"},
                        {"name": "Document Structure", "children": [], "summary": "Organization of content"}
                    ]
                },
                {
                    "name": "Content Analysis",
                    "summary": "Detailed analysis of document content",
                    "children": [
                        {"name": "Core Concepts", "children": [], "summary": "Fundamental ideas presented"},
                        {"name": "Supporting Details", "children": [], "summary": "Evidence and examples"},
                        {"name": "Conclusions", "children": [], "summary": "Main findings and outcomes"}
                    ]
                }
            ]
            if top_keywords:
                term_children = []
                for keyword, freq in top_keywords[:8]:
                    term_children.append({
                        "name": keyword.title(),
                        "children": [],
                        "summary": f"Mentioned {freq} times in document"
                    })
                children.append({
                    "name": "Key Terms and Concepts",
                    "summary": "Important terminology from the document",
                    "children": term_children
                })
            root_name = "Document Analysis"
            if sentences:
                first_sentence = sentences[0]
                if len(first_sentence) < 100:
                    root_name = first_sentence
            return {
                "name": root_name,
                "children": children
            }
        except Exception as e:
            logger.error(f"Error in fallback mindmap: {e}")
            return {
                "name": "Document Analysis",
                "children": [
                    {
                        "name": "Content Overview",
                        "children": [
                            {"name": "Main Information", "children": []},
                            {"name": "Key Details", "children": []},
                            {"name": "Supporting Context", "children": []}
                        ]
                    }
                ]
            }


    def calculate_max_depth(self,node, current_depth=0):
        """Calculate the maximum depth of a mindmap structure."""
        if not isinstance(node, dict) or 'children' not in node:
            return current_depth
        
        children = node.get('children', [])
        if not children:
            return current_depth
        
        max_child_depth = current_depth
        for child in children:
            child_depth = self.calculate_max_depth(child, current_depth + 1)
            max_child_depth = max(max_child_depth, child_depth)
        
        return max_child_depth

    def create_deep_hierarchical_fallback(self,text, batch_number):
        """Create a deep hierarchical fallback mind map for a specific batch."""
        try:
            # Enhanced text analysis for deeper hierarchy
            sentences = [s.strip() for s in text.split('.') if len(s.strip()) > 15][:25]
            paragraphs = [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50][:15]
            
            # Extract keywords with frequency
            words = re.findall(r'\b[A-Za-z]{3,}\b', text.lower())
            word_freq = {}
            stop_words = {'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'was', 'one', 'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'man', 'new', 'now', 'old', 'see', 'two', 'way', 'who', 'boy', 'did', 'its', 'let', 'put', 'say', 'she', 'too', 'use'}
            
            for word in words:
                if word not in stop_words and len(word) > 3:
                    word_freq[word] = word_freq.get(word, 0) + 1
            
            top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:12]
            
            # Create deep hierarchical structure
            children = []
            
            # Group keywords into thematic categories for deeper hierarchy
            if top_keywords:
                # Create 3-4 main categories with deep sub-structures
                keyword_groups = [
                    top_keywords[0:3],   # Primary concepts
                    top_keywords[3:6],   # Secondary concepts  
                    top_keywords[6:9],   # Supporting concepts
                    top_keywords[9:12]   # Additional concepts
                ]
                
                category_names = [
                    "Primary Concepts",
                    "Secondary Concepts", 
                    "Supporting Elements",
                    "Additional Information"
                ]
                
                for i, (category_name, keyword_group) in enumerate(zip(category_names, keyword_groups)):
                    if not keyword_group:
                        continue
                        
                    category_children = []
                    for keyword, freq in keyword_group:
                        # Create deep structure for each keyword
                        keyword_node = {
                            "name": keyword.title().replace('_', ' '),
                            "children": [
                                {
                                    "name": "Definition & Context",
                                    "children": [
                                        {"name": "Core Meaning", "children": [], "summary": f"Primary definition and significance of {keyword}"},
                                        {"name": "Usage Context", "children": [], "summary": f"How {keyword} is used in this section"}
                                    ]
                                },
                                {
                                    "name": "Details & Examples",
                                    "children": [
                                        {"name": "Specific Details", "children": [], "summary": f"Detailed information about {keyword}"},
                                        {"name": "Related Examples", "children": [], "summary": f"Examples and instances of {keyword} mentioned {freq} times"}
                                    ]
                                }
                            ]
                        }
                        category_children.append(keyword_node)
                    
                    if category_children:
                        children.append({
                            "name": category_name,
                            "children": category_children
                        })
            
            # Add content structure analysis
            if sentences:
                content_structure = {
                    "name": "Content Structure",
                    "children": [
                        {
                            "name": "Main Information",
                            "children": [
                                {
                                    "name": "Key Statements",
                                    "children": [
                                        {"name": "Primary Points", "children": [], "summary": "Main ideas expressed in this section"},
                                        {"name": "Supporting Points", "children": [], "summary": "Additional supporting information"}
                                    ]
                                },
                                {
                                    "name": "Detailed Content",
                                    "children": [
                                        {"name": "Explanations", "children": [], "summary": "Detailed explanations and descriptions"},
                                        {"name": "Specifications", "children": [], "summary": "Specific details and technical information"}
                                    ]
                                }
                            ]
                        },
                        {
                            "name": "Organization",
                            "children": [
                                {"name": "Information Flow", "children": [], "summary": "How information is organized in this section"},
                                {"name": "Topic Progression", "children": [], "summary": "Sequence and development of topics"}
                            ]
                        }
                    ]
                }
                children.append(content_structure)
            
            # Determine root name based on content
            root_name = f"Section {batch_number} Analysis"
            if sentences and len(sentences[0]) < 100:
                clean_sentence = re.sub(r'[^\w\s]', '', sentences[0])
                if len(clean_sentence) > 10:
                    root_name = clean_sentence[:70] + "..." if len(clean_sentence) > 70 else clean_sentence
            
            return {
                "name": root_name,
                "children": children if children else [
                    {
                        "name": "Content Overview",
                        "children": [
                            {
                                "name": "Section Information",
                                "children": [
                                    {"name": "Main Content", "children": [], "summary": f"Primary information from section {batch_number}"},
                                    {"name": "Supporting Details", "children": [], "summary": f"Additional details from section {batch_number}"}
                                ]
                            }
                        ]
                    }
                ]
            }
            
        except Exception as e:
            print(f"Error in deep hierarchical fallback for batch {batch_number}: {e}")
            return {
                "name": f"Section {batch_number} Content",
                "children": [
                    {
                        "name": "Document Content",
                        "children": [
                            {
                                "name": "Information",
                                "children": [
                                    {"name": "Content Details", "children": [], "summary": f"Content from batch {batch_number}"}
                                ]
                            }
                        ]
                    }
                ]
            }

    def normalize_topic_name(name):
        """Normalize topic name for comparison - case insensitive, strip whitespace, handle variations."""
        if not name or not isinstance(name, str):
            return ""
        
        # Convert to lowercase and strip whitespace
        normalized = name.strip().lower()
        
        # Remove common variations in formatting
        normalized = re.sub(r'\s+', ' ', normalized)  # Multiple spaces to single space
        normalized = re.sub(r'[^\w\s]', '', normalized)  # Remove special characters
        
        return normalized

    def find_matching_topic(self,topic_name, existing_children):
        """Find if a topic already exists in the children list with fuzzy matching."""
        normalized_new = self.normalize_topic_name(topic_name)
        
        for i, child in enumerate(existing_children):
            if isinstance(child, dict) and 'name' in child:
                normalized_existing = self.normalize_topic_name(child['name'])
                
                # Exact match
                if normalized_new == normalized_existing:
                    return i, child
                
                # Fuzzy match - check if one contains the other (for similar topics)
                if len(normalized_new) > 3 and len(normalized_existing) > 3:
                    if normalized_new in normalized_existing or normalized_existing in normalized_new:
                        # Additional check - ensure they're actually related topics
                        words_new = set(normalized_new.split())
                        words_existing = set(normalized_existing.split())
                        
                        # If they share significant words, consider them a match
                        if len(words_new.intersection(words_existing)) >= min(len(words_new), len(words_existing)) * 0.6:
                            return i, child
        
        return None, None

    def merge_hierarchical_topics(self,existing_node, new_node, depth=0):
        """
        Advanced hierarchical merging that handles deep topic structures:
        topic -> subtopic -> sub-subtopic -> continue...
        
        Logic:
        1. If topic exists, merge subtopics into it
        2. If topic doesn't exist, add it completely
        3. Recursively apply this at all levels
        """
        try:
            if not existing_node or not isinstance(existing_node, dict):
                return new_node if new_node else {}
            
            if not new_node or not isinstance(new_node, dict):
                return existing_node
            
            print(f"{'  ' * depth}Merging at depth {depth}: '{existing_node.get('name', 'Unknown')}' + '{new_node.get('name', 'Unknown')}'")
            
            # Start with existing node as base
            merged_node = existing_node.copy()
            
            # Update summary if new node has better content
            if 'summary' in new_node:
                if 'summary' not in merged_node or not merged_node.get('summary') or len(new_node['summary']) > len(merged_node.get('summary', '')):
                    merged_node['summary'] = new_node['summary']
            
            # Get children from both nodes
            existing_children = existing_node.get('children', [])
            new_children = new_node.get('children', [])
            
            if not new_children:
                return merged_node
            
            # Start with copy of existing children
            merged_children = existing_children.copy()
            
            # Process each new child
            for new_child in new_children:
                if not isinstance(new_child, dict) or 'name' not in new_child:
                    continue
                
                # Try to find matching existing topic
                match_index, matching_child = self.find_matching_topic(new_child['name'], merged_children)
                
                if matching_child is not None:
                    # Topic already exists - merge the subtopics recursively
                    print(f"{'  ' * depth}  ↳ Merging subtopics into existing topic: '{matching_child['name']}'")
                    merged_children[match_index] = self.merge_hierarchical_topics(matching_child, new_child, depth + 1)
                    
                    # Count new subtopics added
                    existing_subtopic_count = len(matching_child.get('children', []))
                    new_subtopic_count = len(merged_children[match_index].get('children', []))
                    if new_subtopic_count > existing_subtopic_count:
                        print(f"{'  ' * depth}    ✓ Added {new_subtopic_count - existing_subtopic_count} new subtopics")
                        
                else:
                    # Topic doesn't exist - add it completely with all its hierarchy
                    merged_children.append(new_child)
                    subtopic_count = self.count_mindmap_nodes(new_child) - 1  # Subtract 1 for the topic itself
                    print(f"{'  ' * depth}  ✓ Added new topic: '{new_child['name']}' (with {subtopic_count} subtopics)")
            
            merged_node['children'] = merged_children
            return merged_node
            
        except Exception as e:
            print(f"Error in hierarchical merging at depth {depth}: {e}")
            return existing_node if existing_node else new_node

    def generate_mindmap_heading(self, selected_documents=None, document_info=None, api_key=None):
        """
        Use GPT-4o API to generate a meaningful mindmap heading based on selected document filenames.
        Uses only the filenames of selected documents to create a descriptive title.
        """
        try:
            import openai
            load_dotenv()
            api_key = os.getenv("OPENAI_API_KEY")
            # Set API key
            if api_key:
                openai.api_key = api_key
            else:
                openai.api_key = OPENAI_API_KEY
            
            # Extract raw filenames - no cleaning, let LLM handle everything
            filenames = []
            if document_info and isinstance(document_info, list):
                filenames = [doc.get('filename', '') for doc in document_info if doc.get('filename')]
            elif selected_documents and isinstance(selected_documents, list):
                filenames = [f"Document_{doc_id}" for doc_id in selected_documents]
            
            if not filenames:
                filenames = ["untitled_document"]
            
            # Just pass raw filenames to LLM
            filenames_text = "\n".join([f"- {filename}" for filename in filenames])
            
            # Let GPT-4o do everything - but MUST stay related to document names
            if len(filenames) == 1:
                prompt = f"""Generate a professional 7-8 word mindmap heading that is DIRECTLY RELATED to this document filename. The heading must reflect the actual content suggested by the filename.
   
    Document filename: {filenames[0]}
   
    CRITICAL REQUIREMENTS:
    - The heading MUST be directly related to what this filename suggests
    - Use keywords and concepts from the actual filename
    - Create EXACTLY 7 or 8 words
    - Make it professional but stay true to the document's apparent topic
    - If filename mentions "sales" - heading must be about sales
    - If filename mentions "marketing" - heading must be about marketing
    - If filename mentions "technical" - heading must be about technical topics
    - DO NOT create generic headings unrelated to the filename
    - Extract the main subject from the filename and build around it
   
    Generate the professional mindmap heading that reflects this specific document:"""
   
            else:
                prompt = f"""Generate a professional 7-8 word mindmap heading that is DIRECTLY RELATED to these specific document filenames. The heading must reflect the actual content suggested by ALL these filenames.
   
    Document filenames:
    {filenames_text}
   
    CRITICAL REQUIREMENTS:
    - The heading MUST be directly related to what these filenames suggest
    - Use actual keywords and concepts from the filenames themselves
    - Find the real common themes from the actual filenames (not generic themes)
    - Create EXACTLY 7 or 8 words
    - If filenames are about "budget" and "finance" - heading must be about financial topics
    - If filenames are about "user guide" and "manual" - heading must be about documentation
    - If filenames are about "research" and "data" - heading must be about research/data
    - DO NOT create generic headings unrelated to the actual filename content
    - The heading should make someone think "yes, this is exactly about those documents"
   
    Generate the professional mindmap heading that reflects these specific documents:"""
   
            # Call GPT-4o API with emphasis on staying document-relevant
            response = openai.chat.completions.create(
                model="gpt-4o",  # Latest GPT-4o model
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert at creating professional titles that are DIRECTLY related to document content. Never create generic titles. Always extract and use the actual subject matter from the filenames provided. The title must make it obvious what the documents are about."
                    },
                    {
                        "role": "user",
                        "content": prompt
                    }
                ],
                max_tokens=50,
                temperature=0.6,  # Slightly lower to stay more focused on actual content
                top_p=0.85,
                frequency_penalty=0.3,
                presence_penalty=0.2
            )
            
            # Extract the generated title - minimal processing
            title = response.choices[0].message.content.strip()
            
            # Basic word count check - if wrong, try once more
            word_count = len(title.split())
            if word_count < 7 or word_count > 8:
                # Retry with more specific instruction about staying document-relevant
                retry_response = openai.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {
                            "role": "system",
                            "content": "Generate exactly 7 words that MUST relate to the specific document content. No generic titles allowed."
                        },
                        {
                            "role": "user",
                            "content": f"Based on these specific filenames: {', '.join(filenames[:3])}, create exactly 7 words that someone would recognize as being about these exact documents. Use actual keywords from the filenames."
                        }
                    ],
                    max_tokens=30,
                    temperature=0.4  # Lower temperature for more document-focused results
                )
                retry_title = retry_response.choices[0].message.content.strip()
                retry_word_count = len(retry_title.split())
                if 7 <= retry_word_count <= 8:
                    title = retry_title
            
            logger.info(f"Generated {len(title.split())}-word heading: '{title}' from {len(filenames)} files")
            return title
            
        except Exception as e:
            logger.error(f"Error generating mindmap heading with GPT-4o: {e}")
            # Simple fallback without any processing
            if len(filenames) == 1:
                return f"Professional Analysis of {filenames[0].split('.')[0]} Research Framework"
            else:
                return f"Integrated Analysis of {len(filenames)} Document Research Collection"


# In your backend views.py, update the MindMapQuestionView:

class MindMapQuestionView(APIView):
    """Generate question for mindmap node - matches your 'mindmap-question/' endpoint"""
    permission_classes = [IsAuthenticated]
   
    def post(self, request):
        try:
            user = request.user
            main_project_id = request.data.get('main_project_id')
            if not main_project_id:
                return Response({'error': 'Main project ID is required'}, status=status.HTTP_400_BAD_REQUEST)
            mindmap_id = request.data.get('mindmap_id')
            topic_name = request.data.get('topic_name', '')
            topic_summary = request.data.get('topic_summary', '')
            node_path = request.data.get('node_path', '')
            selected_documents = request.data.get('selected_documents', [])
            target_user_id = request.data.get('target_user_id')
            force_new_context = request.data.get('force_new_context', False)
            current_timestamp = request.data.get('current_timestamp')
            mindmap_document_sources = request.data.get('mindmap_document_sources', [])

            # Handle admin operation for another user
            if target_user_id and request.user.username == 'admin':
                try:
                    user = User.objects.get(id=target_user_id)
                except User.DoesNotExist:
                    return Response({
                        'error': 'Target user not found'
                    }, status=status.HTTP_404_NOT_FOUND)

            if not all([main_project_id, topic_name]):
                return Response({
                    'error': 'Missing required parameters'
                }, status=status.HTTP_400_BAD_REQUEST)

            # If multiple documents are selected, combine their context for question generation
            if selected_documents:
                valid_docs = Document.objects.filter(
                    id__in=selected_documents,
                    user=user
                ).count()
                if valid_docs != len(selected_documents):
                    valid_doc_ids = list(Document.objects.filter(
                        id__in=selected_documents,
                        user=user
                    ).values_list('id', flat=True))
                    selected_documents = [str(doc_id) for doc_id in valid_doc_ids]

            # Generate specific question for the topic using combined context
            question = self.generate_topic_question(
                topic_name,
                topic_summary,
                node_path,
                mindmap_id=mindmap_id,
                document_context=selected_documents
            )

            return Response({
                'success': True,
                'question': question,
                'topic': topic_name,
                'node_path': node_path,
                'mindmap_id': mindmap_id,
                'selected_documents': selected_documents,
                'timestamp': current_timestamp or int(time.time() * 1000)
            }, status=status.HTTP_200_OK)

        except Exception as e:
            logger.error(f"Error in MindMapQuestionView: {str(e)}")
            return Response({
                'error': f'Internal server error: {str(e)}',
                'success': False
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def generate_topic_question(self, topic_name, topic_summary, node_path, mindmap_id=None, document_context=None):
        """Generate a specific question for the mindmap topic with enhanced context."""
        try:
            time.sleep(1.0)  # Rate limiting
            
            # ✅ NEW: Enhanced prompt with document context awareness
            context_info = ""
            if document_context and len(document_context) > 0:
                context_info = f"The user has selected {len(document_context)} specific document(s) for this question context. "
            
            prompt = f"""
                Generate ONE simple and helpful question about "{topic_name}" that will make it easier for users to understand this topic.

                Topic: {topic_name}  
                Summary: {topic_summary}  
                Context Path: {node_path}
                Mindmap ID: {mindmap_id}
                {context_info}

                Create a question that:
                1. Uses simple language  
                2. Helps someone understand the topic better  
                3. Is 10–25 words long  
                4. Focuses on practical use or basic understanding
                5. {"Should be answerable from the selected documents" if document_context else "Can be answered from general knowledge"}

                Examples of good questions:
                - "How is [topic] used in real life?"  
                - "What is the main idea behind [topic]?"  
                - "Why is [topic] important or useful?"

                Generate ONE question:
            """
           
            response = mindmap_client.models.generate_content(
                model=GEMINI_MODEL,
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.2,
                    max_output_tokens=512,
                    thinking_config=types.ThinkingConfig(thinking_budget=0)
                )
            )
           
            question = response.text.strip()
           
            # Clean the question
            if question.startswith('"') and question.endswith('"'):
                question = question[1:-1]
                
            # ✅ NEW: Log the generated question with context
            logger.info(f"Generated question for mindmap {mindmap_id}: {question}")
            logger.info(f"Question context: {len(document_context or [])} documents")
           
            return question
           
        except Exception as e:
            logger.error(f"Error generating topic question: {e}")
            return f"What are the key aspects and details about '{topic_name}' in this document?"




# Function-based views to match your existing URL structure
@api_view(['GET'])
def get_user_mindmaps(request):
    """Get user's mindmaps for a project - matches your 'user-mindmaps/' endpoint"""
    try:
        user = request.user
        main_project_id = request.GET.get('main_project_id')
        target_user_id = request.GET.get('target_user_id')
        
        # Handle admin viewing mindmaps for another user
        if target_user_id and request.user.username == 'admin':
            try:
                user = User.objects.get(id=target_user_id)
            except User.DoesNotExist:
                return Response({
                    'error': 'Target user not found'
                }, status=status.HTTP_404_NOT_FOUND)
       
        if not main_project_id:
            return Response({
                'error': 'Main project ID is required'
            }, status=status.HTTP_400_BAD_REQUEST)
       
        mindmaps = MindMap.objects.filter(
            user=user,
            main_project_id=main_project_id
        ).order_by('-created_at')
       
        mindmap_list = []
        for mindmap in mindmaps:
            mindmap_list.append({
                'id': mindmap.id,
                'created_at': mindmap.created_at.isoformat(),
                'updated_at': mindmap.updated_at.isoformat(),
                'total_nodes': mindmap.total_nodes,
                'document_sources': mindmap.get_document_sources_list(),
                'title': f"Mindmap - {mindmap.created_at.strftime('%b %d, %Y %H:%M')}",
                'preview': mindmap.data.get('name', 'Untitled Mindmap') if mindmap.data else 'Untitled Mindmap'
            })
       
        return Response({
            'success': True,
            'mindmaps': mindmap_list,
            'count': len(mindmap_list)
        }, status=status.HTTP_200_OK)
       
    except Exception as e:
        logger.error(f"Error fetching mindmaps: {str(e)}")
        return Response({
            'error': f'Failed to fetch mindmaps: {str(e)}',
            'success': False
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET', 'DELETE'])
def get_mindmap_data(request, mindmap_id):
    """Get or delete specific mindmap data - matches your 'mindmap/<id>/' endpoint"""
    try:
        user = request.user
        
        try:
            mindmap = MindMap.objects.get(id=mindmap_id, user=user)
        except MindMap.DoesNotExist:
            return Response({
                'error': 'Mindmap not found',
                'success': False
            }, status=status.HTTP_404_NOT_FOUND)
        
        if request.method == 'GET':
            # CRITICAL: Get document IDs for the mindmap's source documents
            document_ids = []
            document_sources = mindmap.get_document_sources_list()
            
            print(f"🔍 Backend: Processing mindmap {mindmap_id} with sources: {document_sources}")
            
            if document_sources:
                try:
                    # Import Document model at the top of your file if not already imported
                    # from your_app.models import Document
                    
                    documents = Document.objects.filter(
                        user=user,
                        filename__in=document_sources
                    ).values('id', 'filename')
                    
                    document_ids = [str(doc['id']) for doc in documents]
                    
                    print(f"✅ Backend: Found {len(document_ids)} matching documents for mindmap {mindmap_id}")
                    print(f"📄 Backend: Document IDs: {document_ids}")
                    
                except Exception as e:
                    print(f"❌ Backend: Error getting document IDs for mindmap {mindmap_id}: {str(e)}")
                    logger.error(f"Error getting document IDs for mindmap {mindmap_id}: {str(e)}")
            else:
                print(f"⚠️ Backend: No document sources found for mindmap {mindmap_id}")
            
            response_data = {
                'success': True,
                'mindmap': mindmap.data,
                'mindmap_id': mindmap.id,
                'stats': {
                    'total_nodes': mindmap.total_nodes,
                    'created_at': mindmap.created_at.isoformat(),
                    'updated_at': mindmap.updated_at.isoformat(),
                    'document_sources': document_sources,
                    'document_ids': document_ids,  # CRITICAL: Include document IDs
                    'documents_processed': len(document_sources),
                    'mindmap_nodes': mindmap.total_nodes
                }
            }
            
            print(f"📤 Backend: Sending response with {len(document_ids)} document IDs")
            
            return Response(response_data, status=status.HTTP_200_OK)
            
        elif request.method == 'DELETE':
            mindmap.delete()
            logger.info(f"Deleted mindmap {mindmap_id} for user {user.username}")
            
            return Response({
                'success': True,
                'message': 'Mindmap deleted successfully'
            }, status=status.HTTP_200_OK)
            
    except Exception as e:
        print(f"❌ Backend: Error with mindmap {mindmap_id}: {str(e)}")
        logger.error(f"Error with mindmap {mindmap_id}: {str(e)}")
        return Response({
            'error': f'Failed to process mindmap: {str(e)}',
            'success': False
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

from django.db.models import Sum, Count
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.permissions import IsAuthenticated
from django.contrib.auth.models import User
from .models import Document, ChatMessage

class AdminNotebookUserStatsView(APIView):
    permission_classes = [IsAuthenticated]
 
    def get(self, request):
        # Only allow admin user
        if not request.user.username == 'admin':
            return Response({'error': 'Unauthorized'}, status=403)
 
        users = User.objects.all()
        stats = []
        total_short = 0
        total_comprehensive = 0
 
        for user in users:
            docs = Document.objects.filter(user=user)
            doc_stats = []
            for doc in docs:
                from .models import ChatHistory
                conversations = ChatHistory.objects.filter(user=user, documents=doc)
                question_count = ChatMessage.objects.filter(
                    chat_history__in=conversations,
                    role='user'
                ).count()
                # Count assistant messages by response_length
                short_count = ChatMessage.objects.filter(
                    chat_history__in=conversations,
                    role='assistant',
                    response_length='short'
                ).count()
                comprehensive_count = ChatMessage.objects.filter(
                    chat_history__in=conversations,
                    role='assistant',
                    response_length='comprehensive'
                ).count()
                total_short += short_count
                total_comprehensive += comprehensive_count
                doc_stats.append({
                    'document_id': doc.id,
                    'filename': doc.filename,
                    'questions_asked': question_count,
                    'short_responses': short_count,
                    'comprehensive_responses': comprehensive_count,
                })
 
            doc_upload_count = docs.count()
 
            # --- Add aggregate stats for user ---
            # Token and question stats
            from notebook.models import ConversationTransaction, DocumentTransaction
            conv_stats = ConversationTransaction.objects.filter(
                user_transaction__user=user
            ).aggregate(
                total_tokens=Sum('total_tokens_used'),
                total_input_tokens=Sum('input_api_tokens'),
                total_output_tokens=Sum('output_api_tokens'),
                total_questions=Sum('question_count'),
            )
            doc_stats_agg = DocumentTransaction.objects.filter(
                user_transaction__user=user
            ).aggregate(
                total_pages=Sum('no_pages'),
                total_documents=Count('id'),
            )
 
            stats.append({
                'user_id': user.id,
                'username': user.username,
                'document_upload_count': doc_upload_count,
                'documents': doc_stats,
                'total_tokens_used': conv_stats['total_tokens'] or 0,
                'total_input_tokens': conv_stats['total_input_tokens'] or 0,
                'total_output_tokens': conv_stats['total_output_tokens'] or 0,
                'total_questions_asked': conv_stats['total_questions'] or 0,
                'total_pages_processed': doc_stats_agg['total_pages'] or 0,
                'total_documents_uploaded': doc_stats_agg['total_documents'] or 0,
            })
 
        return Response({
            'user_stats': stats,
            'total_short_responses': total_short,
            'total_comprehensive_responses': total_comprehensive
        }, status=200)
    
 
class GenerateIdeaContextView(APIView):
    """
    API endpoint to extract structured idea generation parameters
    from documents or query results using pgvector.
    """
 
    def post(self, request):
        user = request.user
        document_id = request.data.get('document_id')
        query = request.data.get('query')
        main_project_id = request.data.get('main_project_id')
       
        # Validate input - need either document_id or query
        if not document_id and not query:
            return Response({
                'error': 'Either document_id or query parameter is required'
            }, status=status.HTTP_400_BAD_REQUEST)
           
        try:
            # Case 1: Using Document ID - fetch existing parameters or extract new ones
            if document_id:
                document = get_object_or_404(Document, id=document_id, user=user)
               
                # Get document name without extension
                document_name_no_ext = self.remove_file_extension(document.filename)
               
                # Generate a unique project name - handle the case when main_project_id is None
                suggested_project_name = f"Ideas from {document_name_no_ext}"
                if main_project_id:
                    try:
                        suggested_project_name = self.generate_unique_project_name(document_name_no_ext, main_project_id)
                    except Exception as e:
                        # Log the error but continue with the default name
                        print(f"Error generating unique project name: {str(e)}")
               
                try:
                    # Check if we already have parameters stored
                    processed_index = ProcessedIndex.objects.get(document=document)
                   
                    # If parameters exist, return them
                    if processed_index.idea_parameters:
                        return Response({
                            'document_id': document_id,
                            'document_name': document.filename,
                            'document_name_no_ext': document_name_no_ext,
                            'idea_parameters': processed_index.idea_parameters,
                            'suggested_project_name': suggested_project_name
                        })
                   
                    # If no parameters yet, extract them from the document using pgvector
                    # First check if the document has a markdown path (LlamaParse document)
                    if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                        # This is a LlamaParse document, read the markdown content directly
                        try:
                            with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                                full_text = f.read()
                               
                            # Extract parameters from markdown content
                            idea_params = self.extract_idea_parameters(full_text)
                           
                            # Save the parameters for future use
                            processed_index.idea_parameters = idea_params
                            processed_index.save()
                           
                            return Response({
                                'document_id': document_id,
                                'document_name': document.filename,
                                'document_name_no_ext': document_name_no_ext,
                                'idea_parameters': idea_params,
                                'suggested_project_name': suggested_project_name
                            })
                        except Exception as e:
                            print(f"Error reading markdown file: {str(e)}")
                            # Continue with pgvector approach as fallback
                   
                    # Load document content from pgvector embeddings
                    document_embeddings = DocumentEmbedding.objects.filter(
                        document=document
                    ).order_by('chunk_id')
                   
                    if not document_embeddings.exists():
                        return Response({
                            'error': 'No content found in the document. Document may not be processed yet.'
                        }, status=status.HTTP_400_BAD_REQUEST)
                   
                    # Extract parameters from document content
                    full_text = " ".join([embedding.content for embedding in document_embeddings])
                    idea_params = self.extract_idea_parameters(full_text)
                   
                    # Save the parameters for future use
                    processed_index.idea_parameters = idea_params
                    processed_index.save()
 
                    if main_project_id:
                        update_project_timestamp(main_project_id, user)
                   
                    return Response({
                        'document_id': document_id,
                        'document_name': document.filename,
                        'document_name_no_ext': document_name_no_ext,
                        'idea_parameters': idea_params,
                        'suggested_project_name': suggested_project_name
                    })
                   
                except ProcessedIndex.DoesNotExist:
                    return Response({
                        'error': 'Document has not been processed yet'
                    }, status=status.HTTP_404_NOT_FOUND)
           
            # Case 2: Using Query - search across documents and extract from relevant chunks
            else:
                # Get active/selected document
                active_doc_id = request.session.get('active_document_id')
                if not active_doc_id:
                    return Response({
                        'error': 'No active document selected'
                    }, status=status.HTTP_400_BAD_REQUEST)
               
                document = get_object_or_404(Document, id=active_doc_id, user=user)
                processed_index = get_object_or_404(ProcessedIndex, document=document)
               
                # Get document name without extension
                document_name_no_ext = self.remove_file_extension(document.filename)
               
                # Generate a unique project name - handle the case when main_project_id is None
                suggested_project_name = f"Ideas from {document_name_no_ext}"
                if main_project_id:
                    try:
                        suggested_project_name = self.generate_unique_project_name(document_name_no_ext, main_project_id)
                    except Exception as e:
                        # Log the error but continue with the default name
                        print(f"Error generating unique project name: {str(e)}")
               
                # First check if the document has a markdown path (LlamaParse document)
                if processed_index.markdown_path and os.path.exists(processed_index.markdown_path):
                    # This is a LlamaParse document, read the markdown content directly
                    try:
                        with open(processed_index.markdown_path, 'r', encoding='utf-8') as f:
                            full_text = f.read()
                           
                        # Extract parameters from markdown content with query context
                        idea_params = self.extract_idea_parameters(query, None, full_text)
                       
                        return Response({
                            'document_id': active_doc_id,
                            'document_name': document.filename,
                            'document_name_no_ext': document_name_no_ext,
                            'query': query,
                            'idea_parameters': idea_params,
                            'suggested_project_name': suggested_project_name
                        })
                    except Exception as e:
                        print(f"Error reading markdown file: {str(e)}")
                        # Continue with pgvector approach as fallback
               
                # Use pgvector similarity search to find relevant chunks
                relevant_chunks = self.search_pgvector_similarity(document, query, k=5)
               
                if not relevant_chunks:
                    return Response({
                        'error': 'No relevant content found for the query'
                    }, status=status.HTTP_400_BAD_REQUEST)
               
                # Extract parameters from relevant chunks
                idea_params = self.extract_idea_parameters(query, relevant_chunks)
 
                if main_project_id:
                    update_project_timestamp(main_project_id, user)
               
                return Response({
                    'document_id': active_doc_id,
                    'document_name': document.filename,
                    'document_name_no_ext': document_name_no_ext,
                    'query': query,
                    'idea_parameters': idea_params,
                    'suggested_project_name': suggested_project_name
                })
               
        except Exception as e:
            print(f"Error generating idea context: {str(e)}")
            return Response({
                'error': str(e),
                'detail': 'An error occurred while generating idea context'
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
   
    def remove_file_extension(self, filename):
        """
        Remove file extension from filename
        """
        import os
        return os.path.splitext(filename)[0]
   
    def generate_unique_project_name(self, document_name, main_project_id):
        """
        Generate a unique project name based on document name,
        adding (1), (2), etc. if needed to avoid duplicates
        """
        from ideaGen.models import Project
       
        base_name = f"Ideas from {document_name}"
        project_name = base_name
        counter = 1
       
        # Check for existing projects with this name in the main project
        while Project.objects.filter(
            name=project_name,
            main_project_id=main_project_id
        ).exists():
            # Increment counter and update name
            project_name = f"{base_name} ({counter})"
            counter += 1
           
        return project_name
   
    def get_embeddings(self, texts):
        """
        Get embeddings for a list of texts using OpenAI
        This method should match the one used in your document processing
        """
        from openai import OpenAI
        client = OpenAI()
       
        try:
            if not texts:
                return []
           
            # Ensure texts is a list
            if isinstance(texts, str):
                texts = [texts]
           
            response = client.embeddings.create(
                input=texts,
                model="text-embedding-3-small"
            )
           
            return [data.embedding for data in response.data]
           
        except Exception as e:
            print(f"Error getting embeddings: {str(e)}")
            return []
   
    def get_query_embedding(self, query):
        """Get embedding for the query using OpenAI"""
        from openai import OpenAI
        client = OpenAI()  # Initialize the client
       
        try:
            response = client.embeddings.create(
                input=[query],
                model="text-embedding-3-small"
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"Error getting embedding for query: {str(e)}")
            raise
   
    def search_pgvector_similarity(self, document, query, k=10):
        """
        Enhanced search function using pgvector for similarity search
        Based on the QnA search_similar_content implementation
       
        Args:
            document: The document object to search within
            query: The search query
            k: Number of top results to return
           
        Returns:
            list: List of relevant chunks with their content
        """
        print(f"\n🔍 IDEA SEARCH DEBUG: Starting pgvector search for query: '{query}'")
        print(f"🔍 IDEA SEARCH DEBUG: Document: {document.filename}")
 
        # Get embeddings for the query
        query_embedding = self.get_embeddings([query])
        if not query_embedding:
            print("❌ IDEA SEARCH DEBUG: Failed to get query embedding")
            return []
 
        print("✅ IDEA SEARCH DEBUG: Got query embedding")
 
        try:
            # Check if embeddings exist for this document
            embeddings_exist = DocumentEmbedding.objects.filter(document=document).exists()
            if not embeddings_exist:
                print(f"❌ IDEA SEARCH DEBUG: No embeddings found for {document.filename}")
                return []
 
            # Perform pgvector similarity search
            from pgvector.django import CosineDistance
 
            # Query embeddings using pgvector
            similar_embeddings = DocumentEmbedding.objects.filter(
                document=document
            ).annotate(
                distance=CosineDistance('embedding', query_embedding[0])
            ).order_by('distance')[:k]
 
            print(f"✅ IDEA SEARCH DEBUG: pgvector search returned {len(similar_embeddings)} results")
 
            # Process results
            relevant_chunks = []
            seen_content_hashes = set()
           
            for embedding_obj in similar_embeddings:
                distance = float(embedding_obj.distance)
               
                # Filter by distance threshold (cosine distance, lower is better)
                if distance < 0.8:
                    content = embedding_obj.content
                   
                    if content and content.strip():
                        # Duplicate detection
                        content_hash = hash(content[:150])
                       
                        if content_hash not in seen_content_hashes:
                            seen_content_hashes.add(content_hash)
                           
                            relevant_chunks.append({
                                'text': content,
                                'chunk_id': embedding_obj.chunk_id,
                                'source': embedding_obj.source or document.filename,
                                'source_file': embedding_obj.source_file or document.filename,
                                'distance': distance,
                                'document_id': document.id
                            })
                           
                            print(f"   ✅ Added result: distance={distance:.4f}, content='{content[:100]}...'")
 
            print(f"✅ IDEA SEARCH DEBUG: Added {len(relevant_chunks)} relevant chunks")
 
            # If no results with similarity search, try text search fallback
            if not relevant_chunks:
                print(f"🔄 IDEA SEARCH DEBUG: No similarity results, trying text search fallback")
                try:
                    query_lower = query.lower()
                   
                    text_matches = DocumentEmbedding.objects.filter(
                        document=document,
                        content__icontains=query_lower
                    )[:k]
                   
                    print(f"   Found {len(text_matches)} text matches")
                   
                    for embedding_obj in text_matches:
                        content = embedding_obj.content
                        content_hash = hash(content[:150])
                       
                        if content_hash not in seen_content_hashes:
                            seen_content_hashes.add(content_hash)
                           
                            relevant_chunks.append({
                                'text': content,
                                'chunk_id': embedding_obj.chunk_id,
                                'source': embedding_obj.source or document.filename,
                                'source_file': embedding_obj.source_file or document.filename,
                                'distance': 0.5,  # Default relevance score for text matches
                                'document_id': document.id
                            })
                           
                            print(f"   ✅ Added text match: '{content[:100]}...'")
                           
                except Exception as fallback_error:
                    print(f"❌ IDEA SEARCH DEBUG: Fallback search failed: {str(fallback_error)}")
 
            return relevant_chunks
 
        except Exception as e:
            print(f"❌ IDEA SEARCH DEBUG: Error in pgvector similarity search: {str(e)}")
           
            # Final fallback: return first k chunks
            try:
                print(f"🔄 IDEA SEARCH DEBUG: Using final fallback - first {k} chunks")
                all_embeddings = DocumentEmbedding.objects.filter(
                    document=document
                ).order_by('chunk_id')[:k]
               
                fallback_chunks = []
                for emb in all_embeddings:
                    fallback_chunks.append({
                        'text': emb.content,
                        'chunk_id': emb.chunk_id,
                        'source': emb.source or document.filename,
                        'source_file': emb.source_file or document.filename,
                        'distance': 1.0,  # High distance indicates low relevance
                        'document_id': document.id
                    })
               
                print(f"✅ IDEA SEARCH DEBUG: Fallback returned {len(fallback_chunks)} chunks")
                return fallback_chunks
               
            except Exception as fallback_error:
                print(f"❌ IDEA SEARCH DEBUG: Final fallback also failed: {str(fallback_error)}")
                return []
   
    def extract_idea_parameters(self, context, relevant_chunks=None, full_text=None):
        """
        Extract structured parameters for the Idea Generator
       
        Args:
            context (str): Either full document content or query
            relevant_chunks (list, optional): List of relevant chunks from search
            full_text (str, optional): For LlamaParse documents, the full markdown text
           
        Returns:
            dict: Structured parameters for idea generation
        """
        from openai import OpenAI
        client = OpenAI()  # Initialize the client
        import json
       
        try:
            # Determine extraction context source
            if full_text:
                # If we have full text (e.g., from markdown), use that
                extraction_context = full_text
            elif relevant_chunks and len(relevant_chunks) > 0:
                # If relevant chunks are provided, use them for extraction context
                extraction_context = "\n\n".join([chunk.get('text', '') for chunk in relevant_chunks])
            else:
                # If no chunks available, use the context directly
                extraction_context = context
               
            # Create extraction prompt
            extraction_prompt = f"""
            Extract the following key parameters from the provided context.
            If a parameter isn't explicitly mentioned, infer it from context or leave it blank.
           
            Context:
            {extraction_context}
           
            Extract these parameters:
            - Brand_Name: The company or product brand mentioned
            - Category: The product category or market area
            - Concept: The main idea, campaign, or product concept
            - Benefits: Key benefits or value propositions (list up to 3)
            - RTB: Reason to Believe - evidence supporting the benefits (list up to 3)
            - Ingredients: Key ingredients or components (if applicable)
            - Features: Notable product/service features (list up to 3)
            - Theme: Overall theme or tone
            - Demographics: Target audience demographics
           
            Format the response as a valid JSON object with these fields.
            """
           
            # Call LLM to extract parameters
            response = client.chat.completions.create(
                model="gpt-4o",  
                messages=[
                    {"role": "system", "content": "You are a specialist in extracting structured information from documents."},
                    {"role": "user", "content": extraction_prompt}
                ],
                response_format={"type": "json_object"}
            )
           
            # Parse JSON response
            extracted_params = json.loads(response.choices[0].message.content)
           
            return extracted_params
       
        except Exception as e:
            print(f"Error extracting idea parameters: {str(e)}")
            # Return empty structure if extraction fails
            return {
                "Brand_Name": "",
                "Category": "",
                "Concept": "",
                "Benefits": "",
                "RTB": "",
                "Ingredients": "",
                "Features": "",
                "Theme": "",
                "Demographics": ""
            }
        
import base64
import uuid
from datetime import datetime
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.parsers import MultiPartParser
from rest_framework import status
from openai import OpenAI
import logging
 
# Assuming these are your existing imports and models
from .models import ChatHistory, ChatMessage, ConversationMemoryBuffer, Document
 
logger = logging.getLogger(__name__)
 
class GPTImageChatView(APIView):
    parser_classes = [MultiPartParser]
 
 
 
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
 
    if not OPENAI_API_KEY:
        raise ValueError("Missing required API keys in environment variables")
 
 
 
    client = OpenAI(api_key=OPENAI_API_KEY)
 
   
    def post(self, request):
        user = request.user
        try:
            # Extract data from request
            images = request.FILES.getlist("images")
            question = request.data.get("message")
            conversation_id = request.data.get('conversation_id')
            main_project_id = request.data.get('main_project_id')
            context_mode = request.data.get("context_mode", "image")
 
            print(f"Main project ID: {main_project_id}")
            print(f"Conversation ID: {conversation_id}")
            print(f"Images: {images}")
 
            # Validation
            if not images or not question:
                return Response({
                    'error': 'Images and question are required'
                }, status=status.HTTP_400_BAD_REQUEST)
 
            if not main_project_id:
                return Response({
                    'error': 'Main project ID is required'
                }, status=status.HTTP_400_BAD_REQUEST)
 
            # Log the inputs for debugging
            print(f"Image Chat request received")
            print(f"Question: {question}")
            print(f"Number of images: {len(images)}")
 
            conversation_context = ""
            if conversation_id:
                try:
                    # Try to get existing conversation
                    conversation = ChatHistory.objects.get(
                        conversation_id=conversation_id,
                        user=user
                    )
                   
                    # Get last 10 conversation messages (5 user + 5 assistant = 5 complete turns)
                    recent_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('-created_at')[:10]
                   
                    # Format for context - get full messages, not truncated
                    if recent_messages:
                        context_messages = []
                        for msg in reversed(recent_messages):  # Reverse to get chronological order
                            prefix = "User: " if msg.role == 'user' else "Assistant: "
                            # Don't truncate the messages - get full context
                            context_messages.append(f"{prefix}{msg.content}")
                       
                        conversation_context = "\n\n".join(context_messages)
                        print(f"Retrieved conversation context with {len(recent_messages)} messages")
                   
                    # Try to get memory buffer for additional context summary
                    try:
                        memory_buffer = ConversationMemoryBuffer.objects.get(conversation=conversation)
                        if memory_buffer.context_summary:
                            conversation_context += f"\n\n--- Conversation Summary ---\n{memory_buffer.context_summary}"
                    except ConversationMemoryBuffer.DoesNotExist:
                        # No memory buffer yet, that's fine
                        pass
                       
                except ChatHistory.DoesNotExist:
                    # No conversation yet, that's fine
                    pass
 
            # Prepare message content for OpenAI API
            message_content = [{"type": "text", "text": question}]
 
            # Include conversation context if available
            conversation_prompt = ""
            if conversation_context and conversation_context.strip():
                conversation_prompt = f"""
                RECENT CONVERSATION HISTORY:
                {conversation_context}
               
                The above messages are the previous parts of the same ongoing conversation with the user.
                When generating your response:
                - Consider the conversation flow and any references to previous topics
                - If the user is asking follow-up questions or clarifications, refer back to the relevant previous context
                - Maintain consistency with earlier responses and user preferences shown in the conversation
                - If the current question builds upon or relates to previous discussions, acknowledge that connection
                - Avoid repeating information already covered unless specifically asked for clarification
                - Keep the conversation natural and flowing, showing awareness of the ongoing dialogue
               
                Use this conversation history to provide more contextually aware and relevant responses.
                """
                print(f"Including conversation context in prompt (length: {len(conversation_context)} chars)")
            else:
                print("No conversation context available for this response")
 
            message_content.append({"type": "text", "text": conversation_prompt})
            # Process images and add to message content
            for image in images:
                img_bytes = image.read()
                base64_image = base64.b64encode(img_bytes).decode("utf-8")
                message_content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{base64_image}"
                    }
                })
 
            # Set fixed max_tokens
            max_tokens = 1000
 
            # Generate response using OpenAI API
            try:
                response = self.client.chat.completions.create(
                    model="gpt-4o",
                    messages=[{
                        "role": "user",
                        "content": message_content,
                    }],
                    max_tokens=max_tokens,
                )
 
                token_usage = {
                    'input_tokens': response.usage.prompt_tokens if response.usage else 0,
                    'output_tokens': response.usage.completion_tokens if response.usage else 0,
                    'total_tokens': response.usage.total_tokens if response.usage else 0
                }
 
                clean_response = response.choices[0].message.content
                print(f"Generated GPT response: {clean_response[:200]}...")
 
            except Exception as e:
                logger.error(f"OpenAI API error: {str(e)}")
                return Response({
                    'error': f'Error generating response: {str(e)}'
                }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
 
            # Generate follow-up questions based on image chat
            follow_up_questions = self.generate_image_follow_up_questions(question, clean_response)
 
            # Prepare conversation details
            conversation_id = conversation_id or str(uuid.uuid4())
            title = f"Image Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"
 
            # Create or retrieve conversation
            conversation, created = ChatHistory.objects.get_or_create(
                user=user,
                conversation_id=conversation_id,
                main_project_id=main_project_id,
                defaults={
                    'title': title,
                    'is_active': True,
                    'follow_up_questions': follow_up_questions,
                }
            )
 
            if created:
                transaction_result = self.create_conversation_transaction(
                    user, conversation, main_project_id,
                    None, None, None, token_usage # Image chat specific params
                )
               
                # Check if token limit was exceeded
                if isinstance(transaction_result, dict) and transaction_result.get('error'):
                    # Instead of deleting conversation and returning HTTP error,
                    # create user message and error response message
                   
                    # Create user message
                    user_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='user',
                        content=question
                    )
                   
                    # Create AI error response message
                    error_message = transaction_result.get('message', 'Token limit exceeded')
                    ai_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='assistant',
                        content=error_message,
                        sources="System",
                        citations=[],
                    )
                   
                    # Update memory buffer
                    memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                        conversation=conversation
                    )
                   
                    # Get all messages for this conversation
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                   
                    # Update memory with all messages
                    memory_buffer.update_memory(all_messages)
                   
                    # Update conversation details
                    conversation.title = title
                    conversation.follow_up_questions = []  # No follow-up questions for error
                    conversation.save()
                   
                    # Get all messages for response
                    message_list = []
                    for message in all_messages:
                        message_data = {
                            'id': message.id,
                            'role': message.role,
                            'content': message.content,
                            'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': message.citations or [],
                            'sources_info': getattr(message, 'sources', ''),
                            'extracted_urls': getattr(message, 'extracted_urls', []),
                            'context_mode': 'image' if message.role == 'assistant' else None,
                        }
                        message_list.append(message_data)
                   
                    response_data = {
                        'response': error_message,
                        'follow_up_questions': [],
                        'conversation_id': str(conversation.conversation_id),
                        'citations': [],
                        'active_document_id': None,
                        'sources_info': "System",
                        'url_content_used': False,
                        'extracted_urls': [],
                        'messages': message_list,
                        'image_count': len(images),
                        'context_mode': context_mode,
                    }
                   
                    # Return successful response with error message as content
                    return Response(response_data, status=status.HTTP_200_OK)
                   
            else:
                transaction_result = self.update_conversation_transaction(conversation, None, token_usage)
               
                # Check if token limit was exceeded
                if isinstance(transaction_result, dict) and transaction_result.get('error'):
                    # Create user message
                    user_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='user',
                        content=question
                    )
                   
                    # Create AI error response message
                    error_message = transaction_result.get('message', 'Token limit exceeded')
                    ai_message = ChatMessage.objects.create(
                        chat_history=conversation,
                        role='assistant',
                        content=error_message,
                        sources="System",
                        citations=[],
                    )
                   
                    # Update memory buffer
                    memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                        conversation=conversation
                    )
                   
                    # Get all messages for this conversation
                    all_messages = ChatMessage.objects.filter(
                        chat_history=conversation
                    ).order_by('created_at')
                   
                    # Update memory with all messages
                    memory_buffer.update_memory(all_messages)
                   
                    # Update conversation details
                    conversation.title = title
                    conversation.follow_up_questions = []  # No follow-up questions for error
                    conversation.save()
                   
                    # Get all messages for response
                    message_list = []
                    for message in all_messages:
                        message_data = {
                            'id': message.id,
                            'role': message.role,
                            'content': message.content,
                            'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                            'citations': message.citations or [],
                            'sources_info': getattr(message, 'sources', ''),
                            'extracted_urls': getattr(message, 'extracted_urls', []),
                            'context_mode': 'image' if message.role == 'assistant' else None,
                        }
                        message_list.append(message_data)
                   
                    response_data = {
                        'response': error_message,
                        'follow_up_questions': [],
                        'conversation_id': str(conversation.conversation_id),
                        'citations': [],
                        'active_document_id': None,
                        'sources_info': "System",
                        'url_content_used': False,
                        'extracted_urls': [],
                        'messages': message_list,
                        'image_count': len(images),
                        'context_mode': context_mode,
                    }
                   
                    # Return successful response with error message as content
                    return Response(response_data, status=status.HTTP_200_OK)
 
            # Continue with normal flow if no token limit exceeded...
            # Create user message
            user_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='user',
                content=question
            )
 
            # Create user message
            user_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='user',
                content=question
            )
 
            # Create AI response message
            ai_message = ChatMessage.objects.create(
                chat_history=conversation,
                role='assistant',
                content=clean_response,
                sources="Image Analysis",
                citations=[],
            )
 
            # Update or create memory buffer after saving the new messages
            memory_buffer, created = ConversationMemoryBuffer.objects.get_or_create(
                conversation=conversation
            )
 
            # Get all messages for this conversation (including the new ones we just created)
            all_messages = ChatMessage.objects.filter(
                chat_history=conversation
            )
 
            # Update memory with all messages - this will now properly populate the buffer
            memory_buffer.update_memory(all_messages)
            print(f"Updated memory buffer for conversation {conversation_id}")
 
            if main_project_id:
                update_project_timestamp(main_project_id, request.user)
 
            conversation.title = title
            conversation.follow_up_questions = follow_up_questions
            conversation.save()
 
            # Get all messages for response
            message_list = []
            for message in all_messages:
                message_data = {
                    'id': message.id,
                    'role': message.role,
                    'content': message.content,
                    'created_at': message.created_at.strftime('%Y-%m-%d %H:%M'),
                    'citations': message.citations or [],
                    'sources_info': getattr(message, 'sources', ''),
                    'extracted_urls': getattr(message, 'extracted_urls', []),
                    'context_mode': 'image' if message.role == 'assistant' else None,
                }
                message_list.append(message_data)
 
            response_data = {
                'response': clean_response,
                'follow_up_questions': follow_up_questions,
                'conversation_id': str(conversation.conversation_id),
                'citations': [],
                'active_document_id': None,
                'sources_info': "Image Analysis",
                'url_content_used': False,
                'extracted_urls': [],
                'messages': message_list,
                'image_count': len(images),
                'context_mode': context_mode,
            }
 
            print("\n--- Image Chat Interaction Logged ---")
            print(f"User Question: {question}")
            print(f"Number of images: {len(images)}")
            print(f"Assistant Response: {clean_response[:500]}...")
            print(f"Follow-up Questions: {len(follow_up_questions)}")
            print("-----------------------------\n")
 
            return Response(response_data, status=status.HTTP_200_OK)
 
        except Exception as e:
            logger.error(f"Unexpected error in GPTImageChatView: {str(e)}", exc_info=True)
            return Response(
                {'error': f'An unexpected error occurred: {str(e)}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )
 
   
    def generate_image_follow_up_questions(self, question, response):
        """Generate follow-up questions specific to image analysis"""
        base_questions = [
            "Can you analyze any other aspects of the image?",
            "What other details can you identify in the image?",
            "Are there any specific elements you'd like me to focus on?"
        ]
       
        # You can enhance this with more sophisticated logic
        # based on the question and response content
        try:
            if "color" in question.lower():
                base_questions.append("Can you describe the color scheme in more detail?")
            if "text" in question.lower():
                base_questions.append("Is there any other text visible in the image?")
            if "object" in question.lower():
                base_questions.append("What other objects can you identify?")
        except:
            pass
       
        return base_questions[:3]  # Return first 3 questions
 
    def create_conversation_transaction(self, user, conversation, main_project_id, use_web_knowledge, response_format, response_length, token_usage=None):
        """Create transaction record for conversation with token limit validation"""
        try:
            # Get or create user API tokens record
            user_api_tokens, created = UserAPITokens.objects.get_or_create(
                user=user,
                defaults={
                    'page_limit': 20,
                    'token_limit': 1_000_000
                }
            )
           
            # Calculate current token usage for this user
            current_total_tokens = self.get_user_total_token_usage(user)
           
            # Get tokens for this request
            request_tokens = 0
            if token_usage:
                request_tokens = token_usage.get('total_tokens', 0)
                if request_tokens == 0:  # Fallback calculation
                    request_tokens = token_usage.get('input_tokens', 0) + token_usage.get('output_tokens', 0)
           
            # Check if adding these tokens would exceed the limit
            new_total = current_total_tokens + request_tokens
           
            print(f"Token usage check - Current: {current_total_tokens}, Request: {request_tokens}, New Total: {new_total}, Limit: {user_api_tokens.token_limit}")
           
            if new_total > user_api_tokens.token_limit:
                # Return error data instead of creating transaction
                remaining_tokens = max(0, user_api_tokens.token_limit - current_total_tokens)
                error_data = {
                    'error': 'token_limit_exceeded',
                    'message': f"You have reached your conversation token limit of {user_api_tokens.token_limit:,} tokens. Current usage: {current_total_tokens:,} tokens. Remaining: {remaining_tokens:,} tokens. The current message exceeds your remaining token quota and cannot be processed. Please limit your usage or request for reset at - contact@klarifai.ai",
                    'current_usage': current_total_tokens,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': remaining_tokens,
                    'requested_tokens': request_tokens
                }
                print(f"Token limit exceeded for user {user.id}: {new_total} > {user_api_tokens.token_limit}")
                return error_data
           
            # Get the project
            main_project = Project.objects.get(id=main_project_id, user=user)
           
            # Create main transaction record
            user_transaction = UserTransaction.objects.create(
                user=user,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                conversation_title=conversation.title,
                conversation_id=conversation.conversation_id,
                main_project=main_project,
                metadata={
                    'creation_timestamp': timezone.now().isoformat(),
                    'web_knowledge_used': use_web_knowledge,
                    'response_format': response_format,
                    'response_length': response_length,
                    'token_usage': token_usage or {},
                    'token_limit_check': {
                        'previous_total': current_total_tokens,
                        'request_tokens': request_tokens,
                        'new_total': new_total,
                        'limit': user_api_tokens.token_limit
                    }
                }
            )
           
            # Extract token information if available
            input_tokens = token_usage.get('input_tokens', 0) if token_usage else 0
            output_tokens = token_usage.get('output_tokens', 0) if token_usage else 0
           
            # Create detailed conversation transaction
            ConversationTransaction.objects.create(
                user_transaction=user_transaction,
                message_count=1,  # Initial message
                question_count=1,  # First question
                input_api_tokens=input_tokens,
                output_api_tokens=output_tokens,
                web_knowledge_used=True,
                response_format="image",
                response_length="image"
            )
           
            print(f"Transaction recorded for conversation: {conversation.title}")
            print(f"Token usage - Input: {input_tokens}, Output: {output_tokens}, Total: {input_tokens + output_tokens}")
            print(f"User total token usage: {new_total}/{user_api_tokens.token_limit}")
           
            # Return success data
            return {
                'success': True,
                'current_usage': new_total,
                'token_limit': user_api_tokens.token_limit,
                'remaining_tokens': user_api_tokens.token_limit - new_total
            }
           
        except Exception as e:
            print(f"Error creating conversation transaction: {str(e)}")
            return {
                'error': 'transaction_error',
                'message': f'Error creating conversation transaction: {str(e)}'
            }
 
    def get_user_total_token_usage(self, user):
        """Calculate total token usage for a user across all conversations"""
        from django.db.models import Sum,Count
        try:
            from django.db.models import Sum,Count
            # Get all conversation transactions for this user
            total_tokens = ConversationTransaction.objects.filter(
                user_transaction__user=user,
                user_transaction__is_active=True
            ).aggregate(
                total_input=Sum('input_api_tokens'),
                total_output=Sum('output_api_tokens')
            )
           
            input_total = total_tokens.get('total_input') or 0
            output_total = total_tokens.get('total_output') or 0
           
            return input_total + output_total
           
        except Exception as e:
            print(f"Error calculating user token usage: {str(e)}")
            return 0
 
    def update_conversation_transaction(self, conversation, use_web_knowledge=True, token_usage=None):
        """Update existing conversation transaction when new messages are added with token limit validation"""
        try:
            # Get user from conversation
            user = conversation.user
           
            # Get or create user API tokens record
            user_api_tokens, created = UserAPITokens.objects.get_or_create(
                user=user,
                defaults={
                    'page_limit': 20,
                    'token_limit': 1_000_000
                }
            )
           
            # Calculate current token usage for this user
            current_total_tokens = self.get_user_total_token_usage(user)
           
            # Get tokens for this request
            request_tokens = 0
            if token_usage:
                request_tokens = token_usage.get('total_tokens', 0)
                if request_tokens == 0:  # Fallback calculation
                    request_tokens = token_usage.get('input_tokens', 0) + token_usage.get('output_tokens', 0)
           
            # Check if adding these tokens would exceed the limit
            new_total = current_total_tokens + request_tokens
           
            print(f"Token usage update check - Current: {current_total_tokens}, Request: {request_tokens}, New Total: {new_total}, Limit: {user_api_tokens.token_limit}")
           
            if new_total > user_api_tokens.token_limit:
                # Return error data instead of updating transaction
                remaining_tokens = max(0, user_api_tokens.token_limit - current_total_tokens)
                error_data = {
                    'error': 'token_limit_exceeded',
                    'message': f"You have reached your conversation token limit of {user_api_tokens.token_limit:,} tokens. Current usage: {current_total_tokens:,} tokens. Remaining: {remaining_tokens:,} tokens. The current message exceeds your remaining token quota and cannot be processed. Please limit your usage or request for reset at - contact@klarifai.ai",
                    'current_usage': current_total_tokens,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': remaining_tokens,
                    'requested_tokens': request_tokens
                }
                print(f"Token limit exceeded during update for user {user.id}: {new_total} > {user_api_tokens.token_limit}")
                return error_data
           
            # Find existing transaction for this conversation
            user_transaction = UserTransaction.objects.filter(
                conversation_id=conversation.conversation_id,
                transaction_type=TransactionType.CONVERSATION_CREATE,
                is_active=True
            ).first()
           
            if user_transaction and hasattr(user_transaction, 'notebook_conversation_details'):
                # Get the conversation details
                conv_details = user_transaction.notebook_conversation_details
               
                # Update message count and question count
                conv_details.message_count += 1
                conv_details.question_count += 1  # Each update represents a new question
               
                # Update token usage if provided
                if token_usage:
                    conv_details.input_api_tokens += token_usage.get('input_tokens', 0)
                    conv_details.output_api_tokens += token_usage.get('output_tokens', 0)
               
                if use_web_knowledge:
                    conv_details.web_knowledge_used = True
                   
                conv_details.save()
               
                # Update metadata with token usage
                if token_usage:
                    if 'token_history' not in user_transaction.metadata:
                        user_transaction.metadata['token_history'] = []
                    user_transaction.metadata['token_history'].append({
                        'timestamp': timezone.now().isoformat(),
                        'input_tokens': token_usage.get('input_tokens', 0),
                        'output_tokens': token_usage.get('output_tokens', 0),
                        'total_tokens': token_usage.get('total_tokens', 0)
                    })
               
                user_transaction.metadata['last_message_timestamp'] = timezone.now().isoformat()
                user_transaction.metadata['token_limit_check'] = {
                    'previous_total': current_total_tokens,
                    'request_tokens': request_tokens,
                    'new_total': new_total,
                    'limit': user_api_tokens.token_limit
                }
                user_transaction.save()
               
                print(f"Updated conversation transaction - Questions: {conv_details.question_count}, Total tokens: {conv_details.total_tokens_used}")
                print(f"User total token usage: {new_total}/{user_api_tokens.token_limit}")
               
                # Return success data
                return {
                    'success': True,
                    'current_usage': new_total,
                    'token_limit': user_api_tokens.token_limit,
                    'remaining_tokens': user_api_tokens.token_limit - new_total
                }
               
        except Exception as e:
            print(f"Error updating conversation transaction: {str(e)}")
            return {
                'error': 'transaction_error',
                'message': f'Error updating conversation transaction: {str(e)}'
            }