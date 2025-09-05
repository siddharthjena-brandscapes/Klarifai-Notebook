# core/views.py
from django.shortcuts import render, redirect, get_object_or_404
from django.http import JsonResponse
from django.core.exceptions import PermissionDenied
from .models import Project, ProjectModule, DocumentQAModule, IdeaGeneratorModule, Category, UserFeaturePermissions
from chat.models import UserAPITokens, UserUploadPermissions, UserModulePermissions
from django.db import transaction
from rest_framework.decorators import api_view, authentication_classes, permission_classes
from rest_framework.authentication import TokenAuthentication
from rest_framework.permissions import IsAuthenticated
from .utils import update_project_timestamp

import tempfile
import os
import fitz  # PyMuPDF for PDF handling
from pptx import Presentation
import google.generativeai as genai
from django.contrib.auth.models import User
from django.db.models import Q
 
 
@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def archive_project(request, project_id):
    """Archive a project by setting is_active to False"""
    try:
        project = get_object_or_404(Project, id=project_id, user=request.user)
        project.is_active = False
        project.save()
        
        return JsonResponse({
            'status': 'success',
            'message': f'Project "{project.name}" has been archived successfully',
            'project_id': project.id
        })
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)


@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def unarchive_project(request, project_id):
    """Restore a project from archive by setting is_active to True"""
    try:
        # Find the project even if it's archived (is_active=False)
        project = get_object_or_404(Project, id=project_id, user=request.user)
        project.is_active = True
        project.save()
        
        return JsonResponse({
            'status': 'success',
            'message': f'Project "{project.name}" has been restored successfully',
            'project_id': project.id
        })
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)


@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def archived_projects(request):
    """Get a list of all archived projects"""
    projects = Project.objects.filter(user=request.user, is_active=False)
    return JsonResponse({
        'projects': [{
            'id': project.id,
            'name': project.name,
            'description': project.description,
            'category': project.category,
            'created_at': project.created_at.isoformat(),
            'updated_at': project.updated_at.isoformat(),
            'selected_modules': project.selected_modules,
        } for project in projects]
    })
@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def enhance_prompt_with_ai(request):
    user=request.user
    try:
        # Get the user-written prompt from the request
        data = request.data
        user_prompt = data.get('prompt', '')
       
        if not user_prompt or len(user_prompt.strip()) < 10:
            return JsonResponse({
                'status': 'error',
                'message': 'Please provide a more detailed initial description (at least 10 characters).'
            }, status=400)
       
        # Initialize Gemini model
        model = initialize_gemini(user=user)
        if not model:
            return JsonResponse({
                'status': 'error',
                'message': 'Failed to initialize AI model.'
            }, status=500)
       
        # Create prompt for enhancing the user's input
        ai_prompt = f"""
        You are an expert at improving and enhancing text descriptions for AI projects.
        Please enhance and expand the following project description to make it more detailed,
        well-structured, and effective. Maintain the original intent and key points, but
        make it more professional, clear, and comprehensive.
       
        Original description:
        ---
        {user_prompt}
        ---
       
        Please provide an enhanced version that:
        1. Is well-organized with clear structure
        2. Expands on key concepts
        3. Uses professional language
        4. Provides more context where helpful
        5. Is comprehensive but concise
        """
       
        # Generate the enhanced prompt
        try:
            response = model.generate_content(ai_prompt)
            enhanced_prompt = response.text
           
            return JsonResponse({
                'status': 'success',
                'enhanced_prompt': enhanced_prompt
            })
        except Exception as e:
            print(f"Error generating enhanced prompt: {str(e)}")
            return JsonResponse({
                'status': 'error',
                'message': 'Failed to enhance the prompt. Please try again or use the original description.'
            }, status=500)
           
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
 
@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def upload_document_for_prompt(request):
    user=request.user
    try:
        if 'document' not in request.FILES:
            return JsonResponse({
                'status': 'error',
                'message': 'No document file provided'
            }, status=400)
       
        document_file = request.FILES['document']
        file_name = document_file.name.lower()
       
        # Extract text based on file type
        if file_name.endswith('.pdf'):
            document_text = extract_text_from_pdf(document_file)
        elif file_name.endswith(('.pptx', '.ppt')):
            document_text = extract_text_from_ppt(document_file)
        elif file_name.endswith('.txt'):
            document_text = extract_text_from_txt(document_file)
        else:
            return JsonResponse({
                'status': 'error',
                'message': 'Unsupported file format. Please upload PDF, PPTX, or TXT files.'
            }, status=400)
       
        # Generate prompt from document content
        generated_prompt = generate_prompt(document_text,user=user)
       
        return JsonResponse({
            'status': 'success',
            'prompt': generated_prompt
        })
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
# Helper functions for document processing
 
def extract_text_from_pdf(pdf_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
        tmp_file.write(pdf_file.read())
        tmp_path = tmp_file.name
   
    doc = fitz.open(tmp_path)
    text = ""
    for page in doc:
        text += page.get_text()
   
    os.unlink(tmp_path)  # Delete temp file
    return text
 
def extract_text_from_ppt(ppt_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(ppt_file.read())
        tmp_path = tmp_file.name
   
    presentation = Presentation(tmp_path)
    text = ""
   
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
   
    os.unlink(tmp_path)  # Delete temp file
    return text

def extract_text_from_txt(txt_file):
    # For text files, we can simply read the content directly
    text = txt_file.read()
    
    # Decode bytes to string if necessary
    if isinstance(text, bytes):
        text = text.decode('utf-8', errors='replace')
        
    return text
 
def initialize_gemini(user=None):

    user_api_tokens = UserAPITokens.objects.get(user=user)
    api_key = user_api_tokens.gemini_token 
   
    if not api_key:
        return None
   
    genai.configure(api_key=api_key)
    return genai.GenerativeModel('gemini-2.0-flash')
 
def generate_prompt(document_text, user=None):
    model = initialize_gemini(user=user)
    if not model:
        return "Failed to initialize AI model. Please provide a description manually."
   
    prompt = """
    You are an assistant that creates effective prompts based on document content.
    Create a detailed, well-structured prompt that captures the key information from the document.
    The prompt should be designed to help a chatbot provide relevant and accurate responses about the document's content.
    Structure the prompt in a way that guides the chatbot to understand the document's main topics, key points, and important details.
    dont write anything like prmopt for chatbot: (in the heading ), just write the prompt
   
    Create a prompt based on the following document content:
   
    """
   
    prompt += document_text[:8000]
   
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating prompt: {str(e)}")
        return "Failed to generate prompt. Please provide a description manually."



@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def create_project(request):
    try:
        with transaction.atomic():
            # Get data from request.data instead of request.POST
            data = request.data
            
            # Extract fields
            name = data.get('name')
            description = data.get('description')
            category = data.get('category', [])  # ✅ Change this - expect array, default to empty list
            selected_modules = data.get('selected_modules', [])
            
            # Validate required fields
            if not name or not category:  # ✅ Change validation - check if category array is not empty
                return JsonResponse({
                    'status': 'error',
                    'message': 'Missing required fields'
                }, status=400)
            
            # ✅ Add validation for category array
            if not isinstance(category, list) or len(category) == 0:
                return JsonResponse({
                    'status': 'error',
                    'message': 'Please select at least one category'
                }, status=400)
            
            # Check if project with same name already exists for this user
            if Project.objects.filter(user=request.user, name=name).exists():
                return JsonResponse({
                    'status': 'error',
                    'message': 'Project with this name already exists for your account'
                }, status=400)
            
            # Create the project
            project = Project.objects.create(
                name=name,
                description=description,
                category=category,  # ✅ This will now store array
                user=request.user,
                selected_modules=selected_modules
            )
            
            # Create project modules (unchanged)
            for module_type in selected_modules:
                project_module = ProjectModule.objects.create(
                    project=project,
                    module_type=module_type
                )
                
                if module_type == 'document-qa':
                    DocumentQAModule.objects.create(project_module=project_module)
                elif module_type == 'idea-generator':
                    IdeaGeneratorModule.objects.create(project_module=project_module)
            
            return JsonResponse({
                'status': 'success',
                'project': {
                    'id': project.id,
                    'name': project.name,
                    'description': project.description,
                    'category': project.category,  # ✅ This will return array
                    'created_at': project.created_at.isoformat(),
                    'updated_at': project.updated_at.isoformat(),
                    'selected_modules': project.selected_modules
                }
            })
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
    
    
@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def project_list(request):
    projects = Project.objects.filter(user=request.user, is_active=True)
    return JsonResponse({
        'projects': [{
            'id': project.id,
            'name': project.name,
            'description': project.description,
            'category': project.category,
            'created_at': project.created_at.isoformat(),
            'updated_at': project.updated_at.isoformat(),
            'selected_modules': project.selected_modules,
        } for project in projects]
    })

@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def project_detail(request, project_id):
    project = get_object_or_404(Project, id=project_id)
    
    if project.user != request.user:
        raise PermissionDenied
    
    return JsonResponse({
        'project': {
            'id': project.id,
            'name': project.name,
            'description': project.description,
            'category': project.category,
            'created_at': project.created_at.isoformat(),
            'updated_at': project.updated_at.isoformat(),
            'selected_modules': project.selected_modules,
            'modules': [{
                'type': module.module_type,
                'data': get_module_data(module)
            } for module in project.modules.all()]
        }
    })

def get_module_data(module):
    if module.module_type == 'document-qa':
        qa_module = module.document_qa
        return {
            'documents': [{
                'id': doc.id,
                'filename': doc.filename
            } for doc in qa_module.documents.all()],
            'chat_histories': [{
                'id': chat.id,
                'title': chat.title
            } for chat in qa_module.chat_histories.all()]
        }
    elif module.module_type == 'idea-generator':
        idea_module = module.idea_generator
        return {
            'product_ideas': [{
                'id': idea.id,
                'product': idea.product,
                'brand': idea.brand
            } for idea in idea_module.product_ideas.all()]
        }
    return {}

@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def delete_project(request, project_id):
    if request.method == 'POST':
        project = get_object_or_404(Project, id=project_id, user=request.user)
        project_name = project.name  # Store project name before deleting
        project.delete()  # Permanently delete the project
           
        return JsonResponse({'status': 'success'})
   
    return JsonResponse({'status': 'error', 'message': 'Invalid request method'})

from rest_framework import status
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from django.db import transaction
from .models import Project, ProjectModule

@api_view(['PUT'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def update_project(request, project_id):
    try:
        project = get_object_or_404(Project, id=project_id, user=request.user)
        
        # Update project fields
        project.name = request.data.get('name', project.name)
        project.description = request.data.get('description', project.description)
        project.category = request.data.get('category', project.category)
        
        # Update selected_modules if provided
        if 'selected_modules' in request.data:
            project.selected_modules = request.data['selected_modules']
            
        project.save()

        update_project_timestamp(project_id, request.user)
        
        return JsonResponse({
            'status': 'success',
            'project': {
                'id': project.id,
                'name': project.name,
                'description': project.description,
                'category': project.category,
                'created_at': project.created_at.isoformat(),
                'updated_at': project.updated_at.isoformat(),
                'selected_modules': project.selected_modules
            }
        })
        
    except Project.DoesNotExist:
        return JsonResponse({
            'status': 'error',
            'message': 'Project not found'
        }, status=404)
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)

@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def create_category_by_user(request):
    """Authenticated user creates their own personal category"""
    try:
        data = request.data
        name = data.get('name')
 
        if not name:
            return JsonResponse({
                'status': 'error',
                'message': 'Missing required field: name'
            }, status=400)
 
        # Ensure this category name doesn't already exist for the user
        if Category.objects.filter(user=request.user, name=name).exists():
            return JsonResponse({
                'status': 'error',
                'message': f'Category "{name}" already exists for your account'
            }, status=400)
 
        # Create category for current user
        category = Category.objects.create(
            name=name,
            user=request.user
        )
 
        return JsonResponse({
            'status': 'success',
            'message': f'Category "{name}" created successfully',
            'category': {
                'id': category.id,
                'name': category.name,
                'user_id': category.user.id,
                'username': category.user.username,
                'created_at': category.created_at.isoformat(),
            }
        })
 
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)



@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def create_category_for_user(request):
    """Admin creates a category for a specific user or as a global category"""
    try:
        if request.user.username != 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
 
        data = request.data
        name = data.get('name')
        user_id = data.get('user_id')  # Optional
 
        if not name:
            return JsonResponse({
                'status': 'error',
                'message': 'Missing required field: name'
            }, status=400)
 
        target_user = None
        if user_id:
            target_user = get_object_or_404(User, id=user_id)
 
        # Check if category already exists for this user or globally
        if Category.objects.filter(user=target_user, name=name).exists():
            return JsonResponse({
                'status': 'error',
                'message': f'Category \"{name}\" already exists for this user'
            }, status=400)
 
        category = Category.objects.create(name=name, user=target_user)
 
        return JsonResponse({
            'status': 'success',
            'message': f'Category \"{name}\" created successfully for {target_user.username if target_user else "Global"}',
            'category': {
                'id': category.id,
                'name': category.name,
                'user_id': category.user.id if category.user else None,
                'username': category.user.username if category.user else None,
                'created_at': category.created_at.isoformat(),
            }
        })
 
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
 
@api_view(['PUT'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def update_category(request, category_id):
    """Admin updates a category"""
    try:
        if request.user.username != 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
 
        category = get_object_or_404(Category, id=category_id)
        data = request.data
 
        if 'name' in data:
            # Check if a category with the same name already exists for the same user (or global)
            if Category.objects.filter(
                user=category.user,
                name=data['name']
            ).exclude(id=category_id).exists():
                return JsonResponse({
                    'status': 'error',
                    'message': f"Category \"{data['name']}\" already exists for this user"
                }, status=400)
 
            category.name = data['name']
 
        category.save()
 
        return JsonResponse({
            'status': 'success',
            'message': 'Category updated successfully',
            'category': {
                'id': category.id,
                'name': category.name,
                'user_id': category.user.id if category.user else None,
                'username': category.user.username if category.user else 'Global',
                'updated_at': category.updated_at.isoformat(),
                'is_global': category.user is None,
            }
        })
 
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
 
@api_view(['DELETE'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def delete_category(request, category_id):
    """Admin deletes a category"""
    try:
        if request.user.username != 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
 
        category = get_object_or_404(Category, id=category_id)
        category_name = category.name
        username = category.user.username if category.user else 'Global'
 
        # Soft delete
        category.is_active = False
        category.save()
 
        return JsonResponse({
            'status': 'success',
            'message': f'Category \"{category_name}\" deleted successfully for {username}'
        })
 
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
 
@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def get_all_user_categories(request):
    """Admin gets all categories, including global and user-specific"""
    try:
        # Only admin can access this
        if request.user.username != 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
 
        categories = Category.objects.filter(is_active=True).select_related('user').order_by('name')
 
        return JsonResponse({
            'status': 'success',
            'categories': [{
                'id': category.id,
                'name': category.name,
                'user_id': category.user.id if category.user else None,
                'username': category.user.username if category.user else 'Global',
                'user_email': category.user.email if category.user else 'N/A',
                'created_at': category.created_at.isoformat(),
                'updated_at': category.updated_at.isoformat(),
                'is_global': category.user is None,
            } for category in categories]
        })
 
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
 
   
 
@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def get_user_categories(request, user_id=None):
    """Get categories for a specific user or current user, including global ones"""
    try:
        if user_id:
            if not request.user.username == 'admin':
                return JsonResponse({
                    'status': 'error',
                    'message': 'Permission denied'
                }, status=403)
            target_user = get_object_or_404(User, id=user_id)
        else:
            target_user = request.user
 
        categories = Category.objects.filter(
            Q(user=target_user) | Q(user__isnull=True),
            is_active=True
        ).order_by('name')
 
        return JsonResponse({
            'status': 'success',
            'categories': [{
                'id': category.id,
                'name': category.name,
                'is_global': category.user is None,
                'created_at': category.created_at.isoformat(),
                'updated_at': category.updated_at.isoformat(),
            } for category in categories]
        })
 
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
    

@api_view(['PATCH'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def update_user_right_panel_permissions(request, user_id):
    try:
        if not request.user.username == 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
            
        target_user = get_object_or_404(User, id=user_id)
        feature_permissions = UserFeaturePermissions.get_or_create_for_user(target_user)
        
        data = request.data
        print(f"=== UPDATING USER {target_user.username} ===")
        print(f"Received data: {data}")
        
        if 'disabled_features' in data:
            print(f"Before update: {feature_permissions.to_disabled_dict()}")
            feature_permissions.update_from_disabled_dict(data['disabled_features'])
            feature_permissions.created_by = request.user
            feature_permissions.save()
            print(f"After update: {feature_permissions.to_disabled_dict()}")
            
        return JsonResponse({
            'status': 'success',
            'message': 'Right panel permissions updated successfully',
            'data': {
                'success': True,
                'disabled_features': feature_permissions.to_disabled_dict()
            }
        })
        
    except Exception as e:
        import traceback
        print(f"Error in update_user_right_panel_permissions: {traceback.format_exc()}")
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def get_current_user_right_panel_permissions(request):
    """Get current user's RIGHT PANEL permissions only"""
    try:
        # Get only right panel feature permissions
        feature_permissions = UserFeaturePermissions.get_or_create_for_user(request.user)
        
        return JsonResponse({
            'status': 'success',
            'data': {
                'user_id': request.user.id,
                'username': request.user.username,
                'disabled_features': feature_permissions.to_disabled_dict(),
            }
        })
        
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)


@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def get_all_users_admin(request):
    try:
        if not request.user.username == 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
           
        users = User.objects.all()
        users_data = []
       
        for user in users:
            api_tokens, _ = UserAPITokens.objects.get_or_create(user=user)
 
            # Get module permissions from chat app
            try:
                module_permissions = UserModulePermissions.objects.get(user=user)
                disabled_modules = module_permissions.disabled_modules
            except UserModulePermissions.DoesNotExist:
                disabled_modules = {}
           
            # Get upload permissions from chat app
            try:
                upload_permissions = UserUploadPermissions.objects.get(user=user)
                can_upload = upload_permissions.can_upload
            except UserUploadPermissions.DoesNotExist:
                can_upload = True
         
           
            # Get RIGHT PANEL feature permissions
            feature_permissions = UserFeaturePermissions.get_or_create_for_user(user)
            disabled_features_data = feature_permissions.to_disabled_dict()
           
            # 🔥 Make sure to include disabled_features in the response
            users_data.append({
                'id': user.id,
                'username': user.username,
                'email': user.email,
                'api_tokens': {
                    'nebius_token': api_tokens.nebius_token or '',
                    'gemini_token': api_tokens.gemini_token or '',
                    'llama_token': api_tokens.llama_token or '',
                    'token_limit': api_tokens.token_limit or '',
                    'page_limit': api_tokens.page_limit or '',
                },
               'disabled_modules': disabled_modules,  # From chat app
            'upload_permissions': {'can_upload': can_upload},
 
               
                # ✅ ADD THIS LINE - This is what's missing!
                'disabled_features': disabled_features_data,
            })
           
        return JsonResponse(users_data, safe=False)
       
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)
    

    
@api_view(['PATCH'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def update_user_module_permissions(request, user_id):
    try:
        if not request.user.username == 'admin':
            return JsonResponse({
                'status': 'error',
                'message': 'Permission denied'
            }, status=403)
            
        target_user = get_object_or_404(User, id=user_id)
        api_tokens, _ = UserAPITokens.objects.get_or_create(user=target_user)
        
        data = request.data
        if 'disabled_modules' in data:
            api_tokens.disabled_modules = data['disabled_modules']
            api_tokens.save()
            
        return JsonResponse({
            'status': 'success',
            'message': 'Module permissions updated successfully'
        })
        
    except Exception as e:
        return JsonResponse({
            'status': 'error',
            'message': str(e)
        }, status=500)

@api_view(['PATCH'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def update_user_upload_permissions(request, user_id):
    try:
        if not request.user.username == 'admin':
            return JsonResponse({'status': 'error', 'message': 'Permission denied'}, status=403)
            
        target_user = get_object_or_404(User, id=user_id)
        
        # Get or create upload permissions from chat app
        upload_permissions, created = UserUploadPermissions.objects.get_or_create(
            user=target_user,
            defaults={'can_upload': True}
        )
        
        data = request.data
        if 'can_upload' in data:
            upload_permissions.can_upload = data['can_upload']
            upload_permissions.save()
            
            print(f"Admin user {request.user.username} set upload permissions for user {target_user.username} to {data['can_upload']}")
            
        return JsonResponse({'status': 'success', 'data': {'success': True}})
        
    except Exception as e:
        print(f"Error updating upload permissions: {str(e)}")
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)




 
from rest_framework.decorators import api_view
from rest_framework.response import Response
from django.db.models import Min
from .models import Project
from notebook.models import ChatHistory
from ideaGen.models import ProductIdea2
from notebook.models import UserTransaction, TransactionType
@api_view(['GET'])
def get_first_activities(request):
    """Get timestamps of first activities across all modules"""
   
    # Get first project creation
    first_project = Project.objects.select_related('user').order_by('created_at').first()
   
    # Get first notebook question using both transaction tables
    first_conversation = UserTransaction.objects.filter(
        transaction_type=TransactionType.CONVERSATION_CREATE,
    ).select_related(
        'user',
        'notebook_conversation_details'  # Join with ConversationTransaction
    ).order_by('created_at').first()
   
    # Get first idea generation
    first_idea = ProductIdea2.objects.select_related('user').order_by('created_at').first()
 
    return Response({
        'first_project': {
            'username': first_project.user.username if first_project else None,
            'timestamp': first_project.created_at if first_project else None,
            'project_name': first_project.name if first_project else None
        },
        'first_question': {
            'username': first_conversation.user.username if first_conversation else None,
            'timestamp': first_conversation.created_at if first_conversation else None,
            'question_count': first_conversation.notebook_conversation_details.question_count if first_conversation and hasattr(first_conversation, 'notebook_conversation_details') else 0,
            'conversation_title': first_conversation.conversation_title if first_conversation else None
        },
        'first_idea': {
            'username': first_idea.user.username if first_idea else None,
            'timestamp': first_idea.created_at if first_idea else None,
            'idea_name': first_idea.product if first_idea else None
        }
    })

@api_view(['GET'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
def get_all_activities(request):
    """Get all activities chronologically organized"""
   
    # Fetch all projects
    projects = Project.objects.select_related('user').order_by('created_at')
    
    # Fetch all conversations with their details
    conversations = UserTransaction.objects.filter(
        transaction_type=TransactionType.CONVERSATION_CREATE,
    ).select_related(
        'user',
        'notebook_conversation_details'
    ).order_by('created_at')
    
    # Fetch all ideas
    ideas = ProductIdea2.objects.select_related('user', 'project').order_by('created_at')
   
    activities = []
   
    # Add projects
    for project in projects:
        activities.append({
            'type': 'Project Creation',
            'username': project.user.username,
            'details': f"Created project '{project.name}'",
            'timestamp': project.created_at
        })
   
    # Add conversations
    for conv in conversations:
        conv_details = getattr(conv, 'notebook_conversation_details', None)
        activities.append({
            'type': 'Notebook Question',
            'username': conv.user.username,
            'details': f"Started '{conv.conversation_title or 'Untitled'}' with {conv_details.question_count if conv_details else 0} questions",
            'timestamp': conv.created_at,
            'metadata': {
                'question_count': conv_details.question_count if conv_details else 0,
                'message_count': conv_details.message_count if conv_details else 0
            }
        })
   
    # Add ideas
    for idea in ideas:
        try:
            project_name = idea.project.name if idea.project else "Unknown Project"
            activities.append({
                'type': 'Idea Generation',
                'username': idea.user.username,
                'details': f"Created idea session '{project_name}'",
                'timestamp': idea.created_at
            })
        except Exception as e:
            activities.append({
                'type': 'Idea Generation',
                'username': idea.user.username,
                'details': f"Generated idea '{idea.product}'",
                'timestamp': idea.created_at
            })
   
    # Sort all activities by timestamp
    activities.sort(key=lambda x: x['timestamp'], reverse=True)
   
    return Response({
        'status': 'success',
        'activities': activities
    })